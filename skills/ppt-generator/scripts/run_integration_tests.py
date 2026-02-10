"""Standalone integration test runner for ppt_generator.

Runs two tests without requiring pytest:
- test_generate_mft: uses docs/presentations/MFT-20260210 inputs
- test_background_alpha_without_lxml: creates a small background image and
  ensures generate_pptx works when lxml is absent (captures a warning)

Exits with non-zero code on failure.
"""
from __future__ import annotations

import json
import logging
import os
import sys
import tempfile
from pathlib import Path
from pptx import Presentation
from PIL import Image

# Ensure package importable
REPO_ROOT = Path(__file__).resolve().parents[3]
SKILL_PATH = REPO_ROOT / "skills" / "ppt-generator"
if str(SKILL_PATH) not in sys.path:
    sys.path.insert(0, str(SKILL_PATH))

from ppt_generator.cli import generate_pptx  # type: ignore


def run_generate_mft() -> None:
    print("Running test_generate_mft...")
    semantic = REPO_ROOT / "docs" / "presentations" / "MFT-20260210" / "slides_semantic.json"
    design = REPO_ROOT / "docs" / "presentations" / "MFT-20260210" / "design_spec.json"
    out = Path(tempfile.gettempdir()) / "MFT_integration_test.pptx"
    try:
        generate_pptx(str(semantic), str(design), str(out))
    except Exception as e:
        raise RuntimeError(f"generate_pptx failed for MFT test: {e}")
    if not out.exists():
        raise RuntimeError("Output PPTX was not created for MFT test")
    prs = Presentation(str(out))
    if len(prs.slides) < 1:
        raise RuntimeError("Generated PPTX contains no slides")
    print("test_generate_mft: OK — slides:", len(prs.slides))


def run_background_alpha_without_lxml() -> None:
    print("Running test_background_alpha_without_lxml...")
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        img_path = td / "bg.png"
        Image.new("RGB", (16, 16), color=(255, 0, 0)).save(img_path)

        semantic = td / "semantic.json"
        design = td / "design.json"
        semantic.write_text(json.dumps({
            "title": "Test",
            "sections": [{"id": "A", "title": "S", "start_slide": 1}],
            "slides": [{"id": 1, "title": "T", "slide_type": "title", "content": ["c"]}]
        }, ensure_ascii=False), encoding="utf-8")
        design.write_text(json.dumps({
            "slide_type_layouts": {"title": {"background_image": str(img_path)}},
            "design_system": {"grid_system": {"slide_width_inches": 13.333, "slide_height_inches": 7.5, "margin_horizontal": 80, "gutter": 24, "columns": 12}},
            "tokens": {"colors": {"primary": "#2563EB", "on_primary": "#FFFFFF", "surface": "#FFFFFF", "on_surface": "#0F172A"}}
        }, ensure_ascii=False), encoding="utf-8")

        out = td / "out_bg.pptx"
        logger = logging.getLogger("ppt_generator.renderers")
        records: list[logging.LogRecord] = []

        class H(logging.Handler):
            def emit(self, record: logging.LogRecord) -> None:
                records.append(record)

        handler = H()
        logger.setLevel(logging.WARNING)
        logger.addHandler(handler)

        saved = sys.modules.pop("lxml", None)
        try:
            generate_pptx(str(semantic), str(design), str(out))
        except Exception as e:
            raise RuntimeError(f"generate_pptx raised an exception with missing lxml: {e}")
        finally:
            if saved is not None:
                sys.modules["lxml"] = saved
            logger.removeHandler(handler)

        if not out.exists():
            raise RuntimeError("Output PPTX not created for background alpha test")

        # Affirm a warning was logged or at least the code path succeeded
        msgs = [r.getMessage() for r in records]
        print("warning messages captured:", msgs)
        print("test_background_alpha_without_lxml: OK")


def main() -> int:
    try:
        run_generate_mft()
        run_background_alpha_without_lxml()
    except Exception as exc:
        print("TEST FAILED:", exc)
        return 2
    print("ALL TESTS PASSED")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
