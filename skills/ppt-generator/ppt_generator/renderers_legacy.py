"""Renderers for slides and components.

This module reuses helpers and grid to provide all rendering functions
exported by the original monolithic script.
"""
from typing import Any, Dict, List
import json
import os
import re
import logging

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from .helpers import (
    get_color, get_color_hex, get_font_size, hex_to_rgb, px_to_inches
)
from .grid import GridSystem


# For brevity we copy selectively the rendering helpers from the original
# script, focusing on title bar, bottom bar, speaker notes, and the main
# per-slide dispatcher. Additional component renderers (kpis, bullets,
# tables, visuals) are also included as smaller functions.


def get_layout_zones(spec: Dict) -> Dict[str, float]:
    lz = spec.get('layout_zones', {})
    return {
        'title_bar_h': lz.get('title_bar_height_default', 0.55),
        'title_bar_h_narrow': lz.get('title_bar_height_narrow', 0.4),
        'bottom_bar_h': max(lz.get('bottom_bar_height', 0.25), 0.25),
        'content_margin_top': lz.get('content_margin_top', 0.12),
        'content_bottom_margin': lz.get('content_bottom_margin', 0.2),
    }


def get_title_bar_mode(spec: Dict, slide_type: str) -> str:
    if slide_type in ('title', 'section_divider'):
        return 'none'
    stl = spec.get('slide_type_layouts', {})
    entry = stl.get(slide_type, stl.get('default', {}))
    return entry.get('title_bar', 'standard')


def get_slide_override(spec: Dict, slide_id: int = None) -> Dict:
    """Get per-slide layout override from slide_overrides, if any."""
    if slide_id is not None:
        return spec.get('slide_overrides', {}).get(str(slide_id), {})
    return {}


def get_content_fill(spec: Dict, slide_type: str, slide_id: int = None) -> str:
    """Get content fill strategy, checking per-slide override first.

    Priority: slide_overrides[slide_id].content_fill > slide_type_layouts[type].content_fill > 'expand'.
    Returns 'expand', 'center', or 'top-align'.
    """
    override = get_slide_override(spec, slide_id)
    if 'content_fill' in override:
        return override['content_fill']
    stl = spec.get('slide_type_layouts', {})
    entry = stl.get(slide_type, stl.get('default', {}))
    return entry.get('content_fill', 'expand')


def get_max_card_h(spec: Dict, slide_id: int = None) -> float:
    """Get max card height cap from slide_overrides. Returns 0.0 (no cap) if not set."""
    override = get_slide_override(spec, slide_id)
    return override.get('max_card_h', 0.0)


def apply_font_to_run(run: Any, spec: Dict) -> None:
    try:
        font_family = spec.get('typography_system', {}).get('font_family')
        cjk_font = spec.get('typography_system', {}).get('cjk_font_family')
        if font_family:
            run.font.name = font_family
            r_pr = getattr(run._element, 'rPr', None)
            if r_pr is not None:
                try:
                    r_pr.rFonts.set(qn('a:ascii'), font_family)
                    r_pr.rFonts.set(qn('a:hAnsi'), font_family)
                except Exception:
                    pass
        if cjk_font:
            r_pr = getattr(run._element, 'rPr', None)
            if r_pr is not None:
                try:
                    r_pr.rFonts.set(qn('a:ea'), cjk_font)
                except Exception:
                    pass
    except Exception:
        pass


def create_card_shape(slide: Any, spec: Dict, left_in: float, top_in: float, width_in: float, height_in: float):
    try:
        br = spec.get('component_library', {}).get('card', {}).get('border_radius')
    except Exception:
        br = None
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if br is None or br != 0 else MSO_SHAPE.RECTANGLE
    return slide.shapes.add_shape(shape_type, Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))


def add_shadow(shape, blur_pt: int = 6, offset_pt: int = 2) -> None:
    """Add a subtle shadow to a shape if supported by python-pptx.

    This helper is intentionally defensive: if the presentation backend does
    not support setting shadow properties, it silently no-ops to preserve
    compatibility with minimal environments.
    """
    try:
        # python-pptx does not expose a high-level shadow API consistently
        # across backends; attempt to access `.shadow` if available, else no-op.
        sh = getattr(shape, 'shadow', None)
        if sh is not None:
            # If present, attempt non-fatal properties application
            try:
                # these attributes may not exist; ignore failures
                sh.inherit = False
            except Exception:
                pass
    except Exception:
        pass


# ── Icon/Symbol System ──
# Maps keywords in labels/text to decorative Unicode symbols for visual variety
_ICON_MAP = {
    # Performance & metrics
    '效率': '⚡', 'efficiency': '⚡', '性能': '⚡', 'performance': '⚡',
    '温升': '🌡', '温度': '🌡', 'temperature': '🌡', '热': '🔥', 'thermal': '🔥',
    '频率': '〜', 'frequency': '〜', '频段': '〜',
    # Quality & reliability
    '质量': '✦', 'quality': '✦', '可靠': '✦', 'reliability': '✦',
    '可用率': '✦', 'availability': '✦', 'cpk': '✦',
    'pd': '⬡', '局放': '⬡', 'fpy': '✓', '放行率': '✓',
    # Risk & safety
    '风险': '⚠', 'risk': '⚠', '安全': '🛡', 'safety': '🛡',
    '短缺': '⚠', 'shortage': '⚠', '缓解': '↻',
    # Materials & components
    '材料': '◆', 'material': '◆', '磁芯': '◇', '纳米晶': '◈', '非晶': '◈', '粉末': '◈',
    '绕组': '⊕', 'winding': '⊕', '绝缘': '▣', 'insulation': '▣',
    # Process & manufacturing
    '制造': '⚙', 'manufacturing': '⚙', '工艺': '⚙', 'process': '⚙',
    '测试': '✎', 'test': '✎', '验证': '✔', 'validation': '✔',
    '仿真': '▷', 'simulation': '▷',
    # Commercial & business
    '成本': '¥', 'cost': '$', '预算': '¥', 'budget': '$',
    '收入': '↗', 'revenue': '↗', '商业': '◉', 'business': '◉',
    '市场': '◎', 'market': '◎', 'roi': '↗',
    # Infrastructure & systems
    '冷却': '❄', '风冷': '❄', '液冷': '❄', '被动冷却': '❄', 'cooling': '❄',
    'emc': '⊘', '电磁': '⊘', '屏蔽': '⊘',
    '充电': '⚡', 'ev': '⚡', '牵引': '⇌', '数据中心': '▦', 'sst': '⇌',
    # Organization & governance
    '组织': '⊞', '治理': '⊞', 'governance': '⊞',
    '标准': '☰', 'standard': '☰', '合规': '☰',
    '决策': '◈', 'decision': '◈',
    # Time & roadmap
    '里程碑': '◆', 'milestone': '◆', '路线图': '▸', 'roadmap': '▸',
    '运维': '⟳', 'maintenance': '⟳', '监测': '◉', 'monitoring': '◉',
    # Default fallback icons by index
    '_idx_0': '●', '_idx_1': '◆', '_idx_2': '■', '_idx_3': '▲',
    '_idx_4': '◉', '_idx_5': '⬟',
}


def get_icon_for_text(text: str, index: int = 0) -> str:
    """Find the best matching icon for given text content."""
    text_lower = text.lower()
    for keyword, icon in _ICON_MAP.items():
        if keyword.startswith('_idx_'):
            continue
        if keyword in text_lower:
            return icon
    # Fallback: use index-based icon for variety
    return _ICON_MAP.get(f'_idx_{index % 6}', '●')


def render_icon_circle(slide: Any, spec: Dict, icon: str, cx: float, cy: float, size: float = 0.40,
                       bg_token: str = 'primary', fg_token: str = 'on_primary') -> None:
    """Render a decorative circle with an icon/symbol inside."""
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx), Inches(cy),
        Inches(size), Inches(size)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = get_color(spec, bg_token)
    circle.line.fill.background()
    # Put icon text directly on the oval shape for native centering
    # Use margin_top offset to compensate for Unicode glyph ascender bias
    icon_font_size = int(size * 32)  # slightly smaller for better centering
    glyph_offset = Pt(max(1, int(size * 5)))  # push glyph down ~5% of size
    circle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    circle.text_frame.margin_top = glyph_offset
    circle.text_frame.margin_bottom = Emu(1)
    circle.text_frame.margin_left = Emu(1)
    circle.text_frame.margin_right = Emu(1)
    p = circle.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = icon
    run.font.size = Pt(icon_font_size)
    run.font.bold = True
    run.font.color.rgb = get_color(spec, fg_token)


def render_title_bar(slide: Any, spec: Dict, grid: GridSystem, title: str, slide_num: int, _total_slides: int,
                     section_label: str = '', accent_color_token: str = 'primary', mode: str = 'standard') -> float:
    if mode == 'none':
        return 0.0
    lz = get_layout_zones(spec)
    bar_h = lz['title_bar_h_narrow'] if mode == 'narrow' else lz['title_bar_h']
    title_font_pt = get_font_size(spec, 'slide_title')
    title_text_h = title_font_pt / 72.0 + 0.08
    if section_label:
        min_bar_h = 0.28 + title_text_h
    else:
        min_bar_h = 0.10 + title_text_h
    bar_h = max(bar_h, min_bar_h)

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(grid.slide_w), Inches(bar_h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = get_color(spec, accent_color_token)
    bar.line.fill.background()

    if section_label:
        tb = slide.shapes.add_textbox(
            Inches(grid.margin_h), Inches(0.06), Inches(4), Inches(0.22)
        )
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = section_label
        run.font.size = Pt(get_font_size(spec, 'section_label'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_primary')
        apply_font_to_run(run, spec)

    title_top = 0.22 if section_label else 0.10
    tb = slide.shapes.add_textbox(
        Inches(grid.margin_h), Inches(title_top),
        Inches(grid.slide_w - 2 * grid.margin_h - 1.0), Inches(bar_h - title_top - 0.05)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(get_font_size(spec, 'slide_title'))
    run.font.bold = True
    run.font.color.rgb = get_color(spec, 'on_primary')
    apply_font_to_run(run, spec)

    tb_num = slide.shapes.add_textbox(
        Inches(grid.slide_w - 1.0), Inches(title_top), Inches(0.7), Inches(0.35)
    )
    tf_num = tb_num.text_frame
    p = tf_num.paragraphs[0]
    run = p.add_run()
    run.text = f"{slide_num}"
    run.font.size = Pt(get_font_size(spec, 'page_number'))
    run.font.bold = True
    run.font.color.rgb = get_color(spec, 'on_primary')
    apply_font_to_run(run, spec)
    p.alignment = PP_ALIGN.RIGHT

    return bar_h


def render_bottom_bar(slide: Any, spec: Dict, grid: GridSystem, section_name: str, slide_num: int, total_slides: int,
                      section_index: int = 0, total_sections: int = 6, accent_token: str = 'primary') -> None:
    lz = get_layout_zones(spec)
    bar_h = lz['bottom_bar_h']
    bar_top = grid.slide_h - bar_h

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(bar_top), Inches(grid.slide_w), Inches(bar_h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = get_color(spec, 'surface_variant')
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(grid.margin_h), Inches(bar_top + 0.03), Inches(3), Inches(bar_h - 0.06))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = section_name
    run.font.size = Pt(get_font_size(spec, 'label'))
    run.font.bold = True
    run.font.color.rgb = get_color(spec, accent_token)

    prog_left = grid.slide_w / 2 - 1.5
    seg_w = 3.0 / max(total_sections, 1)
    for i in range(total_sections):
        seg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(prog_left + i * seg_w + 0.02), Inches(bar_top + bar_h / 2 - 0.03),
            Inches(seg_w - 0.04), Inches(0.06)
        )
        seg.fill.solid()
        if i <= section_index:
            seg.fill.fore_color.rgb = get_color(spec, accent_token)
        else:
            seg.fill.fore_color.rgb = get_color(spec, 'outline')
        seg.line.fill.background()

    tb = slide.shapes.add_textbox(
        Inches(grid.slide_w - grid.margin_h - 1.5), Inches(bar_top + 0.03), Inches(1.5), Inches(bar_h - 0.06)
    )
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{slide_num} / {total_slides}"
    run.font.size = Pt(get_font_size(spec, 'page_number'))
    run.font.color.rgb = get_color(spec, 'on_surface_variant')


def render_speaker_notes(slide: Any, notes_text: Any) -> None:
    if not notes_text:
        return
    text = notes_text if isinstance(notes_text, str) else json.dumps(notes_text, ensure_ascii=False)
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# Minimal component renderers used by slide renderers

def render_kpis(slide: Any, kpis: List[Dict], spec: Dict, _grid: GridSystem, left: float, top: float, width: float, avail_h: float = None) -> float:
    """Render KPI cards with value-first visual hierarchy.

    Layout (per card):
    ┌───────────────────┐
    │ ● label           │  ← Row 1: icon badge + label (small, muted — the "what")
    │                   │
    │   >= 98%          │  ← Row 2: VALUE as hero (large, bold — the "answer")
    └───────────────────┘
    """
    if not kpis:
        return 0
    n = len(kpis)
    card_gap = 0.15
    card_w = (width - card_gap * (n - 1)) / max(n, 1)

    # ── Card height: content-driven ──
    # label_row(~0.22) + gap(0.06) + value_row(~0.40) + padding(0.30) ≈ 0.98
    base_card_h = 1.00 if n >= 4 else 1.10
    KPI_H_CAP = 1.30
    if avail_h and avail_h > base_card_h + 1.5:
        card_h = min(base_card_h + 0.15, KPI_H_CAP)
    else:
        card_h = base_card_h

    # ── Font sizing adaptive to card width ──
    val_font_pt = get_font_size(spec, 'kpi_value')  # default 20
    if card_w < 2.5:
        val_font_pt = min(val_font_pt, 16)
    elif card_w < 3.5:
        val_font_pt = min(val_font_pt, 18)

    for i, kpi in enumerate(kpis):
        cx = left + i * (card_w + card_gap)
        card = create_card_shape(slide, spec, cx, top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = get_color(spec, 'primary_container')
        card.line.fill.background()

        kpi_label = kpi.get('label', '')
        kpi_icon = get_icon_for_text(kpi_label, i)
        val = kpi.get('value', '')
        trend = kpi.get('trend', '')
        trend_arrow = {'up': ' ↑', 'down': ' ↓', 'stable': ' →'}.get(trend, '')

        pad_l = 0.14  # left padding inside card
        pad_r = 0.12  # right padding

        # ── Row 1: label line with small color bar (the "what") ──
        label_y = top + 0.12
        label_h = 0.22
        # Small color accent bar (left edge indicator)
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(cx + pad_l), Inches(label_y + 0.02),
            Inches(0.04), Inches(label_h - 0.04)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = get_color(spec, 'accent_2')
        accent_bar.line.fill.background()
        # Label text — right of accent bar
        label_x = cx + pad_l + 0.10
        label_w = card_w - pad_l - pad_r - 0.10
        tb_label = slide.shapes.add_textbox(
            Inches(label_x), Inches(label_y),
            Inches(max(label_w, 0.5)), Inches(label_h)
        )
        tf_label = tb_label.text_frame
        tf_label.word_wrap = True
        tf_label.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_label = tf_label.paragraphs[0]
        p_label.alignment = PP_ALIGN.LEFT
        run_label = p_label.add_run()
        run_label.text = kpi_label
        run_label.font.size = Pt(get_font_size(spec, 'kpi_label'))
        run_label.font.color.rgb = get_color(spec, 'on_surface_variant')

        # ── Row 2: VALUE as hero (the "answer") — dominant element ──
        val_y = label_y + label_h + 0.04
        val_h = card_h - (val_y - top) - 0.08
        tb_val = slide.shapes.add_textbox(
            Inches(cx + pad_l), Inches(val_y),
            Inches(card_w - pad_l - pad_r), Inches(max(val_h, 0.35))
        )
        tf_val = tb_val.text_frame
        tf_val.word_wrap = True
        tf_val.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_val = tf_val.paragraphs[0]
        p_val.alignment = PP_ALIGN.LEFT
        run_val = p_val.add_run()
        run_val.text = f"{val}{trend_arrow}"
        # Hero value: larger font, bold, strong color
        hero_font_pt = val_font_pt + 4  # boost over base
        run_val.font.size = Pt(hero_font_pt)
        run_val.font.bold = True
        run_val.font.color.rgb = get_color(spec, 'primary')
        apply_font_to_run(run_val, spec)
    return card_h + 0.15


# Visual/placeholder renderers (plugin-based architecture)

def render_visual(slide: Any, visual: Dict, spec: Dict, _grid: GridSystem, left: float, top: float, width: float, height: float) -> None:
    """Render visual using plugin-based renderer registry.
    
    New architecture:
    1. Try plugin renderers first (native, mermaid, etc.)
    2. Fall back to legacy placeholder renderers if no plugin available
    """
    if not visual or visual.get('type') in (None, 'none'):
        return
    
    # Import registry (lazy import to avoid circular dependency)
    try:
        from .core import get_registry
        from .protocols import VisualDataProtocol
        
        # Attempt to construct VDP from visual dict
        try:
            # Try new format (with structured 'data' field)
            visual_data = VisualDataProtocol(**visual)
        except Exception:
            # Legacy format: convert placeholder_data to structured data
            pd = visual.get('placeholder_data', {})
            gantt_data = pd.get('gantt_data')
            mermaid_code = pd.get('mermaid_code')
            
            if gantt_data:
                # Has structured gantt data
                visual_data = VisualDataProtocol(
                    type=visual.get('type'),
                    title=visual.get('title'),
                    data=gantt_data,
                    placeholder_data=pd
                )
            elif mermaid_code:
                # Has mermaid code — always use mermaid type for renderer lookup
                visual_data = VisualDataProtocol(
                    type='mermaid',  # Override type to ensure mermaid renderer picks it up
                    title=visual.get('title'),
                    data={},  # Empty dict for mermaid-based visuals
                    placeholder_data=pd
                )
            else:
                # No structured data, use legacy renderers
                visual_data = None
        
        # Try plugin renderers
        if visual_data:
            registry = get_registry()
            if registry.render_with_fallback(slide, visual_data, spec, left, top, width, height):
                return  # Success!
            # If plugin failed, fall through to legacy renderers
    
    except ImportError:
        # Registry not available, use legacy renderers
        pass
    except Exception as e:
        import logging
        logging.warning(f"Plugin renderer failed, falling back to legacy: {e}")
    
    # Legacy renderers (backward compatibility)
    pd = visual.get('placeholder_data', {})
    if pd.get('chart_config'):
        render_chart_table(slide, visual, spec, left, top, width, height)
    elif pd.get('mermaid_code'):
        render_mermaid_placeholder(slide, visual, spec, left, top, width, height)
    else:
        render_visual_placeholder(slide, visual, spec, left, top, width, height)


def render_chart_table(slide: Any, visual: Dict, spec: Dict, left: float, top: float, width: float, height: float) -> None:
    """Render data as a proper PowerPoint table with clear grid lines."""
    config = visual.get('placeholder_data', {}).get('chart_config', {})
    labels = config.get('labels', [])
    series = config.get('series', [])
    rows_data = config.get('rows', [])
    if not labels:
        return

    title_offset = 0.0
    if visual.get('title'):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top - 0.30), Inches(width), Inches(0.28))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = visual['title']
        run.font.size = Pt(get_font_size(spec, 'table_header'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_surface')
        title_offset = 0.35

    table_top = top + title_offset

    # Determine table structure
    if rows_data:
        # Matrix format: labels are column headers, rows_data are data rows
        n_table_cols = len(labels)
        n_table_rows = len(rows_data) + 1  # +1 for header
        header_labels = labels
        body_rows = rows_data
    elif series:
        # Series format: check if series have names
        has_names = any(s.get('name') for s in series)
        if has_names:
            # Detect if first label is a row-header category
            first_data_len = len(series[0].get('data', []))
            if first_data_len < len(labels):
                # First label is row header (e.g. "方案"), rest are data columns
                header_labels = labels  # keep all labels
                n_table_cols = len(labels)
                n_table_rows = len(series) + 1
                body_rows = []
                for s in series:
                    row = [s.get('name', '')] + [str(v) for v in s.get('data', [])]
                    body_rows.append(row)
            else:
                # Series names become first column, labels are data columns
                header_labels = [''] + labels
                n_table_cols = len(labels) + 1
                n_table_rows = len(series) + 1
                body_rows = []
                for s in series:
                    row = [s.get('name', '')] + [str(v) for v in s.get('data', [])]
                    body_rows.append(row)
        else:
            # No names: labels are column headers, data rows below
            header_labels = labels
            n_table_cols = len(labels)
            n_table_rows = len(series) + 1
            body_rows = [[str(v) for v in s.get('data', [])] for s in series]
    else:
        return

    # Create PowerPoint table
    usable_h = height - title_offset
    tbl_shape = slide.shapes.add_table(
        n_table_rows, n_table_cols,
        Inches(left), Inches(table_top),
        Inches(width), Inches(min(usable_h, n_table_rows * 0.55))
    )
    table = tbl_shape.table

    # Style settings from design spec
    hdr_bg = get_color(spec, 'primary')
    hdr_fg = get_color(spec, 'on_primary')
    alt_bg = get_color(spec, 'surface_dim')
    border_color = get_color(spec, 'outline')
    cell_fg = get_color(spec, 'on_surface')
    hdr_font_pt = get_font_size(spec, 'table_header')
    cell_font_pt = get_font_size(spec, 'table_cell')

    def _style_cell(cell, text, is_header=False, is_row_label=False):
        """Apply consistent styling to a table cell."""
        cell.text = ''
        tf = cell.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(text)
        if is_header:
            run.font.size = Pt(hdr_font_pt)
            run.font.bold = True
            run.font.color.rgb = hdr_fg
            cell.fill.solid()
            cell.fill.fore_color.rgb = hdr_bg
        elif is_row_label:
            run.font.size = Pt(cell_font_pt)
            run.font.bold = True
            run.font.color.rgb = cell_fg
        else:
            run.font.size = Pt(cell_font_pt)
            run.font.color.rgb = cell_fg

    # Header row
    for j, lbl in enumerate(header_labels):
        _style_cell(table.cell(0, j), lbl, is_header=True)

    # Data rows
    for r, row in enumerate(body_rows):
        # Alternating row background
        for j in range(n_table_cols):
            cell = table.cell(r + 1, j)
            if r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_bg
            else:
                cell.fill.background()
            val = row[j] if j < len(row) else ''
            is_row_label = (j == 0 and any(s.get('name') for s in series)) if series else False
            _style_cell(cell, val, is_row_label=is_row_label)

    # Set borders on all cells
    try:
        from pptx.oxml.ns import qn
        for r_idx in range(n_table_rows):
            for c_idx in range(n_table_cols):
                cell = table.cell(r_idx, c_idx)
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
                    ln = tcPr.find(qn(f'a:{border_name}'))
                    if ln is not None:
                        tcPr.remove(ln)
                    ln = tcPr.makeelement(qn(f'a:{border_name}'), {})
                    ln.set('w', '19050')  # 1.5pt for clear visibility
                    solidFill = ln.makeelement(qn('a:solidFill'), {})
                    srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': str(border_color)})
                    solidFill.append(srgb)
                    ln.append(solidFill)
                    tcPr.append(ln)
    except Exception:
        pass  # Borders are nice-to-have; table is readable without them


def render_mermaid_placeholder(slide: Any, visual: Dict, spec: Dict, left: float, top: float, width: float, height: float) -> None:
    mermaid = visual.get('placeholder_data', {}).get('mermaid_code', '')
    
    # Title setup
    has_title = bool(visual.get('title'))
    
    # Preview content preparation
    preview_lines = mermaid.strip().split('\n')
    max_lines = 8
    display_lines = preview_lines[:max_lines]
    preview_text = '\n'.join(display_lines)
    if len(preview_lines) > max_lines:
        preview_text += '\n...'
        display_line_count = max_lines + 1
    else:
        display_line_count = len(display_lines)

    # Dynamic Height Calculation (Shrink-to-Fit)
    # Estimate needed height: Title space + Text space + Padding
    header_space = 0.5 if has_title else 0.0
    text_space = display_line_count * 0.35  # Approx 0.35" per line at 14pt
    padding = 0.4
    
    ideal_h = header_space + text_space + padding
    # Clamp height: Min 1.2", Max allocated 'height', Max absolute 3.8"
    card_h = max(1.2, min(height, ideal_h, 3.8))
    
    # Vertical Centering in Allocated Slot
    # If allocated height > card_h, move card down to center it
    v_offset = (height - card_h) / 2
    render_top = top + v_offset
    
    # Draw Card
    card = create_card_shape(slide, spec, left, render_top, width, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = get_color(spec, 'surface_variant')
    card.line.color.rgb = get_color(spec, 'outline')
    card.line.width = Pt(1)
    
    # 1. Render Title
    if has_title:
        # Title is relative to render_top
        tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(render_top + 0.12), Inches(width - 0.4), Inches(0.28))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"📊 {visual['title']}"
        run.font.size = Pt(get_font_size(spec, 'table_header'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_surface')
        p.alignment = PP_ALIGN.CENTER
        apply_font_to_run(run, spec)

    # 2. Render content
    # Content area starts below title
    content_top = render_top + 0.12 + (0.35 if has_title else 0)
    content_h = max(0.2, card_h - (0.35 if has_title else 0) - 0.2)
    
    # Textbox for code
    tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(content_top), Inches(width - 0.4), Inches(content_h))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    
    run = p.add_run()
    run.font.size = Pt(14)
    run.font.color.rgb = get_color(spec, 'on_surface_variant')
    run.text = preview_text
    apply_font_to_run(run, spec)


def render_visual_placeholder(slide: Any, visual: Dict, spec: Dict, left: float, top: float, width: float, height: float) -> None:
    card_h = min(height, 2.5)
    card = create_card_shape(slide, spec, left, top, width, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = get_color(spec, 'surface_variant')
    card.line.color.rgb = get_color(spec, 'outline')
    card.line.width = Pt(1)
    label = visual.get('title', visual.get('type', 'Visual'))
    reqs = visual.get('content_requirements', [])
    text = f"[{label}]"
    if reqs:
        text += '\n' + '\n'.join(f"  • {r}" for r in reqs[:3])
    inner_h = min(card_h - 0.6, 2.0)
    inner_h = max(inner_h, 0.4)
    inner_top = top + max(0.12, (card_h - inner_h) / 2)
    tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(inner_top), Inches(width - 0.4), Inches(inner_h))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(get_font_size(spec, 'table_header'))
    run.font.color.rgb = get_color(spec, 'outline')
    apply_font_to_run(run, spec)

# Compatibility wrapper: provide the full-featured name expected by older renderers
def render_comparison_items(slide: Any, items: List[Dict], spec: Dict, grid: GridSystem, left: float, top: float, width: float, avail_h: float = None, slide_id: int = None) -> float:
    return render_comparison_items_simple(slide, items, spec, grid, left, top, width, avail_h, slide_id=slide_id)



# ════════════════════════════════════════════════════════════════
# §7  Per-Slide-Type Renderers (restored)
# ════════════════════════════════════════════════════════════════

def render_slide_title(slide: Any, sd: Dict, spec: Dict, grid: GridSystem, **ctx) -> None:
    """Title slide — full color bg, centered title, KPIs, author/date.

    Implements improved dynamic title height calculation for CJK/long titles
    and applies configured fonts.
    """
    title_text = sd.get('_deck_title', '') or sd.get('title', '')
    # Content layout anchors
    content_top = grid.slide_h * 0.30
    margin = grid.margin_h + 1.0
    cw = grid.slide_w - 2 * margin

    def is_cjk(ch: str) -> bool:
        return any([
            '\u4e00' <= ch <= '\u9fff',  # CJK Unified Ideographs
            '\u3400' <= ch <= '\u4dbf',  # CJK Unified Ideographs Extension A
            '\u3000' <= ch <= '\u303f'   # punctuation
        ])

    def estimate_text_height(text: str, role: str, width_in: float):
        font_pt = get_font_size(spec, role)
        total_w = 0.0
        for ch in text:
            if is_cjk(ch):
                total_w += (font_pt / 72.0) * 0.95
            elif ch.isspace():
                total_w += (font_pt / 72.0) * 0.35
            else:
                total_w += (font_pt / 72.0) * 0.55
        chars_per_line = max(1, int(width_in / max(0.0001, (total_w / max(1, len(text))))))
        if len(text) <= chars_per_line:
            lines = 1
        else:
            lines = max(1, int((total_w / width_in) + 0.999))
        line_h = font_pt / 72.0 + 0.08
        return lines * line_h, font_pt

    title_h, title_pt = estimate_text_height(title_text, 'display_large', cw)
    subtitles = sd.get('content', [])
    subtitle_line_h = get_font_size(spec, 'slide_subtitle') / 72.0 + 0.06
    subtitle_h = len(subtitles) * subtitle_line_h

    comps = sd.get('components', {})
    kpi_h = 0
    if comps.get('kpis'):
        kpi_h = 0.85 + 0.15

    available = max(1.0, grid.slide_h * 0.55 - content_top - 0.3)

    adjusted_pt = title_pt
    adjusted_title_h = title_h
    while adjusted_title_h + subtitle_h + kpi_h > available and adjusted_pt > 16:
        adjusted_pt -= 2
        total_w = 0.0
        for ch in title_text:
            if is_cjk(ch):
                total_w += (adjusted_pt / 72.0) * 0.95
            elif ch.isspace():
                total_w += (adjusted_pt / 72.0) * 0.35
            else:
                total_w += (adjusted_pt / 72.0) * 0.55
        adjusted_lines = max(1, int((total_w / cw) + 0.999))
        adjusted_title_h = adjusted_lines * (adjusted_pt / 72.0 + 0.08)

    max_title_lines = 2
    def compute_lines_for_pt(pt: int) -> int:
        total_w = 0.0
        for ch in title_text:
            if is_cjk(ch):
                total_w += (pt / 72.0) * 0.95
            elif ch.isspace():
                total_w += (pt / 72.0) * 0.35
            else:
                total_w += (pt / 72.0) * 0.55
        return max(1, int((total_w / cw) + 0.999))

    lines_now = compute_lines_for_pt(adjusted_pt)
    while lines_now > max_title_lines and adjusted_pt > 14:
        adjusted_pt -= 1
        lines_now = compute_lines_for_pt(adjusted_pt)
    if lines_now > max_title_lines:
        mid = len(title_text) // 2
        left = title_text.rfind(' ', 0, mid)
        right = title_text.find(' ', mid)
        if left != -1:
            split_pos = left
        elif right != -1:
            split_pos = right
        else:
            split_pos = mid
        title_text = title_text[:split_pos].strip() + '\n' + title_text[split_pos:].strip()
        lines_now = compute_lines_for_pt(adjusted_pt)
        adjusted_title_h = lines_now * (adjusted_pt / 72.0 + 0.08)

    if any(is_cjk(ch) for ch in title_text):
        cjk_count = sum(1 for ch in title_text if is_cjk(ch))
        txt_len = len(title_text.replace(' ', ''))
        if txt_len > 24:
            cap = 24
        elif txt_len > 16:
            cap = 28
        else:
            cap = 32
        if adjusted_pt > cap:
            adjusted_pt = cap
            total_w = 0.0
            for ch in title_text:
                if is_cjk(ch):
                    total_w += (adjusted_pt / 72.0) * 0.95
                elif ch.isspace():
                    total_w += (adjusted_pt / 72.0) * 0.35
                else:
                    total_w += (adjusted_pt / 72.0) * 0.55
            chars_per_line = max(1, int(cw / max(0.0001, (total_w / max(1, len(title_text))))))
            if len(title_text) <= chars_per_line:
                lines_now = 1
            else:
                lines_now = max(1, int((total_w / cw) + 0.999))
            adjusted_title_h = lines_now * (adjusted_pt / 72.0 + 0.08)

    # === Actual rendering ===
    # Title textbox
    title_tb = slide.shapes.add_textbox(
        Inches(margin), Inches(content_top),
        Inches(cw), Inches(adjusted_title_h + 0.15)
    )
    title_tf = title_tb.text_frame
    title_tf.word_wrap = True
    p = title_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(adjusted_pt)
    run.font.bold = True
    run.font.color.rgb = get_color(spec, 'on_primary')
    apply_font_to_run(run, spec)
    cursor_y = content_top + adjusted_title_h + 0.20

    # Decorative line
    dec_line_w = 2.5
    dec_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches((grid.slide_w - dec_line_w) / 2), Inches(cursor_y),
        Inches(dec_line_w), Inches(0.04)
    )
    dec_line.fill.solid()
    dec_line.fill.fore_color.rgb = get_color(spec, 'primary_container')
    dec_line.line.fill.background()
    cursor_y += 0.25

    # Subtitles from content[]
    for sub in subtitles:
        sub_tb = slide.shapes.add_textbox(
            Inches(margin + 0.5), Inches(cursor_y),
            Inches(cw - 1.0), Inches(subtitle_line_h + 0.10)
        )
        sub_tf = sub_tb.text_frame
        sub_tf.word_wrap = True
        p2 = sub_tf.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = sub
        run2.font.size = Pt(get_font_size(spec, 'slide_subtitle'))
        run2.font.color.rgb = get_color(spec, 'on_primary')
        apply_font_to_run(run2, spec)
        cursor_y += subtitle_line_h + 0.12

    # KPIs (if present)
    if comps.get('kpis'):
        kpi_top = max(cursor_y + 0.30, grid.slide_h * 0.68)
        render_kpis(slide, comps['kpis'], spec, grid, margin, kpi_top, cw)


def render_callouts(slide: Any, callouts: List[Dict], spec: Dict, grid: GridSystem, left: float, top: float, width: float) -> float:
    """Render callout boxes."""
    if not callouts:
        return 0
    y = top
    for co in callouts:
        h = 0.55
        bg = create_card_shape(slide, spec, left, y, width, h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = get_color(spec, 'primary_container')
        bg.line.fill.background()
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(y), Inches(0.05), Inches(h)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = get_color(spec, 'primary')
        bar.line.fill.background()
        tb = slide.shapes.add_textbox(
            Inches(left + 0.15), Inches(y + 0.08),
            Inches(width - 0.25), Inches(h - 0.16)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        text = co if isinstance(co, str) else co.get('text', '') or co.get('description', '') or co.get('title', '')
        title = co.get('title', '') if isinstance(co, dict) else ''
        if title and text and text != title:
            run.text = f"{title}: {text}"
        else:
            run.text = text or title
        run.font.size = Pt(get_font_size(spec, 'callout_text'))
        run.font.color.rgb = get_color(spec, 'on_primary_container')
        apply_font_to_run(run, spec)
        y += h + 0.08
    return y - top


def render_slide_comparison(slide: Any, sd: Dict, spec: Dict, grid: GridSystem, **ctx) -> None:
    """Comparison — cards + optional visual hybrid layout.

    When both comparison_items AND a visual (chart_config/mermaid) are present,
    uses a split layout: left 7 cols for cards, right 5 cols for chart.
    When only comparison_items are present, uses full-width card layout.
    Cards expand to fill the available content zone height.
    """
    bar_h = ctx.get('bar_h', 0.55)
    slide_id = ctx.get('slide_id')
    top, avail_h = _content_zone(grid, bar_h)
    comps = sd.get('components', {})
    fill_mode = get_content_fill(spec, sd.get('slide_type', 'comparison'), slide_id=slide_id)
    has_visual = (sd.get('visual') and
                  sd['visual'].get('type') not in (None, 'none') and
                  sd['visual'].get('placeholder_data', {}))

    # Normalize: treat 'risks' as comparison_items if comparison_items is absent
    comp_items = comps.get('comparison_items') or comps.get('risks')
    # Also check visual.placeholder_data for misplaced comparison_items
    comp_items_from_visual = False
    if not comp_items:
        pd = sd.get('visual', {}).get('placeholder_data', {})
        if isinstance(pd, dict) and pd.get('comparison_items'):
            comp_items = pd['comparison_items']
            comp_items_from_visual = True  # Mark: data came from visual, skip re-rendering visual

    if comp_items:
        # Check for duplication: if chart_config labels overlap with card labels, skip visual
        skip_visual = False
        if comp_items_from_visual:
            skip_visual = True  # Visual data already consumed as cards
        elif has_visual:
            pd = sd.get('visual', {}).get('placeholder_data', {})
            chart_cfg = pd.get('chart_config', {}) if isinstance(pd, dict) else {}
            chart_labels = set(chart_cfg.get('labels', []))
            card_labels = {item.get('label', '') for item in comp_items}
            if chart_labels and card_labels:
                overlap_count = len(chart_labels & card_labels)
                if overlap_count >= len(card_labels) * 0.8:
                    skip_visual = True  # Too much label overlap — skip chart

        # Determine layout: hybrid (cards + chart) or full-width cards
        if has_visual and not skip_visual:
            # Hybrid: left 7 cols for cards, right 5 cols for visual
            cards_left, cards_w = grid.col_span(7, 0)
            vis_left, vis_w = grid.col_span(5, 7)
        else:
            # Full-width cards (or dedup/consumed visual)
            cards_left = grid.margin_h
            cards_w = grid.usable_w
            vis_left = vis_w = 0

        content = sd.get('content', [])
        n_bullets = min(len(content), 4)
        bullet_h = n_bullets * 0.40 if n_bullets else 0
        callout_h = 0.8 if comps.get('callouts') else 0
        # Estimate total content height for vertical centering
        skip_keys = {'comparison_items', 'risks', 'callouts'}
        if content:
            skip_keys.add('bullets')
        extra_comps = {k: v for k, v in comps.items()
                       if k not in skip_keys and v}
        est_cards_h = _estimate_components_height({'comparison_items': comp_items}, spec, cards_w)
        est_extra_h = _estimate_components_height(extra_comps, spec, cards_w) if extra_comps else 0
        est_total = est_cards_h + bullet_h + est_extra_h + callout_h
        # Vertical centering: only when NOT expanding (expand uses full space from top)
        if fill_mode == 'expand':
            v_offset = 0.0
        elif est_total < avail_h - 0.5:
            v_offset = max(0, (avail_h - est_total) / 2 - 0.15)
        else:
            v_offset = 0.0
        # Cards fill available space minus bullets and callouts
        if fill_mode == 'expand':
            cards_avail = avail_h - 0.15 - bullet_h - callout_h - v_offset
        else:
            cards_avail = None  # use default card height
        h = render_comparison_items(slide, comp_items, spec, grid,
                                    cards_left, top + 0.15 + v_offset, cards_w,
                                    avail_h=cards_avail, slide_id=slide_id)
        cursor = top + 0.15 + v_offset + h
        # Content bullets below cards (complementary context)
        if content:
            bullet_spacing = 0.40
            for i, bullet in enumerate(content[:4]):
                tb = slide.shapes.add_textbox(
                    Inches(cards_left + 0.1), Inches(cursor + i * bullet_spacing),
                    Inches(cards_w - 0.2), Inches(0.35)
                )
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = f"• {bullet}"
                run.font.size = Pt(get_font_size(spec, 'bullet_text'))
                run.font.color.rgb = get_color(spec, 'on_surface')
            cursor += n_bullets * bullet_spacing
        # Render kpis/other components
        if extra_comps:
            h2 = render_components(slide, extra_comps, spec, grid,
                                   cards_left, cursor, cards_w,
                                   avail_h - (cursor - top))
            cursor += h2
        # Callouts below
        if comps.get('callouts'):
            render_callouts(slide, comps['callouts'], spec, grid,
                            cards_left, cursor, cards_w)
        # Visual on right (hybrid layout, only if not skipped)
        if has_visual and vis_w > 0 and not skip_visual:
            render_visual(slide, sd['visual'], spec, grid,
                          vis_left, top + 0.15 + v_offset, vis_w, avail_h - 0.3 - v_offset)
    else:
        # No comparison_items — render components + visual
        non_vis_comps = {k: v for k, v in comps.items() if v}
        has_comps_data = bool(non_vis_comps)
        if has_visual and has_comps_data:
            # Stacked vertical layout: visual on top (full width), components below
            # This avoids the left-right split that creates visual clutter
            vis_est_h = 1.8  # estimate for chart table
            pd_check = sd.get('visual', {}).get('placeholder_data', {})
            chart_cfg = pd_check.get('chart_config', {}) if isinstance(pd_check, dict) else {}
            n_series = len(chart_cfg.get('series', []))
            if chart_cfg.get('labels'):
                vis_est_h = 0.50 + (n_series + 1) * 0.45  # header row + data rows
            comp_est_h = _estimate_components_height(non_vis_comps, spec, grid.usable_w)
            total_h = vis_est_h + 0.25 + comp_est_h
            # Center the combined block vertically
            v_off = max(0, (avail_h - total_h) / 2 - 0.10) if total_h < avail_h - 0.3 else 0.0
            cursor = top + 0.15 + v_off
            # Render visual (chart/table) at top, full width
            render_visual(slide, sd['visual'], spec, grid,
                          grid.margin_h, cursor, grid.usable_w, vis_est_h)
            cursor += vis_est_h + 0.25
            # Render components (kpis, bullets) below, full width
            render_components(slide, non_vis_comps, spec, grid,
                              grid.margin_h, cursor, grid.usable_w,
                              avail_h - (cursor - top))
        elif has_visual:
            render_visual(slide, sd['visual'], spec, grid,
                          grid.margin_h, top + 0.15, grid.usable_w, avail_h - 0.3)
        elif has_comps_data:
            # Components only — full width, vertically centered
            comp_est_h = _estimate_components_height(non_vis_comps, spec, grid.usable_w)
            v_off = max(0, (avail_h - comp_est_h) / 2 - 0.15) if comp_est_h < avail_h * 0.6 else 0.0
            cursor = top + 0.15 + v_off
            render_components(slide, non_vis_comps, spec, grid,
                              grid.margin_h, cursor, grid.usable_w, avail_h - 0.3)
        else:
            render_slide_two_column(slide, sd, spec, grid, **ctx)

# Dispatcher mapping will be defined after per-slide-type renderers are implemented.
# See bottom of file for RENDERERS mapping.


def render_components(slide: Any, comps: Dict, spec: Dict, grid: GridSystem, left: float, top: float, width: float, avail_h: float = None) -> float:
    """Compatibility wrapper for older API names."""
    return render_components_simple(slide, comps, spec, grid, left, top, width, avail_h or 0)


def render_decisions(slide: Any, decisions: List[Dict], spec: Dict, grid: GridSystem, left: float, top: float, width: float, avail_h: float = None) -> float:
    """Compatibility wrapper for older API names."""
    return render_decisions_simple(slide, decisions, spec, grid, left, top, width, avail_h=avail_h)


def render_slide_decision(slide: Any, sd: Dict, spec: Dict, grid: GridSystem, **ctx) -> None:
    """Decision slide — left decisions + right visual/chart.

    Decision cards and the right-side visual/table are expanded to fill
    the full content zone height, preventing large whitespace areas.
    """
    bar_h = ctx.get('bar_h', 0.55)
    slide_id = ctx.get('slide_id')
    top, avail_h = _content_zone(grid, bar_h)
    comps = sd.get('components', {})
    fill_mode = get_content_fill(spec, sd.get('slide_type', 'decision'), slide_id=slide_id)

    has_visual = (sd.get('visual') and sd['visual'].get('type') not in (None, 'none'))
    non_decision_comps = {k: v for k, v in comps.items() if k != 'decisions' and v}
    has_right_content = has_visual or bool(non_decision_comps)

    if has_right_content:
        l_left, l_w = grid.col_span(5, 0)
        r_left, r_w = grid.col_span(7, 5)
    else:
        # No visual or other components — full width for decisions
        l_left = grid.margin_h
        l_w = grid.usable_w
        r_left = r_w = 0

    # Decision cards
    cursor = top + 0.15
    if comps.get('decisions'):
        decision_labels = {(d.get('title') or d.get('label') or '').strip()
                           for d in comps['decisions']}
        extra_bullets = [b for b in sd.get('content', [])
                         if b.strip() not in decision_labels]
        bullet_reserve = len(extra_bullets[:3]) * 0.50
        if fill_mode == 'expand':
            decisions_avail = avail_h - 0.15 - bullet_reserve
        else:
            decisions_avail = None
        h = render_decisions(slide, comps['decisions'], spec, grid,
                             l_left, cursor, l_w, avail_h=decisions_avail)
        cursor += h
        # Extra bullets (not duplicating decision titles)
        for i, bullet in enumerate(extra_bullets[:3]):
            tb = slide.shapes.add_textbox(
                Inches(l_left + 0.1), Inches(cursor + i * 0.48), Inches(l_w - 0.2), Inches(0.45)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"• {bullet}"
            run.font.size = Pt(get_font_size(spec, 'bullet_text'))
            run.font.color.rgb = get_color(spec, 'on_surface')
        cursor += len(extra_bullets[:3]) * 0.48

    # Right column: visual or non-decision components (prevents duplication)
    if has_visual:
        render_visual(slide, sd['visual'], spec, grid,
                      r_left, top + 0.15, r_w, avail_h - 0.3)
    elif non_decision_comps and r_w > 0:
        render_components(slide, non_decision_comps, spec, grid,
                          r_left, top + 0.15, r_w, avail_h - 0.3)


def _resolve_bg_image(spec: Dict, slide_type: str, sd: Dict) -> str:
    img = sd.get('background_image', '')
    if img:
        return img
    stl = spec.get('slide_type_layouts', {})
    entry = stl.get(slide_type, stl.get('default', {}))
    return entry.get('background_image', '')


def _content_zone(grid: GridSystem, bar_h: float):
    """Calculate content zone below title bar."""
    lz_top = bar_h + 0.12
    lz_h = grid.slide_h - lz_top - 0.35  # leave room for bottom bar
    return lz_top, lz_h


def _apply_vertical_offset(content_h: float, avail_h: float, fill_mode: str) -> float:
    """Calculate vertical offset to apply based on content_fill mode.

    Returns additional top offset (inches) to shift content down.
    - 'center': center content vertically in available space
    - 'expand' / 'top-align' / 'top': no offset (content starts at top)
    """
    if fill_mode == 'center' and content_h < avail_h:
        return (avail_h - content_h) / 2
    return 0.0


def _effective_fill(fill_mode: str, has_visual: bool, is_split: bool = False, has_text: bool = False) -> str:
    """Context-aware fill mode resolver.

    Behavior when `fill_mode=='center'` and a visual is present:
    - Split layout (side-by-side): return 'split-center' → center each column independently
    - Stacked layout with text+visual: return 'combined-center' → center the combined block (text+visual)
    - Visual-only (no text): keep 'center' so the visual itself can center
    - Otherwise: return the original fill_mode
    """
    if not has_visual or fill_mode != 'center':
        return fill_mode
    if is_split:
        return 'split-center'
    if not has_text:
        return 'center'
    return 'combined-center'


def _estimate_components_height(comps: Dict, spec: Dict, width: float) -> float:
    """Estimate total height of components block (KPIs, bullets, content, etc.) for centering."""
    h = 0.0
    if comps.get('kpis'):
        n_kpis = len(comps['kpis'])
        per_kpi = 0.95 + 0.15  # card_h + gap (using new compact sizing)
        h += per_kpi
    if comps.get('comparison_items'):
        h += 1.8 + 0.15  # min_card_h + gap
    if comps.get('decisions'):
        h += len(comps['decisions']) * (0.65 + 0.12)
    if comps.get('bullets'):
        h += min(len(comps['bullets']), 8) * 0.48
    if comps.get('callouts'):
        h += len(comps['callouts']) * 0.63  # 0.55 height + 0.08 gap
    if comps.get('table_data'):
        h += 2.0  # rough table estimate
    return h


def _estimate_visual_height(visual: Dict, spec: Dict, avail_h: float, reserved: float = 0.5) -> float:
    """Conservative estimate of visual height (in inches).

    - Use a type-driven default heuristic, but cap to available space minus reserved.
    - Ensures combined-centering uses a realistic visual height and avoids overflow.
    """
    if not visual or visual.get('type') in (None, 'none'):
        return 0.0
    vtype = (visual.get('type') or '').lower()
    # defaults by visual type (inches)
    defaults = {
        'gantt': 2.8,
        'mermaid': 2.8,
        'flowchart': 2.5,
        'sequence': 2.2,
        'chart': 2.2,
        'table': 2.0,
    }
    default_h = defaults.get(vtype, 2.5)
    max_allowed = max(0.5, avail_h - reserved)
    return min(default_h, max_allowed)


def _fit_visual_height_for_combined(visual: Dict, spec: Dict, avail_h: float, text_h: float, reserved: float = 0.2, min_h: float = 1.0) -> (float, bool):
    """Try to fit the visual into the remaining space when centering combined blocks.

    Returns (chosen_height_in_inches, fits_bool).
    If the visual can be reduced to fit the available area (>= min_h), returns chosen height and True.
    Otherwise returns a conservative min height and False.
    """
    if not visual or visual.get('type') in (None, 'none'):
        return 0.0, True
    avail_for_vis = max(0.0, avail_h - text_h - reserved)
    # get default estimate
    vtype = (visual.get('type') or '').lower()
    defaults = {
        'gantt': 2.8,
        'mermaid': 2.8,
        'flowchart': 2.5,
        'sequence': 2.2,
        'chart': 2.2,
        'table': 2.0,
    }
    default_h = defaults.get(vtype, 2.5)
    if avail_for_vis >= min_h:
        chosen = min(default_h, avail_for_vis)
        return chosen, True
    # cannot fit to min_h
    chosen = max(min_h, min(default_h, avail_for_vis if avail_for_vis > 0 else min_h))
    return chosen, False


def render_slide_data_heavy(slide: Any, sd: Dict, spec: Dict, grid: GridSystem, **ctx) -> None:
    """Data-heavy: components first (KPIs, decisions, etc), then visual."""
    bar_h = ctx.get('bar_h', 0.55)
    slide_id = ctx.get('slide_id')
    top, avail_h = _content_zone(grid, bar_h)
    comps = sd.get('components', {})
    fill_mode = get_content_fill(spec, sd.get('slide_type', 'data-heavy'), slide_id=slide_id)
    has_visual = (sd.get('visual') and
                  sd['visual'].get('type') not in (None, 'none'))

    # If both components AND visual exist, use STACKED layout (KPIs above, visual below)
    # to prevent overlap. Only use split for heavy text content.
    has_comps_check = any(comps.get(k) for k in comps)
    comp_keys = [k for k in comps if comps.get(k)]
    is_kpi_only = comp_keys == ['kpis'] or set(comp_keys) <= {'kpis', 'bullets'}
    if has_visual and has_comps_check and not is_kpi_only:
        comp_left, comp_w = grid.col_span(6, 0)
        vis_left, vis_w = grid.col_span(6, 6)
    elif has_visual and has_comps_check and is_kpi_only:
        # Stacked: KPIs on top full-width, visual below full-width
        comp_left = grid.margin_h
        comp_w = grid.usable_w
        vis_left = grid.margin_h
        vis_w = grid.usable_w
    elif has_visual:
        comp_left = comp_w = 0
        vis_left = grid.margin_h
        vis_w = grid.usable_w
    else:
        comp_left = grid.margin_h
        comp_w = grid.usable_w
        vis_left = vis_w = 0

    # Estimate content height for vertical centering
    has_comps = any(comps.get(k) for k in comps)
    content_h = _estimate_components_height(comps, spec, comp_w) if has_comps else 0
    # In split layout, always center both columns independently for visual balance
    is_split = has_visual and has_comps and not is_kpi_only
    # For stacked layout, include visual height in total content estimate for centering
    is_stacked = has_visual and has_comps and is_kpi_only
    if is_stacked:
        vis_est_for_center = _estimate_visual_height(sd.get('visual'), spec, avail_h, reserved=0.2)
        content_h += vis_est_for_center + 0.20  # +gap between components and visual
    fill_mode = _effective_fill(fill_mode, has_visual, is_split=is_split, has_text=has_comps)
    # Force split-center for split layouts to avoid asymmetric alignment
    if is_split and fill_mode not in ('split-center',):
        fill_mode = 'split-center'
    # For 'split-center', center components within the left column and center visual independently in right column
    if is_split and fill_mode == 'split-center' and has_comps:
        comp_v_offset = _apply_vertical_offset(content_h, avail_h, 'center')
    else:
        comp_v_offset = _apply_vertical_offset(content_h, avail_h, fill_mode) if has_comps and fill_mode == 'center' else 0.0

    # For split layout, compute visual offset when centering its column
    if is_split and fill_mode == 'split-center':
        vis_est_h = _estimate_visual_height(sd.get('visual'), spec, avail_h, reserved=0.2)
        vis_v_offset = _apply_vertical_offset(vis_est_h, avail_h, 'center')
    else:
        vis_v_offset = 0.0

    cursor = top + 0.10 + comp_v_offset

    # Render content bullets (sd.content) if different from components.bullets
    content_bullets = sd.get('content', [])
    comp_bullet_texts = set()
    for b in comps.get('bullets', []):
        comp_bullet_texts.add(b.get('text', '') if isinstance(b, dict) else str(b))
    unique_content = [b for b in content_bullets if b not in comp_bullet_texts]
    if unique_content and comp_w > 0:
        n_items = min(len(unique_content), 4)
        is_content_sparse = n_items <= 4 and not has_visual
        if is_content_sparse and n_items >= 2:
            # Numbered takeaway cards — horizontal layout for sparse key points
            card_gap = 0.12
            card_w = (comp_w - card_gap * (n_items - 1)) / n_items
            card_h = 0.70
            for i, bullet in enumerate(unique_content[:n_items]):
                cx = comp_left + i * (card_w + card_gap)
                # Card background
                card_bg = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(cx), Inches(cursor),
                    Inches(card_w), Inches(card_h)
                )
                card_bg.fill.solid()
                card_bg.fill.fore_color.rgb = get_color(spec, 'primary_container')
                card_bg.line.fill.background()
                # Accent left bar
                accent_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(cx), Inches(cursor),
                    Inches(0.04), Inches(card_h)
                )
                accent_bar.fill.solid()
                accent_bar.fill.fore_color.rgb = get_color(spec, 'primary')
                accent_bar.line.fill.background()
                # Number badge
                num_tb = slide.shapes.add_textbox(
                    Inches(cx + 0.10), Inches(cursor + 0.06),
                    Inches(0.25), Inches(0.25)
                )
                num_tf = num_tb.text_frame
                num_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                num_p = num_tf.paragraphs[0]
                num_run = num_p.add_run()
                num_run.text = str(i + 1)
                num_run.font.size = Pt(16)
                num_run.font.bold = True
                num_run.font.color.rgb = get_color(spec, 'primary')
                # Text content
                txt_tb = slide.shapes.add_textbox(
                    Inches(cx + 0.10), Inches(cursor + 0.28),
                    Inches(card_w - 0.20), Inches(card_h - 0.32)
                )
                txt_tf = txt_tb.text_frame
                txt_tf.word_wrap = True
                txt_tf.vertical_anchor = MSO_ANCHOR.TOP
                txt_p = txt_tf.paragraphs[0]
                txt_run = txt_p.add_run()
                txt_run.text = bullet
                txt_run.font.size = Pt(get_font_size(spec, 'body'))
                txt_run.font.color.rgb = get_color(spec, 'on_primary_container')
                txt_p.line_spacing = 1.3
            cursor += card_h + 0.15
        else:
            # Standard bullet list for many items or non-sparse content
            c_font = get_font_size(spec, 'bullet_text')
            c_spacing = 0.48
            for i, bullet in enumerate(unique_content[:4]):
                tb = slide.shapes.add_textbox(
                    Inches(comp_left + 0.1), Inches(cursor + i * c_spacing),
                    Inches(comp_w - 0.2), Inches(0.38)
                )
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = f"• {bullet}"
                run.font.size = Pt(c_font)
                run.font.color.rgb = get_color(spec, 'on_surface')
                p.line_spacing = 1.5
            cursor += min(len(unique_content), 4) * c_spacing + 0.15

    # Render all components
    if has_comps and comp_w > 0:
        h = render_components(slide, comps, spec, grid,
                              comp_left, cursor, comp_w,
                              avail_h - 0.10)
        cursor += h

    # Visual: layout depends on split vs stacked vs visual-only
    if has_visual:
        if is_split and not is_kpi_only:
            # Split layout: visual fills right column independently
            vis_render_top = top + 0.10 + vis_v_offset
            render_visual(slide, sd['visual'], spec, grid,
                          vis_left, vis_render_top, vis_w, avail_h - 0.10 - vis_v_offset)
        elif is_kpi_only and has_comps:
            # Stacked layout: visual goes BELOW KPIs with gap
            vis_gap = 0.20
            vis_render_top = cursor + vis_gap
            vis_available = max(1.5, avail_h - (vis_render_top - top) - 0.10)
            render_visual(slide, sd['visual'], spec, grid,
                          vis_left, vis_render_top, vis_w, vis_available)
        else:
            # Visual-only: center if requested
            vis_v_offset = _apply_vertical_offset(avail_h * 0.5, avail_h, fill_mode)
            render_visual(slide, sd['visual'], spec, grid,
                          vis_left, top + 0.10 + vis_v_offset, vis_w, avail_h - 0.10 - vis_v_offset)


def render_slide_matrix(slide: Any, sd: Dict, spec: Dict, grid: GridSystem, **ctx) -> None:
    """Matrix layout — for risk/evaluation grids."""
    bar_h = ctx.get('bar_h', 0.55)
    slide_id = ctx.get('slide_id')
    top, avail_h = _content_zone(grid, bar_h)
    comps = sd.get('components', {})
    fill_mode = get_content_fill(spec, sd.get('slide_type', 'matrix'), slide_id=slide_id)
    has_visual = sd.get('visual') and sd['visual'].get('type') not in (None, 'none')
    fill_mode = _effective_fill(fill_mode, has_visual)

    # Estimate content height
    bullets = sd.get('content', [])[:4]
    content_h = len(bullets) * 0.42 + 0.15 + _estimate_components_height(comps, spec, grid.usable_w)
    v_offset = _apply_vertical_offset(content_h, avail_h, fill_mode)

    cursor = top + 0.10 + v_offset

    # Bullets
    for i, bullet in enumerate(bullets):
        tb = slide.shapes.add_textbox(
            Inches(grid.margin_h + 0.1), Inches(cursor + i * 0.42),
            Inches(grid.usable_w - 0.2), Inches(0.38)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.size = Pt(get_font_size(spec, 'bullet_text'))
        run.font.color.rgb = get_color(spec, 'on_surface')
    cursor += len(bullets) * 0.42 + 0.15

    # Components (risks, table_data, etc)
    if any(comps.get(k) for k in comps):
        render_components(slide, comps, spec, grid,
                          grid.margin_h, cursor, grid.usable_w,
                          avail_h - (cursor - top))
    elif sd.get('visual'):
        render_visual(slide, sd['visual'], spec, grid,
                      grid.margin_h, cursor, grid.usable_w,
                      avail_h - (cursor - top))


def render_timeline_items(slide: Any, items: List[Dict], spec: Dict, grid: GridSystem, left: float, top: float, width: float) -> float:
    """Render timeline as horizontal milestones."""
    if not items:
        return 0
    n = len(items)
    seg_w = width / max(n, 1)
    dot_r = 0.10
    line_y = top + dot_r
    # Connector line
    if n > 1:
        connector = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left + seg_w / 2), Inches(line_y - 0.015),
            Inches(width - seg_w), Inches(0.03)
        )
        connector.fill.solid()
        connector.fill.fore_color.rgb = get_color(spec, 'outline')
        connector.line.fill.background()
    status_colors = {
        'completed': 'secondary', 'done': 'secondary',
        'in_progress': 'primary', 'active': 'primary',
        'planned': 'outline', 'pending': 'outline',
    }
    for i, item in enumerate(items):
        cx = left + i * seg_w + seg_w / 2 - dot_r
        status = item.get('status', 'planned')
        color_token = status_colors.get(status, 'outline')
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(cx), Inches(top), Inches(dot_r * 2), Inches(dot_r * 2)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = get_color(spec, color_token)
        dot.line.fill.background()
        # Milestone text
        tb = slide.shapes.add_textbox(
            Inches(left + i * seg_w), Inches(top + dot_r * 2 + 0.08),
            Inches(seg_w), Inches(0.30)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = item.get('label', item.get('milestone', item.get('phase', item.get('title', ''))))
        run.font.size = Pt(get_font_size(spec, 'body_text'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_surface')
        # Date
        date_val = item.get('date', item.get('period', item.get('time', item.get('description', ''))))
        if date_val:
            tb2 = slide.shapes.add_textbox(
                Inches(left + i * seg_w), Inches(top + dot_r * 2 + 0.38),
                Inches(seg_w), Inches(0.22)
            )
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = date_val
            run2.font.size = Pt(get_font_size(spec, 'caption'))
            run2.font.color.rgb = get_color(spec, 'on_surface')
    return max(0.32 + 0.30, dot_r * 2 + 0.8)


def render_slide_timeline(slide: Any, sd: Dict, spec: Dict, grid: GridSystem, **ctx) -> None:
    """Timeline/gantt layout."""
    bar_h = ctx.get('bar_h', 0.55)
    slide_id = ctx.get('slide_id')
    top, avail_h = _content_zone(grid, bar_h)
    comps = sd.get('components', {})
    fill_mode = get_content_fill(spec, sd.get('slide_type', 'timeline'), slide_id=slide_id)

    # Estimate content height for centering (timeline + visual below)
    has_visual = sd.get('visual') and sd['visual'].get('type') not in (None, 'none')
    # For timeline, when both timeline and visual exist treat as combined block
    fill_mode = _effective_fill(fill_mode, has_visual, is_split=False, has_text=bool(comps.get('timeline_items') or sd.get('content')))
    timeline_h = 1.2 if comps.get('timeline_items') else len(sd.get('content', [])) * 0.45
    vis_est_h = _estimate_visual_height(sd.get('visual'), spec, avail_h, reserved=0.2) if has_visual else 0
    content_h = timeline_h + vis_est_h + 0.10
    # Gantt charts should use full available height, skip centering
    visual_type = (sd.get('visual', {}).get('type', '') or '').lower()
    if visual_type == 'gantt':
        v_offset = 0.0  # Full height for gantt charts
    elif fill_mode in ('combined-center', 'center') and content_h < avail_h - 0.05:
        v_offset = _apply_vertical_offset(content_h, avail_h, 'center')
    else:
        v_offset = 0.0

    cursor = top + 0.15 + v_offset

    if comps.get('timeline_items'):
        h = render_timeline_items(slide, comps['timeline_items'], spec, grid,
                                  grid.margin_h, cursor, grid.usable_w)
        cursor += h
    else:
        # Render bullets (content[] already populated by normalize_slide_data)
        for i, bullet in enumerate(sd.get('content', [])):
            tb = slide.shapes.add_textbox(
                Inches(grid.margin_h + 0.1), Inches(cursor + i * 0.45),
                Inches(grid.usable_w - 0.2), Inches(0.40)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"• {bullet}"
            run.font.size = Pt(get_font_size(spec, 'bullet_text'))
            run.font.color.rgb = get_color(spec, 'on_surface')
            cursor += 0.45

    # Remaining components
    # Fix: Exclude 'bullets' because they are already handled in the 'else' block via sd['content']
    # This prevents double rendering of bullet points
    remaining_keys = [k for k in comps if k not in ('timeline_items', 'bullets') and comps[k]]
    remaining = {k: comps[k] for k in remaining_keys}
    if remaining:
        render_components(slide, remaining, spec, grid,
                          grid.margin_h, cursor, grid.usable_w,
                          avail_h - (cursor - top))

    # Visual
    if sd.get('visual') and sd['visual'].get('type') not in (None, 'none'):
        vis_h_to_use = chosen_vis_h if 'chosen_vis_h' in locals() and chosen_vis_h else (avail_h - (cursor - top) - 0.2)
        render_visual(slide, sd['visual'], spec, grid,
                      grid.margin_h, cursor + 0.10, grid.usable_w,
                      vis_h_to_use)


def render_slide_flowchart(slide: Any, sd: Dict, spec: Dict, grid: GridSystem, **ctx) -> None:
    """Flowchart layout — bullets + full-width visual."""
    bar_h = ctx.get('bar_h', 0.55)
    slide_id = ctx.get('slide_id')
    top, avail_h = _content_zone(grid, bar_h)
    comps = sd.get('components', {})
    fill_mode = get_content_fill(spec, sd.get('slide_type', 'flowchart'), slide_id=slide_id)
    has_visual = sd.get('visual') and sd['visual'].get('type') not in (None, 'none')
    content_bullets = sd.get('content', [])
    # use combined-center when text+visual exist to center overall block
    fill_mode = _effective_fill(fill_mode, has_visual, is_split=False, has_text=bool(content_bullets or comps))

    # Estimate content height for centering
    if not content_bullets and comps.get('bullets'):
        content_bullets = [(b.get('text', '') if isinstance(b, dict) else str(b)) for b in comps['bullets']]
    bullet_h = min(len(content_bullets), 3) * 0.40 + 0.15
    comp_h = _estimate_components_height({k: v for k, v in comps.items() if k != 'bullets' and v}, spec, grid.usable_w)
    vis_h = 2.0 if (sd.get('visual') and sd['visual'].get('type') not in (None, 'none')) else 0
    total_content_h = bullet_h + comp_h + vis_h
    # Attempt fit-and-center for combined-center: shrink visual if necessary
    if fill_mode in ('combined-center', 'center') and has_visual:
        vis_est_h, fits = _fit_visual_height_for_combined(sd.get('visual'), spec, avail_h, total_content_h - vis_h, reserved=0.2)
        if fits:
            combined_h = (total_content_h - vis_h) + vis_est_h + 0.10
            v_offset = _apply_vertical_offset(combined_h, avail_h, 'center')
            chosen_vis_h = vis_est_h
        else:
            v_offset = 0.0
            chosen_vis_h = None
    else:
        v_offset = 0.0
        chosen_vis_h = None

    cursor = top + 0.15 + v_offset

    # Resolve bullet text
    content_bullets = sd.get('content', [])
    if not content_bullets and comps.get('bullets'):
        content_bullets = [
            (b.get('text', '') if isinstance(b, dict) else str(b))
            for b in comps['bullets']
        ]

    # Bullets
    for i, bullet in enumerate(content_bullets[:3]):
        tb = slide.shapes.add_textbox(
            Inches(grid.margin_h + 0.1), Inches(cursor + i * 0.40),
            Inches(grid.usable_w - 0.2), Inches(0.38)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.size = Pt(get_font_size(spec, 'bullet_text'))
        run.font.color.rgb = get_color(spec, 'on_surface')
    cursor += min(len(content_bullets), 3) * 0.40 + 0.15

    # Components (skip bullets since already rendered)
    remaining_comps = {k: v for k, v in comps.items() if k != 'bullets' and v}
    if remaining_comps:
        h = render_components(slide, remaining_comps, spec, grid,
                              grid.margin_h, cursor, grid.usable_w,
                              avail_h - (cursor - top))
        cursor += h

    # Visual
    if sd.get('visual'):
        vis_h_available = avail_h - (cursor - top)
        
        # Smart height adjustment for horizontal flowcharts
        # Horizontal layouts (LR) with few nodes should use less height to avoid stretching
        visual_data = sd.get('visual', {})
        mermaid_code = visual_data.get('placeholder_data', {}).get('mermaid_code', '')
        is_horizontal = 'LR' in mermaid_code or 'RL' in mermaid_code
        node_count = mermaid_code.count('[') if mermaid_code else 0
        
        if is_horizontal and node_count <= 5:
            # Simple horizontal flowchart: limit height to avoid stretching
            vis_h_to_use = min(3.2, vis_h_available)
        else:
            # Use available height for complex diagrams
            vis_h_to_use = vis_h_available
        
        render_visual(slide, sd['visual'], spec, grid,
                      grid.margin_h, cursor, grid.usable_w, vis_h_to_use)


def find_section_for_slide(slide_id: int, sections: List[Dict]) -> Dict:
    result = {}
    for sec in sorted(sections, key=lambda s: s.get('start_slide', 0)):
        if slide_id >= sec.get('start_slide', 999):
            result = sec
    return result


def _render_background(slide: Any, sd: Dict, spec: Dict, grid: GridSystem, stype: str) -> None:
    """Render the slide background image or solid fill."""
    bg_token = spec.get('slide_type_layouts', {}).get(stype, spec.get('slide_type_layouts', {}).get('default', {})).get('background', 'surface')
    bg_image_path = _resolve_bg_image(spec, stype, sd)
    if bg_image_path and os.path.isfile(bg_image_path):
        pic = slide.shapes.add_picture(bg_image_path, Emu(0), Emu(0), Inches(grid.slide_w), Inches(grid.slide_h))
        sp_tree = slide.shapes._spTree
        sp_tree.remove(pic._element)
        sp_tree.insert(2, pic._element)
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Inches(grid.slide_w), Inches(grid.slide_h))
        fill = overlay.fill
        fill.solid()
        fill.fore_color.rgb = get_color(spec, bg_token)
        try:
            from pptx.oxml.ns import qn as _qn
            # Defensive: not all fill implementations expose an lxml-like API
            solid_fill_elem = getattr(fill, '_fill', None)
            if solid_fill_elem is None or not hasattr(solid_fill_elem, 'find'):
                raise AttributeError('fill._fill is missing or has no find()')
            solid = solid_fill_elem.find(_qn('a:solidFill'))
            if solid is None:
                raise AttributeError('no solidFill element')
            srgb = solid.find(_qn('a:srgbClr'))
            if srgb is None:
                srgb = solid[0] if len(solid) else None
            if srgb is None:
                raise AttributeError('no srgbClr element')
            try:
                from lxml import etree
                alpha = etree.SubElement(srgb, _qn('a:alpha'))
                alpha.set('val', '40000')
            except Exception:
                logging.getLogger(__name__).warning(
                    'lxml not available or failed to set overlay alpha; skipping alpha modification'
                )
        except Exception:
            logging.getLogger(__name__).warning(
                'Could not set overlay alpha (XML API missing); skipping alpha modification'
            )
        overlay.line.fill.background()
        sp_tree = slide.shapes._spTree
        sp_tree.remove(overlay._element)
        sp_tree.insert(3, overlay._element)
    else:
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = get_color(spec, bg_token)


def _render_bullets_fallback(slide: Any, sd: Dict, spec: Dict, grid: GridSystem, bar_h: float) -> None:
    """Render simple bullet-list fallback content (from sd.content or components.bullets)."""
    top = bar_h + 0.12
    content_bullets = sd.get('content', [])
    # components.bullets format: list of {"text": "..."}
    comps = sd.get('components', {})
    if not content_bullets and comps.get('bullets'):
        content_bullets = [b.get('text', '') if isinstance(b, dict) else str(b) for b in comps['bullets']]
    bullet_y = top + 0.15
    for i, bullet in enumerate(content_bullets[:8]):
        tb = slide.shapes.add_textbox(Inches(grid.margin_h + 0.1), Inches(bullet_y + i * 0.48), Inches(grid.usable_w - 0.2), Inches(0.45))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.size = Pt(get_font_size(spec, 'bullet_text'))
        run.font.color.rgb = get_color(spec, 'on_surface')
        p.line_spacing = 1.5


def render_section_divider(slide: Any, sd: Dict, spec: Dict, grid: GridSystem, **ctx) -> None:
    """Section divider — full-color bg, centered title, progress."""
    meta = sd.get('metadata', {})
    sec_idx = meta.get('section_index', ctx.get('section_index', 0))
    total_sec = meta.get('total_sections', ctx.get('total_sections', 6))
    tb = slide.shapes.add_textbox(
        Inches(grid.margin_h + 1), Inches(grid.slide_h * 0.32),
        Inches(grid.usable_w - 2), Inches(1.0)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = sd.get('title', '')
    run.font.size = Pt(get_font_size(spec, 'headline_large'))
    run.font.bold = True
    run.font.color.rgb = get_color(spec, 'on_primary')
    apply_font_to_run(run, spec)
    content_list = sd.get('content', [])
    for b_idx, bullet in enumerate(content_list):
        tb2 = slide.shapes.add_textbox(
            Inches(grid.margin_h + 1.5), Inches(grid.slide_h * 0.48 + b_idx * 0.35),
            Inches(grid.usable_w - 3), Inches(0.32)
        )
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = bullet
        run2.font.size = Pt(get_font_size(spec, 'slide_subtitle'))
        run2.font.color.rgb = get_color(spec, 'on_primary')
    prog_w = 3.0
    prog_left = (grid.slide_w - prog_w) / 2
    seg_w = prog_w / max(total_sec, 1)
    prog_top = grid.slide_h * 0.72
    for i in range(total_sec):
        seg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(prog_left + i * seg_w + 0.03), Inches(prog_top),
            Inches(seg_w - 0.06), Inches(0.08)
        )
        seg.fill.solid()
        if i <= sec_idx:
            seg.fill.fore_color.rgb = get_color(spec, 'on_primary')
        else:
            seg.fill.fore_color.rgb = get_color(spec, 'surface_variant')
        seg.line.fill.background()


def render_slide_bullets(slide: Any, sd: Dict, spec: Dict, grid: GridSystem, **ctx) -> None:
    """Bullet-list slide with optional visual."""
    bar_h = ctx.get('bar_h', 0.55)
    slide_id = ctx.get('slide_id')
    top, avail_h = _content_zone(grid, bar_h)
    fill_mode = get_content_fill(spec, sd.get('slide_type', 'bullet-list'), slide_id=slide_id)
    has_visual = (sd.get('visual') and sd['visual'].get('type') not in (None, 'none'))
    comps = sd.get('components', {})
    has_comps = any(comps.get(k) for k in comps)

    # Detect complex visual types that need vertical layout (full-width)
    visual_type = sd.get('visual', {}).get('type', '') if has_visual else ''
    use_vertical_layout = visual_type in ('flowchart', 'flow', 'sequence', 'process')

    # layout columns
    if has_visual and not use_vertical_layout:
        # Horizontal layout: text left, visual right
        text_left, text_w = grid.col_span(7, 0)
        vis_left, vis_w = grid.col_span(5, 7)
    else:
        # Vertical layout or no visual: text full-width
        text_left, text_w = grid.col_span(12, 0)
        vis_left, vis_w = grid.col_span(12, 0) if use_vertical_layout else (0, 0)

    content_bullets = sd.get('content', [])
    if not content_bullets and comps.get('bullets'):
        content_bullets = [
            (b.get('text', '') if isinstance(b, dict) else str(b))
            for b in comps['bullets']
        ]

    # Estimate content height for vertical centering
    n_bullets = min(len(content_bullets), 8)
    remaining_comps = {k: v for k, v in comps.items() if k != 'bullets' and v}
    bullet_block_h = n_bullets * 0.48
    comp_block_h = _estimate_components_height(remaining_comps, spec, text_w - 0.2) if remaining_comps else 0
    total_content_h = bullet_block_h + (0.15 + comp_block_h if remaining_comps else 0)

    # Resolve effective fill mode with split/stack context
    is_split = bool(has_visual and text_w > 0)
    fill_mode = _effective_fill(get_content_fill(spec, sd.get('slide_type', 'bullet-list'), slide_id=slide_id), has_visual, is_split=is_split, has_text=(total_content_h > 0))

    # Compute positions based on resolved mode
    # For vertical layout, force top alignment to avoid complex calculations
    if use_vertical_layout:
        bullet_y = top + 0.15
    elif fill_mode == 'split-center' and is_split:
        # center text block and visual independently (use conservative visual estimate)
        vis_est_h = _estimate_visual_height(sd.get('visual'), spec, avail_h, reserved=0.2)
        bullet_v_offset = _apply_vertical_offset(total_content_h, avail_h, 'center')
        vis_v_offset = _apply_vertical_offset(vis_est_h, avail_h, 'center')
        bullet_y = top + 0.15 + bullet_v_offset
        vis_top = top + 0.15 + vis_v_offset
    elif fill_mode == 'combined-center' and has_visual:
        # Treat text + visual as a single block and attempt to fit the visual into remaining space
        vis_est_h, fits = _fit_visual_height_for_combined(sd.get('visual'), spec, avail_h, total_content_h, reserved=0.2)
        combined_h = total_content_h + vis_est_h + 0.10
        if fits:
            combined_v_offset = _apply_vertical_offset(combined_h, avail_h, 'center')
            bullet_y = top + 0.10 + combined_v_offset
            vis_top = bullet_y + total_content_h + 0.10
            # store chosen visual height in ctx so we can pass it later
            chosen_vis_h = vis_est_h
        else:
            # cannot fit: fall back to top but reduce top offset to avoid excess whitespace
            bullet_y = top + 0.10
            vis_top = bullet_y + total_content_h + 0.08
            chosen_vis_h = None
    else:
        # top / expand / center-without-visual
        v_offset = _apply_vertical_offset(total_content_h, avail_h, fill_mode)
        bullet_y = top + 0.15 + v_offset
        vis_top = top + 0.15 + v_offset

    # Adjust bullet spacing and font based on layout type and density
    is_sparse = n_bullets <= 3 and not has_visual and not use_vertical_layout
    # Alternate between two sparse styles for consecutive same-type slides
    sparse_style = 'B' if (is_sparse and slide_id and slide_id % 2 == 1) else 'A'
    if use_vertical_layout:
        bullet_spacing = 0.45
        bullet_limit = min(len(content_bullets), 5)
        bullet_font_size = get_font_size(spec, 'bullet_text')
        bullet_text_h = 0.38
    elif is_sparse and sparse_style == 'B':
        # Style B: centered blocks with top accent line, no card background
        bullet_spacing = 1.5
        bullet_limit = n_bullets
        bullet_font_size = 20
        bullet_text_h = 0.90
        sparse_h = n_bullets * bullet_spacing + (0.15 + comp_block_h if remaining_comps else 0)
        bullet_y = top + max(0.15, (avail_h - sparse_h) / 2)
    elif is_sparse:
        # Style A: decorative cards with numbered circles + accent bars
        bullet_spacing = 1.2  # Much wider spacing with decorative cards
        bullet_limit = n_bullets
        bullet_font_size = 20  # Larger font for sparse slides
        bullet_text_h = 0.80
        # Re-center vertically for sparse content
        sparse_h = n_bullets * bullet_spacing + (0.15 + comp_block_h if remaining_comps else 0)
        bullet_y = top + max(0.15, (avail_h - sparse_h) / 2)
    else:
        bullet_spacing = 0.48
        bullet_limit = min(len(content_bullets), 8)
        bullet_font_size = get_font_size(spec, 'bullet_text')
        bullet_text_h = 0.45
    
    for i, bullet in enumerate(content_bullets[:bullet_limit]):
        cur_y = bullet_y + i * bullet_spacing
        if is_sparse and sparse_style == 'B':
            # Style B: centered text blocks with top accent line — visually distinct from Style A
            center_x = text_left + text_w / 2
            accent_w = 0.80
            # Top accent line
            accent_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(center_x - accent_w / 2), Inches(cur_y),
                Inches(accent_w), Inches(0.04)
            )
            accent_line.fill.solid()
            accent_line.fill.fore_color.rgb = get_color(spec, 'primary_container')
            accent_line.line.fill.background()
            # Step label with contextual icon (e.g., "⚡", "◆")
            step_icon = get_icon_for_text(bullet, i)
            step_tb = slide.shapes.add_textbox(
                Inches(text_left + 0.5), Inches(cur_y + 0.10),
                Inches(text_w - 1.0), Inches(0.30)
            )
            step_tf = step_tb.text_frame
            step_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            step_p = step_tf.paragraphs[0]
            step_p.alignment = PP_ALIGN.CENTER
            step_run = step_p.add_run()
            step_run.text = step_icon
            step_run.font.size = Pt(18)
            step_run.font.bold = True
            step_run.font.color.rgb = get_color(spec, 'primary')
            # Main text — centered, larger
            tb = slide.shapes.add_textbox(
                Inches(text_left + 0.5), Inches(cur_y + 0.40),
                Inches(text_w - 1.0), Inches(bullet_text_h - 0.45)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = bullet
            run.font.size = Pt(bullet_font_size)
            run.font.color.rgb = get_color(spec, 'on_surface')
            p.line_spacing = 1.4
        elif is_sparse:
            # Decorative card background for sparse content
            card_left = text_left + 0.3
            card_w = text_w - 0.6
            card_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(card_left), Inches(cur_y),
                Inches(card_w), Inches(bullet_text_h)
            )
            card_bg.fill.solid()
            card_bg.fill.fore_color.rgb = get_color(spec, 'surface_variant')
            card_bg.line.fill.background()
            # Accent bar on left
            accent_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(card_left), Inches(cur_y),
                Inches(0.06), Inches(bullet_text_h)
            )
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = get_color(spec, 'primary')
            accent_bar.line.fill.background()
            # Numbered circle with contextual icon
            circle_size = 0.40
            bullet_icon = get_icon_for_text(bullet, i)
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(card_left + 0.25), Inches(cur_y + (bullet_text_h - circle_size) / 2),
                Inches(circle_size), Inches(circle_size)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = get_color(spec, 'primary')
            circle.line.fill.background()
            num_tb = slide.shapes.add_textbox(
                Inches(card_left + 0.25), Inches(cur_y + (bullet_text_h - circle_size) / 2),
                Inches(circle_size), Inches(circle_size)
            )
            num_tf = num_tb.text_frame
            num_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            num_p = num_tf.paragraphs[0]
            num_p.alignment = PP_ALIGN.CENTER
            num_run = num_p.add_run()
            num_run.text = bullet_icon
            num_run.font.size = Pt(14)
            num_run.font.bold = True
            num_run.font.color.rgb = get_color(spec, 'on_primary')
            # Text (offset past circle)
            tb = slide.shapes.add_textbox(
                Inches(card_left + 0.80), Inches(cur_y + 0.10),
                Inches(card_w - 1.10), Inches(bullet_text_h - 0.20)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = bullet
            run.font.size = Pt(bullet_font_size)
            run.font.color.rgb = get_color(spec, 'on_surface')
            p.line_spacing = 1.5
        else:
            tb = slide.shapes.add_textbox(
                Inches(text_left + 0.1), Inches(cur_y),
                Inches(text_w - 0.2), Inches(bullet_text_h)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"• {bullet}"
            run.font.size = Pt(bullet_font_size)
            run.font.color.rgb = get_color(spec, 'on_surface')
            p.line_spacing = 1.5

    if remaining_comps:
        comp_spacing = 0.10 if use_vertical_layout else 0.15
        comp_top = bullet_y + bullet_limit * bullet_spacing + comp_spacing
        # When visual is below (vertical layout), limit component height to leave room for visual
        if use_vertical_layout and has_visual:
            max_comp_h = min(avail_h * 0.35, 2.0)  # Max 35% of available or 2in for components
        else:
            max_comp_h = avail_h - (comp_top - top)
        render_components_simple(slide, remaining_comps, spec, grid, text_left + 0.1, comp_top, text_w - 0.2, max_comp_h)

    if has_visual:
        # Choose visual placement based on layout mode
        if use_vertical_layout:
            # Vertical layout: visual below text (full-width)
            # Calculate actual bottom position of rendered content
            content_bottom = bullet_y + bullet_limit * bullet_spacing
            
            # Use accurate component height estimate
            if remaining_comps:
                comp_h_est = _estimate_components_height(remaining_comps, spec, text_w - 0.2)
                content_bottom += 0.10 + comp_h_est
            
            # Add proper gap before visual (0.20in for clarity)
            vis_gap = 0.20
            vis_top = content_bottom + vis_gap
            vis_h_available = avail_h - (vis_top - top) - 0.15
            
            # Smart height adjustment for horizontal flowcharts
            # Horizontal layouts (LR) with few nodes should use less height to avoid stretching
            visual_data = sd.get('visual', {})
            mermaid_code = visual_data.get('placeholder_data', {}).get('mermaid_code', '')
            is_horizontal = 'LR' in mermaid_code or 'RL' in mermaid_code
            node_count = mermaid_code.count('[') if mermaid_code else 0
            
            if is_horizontal and node_count <= 5:
                # Simple horizontal flowchart: limit height to avoid stretching
                vis_h_to_use = min(3.2, vis_h_available)
            else:
                # Complex diagrams: ensure minimum height
                vis_h_to_use = max(2.5, vis_h_available)
            
            render_visual(slide, sd['visual'], spec, grid,
                          vis_left, vis_top, vis_w, vis_h_to_use)
        elif is_split:
            # Split layout: either center visual independently or place at top
            if fill_mode == 'split-center':
                vis_h_to_use = min(vis_est_h, max(0.5, avail_h - 0.10 - vis_v_offset))
                render_visual(slide, sd['visual'], spec, grid,
                              vis_left, top + 0.15 + vis_v_offset, vis_w, vis_h_to_use)
            else:
                # Top alignment (default/expand)
                # Constrain visual container height to estimation to prevents excessive centering "gap" 
                # inside the placeholder renderer.
                safe_est = _estimate_visual_height(sd.get('visual'), spec, avail_h, reserved=0.2)
                vis_h_to_use = min(safe_est, max(0.5, avail_h - 0.10))
                render_visual(slide, sd['visual'], spec, grid,
                              vis_left, top + 0.15, vis_w, vis_h_to_use)
        else:
            # Visual-only or combined-center in non-split: center if requested (fill_mode may be 'center')
            vis_v_offset = _apply_vertical_offset(avail_h * 0.5, avail_h, fill_mode)
            vis_h_to_use = chosen_vis_h if ('chosen_vis_h' in locals() and chosen_vis_h) else max(0.5, avail_h - 0.10 - vis_v_offset)
            render_visual(slide, sd['visual'], spec, grid,
                          vis_left, top + 0.15 + vis_v_offset, vis_w, vis_h_to_use)


def render_slide_two_column(slide: Any, sd: Dict, spec: Dict, grid: GridSystem, **ctx) -> None:
    """Two-column or comparison layout."""
    bar_h = ctx.get('bar_h', 0.55)


def render_decisions_simple(slide: Any, decisions: List[Dict], spec: Dict, grid: GridSystem, left: float, top: float, width: float, avail_h: float = None) -> float:
    """Render a simple vertical list of decision cards. Returns used height."""
    if not decisions:
        return 0
    n = len(decisions)
    gap = 0.12
    base_h = 0.65
    # Dynamic expansion: distribute available height among decision cards
    if avail_h and avail_h > (base_h + gap) * n:
        h = min(1.6, (avail_h - gap * (n - 1)) / n)  # Expand up to 1.6in per card
    else:
        h = base_h
    # Dynamic positions based on card height
    title_h = min(0.35, h * 0.40)
    desc_top_offset = title_h + 0.10
    y = top
    for dec in decisions:
        card = create_card_shape(slide, spec, left, y, width, h)
        card.fill.solid()
        card.fill.fore_color.rgb = get_color(spec, 'primary_container')
        card.line.fill.background()
        # Title
        tb = slide.shapes.add_textbox(Inches(left + 0.12), Inches(y + 0.08), Inches(width - 0.24), Inches(title_h))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = dec.get('title', dec.get('label', ''))
        run.font.size = Pt(get_font_size(spec, 'body_text'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_primary_container')
        apply_font_to_run(run, spec)
        # Description
        desc = dec.get('description') or dec.get('rationale') or ''
        if desc:
            tb2 = slide.shapes.add_textbox(Inches(left + 0.12), Inches(y + desc_top_offset), Inches(width - 0.24), Inches(h - desc_top_offset - 0.08))
            p2 = tb2.text_frame.paragraphs[0]
            run2 = p2.add_run()
            run2.text = desc
            run2.font.size = Pt(get_font_size(spec, 'label'))
            run2.font.color.rgb = get_color(spec, 'on_surface_variant')
        y += h + gap
    return y - top


def render_comparison_items_simple(slide: Any, items: List[Dict], spec: Dict, grid: GridSystem, left: float, top: float, width: float, avail_h: float = None, slide_id: int = None) -> float:
    """Render comparison items as side-by-side cards with a mixed sizing strategy.

    Strategy (mixed):
    1. Compute content-driven height for each card (based on header + estimated body lines).
    2. Use max content height across cards as baseline.
    3. If `avail_h` is provided and `content_fill` == 'expand', expand cards to fill
       available height, capped by `max_card_h` from slide_overrides if set.
    Returns used vertical space (card_h + bottom gap).
    """
    if not items:
        return 0
    import math

    n = len(items)
    gap = 0.15

    header_h = 0.35
    sep_h = 0.02
    padding_v = 0.24  # top+bottom combined padding approximation
    min_card_h = 1.8

    def _is_cjk(ch: str) -> bool:
        return any([
            '\u4e00' <= ch <= '\u9fff',
            '\u3400' <= ch <= '\u4dbf',
            '\u3000' <= ch <= '\u303f',
            '\uff00' <= ch <= '\uffef',
        ])

    def _estimate_text_width(text: str, font_pt: float) -> float:
        """Estimate total rendered width in inches, CJK-aware."""
        w = 0.0
        for ch in text:
            if _is_cjk(ch):
                w += (font_pt / 72.0) * 0.95  # CJK ~square
            elif ch == '\n':
                pass
            elif ch.isspace():
                w += (font_pt / 72.0) * 0.35
            else:
                w += (font_pt / 72.0) * 0.50  # Latin
        return w

    def _collect_details(item: Dict) -> list:
        details = []
        attrs = item.get('attributes', {})
        if isinstance(attrs, dict):
            for k, v in attrs.items():
                if v is not None:
                    details.append(f"{k}: {v}")
        for key in ('advantage', 'risk', 'recommendation', 'rationale', 'description', 'short_action', 'action', 'next_step'):
            if item.get(key):
                details.append(str(item.get(key)))
        return details

    def _estimate_card_height(items: List[Dict], card_w: float, font_pt: float) -> float:
        """Estimate max card height across all items at given font size."""
        text_area_w = card_w - 0.40  # match actual rendering: pad_side*2 + body internal margins
        line_h = font_pt / 72.0 + 0.08
        max_h = min_card_h
        for item in items:
            details = _collect_details(item)
            total_lines = 0
            for detail in details:
                tw = _estimate_text_width(detail, font_pt)
                total_lines += max(1, math.ceil(tw / max(0.5, text_area_w)))
            body_h = total_lines * line_h
            ch = header_h + sep_h + body_h + padding_v
            max_h = max(max_h, ch)
        return max_h

    # ── Overflow degradation chain ──
    # Try horizontal layout first; if content overflows, try:
    # 1. Reduce font size (body_text → min 10pt, step -2)
    # 2. Switch to 2×2 grid layout (if n >= 3)
    base_font_pt = get_font_size(spec, 'body_text')
    effective_avail = avail_h - 0.15 if avail_h and avail_h > 0 else 999

    # Determine layout: single row vs 2×2 grid
    use_grid = False
    rows, cols = 1, n
    body_font_pt = base_font_pt

    # Phase 1: try single row, reducing font
    card_w = (width - gap * (n - 1)) / max(n, 1)
    for try_pt in range(int(base_font_pt), 9, -2):
        est_h = _estimate_card_height(items, card_w, try_pt)
        if est_h <= effective_avail:
            body_font_pt = try_pt
            break
        body_font_pt = try_pt
    else:
        # Phase 2: switch to 2×2 grid if n >= 3
        if n >= 3:
            use_grid = True
            cols = math.ceil(n / 2)
            rows = math.ceil(n / cols)
            card_w = (width - gap * (cols - 1)) / max(cols, 1)
            body_font_pt = base_font_pt  # reset font for grid layout
            for try_pt in range(int(base_font_pt), 9, -2):
                est_h = _estimate_card_height(items, card_w, try_pt)
                total_h = est_h * rows + gap * (rows - 1)
                if total_h <= effective_avail:
                    body_font_pt = try_pt
                    break
                body_font_pt = try_pt

    # Compute final card height
    base_h = _estimate_card_height(items, card_w, body_font_pt)

    # Apply expand/cap logic — NEVER expand beyond 2x content height or 3.5in
    CARD_H_CAP = 2.5  # Absolute max card height to prevent spanning entire slide
    card_h = base_h
    fill_mode = get_content_fill(spec, 'comparison', slide_id=slide_id)
    if not use_grid and fill_mode == 'expand' and avail_h and avail_h > 0:
        expanded_h = max(min_card_h, avail_h - 0.15)
        max_cap = get_max_card_h(spec, slide_id)
        if max_cap > 0:
            expanded_h = min(expanded_h, max_cap)
        # Cap at 2x content or absolute max
        expanded_h = min(expanded_h, base_h * 2.0, CARD_H_CAP)
        card_h = max(base_h, expanded_h)
    elif use_grid and fill_mode == 'expand' and avail_h and avail_h > 0:
        per_row_avail = (avail_h - gap * (rows - 1)) / rows
        card_h = max(base_h, per_row_avail - 0.15)
        card_h = min(card_h, CARD_H_CAP)  # Cap
        max_cap = get_max_card_h(spec, slide_id)
        if max_cap > 0:
            card_h = min(card_h, max_cap)

    # ── Render cards ──
    def _render_card(item: Dict, cx: float, cy: float, cw: float, ch: float, card_idx: int = 0) -> None:
        card = create_card_shape(slide, spec, cx, cy, cw, ch)
        card.fill.solid()
        card.fill.fore_color.rgb = get_color(spec, 'surface_variant')
        card.line.fill.background()
        add_shadow(card)
        # Icon circle in header area
        card_label = item.get('label', item.get('name', item.get('title', '')))
        card_icon = get_icon_for_text(card_label, card_idx)
        icon_sz = 0.28
        icon_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(cx + 0.10), Inches(cy + 0.08),
            Inches(icon_sz), Inches(icon_sz)
        )
        icon_circle.fill.solid()
        icon_circle.fill.fore_color.rgb = get_color(spec, 'primary')
        icon_circle.line.fill.background()
        # Put icon text directly on the oval shape for native centering
        # Smaller font + margin_top offset to compensate for Unicode ascender bias
        icon_font_pt = int(icon_sz * 32)  # ~9pt for 0.28in circle (was 36)
        glyph_offset = Pt(max(1, int(icon_sz * 5)))  # push glyph down
        icon_circle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        icon_circle.text_frame.margin_top = glyph_offset
        icon_circle.text_frame.margin_bottom = Emu(1)
        icon_circle.text_frame.margin_left = Emu(1)
        icon_circle.text_frame.margin_right = Emu(1)
        icon_p = icon_circle.text_frame.paragraphs[0]
        icon_p.alignment = PP_ALIGN.CENTER
        icon_p.space_before = Pt(0)
        icon_p.space_after = Pt(0)
        icon_run = icon_p.add_run()
        icon_run.text = card_icon
        icon_run.font.size = Pt(icon_font_pt)
        icon_run.font.bold = True
        icon_run.font.color.rgb = get_color(spec, 'on_primary')
        # Header (shifted right to accommodate icon)
        pad_side = 0.16  # side padding for body & separator
        header_left = cx + 0.10 + icon_sz + 0.06
        header_w = cw - pad_side * 2 - icon_sz - 0.06 + (pad_side - 0.10)
        tb = slide.shapes.add_textbox(Inches(header_left), Inches(cy + 0.10), Inches(max(header_w, 0.5)), Inches(0.35))
        tf_hdr = tb.text_frame
        tf_hdr.margin_left = Inches(0)
        tf_hdr.margin_right = Inches(0)
        p = tf_hdr.paragraphs[0]
        run = p.add_run()
        run.text = item.get('label', item.get('name', item.get('title', '')))
        run.font.size = Pt(get_font_size(spec, 'slide_subtitle'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_surface')
        apply_font_to_run(run, spec)
        # Separator
        s = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(cx + pad_side), Inches(cy + 0.48),
            Inches(cw - pad_side * 2), Inches(sep_h)
        )
        s.fill.solid()
        s.fill.fore_color.rgb = get_color(spec, 'outline')
        s.line.fill.background()
        # Body — generous margins to prevent text overflow
        body_top = cy + 0.55
        body_h = ch - 0.55 - 0.16  # 0.16in bottom margin
        details = _collect_details(item)
        detail_text = '\n'.join(details)
        if detail_text:
            tb2 = slide.shapes.add_textbox(
                Inches(cx + pad_side), Inches(body_top),
                Inches(cw - pad_side * 2), Inches(max(body_h, 0.5))
            )
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            tf2.margin_left = Inches(0.04)
            tf2.margin_right = Inches(0.04)
            tf2.margin_top = Inches(0.02)
            tf2.margin_bottom = Inches(0.02)
            p2 = tf2.paragraphs[0]
            run2 = p2.add_run()
            run2.text = detail_text
            run2.font.size = Pt(body_font_pt)
            run2.font.color.rgb = get_color(spec, 'on_surface')
            p2.line_spacing = 1.4

    if use_grid:
        # 2×2 (or 2×N) grid layout
        total_h = 0.0
        for row in range(rows):
            for col in range(cols):
                idx = row * cols + col
                if idx >= n:
                    break
                cx = left + col * (card_w + gap)
                cy = top + row * (card_h + gap)
                _render_card(items[idx], cx, cy, card_w, card_h, card_idx=idx)
            total_h = (row + 1) * (card_h + gap)
        return total_h
    else:
        # Single row layout
        for i, item in enumerate(items):
            cx = left + i * (card_w + gap)
            _render_card(item, cx, top, card_w, card_h, card_idx=i)
        return card_h + 0.15


def render_components_simple(slide: Any, comps: Dict, spec: Dict, grid: GridSystem, left: float, top: float, width: float, avail_h: float) -> float:
    """Render a subset of components with simple renderers. Returns used height."""
    cursor = 0.0
    # KPIs
    if comps.get('kpis'):
        h = render_kpis(slide, comps['kpis'], spec, grid, left, top + cursor, width, avail_h=avail_h - cursor if avail_h else None)
        cursor += h
    # Comparison items
    if comps.get('comparison_items'):
        h = render_comparison_items_simple(slide, comps['comparison_items'], spec, grid, left, top + cursor, width)
        cursor += h
    # Decisions
    if comps.get('decisions'):
        h = render_decisions_simple(slide, comps['decisions'], spec, grid, left, top + cursor, width, avail_h=avail_h - cursor if avail_h else None)
        cursor += h
    # Bullets (before callouts — bullets provide main content, callouts emphasize)
    if comps.get('bullets') and cursor < avail_h:
        # render remaining bullets — handle dict format with label/value or text
        def _bullet_text(b):
            if isinstance(b, dict):
                t = b.get('text', '')
                if t:
                    return t
                label = b.get('label', '')
                value = b.get('value', '')
                return f"{label}: {value}" if value else label
            return str(b)
        bullets = [_bullet_text(b) for b in comps['bullets']]
        n_items = min(len(bullets), 4)
        remaining_h = avail_h - cursor

        # Use numbered takeaway cards when 2-4 sparse bullets fit horizontally
        # For 4 items on full-width slides, use 2×2 grid to avoid cramped text
        use_grid_2x2 = (n_items == 4 and width > 8.0)
        if n_items >= 2 and n_items <= 4 and remaining_h >= 0.70:
            card_gap = 0.12
            if use_grid_2x2:
                cols = 2
                card_w = (width - card_gap) / cols
                card_h = 0.70
                row_gap = 0.12
            else:
                cols = n_items
                card_w = (width - card_gap * (n_items - 1)) / n_items
                card_h = 0.55
                row_gap = 0
            num_w = 0.32  # width reserved for number column
            for i, b in enumerate(bullets[:n_items]):
                col_idx = i % cols
                row_idx = i // cols
                cx = left + col_idx * (card_w + card_gap)
                cy = top + cursor + row_idx * (card_h + row_gap)
                # Card background
                card_bg = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(cx), Inches(cy),
                    Inches(card_w), Inches(card_h)
                )
                card_bg.fill.solid()
                card_bg.fill.fore_color.rgb = get_color(spec, 'primary_container')
                card_bg.line.fill.background()
                # Accent left bar
                accent_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(cx), Inches(cy),
                    Inches(0.04), Inches(card_h)
                )
                accent_bar.fill.solid()
                accent_bar.fill.fore_color.rgb = get_color(spec, 'primary')
                accent_bar.line.fill.background()
                # Number — large, vertically centered on left side
                num_tb = slide.shapes.add_textbox(
                    Inches(cx + 0.08), Inches(cy),
                    Inches(num_w), Inches(card_h)
                )
                num_tf = num_tb.text_frame
                num_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                num_tf.margin_left = Emu(1)
                num_tf.margin_right = Emu(1)
                num_p = num_tf.paragraphs[0]
                num_p.alignment = PP_ALIGN.CENTER
                num_run = num_p.add_run()
                num_run.text = str(i + 1)
                num_run.font.size = Pt(22)
                num_run.font.bold = True
                num_run.font.color.rgb = get_color(spec, 'primary')
                # Text — right of number, vertically centered
                txt_left = cx + 0.08 + num_w
                txt_w = card_w - 0.08 - num_w - 0.08
                txt_tb = slide.shapes.add_textbox(
                    Inches(txt_left), Inches(cy),
                    Inches(txt_w), Inches(card_h)
                )
                txt_tf = txt_tb.text_frame
                txt_tf.word_wrap = True
                txt_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                txt_tf.margin_left = Inches(0.04)
                txt_tf.margin_right = Inches(0.04)
                txt_p = txt_tf.paragraphs[0]
                txt_run = txt_p.add_run()
                txt_run.text = b
                txt_run.font.size = Pt(get_font_size(spec, 'body') - 1)
                txt_run.font.color.rgb = get_color(spec, 'on_primary_container')
                txt_p.line_spacing = 1.2
            n_rows = (n_items + cols - 1) // cols
            cursor += n_rows * card_h + (n_rows - 1) * row_gap + 0.15
        else:
            # Standard bullet list for many items
            for i, b in enumerate(bullets[:8]):
                if cursor > avail_h - 0.5:
                    break
                tb = slide.shapes.add_textbox(Inches(left + 0.05), Inches(top + cursor), Inches(width - 0.10), Inches(0.45))
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = f"• {b}"
                run.font.size = Pt(get_font_size(spec, 'bullet_text'))
                run.font.color.rgb = get_color(spec, 'on_surface')
                p.line_spacing = 1.5
                cursor += 0.48
    # Callouts (after bullets — callouts emphasize key takeaways)
    if comps.get('callouts') and cursor < avail_h:
        h = render_callouts(slide, comps['callouts'], spec, grid, left, top + cursor, width)
        cursor += h
    return cursor


def render_slide_cta(slide: Any, sd: Dict, spec: Dict, grid: GridSystem, **ctx) -> None:
    """Call-to-action — centered prominent card."""
    bar_h = ctx.get('bar_h', 0.55)
    slide_id = ctx.get('slide_id')
    top, avail_h = _content_zone(grid, bar_h)
    fill_mode = get_content_fill(spec, sd.get('slide_type', 'cta'), slide_id=slide_id)
    comps = sd.get('components', {})

    cl, cw = grid.col_span(8, 2)

    # Estimate card content height
    n_bullets = min(len(sd.get('content', [])), 6)
    est_card_content = n_bullets * 0.55 + 0.45 + _estimate_components_height(comps, spec, cw - 0.70)
    est_card_h = est_card_content + 0.50  # top/bottom padding

    if fill_mode == 'center':
        card_h = min(est_card_h, avail_h - 0.40)
        card_top = top + (avail_h - card_h) / 2
    else:
        card_top = top + 0.40
        card_h = avail_h - 0.80

    # Card background
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(cl), Inches(card_top), Inches(cw), Inches(card_h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = get_color(spec, 'primary_container')
    card.line.fill.background()
    add_shadow(card)

    # Content — use content[] first, fall back to components.bullets
    cta_items = sd.get('content', [])
    if not cta_items and comps.get('bullets'):
        cta_items = []
        for b in comps['bullets']:
            if isinstance(b, dict):
                label = b.get('label', b.get('text', ''))
                value = b.get('value', '')
                cta_items.append(f"{label}: {value}" if value else label)
            else:
                cta_items.append(str(b))

    cursor = card_top + 0.30
    for i, bullet in enumerate(cta_items[:6]):
        tb = slide.shapes.add_textbox(
            Inches(cl + 0.35), Inches(cursor + i * 0.55),
            Inches(cw - 0.70), Inches(0.50)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = bullet
        run.font.size = Pt(get_font_size(spec, 'slide_subtitle'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_primary_container')
    cursor += len(cta_items[:6]) * 0.55 + 0.15

    # Components inside card (skip bullets already rendered)
    remaining_comps = {k: v for k, v in comps.items() if k != 'bullets' and v}
    if remaining_comps:
        render_components(slide, remaining_comps, spec, grid,
                          cl + 0.35, cursor, cw - 0.70,
                          card_h - (cursor - card_top) - 0.20)


def _render_footer_and_notes(slide: Any, sd: Dict, spec: Dict, grid: GridSystem, stype: str, sec_title: str, sec_accent_token: str, sec_index: int, sections_len: int, slide_num: int, total_slides: int) -> None:
    """Render bottom bar and speaker notes."""
    if stype not in ('title', 'section_divider'):
        render_bottom_bar(
            slide,
            spec,
            grid,
            sec_title,
            slide_num,
            total_slides,
            section_index=sec_index,
            total_sections=sections_len,
            accent_token=sec_accent_token,
        )
    render_speaker_notes(slide, sd.get('speaker_notes'))


def normalize_slide_data(sd: Dict) -> Dict:
    """Normalize semantic slide data to canonical keys.

    Unifies equivalent keys so renderers only need to check one name.
    This is the single source of truth for key mapping — add new aliases here.
    Mutations are in-place for efficiency.
    """
    comps = sd.setdefault('components', {})

    # ── 1. Reconcile content[] ↔ components.bullets ──
    # If content is empty but components.bullets exists, populate content
    if not sd.get('content') and comps.get('bullets'):
        sd['content'] = [
            b.get('text', b.get('label', '')) if isinstance(b, dict) else str(b)
            for b in comps['bullets']
        ]

    # ── 2. Normalize timeline_items keys ──
    if comps.get('timeline_items'):
        for item in comps['timeline_items']:
            # phase/milestone/title → label (canonical)
            if not item.get('label'):
                item['label'] = item.get('phase', item.get('milestone', item.get('title', '')))
            # period/date/time/description → date (canonical)
            if not item.get('date'):
                item['date'] = item.get('period', item.get('time', item.get('description', '')))

    # ── 3. Normalize comparison_items keys ──
    if comps.get('comparison_items'):
        for item in comps['comparison_items']:
            if not item.get('label'):
                item['label'] = item.get('name', item.get('title', ''))

    # ── 4. Normalize decisions keys ──
    if comps.get('decisions'):
        for dec in comps['decisions']:
            if not dec.get('title'):
                dec['title'] = dec.get('label', dec.get('name', ''))

    # ── 5. Normalize KPI keys ──
    if comps.get('kpis'):
        for kpi in comps['kpis']:
            if not kpi.get('label'):
                kpi['label'] = kpi.get('name', kpi.get('title', ''))
            if not kpi.get('value'):
                kpi['value'] = kpi.get('metric', kpi.get('data', ''))

    return sd


def render_slide(prs: Presentation, sd: Dict, spec: Dict, grid: GridSystem, sections: List[Dict], slide_num: int, total_slides: int) -> None:
    """Render a single slide with full styling (delegates to helpers)."""
    # Normalize all keys before any rendering
    normalize_slide_data(sd)

    stype = sd.get('slide_type', 'bullet-list')
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Section context — semantic JSON uses 'id', fallback to slide_num
    slide_id = sd.get('id', sd.get('slide_id', slide_num))
    section = find_section_for_slide(slide_id, sections)
    sec_id = section.get('id', 'A')
    sec_title = section.get('title', '')
    sec_accent_token = spec.get('section_accents', {}).get(sec_id, 'primary')
    sec_index = next((i for i, s in enumerate(sections) if s.get('id') == sec_id), 0)

    # Background
    _render_background(slide, sd, spec, grid, stype)

    # Title bar
    tb_mode = get_title_bar_mode(spec, stype)
    bar_h = 0.0
    if tb_mode != 'none':
        section_label = f"Section {sec_id} · {sec_title}" if sec_title else ''
        bar_h = render_title_bar(slide, spec, grid, sd.get('title', ''), slide_num, total_slides, section_label=section_label, accent_color_token=sec_accent_token, mode=tb_mode)

    # Content dispatch by slide type
    renderer = RENDERERS.get(stype, render_slide_bullets)
    renderer(slide, sd, spec, grid, bar_h=bar_h, section_index=sec_index, total_sections=len(sections), slide_id=slide_id)

    # Footer & notes
    _render_footer_and_notes(slide, sd, spec, grid, stype, sec_title, sec_accent_token, sec_index, len(sections), slide_num, total_slides)


# Final RENDERERS mapping
RENDERERS = {
    'title': render_slide_title,
    'section_divider': render_section_divider,
    'bullet-list': render_slide_bullets,
    'two-column': render_slide_two_column,
    'comparison': render_slide_comparison,
    'decision': render_slide_decision,
    'data-heavy': render_slide_data_heavy,
    'data-summary': render_slide_data_heavy,
    'data-process': render_slide_bullets,  # Use bullets layout with vertical visual option
    'matrix': render_slide_matrix,
    'timeline': render_slide_timeline,
    'gantt': render_slide_timeline,
    'flowchart': render_slide_flowchart,
    'sequence': render_slide_flowchart,
    'call_to_action': render_slide_cta,
    'cta': render_slide_cta,
    'next-steps': render_slide_cta,
    'technical': render_slide_two_column,
    'process': render_slide_flowchart,
    'table': render_slide_data_heavy,
    'waterfall': render_slide_data_heavy,
    'kpi_dashboard': render_slide_data_heavy,
    'default': render_slide_bullets,
}
