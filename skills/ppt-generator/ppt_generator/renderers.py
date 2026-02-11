"""Renderers for slides and components.

This module reuses helpers and grid to provide all rendering functions
exported by the original monolithic script.
"""
from typing import Any, Dict, List
import json
import os
import re
import logging

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from .helpers import (
    get_color, get_color_hex, get_font_size, hex_to_rgb, px_to_inches
)
from .grid import GridSystem


# For brevity we copy selectively the rendering helpers from the original
# script, focusing on title bar, bottom bar, speaker notes, and the main
# per-slide dispatcher. Additional component renderers (kpis, bullets,
# tables, visuals) are also included as smaller functions.


def get_layout_zones(spec: Dict) -> Dict[str, float]:
    lz = spec.get('layout_zones', {})
    return {
        'title_bar_h': lz.get('title_bar_height_default', 0.55),
        'title_bar_h_narrow': lz.get('title_bar_height_narrow', 0.4),
        'bottom_bar_h': max(lz.get('bottom_bar_height', 0.25), 0.25),
        'content_margin_top': lz.get('content_margin_top', 0.12),
        'content_bottom_margin': lz.get('content_bottom_margin', 0.2),
    }


def detect_schema_version(slide: Dict) -> int:
    """Detects whether a slide is v1 or v2 schema.

    Returns 2 if `layout_intent.regions` key exists (even if empty), otherwise 1.
    """
    try:
        if isinstance(slide, dict) and 'layout_intent' in slide and isinstance(slide.get('layout_intent'), dict) and 'regions' in slide['layout_intent']:
            return 2
    except Exception:
        pass
    return 1


def resolve_data_source(slide: Dict, path: str):
    """Resolve a dot-separated data source path from a slide dict.

    Examples:
      - 'components.kpis' -> slide['components']['kpis']
      - 'visual' -> slide['visual']
      - 'content' -> slide['content']

    Returns None if the path cannot be resolved.
    """
    if not path or not isinstance(path, str):
        return None
    cur = slide
    parts = path.split('.')
    try:
        for p in parts:
            if isinstance(cur, dict) and p in cur:
                cur = cur[p]
            else:
                # Not found; return None
                return None
        return cur
    except Exception:
        return None


def compute_region_bounds(position: str, grid: GridSystem, bar_h: float, bottom_reserve: float = 0.5) -> tuple:
    """Compute (left, top, width, height) in inches for a given position marker.

    Supported markers:
      - 'full' -> full usable width
      - 'left-<pct>' / 'right-<pct>' -> left or right region with percentage of usable width
      - 'center-<pct>' -> centered region with percentage width
      - 'col-<start>-<span>' -> column-based region using grid.col_span
      - 'top-<pct>' -> full-width but height is pct of available content height

    Args:
        bottom_reserve: Space in inches to reserve at the bottom for footer/insight bar.
                        Default 0.5" (footer only). Use ~0.95" when insight bar is present.

    Best-effort: if parsing fails, returns full area bounds.
    """
    left = grid.margin_h
    top = bar_h + 0.12
    usable_w = grid.usable_w
    slide_h = grid.slide_h

    # Default full area height: reserve bottom space for footer / insight bar
    default_available_h = max(1.0, slide_h - bar_h - bottom_reserve)
    width = usable_w
    height = default_available_h

    if not position or not isinstance(position, str):
        return (left, top, width, height)

    pos = position.strip().lower()
    if pos == 'full':
        return (left, top, width, height)

    import re
    m = re.match(r'^(left|right|center)-(\d{1,3})$', pos)
    if m:
        side, val = m.group(1), int(m.group(2))
        pct = max(1, min(100, val)) / 100.0
        w = usable_w * pct
        if side == 'left':
            return (left, top, w, height)
        elif side == 'right':
            l = left + (usable_w - w)
            return (l, top, w, height)
        else:  # center
            l = left + (usable_w - w) / 2
            return (l, top, w, height)

    m2 = re.match(r'^top-(\d{1,3})$', pos)
    if m2:
        pct = max(1, min(100, int(m2.group(1)))) / 100.0
        h = default_available_h * pct
        return (left, top, width, h)

    m2b = re.match(r'^bottom-(\d{1,3})$', pos)
    if m2b:
        pct = max(1, min(100, int(m2b.group(1)))) / 100.0
        h = default_available_h * pct
        t = top + (default_available_h - h)
        return (left, t, width, h)

    m3 = re.match(r'^col-(\d+)-(\d+)$', pos)
    if m3:
        start = int(m3.group(1))
        span = int(m3.group(2))
        start = max(0, min(grid.columns - 1, start))
        span = max(1, min(grid.columns - start, span))
        l, w = grid.col_span(span, start)
        return (l, top, w, height)

    # fallback: return full
    return (left, top, width, height)


def get_title_bar_mode(spec: Dict, slide_type: str) -> str:
    if slide_type in ('title', 'section_divider'):
        return 'none'
    stl = spec.get('slide_type_layouts', {})
    entry = stl.get(slide_type, stl.get('default', {}))
    return entry.get('title_bar', 'standard')


def apply_font_to_run(run: Any, spec: Dict) -> None:
    try:
        font_family = spec.get('typography_system', {}).get('font_family')
        cjk_font = spec.get('typography_system', {}).get('cjk_font_family')
        if font_family:
            run.font.name = font_family
            r_pr = getattr(run._element, 'rPr', None)
            if r_pr is not None:
                try:
                    r_pr.rFonts.set(qn('a:ascii'), font_family)
                    r_pr.rFonts.set(qn('a:hAnsi'), font_family)
                except Exception:
                    pass
        if cjk_font:
            r_pr = getattr(run._element, 'rPr', None)
            if r_pr is not None:
                try:
                    r_pr.rFonts.set(qn('a:ea'), cjk_font)
                except Exception:
                    pass
    except Exception:
        pass


def create_card_shape(slide: Any, spec: Dict, left_in: float, top_in: float, width_in: float, height_in: float):
    try:
        br = spec.get('component_library', {}).get('card', {}).get('border_radius')
    except Exception:
        br = None
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if br is None or br != 0 else MSO_SHAPE.RECTANGLE
    return slide.shapes.add_shape(shape_type, Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))


def render_assertion_title(slide: Any, spec: Dict, grid: GridSystem, assertion: str, title: str, slide_num: int, _total_slides: int,
                         section_label: str = '', accent_color_token: str = 'primary', mode: str = 'standard') -> float:
    """Render an assertion-style title bar (assertion headline + subtitle).

    Assertion: 16pt bold white (left-aligned)
    Subtitle (original title): 10pt, on_surface_variant (light)
    Minimum bar height: 0.85" (best-effort)
    """
    if mode == 'none':
        return 0.0
    lz = get_layout_zones(spec)
    # Prefer explicit tall bar for assertion headlines
    bar_h = max(lz['title_bar_h'], 0.85)
    if section_label:
        bar_h = max(bar_h, 1.05)  # taller to fit section label + assertion

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(grid.slide_w), Inches(bar_h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = get_color(spec, accent_color_token)
    bar.line.fill.background()

    # Determine assertion text top based on whether section label is present
    assertion_top = 0.12
    if section_label:
        tb = slide.shapes.add_textbox(
            Inches(grid.margin_h), Inches(0.06), Inches(4), Inches(0.18)
        )
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = section_label
        run.font.size = Pt(get_font_size(spec, 'section_label'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_primary')
        apply_font_to_run(run, spec)
        assertion_top = 0.28  # below section label

    # Assertion headline
    tb = slide.shapes.add_textbox(
        Inches(grid.margin_h), Inches(assertion_top), Inches(grid.slide_w - 2 * grid.margin_h - 1.0), Inches(bar_h - assertion_top - 0.03)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = assertion
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = get_color(spec, 'on_primary')
    apply_font_to_run(run, spec)

    # Subtitle (original title)
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = title or ''
    run2.font.size = Pt(10)
    run2.font.color.rgb = get_color(spec, 'on_surface_variant')
    apply_font_to_run(run2, spec)

    # Page number
    tb_num = slide.shapes.add_textbox(
        Inches(grid.slide_w - 1.0), Inches(0.12), Inches(0.7), Inches(0.35)
    )
    tf_num = tb_num.text_frame
    p_num = tf_num.paragraphs[0]
    run_num = p_num.add_run()
    run_num.text = f"{slide_num}"
    run_num.font.size = Pt(get_font_size(spec, 'page_number'))
    run_num.font.bold = True
    run_num.font.color.rgb = get_color(spec, 'on_primary')
    apply_font_to_run(run_num, spec)

    return bar_h


def render_title_bar(slide: Any, spec: Dict, grid: GridSystem, title: str, slide_num: int, _total_slides: int,
                     section_label: str = '', accent_color_token: str = 'primary', mode: str = 'standard') -> float:
    if mode == 'none':
        return 0.0
    lz = get_layout_zones(spec)
    bar_h = lz['title_bar_h_narrow'] if mode == 'narrow' else lz['title_bar_h']
    title_font_pt = get_font_size(spec, 'slide_title')
    title_text_h = title_font_pt / 72.0 + 0.08
    if section_label:
        min_bar_h = 0.28 + title_text_h
    else:
        min_bar_h = 0.10 + title_text_h
    bar_h = max(bar_h, min_bar_h)

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(grid.slide_w), Inches(bar_h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = get_color(spec, accent_color_token)
    bar.line.fill.background()

    if section_label:
        tb = slide.shapes.add_textbox(
            Inches(grid.margin_h), Inches(0.06), Inches(4), Inches(0.22)
        )
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = section_label
        run.font.size = Pt(get_font_size(spec, 'section_label'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_primary')
        apply_font_to_run(run, spec)

    title_top = 0.22 if section_label else 0.10
    tb = slide.shapes.add_textbox(
        Inches(grid.margin_h), Inches(title_top),
        Inches(grid.slide_w - 2 * grid.margin_h - 1.0), Inches(bar_h - title_top - 0.05)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(get_font_size(spec, 'slide_title'))
    run.font.bold = True
    run.font.color.rgb = get_color(spec, 'on_primary')
    apply_font_to_run(run, spec)

    tb_num = slide.shapes.add_textbox(
        Inches(grid.slide_w - 1.0), Inches(title_top), Inches(0.7), Inches(0.35)
    )
    tf_num = tb_num.text_frame
    p = tf_num.paragraphs[0]
    run = p.add_run()
    run.text = f"{slide_num}"
    run.font.size = Pt(get_font_size(spec, 'page_number'))
    run.font.bold = True

    return bar_h


# New: render title slide content for full-bleed / title slides
def render_title_slide(slide: Any, sd: Dict, spec: Dict, grid: GridSystem, slide_num: int, total_slides: int) -> None:
    """Ensure title slides render visible title, subtitle/tagline, and metadata even when title_bar is 'none'.

    Adds three text frames: title (prominent), subtitle/tagline (optional), metadata (date/author) at bottom.
    """
    title_text = sd.get('title') or sd.get('_deck_title') or ''
    subtitle = sd.get('subtitle') or (sd.get('speaker_notes') or '').split('\n', 1)[0][:140]
    # Title (centered area)
    tb_title = slide.shapes.add_textbox(
        Inches(grid.margin_h + 0.5), Inches(grid.slide_h * 0.22), Inches(grid.usable_w - 1.0), Inches(1.2)
    )
    tf = tb_title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(get_font_size(spec, 'display_large'))
    run.font.bold = True
    run.font.color.rgb = get_color(spec, 'on_primary') if get_color_hex(spec, 'primary') == get_color_hex(spec, 'primary') else get_color(spec, 'on_surface')
    apply_font_to_run(run, spec)

    # Subtitle/tagline (if present)
    if subtitle:
        tb_sub = slide.shapes.add_textbox(
            Inches(grid.margin_h + 0.8), Inches(grid.slide_h * 0.38), Inches(grid.usable_w - 1.6), Inches(0.6)
        )
        tf2 = tb_sub.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(get_font_size(spec, 'slide_subtitle'))
        run2.font.color.rgb = get_color(spec, 'on_primary')
        apply_font_to_run(run2, spec)

    # Metadata (bottom-left): date or author if provided
    meta = ''
    if sd.get('author'):
        meta = sd.get('author')
    elif sd.get('date'):
        meta = sd.get('date')
    else:
        # default to generation date to ensure metadata is present
        from datetime import datetime
        meta = datetime.utcnow().strftime('%Y-%m-%d')
    tb_meta = slide.shapes.add_textbox(Inches(grid.margin_h), Inches(grid.slide_h - 0.6), Inches(3.0), Inches(0.35))
    tfm = tb_meta.text_frame
    p3 = tfm.paragraphs[0]
    run3 = p3.add_run()
    run3.text = meta
    run3.font.size = Pt(get_font_size(spec, 'caption'))
    run3.font.color.rgb = get_color(spec, 'on_primary')
    apply_font_to_run(run3, spec)

    run.font.color.rgb = get_color(spec, 'on_primary')
    apply_font_to_run(run, spec)
    p.alignment = PP_ALIGN.RIGHT

    return bar_h


def render_bottom_bar(slide: Any, spec: Dict, grid: GridSystem, section_name: str, slide_num: int, total_slides: int,
                      section_index: int = 0, total_sections: int = 6, accent_token: str = 'primary') -> None:
    lz = get_layout_zones(spec)
    bar_h = lz['bottom_bar_h']
    bar_top = grid.slide_h - bar_h

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(bar_top), Inches(grid.slide_w), Inches(bar_h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = get_color(spec, 'surface_variant')
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(grid.margin_h), Inches(bar_top + 0.03), Inches(3), Inches(bar_h - 0.06))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = section_name
    run.font.size = Pt(get_font_size(spec, 'label'))
    run.font.bold = True
    run.font.color.rgb = get_color(spec, accent_token)

    prog_left = grid.slide_w / 2 - 1.5
    seg_w = 3.0 / max(total_sections, 1)
    for i in range(total_sections):
        seg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(prog_left + i * seg_w + 0.02), Inches(bar_top + bar_h / 2 - 0.03),
            Inches(seg_w - 0.04), Inches(0.06)
        )
        seg.fill.solid()
        if i <= section_index:
            seg.fill.fore_color.rgb = get_color(spec, accent_token)
        else:
            seg.fill.fore_color.rgb = get_color(spec, 'outline')
        seg.line.fill.background()

    tb = slide.shapes.add_textbox(
        Inches(grid.slide_w - grid.margin_h - 1.5), Inches(bar_top + 0.03), Inches(1.5), Inches(bar_h - 0.06)
    )
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{slide_num} / {total_slides}"
    run.font.size = Pt(get_font_size(spec, 'page_number'))
    run.font.color.rgb = get_color(spec, 'on_surface_variant')


def render_insight_bar(slide: Any, insight_text: str, spec: Dict, grid: GridSystem, accent_token: str = 'primary') -> None:
    """Render a compact insight bar above the bottom bar.

    Layout:
      - Height: 0.40"
      - Y position: slide_h - bottom_bar_h - 0.40"
      - Left padding: grid.margin_h
      - Text: white 10pt with a leading 💡 emoji
    """
    try:
        lz = get_layout_zones(spec)
        bar_h = 0.40
        bar_top = grid.slide_h - lz['bottom_bar_h'] - bar_h
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(bar_top), Inches(grid.slide_w), Inches(bar_h)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = get_color(spec, accent_token)
        bar.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(grid.margin_h), Inches(bar_top + 0.06), Inches(grid.usable_w - 0.5), Inches(bar_h - 0.12))
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"💡 {insight_text}"
        run.font.size = Pt(10)
        run.font.color.rgb = get_color(spec, 'on_primary')
        apply_font_to_run(run, spec)
    except Exception:
        # Best-effort; do not raise during slide rendering
        pass


def render_speaker_notes(slide: Any, notes_text: Any) -> None:
    if not notes_text:
        return
    text = notes_text if isinstance(notes_text, str) else json.dumps(notes_text, ensure_ascii=False)
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# Minimal component renderers used by slide renderers

def render_kpis(slide: Any, kpis: List[Dict], spec: Dict, _grid: GridSystem, left: float, top: float, width: float) -> float:
    if not kpis:
        return 0
    n = len(kpis)
    card_gap = 0.15
    card_w = (width - card_gap * (n - 1)) / max(n, 1)
    card_h = 0.85
    for i, kpi in enumerate(kpis):
        cx = left + i * (card_w + card_gap)
        card = create_card_shape(slide, spec, cx, top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = get_color(spec, 'primary_container')
        card.line.fill.background()
        # Value + trend arrow (separate colored run)
        tb = slide.shapes.add_textbox(Inches(cx + 0.12), Inches(top + 0.10), Inches(card_w - 0.24), Inches(0.40))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        # Value run
        run_val = p.add_run()
        run_val.text = str(kpi.get('value', ''))
        run_val.font.size = Pt(get_font_size(spec, 'kpi_value'))
        run_val.font.bold = True
        run_val.font.color.rgb = get_color(spec, 'on_primary_container')
        apply_font_to_run(run_val, spec)
        # Trend arrow run (colored: green ↑, red ↓, gray →)
        trend = kpi.get('trend', '')
        trend_info = {'up': (' ↑', 'success'), 'down': (' ↓', 'error'), 'stable': (' →', 'on_surface_variant')}
        if trend in trend_info:
            arrow_text, arrow_token = trend_info[trend]
            run_arrow = p.add_run()
            run_arrow.text = arrow_text
            run_arrow.font.size = Pt(get_font_size(spec, 'kpi_value'))
            run_arrow.font.bold = True
            run_arrow.font.color.rgb = get_color(spec, arrow_token)
            apply_font_to_run(run_arrow, spec)
        # Label
        tb2 = slide.shapes.add_textbox(Inches(cx + 0.12), Inches(top + 0.52), Inches(card_w - 0.24), Inches(0.25))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = kpi.get('label', '')
        run2.font.size = Pt(get_font_size(spec, 'kpi_label'))
        run2.font.color.rgb = get_color(spec, 'on_primary_container')
    return card_h + 0.15


# Visual/placeholder renderers (simplified)

# Optional: python-pptx native charts
try:
    from pptx.chart.data import CategoryChartData, XyChartData
    from pptx.enum.chart import XL_CHART_TYPE
    HAS_PPTX_CHARTS = True
except Exception:
    HAS_PPTX_CHARTS = False
    CategoryChartData = None
    XyChartData = None
    XL_CHART_TYPE = None

# Map visual types to XL chart types
NATIVE_CHART_TYPE_MAP = {
    'bar_chart': XL_CHART_TYPE.COLUMN_CLUSTERED if XL_CHART_TYPE is not None else None,
    'column_chart': XL_CHART_TYPE.COLUMN_CLUSTERED if XL_CHART_TYPE is not None else None,
    'horizontal_bar': XL_CHART_TYPE.BAR_CLUSTERED if XL_CHART_TYPE is not None else None,
    'line_chart': XL_CHART_TYPE.LINE_MARKERS if XL_CHART_TYPE is not None else None,
    'pie_chart': XL_CHART_TYPE.PIE if XL_CHART_TYPE is not None else None,
    'doughnut_chart': XL_CHART_TYPE.DOUGHNUT if XL_CHART_TYPE is not None else None,
    'radar_chart': XL_CHART_TYPE.RADAR if XL_CHART_TYPE is not None else None,
    'scatter_chart': XL_CHART_TYPE.XY_SCATTER if XL_CHART_TYPE is not None else None,
}


def _log_render_native(*args):
    try:
        with open('/tmp/render_native_debug.log', 'a', encoding='utf-8') as f:
            f.write(' '.join(str(a) for a in args) + '\n')
    except Exception:
        pass


def apply_chart_theme(chart, spec, accent_token='primary'):
    """Theme application: series colors + Y-axis number formatting."""
    try:
        palette = spec.get('md3_palette') or []
        if isinstance(palette, str):
            palette = [palette]
        if not palette:
            palette = ['#2563EB', '#10B981', '#F59E0B', '#A78BFA']
        for i, ser in enumerate(chart.series):
            color = palette[i % len(palette)]
            try:
                ser.format.fill.solid()
                ser.format.fill.fore_color.rgb = hex_to_rgb(color)
            except Exception:
                pass
    except Exception:
        pass

    # Y-axis number formatting for large values
    try:
        va = chart.value_axis
        # Detect max value across all series to choose format
        max_val = 0
        for ser in chart.series:
            try:
                for pt in ser.values:
                    if pt is not None and pt > max_val:
                        max_val = pt
            except Exception:
                pass
        if max_val >= 1_000_000:
            va.tick_labels.number_format = '#,##0,,"M"'
            va.tick_labels.number_format_is_linked = False
            # Set major_unit to avoid duplicate integer labels
            import math
            magnitude = 10 ** int(math.log10(max_val))
            va.major_unit = max(magnitude // 2, 1_000_000)
        elif max_val >= 10_000:
            va.tick_labels.number_format = '#,##0,"K"'
            va.tick_labels.number_format_is_linked = False
            import math
            magnitude = 10 ** int(math.log10(max_val))
            va.major_unit = max(magnitude // 5, 10_000)
    except Exception:
        pass


def render_native_chart(slide: Any, visual: Dict, spec: Dict, left: float, top: float, width: float, height: float, accent_token: str = 'primary') -> bool:
    """Attempt to render a chart using python-pptx native APIs. Returns True on success."""
    if not HAS_PPTX_CHARTS:
        return False
    pd = visual.get('placeholder_data', {})
    config = pd.get('chart_config', {})
    if not config:
        _log_render_native('no config', visual)
        return False

    # Handle composite charts where chart_config contains named sub-charts
    if not config.get('labels') and any(isinstance(v, dict) for v in config.values()):
        # Pick the first sub-chart as a best-effort render target
        first_key = next(iter(config.keys()))
        config = config[first_key]

    labels = config.get('labels') or config.get('x') or []
    series = config.get('series', [])
    if not series:
        _log_render_native('no series', visual)
        return False
    chart_type = visual.get('type', 'line_chart')

    # Map some composite/alias types to supported native types (best-effort)
    ALIAS_TYPE_MAP = {
        'bar_line_chart': 'column_chart',
        'composite_charts': 'line_chart',
    }
    if chart_type in ALIAS_TYPE_MAP:
        chart_type = ALIAS_TYPE_MAP[chart_type]

    if chart_type != 'scatter_chart' and not labels:
        _log_render_native('no labels', visual, config)
        return False

    xl_type = NATIVE_CHART_TYPE_MAP.get(chart_type)
    if not xl_type:
        _log_render_native('unsupported xl_type for chart_type', chart_type)
        return False
    try:
        if xl_type == XL_CHART_TYPE.XY_SCATTER:
            try:
                data = XyChartData()
                all_x = None
                for s in series:
                    xy_x = s.get('x', [])
                    xy_y = s.get('y', [])
                    if not xy_x or not xy_y:
                        return False
                    if all_x is None:
                        all_x = xy_x
                    if len(xy_x) != len(all_x):
                        return False
                    data.add_series(s.get('name', 'Series'), list(zip(xy_x, xy_y)))
            except Exception:
                data = CategoryChartData()
                first_x = series[0].get('x', [])
                data.categories = [str(x) for x in first_x]
                for s in series:
                    data.add_series(s.get('name', 'Series'), s.get('y', []))
                xl_type = XL_CHART_TYPE.LINE_MARKERS
        else:
            data = CategoryChartData()
            data.categories = [str(l) for l in labels]
            for s in series:
                data.add_series(s.get('name', 'Series'), s.get('data', []))

        shape = slide.shapes.add_chart(xl_type, Inches(left), Inches(top), Inches(width), Inches(height), data)
        chart = shape.chart
        # Hide legend when only one series (redundant with chart title)
        if len(series) <= 1:
            chart.has_legend = False
        apply_chart_theme(chart, spec, accent_token)
        return True
    except Exception:
        return False


def render_visual(slide: Any, visual: Dict, spec: Dict, _grid: GridSystem, left: float, top: float, width: float, height: float) -> None:
    if not visual or visual.get('type') in (None, 'none'):
        return
    pd = visual.get('placeholder_data', {})
    # Prefer native chart rendering when possible
    if pd.get('chart_config'):
        try:
            if HAS_PPTX_CHARTS:
                ok = render_native_chart(slide, visual, spec, left, top, width, height)
                if ok:
                    return
        except Exception:
            pass
        # Fallback: render as data table
        render_chart_table(slide, visual, spec, left, top, width, height)
    elif pd.get('mermaid_code'):
        render_mermaid_placeholder(slide, visual, spec, left, top, width, height)
    else:
        render_visual_placeholder(slide, visual, spec, left, top, width, height)


def render_chart_table(slide: Any, visual: Dict, spec: Dict, left: float, top: float, width: float, height: float) -> None:
    config = visual.get('placeholder_data', {}).get('chart_config', {})
    labels = config.get('labels', [])
    series = config.get('series', [])
    if not labels:
        return
    n_cols = len(labels)
    col_w = width / max(n_cols, 1)
    n_rows = len(series) + 1
    title_offset = 0.35 if visual.get('title') else 0
    usable_h = height - title_offset
    row_h = max(0.42, min(1.2, usable_h / max(n_rows, 1)))
    if visual.get('title'):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top - 0.30), Inches(width), Inches(0.28))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = visual['title']
        run.font.size = Pt(get_font_size(spec, 'table_header'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_surface')
    for j, label in enumerate(labels):
        tb = slide.shapes.add_textbox(Inches(left + j * col_w), Inches(top), Inches(col_w), Inches(row_h))
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = str(label)
        run.font.size = Pt(get_font_size(spec, 'table_header'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_surface')
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top + row_h - 0.02), Inches(width), Pt(2))
    sep.fill.solid()
    sep.fill.fore_color.rgb = get_color(spec, 'primary')
    sep.line.fill.background()
    for r, s in enumerate(series):
        ry = top + (r + 1) * row_h
        if r % 2 == 1:
            stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(ry), Inches(width), Inches(row_h))
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = get_color(spec, 'surface_variant')
            stripe.line.fill.background()
        data = s.get('data', [])
        for j, val in enumerate(data):
            if j >= n_cols:
                break
            tb = slide.shapes.add_textbox(Inches(left + j * col_w), Inches(ry), Inches(col_w), Inches(row_h))
            tf = tb.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(get_font_size(spec, 'table_cell'))
            run.font.color.rgb = get_color(spec, 'on_surface')
            if isinstance(val, (int, float)):
                p.alignment = PP_ALIGN.RIGHT


def render_mermaid_placeholder(slide: Any, visual: Dict, spec: Dict, left: float, top: float, width: float, height: float) -> None:
    mermaid = visual.get('placeholder_data', {}).get('mermaid_code', '')
    card = create_card_shape(slide, spec, left, top, width, min(height, 3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = get_color(spec, 'surface_variant')
    card.line.color.rgb = get_color(spec, 'outline')
    card.line.width = Pt(1)
    if visual.get('title'):
        tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.12), Inches(width - 0.4), Inches(0.28))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"📊 {visual['title']}"
        run.font.size = Pt(get_font_size(spec, 'table_header'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_surface')
        apply_font_to_run(run, spec)
    preview = '\n'.join(mermaid.strip().split('\n')[:8])
    if len(mermaid.strip().split('\n')) > 8:
        preview += '\n  ...'
    tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.50), Inches(width - 0.4), Inches(min(height - 0.7, 2.5)))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.font.size = Pt(11)
    run.font.color.rgb = get_color(spec, 'outline')
    apply_font_to_run(run, spec)
    run.text = preview


def render_visual_placeholder(slide: Any, visual: Dict, spec: Dict, left: float, top: float, width: float, height: float) -> None:
    card = create_card_shape(slide, spec, left, top, width, min(height, 2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = get_color(spec, 'surface_variant')
    card.line.color.rgb = get_color(spec, 'outline')
    card.line.width = Pt(1)
    label = visual.get('title', visual.get('type', 'Visual'))
    reqs = visual.get('content_requirements', [])
    text = f"[{label}]"
    if reqs:
        text += '\n' + '\n'.join(f"  • {r}" for r in reqs[:3])
    tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.3), Inches(width - 0.4), Inches(min(height - 0.6, 2.0)))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(get_font_size(spec, 'table_header'))
    run.font.color.rgb = get_color(spec, 'outline')
    apply_font_to_run(run, spec)


# --- Region renderers (Task 4.3) --------------------------------------------

def render_region_chart(pptx_slide: Any, data: Any, bounds: tuple, spec: Dict) -> None:
    """Render a chart region. `data` may be the visual dict or chart_config path-resolved data."""
    left, top, width, height = bounds
    # If data is a visual dict, use render_visual to honor native/chart fallback
    if isinstance(data, dict) and data.get('type'):
        try:
            render_visual(pptx_slide, data, spec, GridSystem(spec), left, top, width, height)
            return
        except Exception:
            pass
    # Otherwise, if data looks like chart_config
    if isinstance(data, dict) and data.get('labels'):
        visual = {'type': 'bar_chart', 'placeholder_data': {'chart_config': data}}
        render_visual(pptx_slide, visual, spec, GridSystem(spec), left, top, width, height)


def _render_comparison_standard(pptx_slide: Any, items: list, attrs_keys: list, bounds: tuple, spec: Dict) -> None:
    """Render comparison table in standard (items=columns) orientation."""
    left, top, width, height = bounds
    n = len(items)
    col_w = width / max(n, 1)
    n_rows = len(attrs_keys)

    header_h = 0.45
    remaining = height - header_h
    row_h = max(0.45, remaining / max(n_rows, 1)) if n_rows else 0.5
    table_h = header_h + n_rows * row_h

    for i, item in enumerate(items):
        hdr_bg = pptx_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left + i * col_w), Inches(top),
            Inches(col_w), Inches(header_h)
        )
        hdr_bg.fill.solid()
        hdr_bg.fill.fore_color.rgb = get_color(spec, 'primary')
        hdr_bg.line.fill.background()

        tb = pptx_slide.shapes.add_textbox(
            Inches(left + i * col_w + 0.08), Inches(top + 0.04),
            Inches(col_w - 0.16), Inches(header_h - 0.08)
        )
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = item.get('label', '')
        run.font.size = Pt(get_font_size(spec, 'table_header'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_primary')
        apply_font_to_run(run, spec)

    for r, key in enumerate(attrs_keys):
        row_top = top + header_h + r * row_h
        if r % 2 == 0:
            for i in range(n):
                stripe = pptx_slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(left + i * col_w), Inches(row_top),
                    Inches(col_w), Inches(row_h)
                )
                stripe.fill.solid()
                stripe.fill.fore_color.rgb = get_color(spec, 'surface_variant')
                stripe.line.fill.background()

        for i, item in enumerate(items):
            val = (item.get('attributes') or {}).get(key, '')
            display_val = str(val) if val else '\u2014'  # em-dash for empty
            tb = pptx_slide.shapes.add_textbox(
                Inches(left + i * col_w + 0.08), Inches(row_top + 0.06),
                Inches(col_w - 0.16), Inches(row_h - 0.12)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = display_val
            run.font.size = Pt(get_font_size(spec, 'table_cell'))
            run.font.color.rgb = get_color(spec, 'on_surface')
            apply_font_to_run(run, spec)

    outline_color = get_color(spec, 'outline')
    for i in range(1, n):
        line_x = left + i * col_w
        line = pptx_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(line_x - 0.005), Inches(top),
            Inches(0.01), Inches(table_h)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = outline_color
        line.line.fill.background()

    h_line = pptx_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top + header_h - 0.005),
        Inches(width), Inches(0.01)
    )
    h_line.fill.solid()
    h_line.fill.fore_color.rgb = outline_color
    h_line.line.fill.background()


def _render_comparison_transposed(pptx_slide: Any, items: list, attrs_keys: list, bounds: tuple, spec: Dict) -> None:
    """Render comparison table transposed: items as rows, attributes as columns.

    Used when there are too many items (>4) to fit as columns.
    Layout: first column = item labels, remaining columns = attribute values.
    """
    left, top, width, height = bounds
    n_items = len(items)
    # Cap attribute columns to 5 to avoid overflow
    attrs_keys = attrs_keys[:5]
    n_attrs = len(attrs_keys)
    # Columns: 1 (label) + n_attrs (attribute values)
    n_cols = 1 + n_attrs
    label_col_w = min(1.5, width * 0.18)
    attr_col_w = (width - label_col_w) / max(n_attrs, 1)

    header_h = 0.40
    remaining = height - header_h
    row_h = max(0.38, remaining / max(n_items, 1))
    table_h = header_h + n_items * row_h

    outline_color = get_color(spec, 'outline')

    # --- Header row: attribute key names ---
    # First header cell (empty label column header)
    hdr_bg = pptx_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top),
        Inches(label_col_w), Inches(header_h)
    )
    hdr_bg.fill.solid()
    hdr_bg.fill.fore_color.rgb = get_color(spec, 'primary')
    hdr_bg.line.fill.background()

    for c, key in enumerate(attrs_keys):
        hdr_bg = pptx_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left + label_col_w + c * attr_col_w), Inches(top),
            Inches(attr_col_w), Inches(header_h)
        )
        hdr_bg.fill.solid()
        hdr_bg.fill.fore_color.rgb = get_color(spec, 'primary')
        hdr_bg.line.fill.background()

        tb = pptx_slide.shapes.add_textbox(
            Inches(left + label_col_w + c * attr_col_w + 0.05), Inches(top + 0.04),
            Inches(attr_col_w - 0.10), Inches(header_h - 0.08)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = key
        run.font.size = Pt(get_font_size(spec, 'table_header'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_primary')
        apply_font_to_run(run, spec)

    # --- Data rows: one row per item ---
    for r, item in enumerate(items):
        row_top = top + header_h + r * row_h

        # Zebra stripe
        if r % 2 == 0:
            stripe = pptx_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left), Inches(row_top),
                Inches(width), Inches(row_h)
            )
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = get_color(spec, 'surface_variant')
            stripe.line.fill.background()

        # Label cell (left column, bold)
        tb = pptx_slide.shapes.add_textbox(
            Inches(left + 0.05), Inches(row_top + 0.04),
            Inches(label_col_w - 0.10), Inches(row_h - 0.08)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = item.get('label', '')
        run.font.size = Pt(get_font_size(spec, 'table_cell'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_surface')
        apply_font_to_run(run, spec)

        # Attribute value cells
        for c, key in enumerate(attrs_keys):
            val = (item.get('attributes') or {}).get(key, '')
            display_val = str(val) if val else '\u2014'
            tb = pptx_slide.shapes.add_textbox(
                Inches(left + label_col_w + c * attr_col_w + 0.05), Inches(row_top + 0.04),
                Inches(attr_col_w - 0.10), Inches(row_h - 0.08)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = display_val
            run.font.size = Pt(get_font_size(spec, 'table_cell'))
            run.font.color.rgb = get_color(spec, 'on_surface')
            apply_font_to_run(run, spec)

    # Vertical grid lines
    # Label column separator
    line = pptx_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left + label_col_w - 0.005), Inches(top),
        Inches(0.01), Inches(table_h)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = outline_color
    line.line.fill.background()
    for c in range(1, n_attrs):
        line_x = left + label_col_w + c * attr_col_w
        line = pptx_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(line_x - 0.005), Inches(top),
            Inches(0.01), Inches(table_h)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = outline_color
        line.line.fill.background()

    # Horizontal separator between header and data
    h_line = pptx_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top + header_h - 0.005),
        Inches(width), Inches(0.01)
    )
    h_line.fill.solid()
    h_line.fill.fore_color.rgb = outline_color
    h_line.line.fill.background()


def render_region_comparison(pptx_slide: Any, data: Any, bounds: tuple, spec: Dict) -> None:
    """Render a comparison table with header background, zebra rows, and grid lines.

    `data` expected: list of comparison items (label, attributes dict).
    Auto-transposes when items > 4 to avoid overly narrow columns.
    """
    left, top, width, height = bounds
    items = data or []
    if not isinstance(items, list) or not items:
        return
    n = len(items)

    # Collect all attribute keys across items (up to 6)
    # Filter: only keep keys that appear in >= 40% of items
    key_counts: dict = {}
    for it in items:
        for k in (it.get('attributes') or {}).keys():
            key_counts[k] = key_counts.get(k, 0) + 1
    threshold = max(1, n * 0.35)
    attrs_keys = [k for k in dict.fromkeys(
        k for it in items for k in (it.get('attributes') or {}).keys()
    ) if key_counts.get(k, 0) >= threshold][:6]

    # Transpose when too many columns OR when columns would be too narrow
    col_w = width / max(n, 1)
    if n > 4 or col_w < 1.5:
        _render_comparison_transposed(pptx_slide, items, attrs_keys, bounds, spec)
    else:
        _render_comparison_standard(pptx_slide, items, attrs_keys, bounds, spec)


def render_region_kpi(pptx_slide: Any, data: Any, bounds: tuple, spec: Dict) -> None:
    """Render KPIs horizontally inside bounds. `data` expected: list of kpi items."""
    left, top, width, height = bounds
    kpis = data or []
    if not isinstance(kpis, list) or not kpis:
        return
    # reuse render_kpis helper with adjusted width
    render_kpis(pptx_slide, kpis, spec, GridSystem(spec), left, top, width)


def render_region_callout(pptx_slide: Any, data: Any, bounds: tuple, spec: Dict) -> None:
    """Render stacked callouts (components.callouts)."""
    left, top, width, height = bounds
    callouts = data or []
    if not isinstance(callouts, list) or not callouts:
        return
    # Stack vertical callouts with label + accent bar
    gap = 0.12
    card_h = min(0.75, max(0.35, (height - (len(callouts)-1)*gap) / len(callouts)))
    accent_bar_w = 0.04  # left accent bar width
    for i, c in enumerate(callouts[:6]):
        cy = top + i * (card_h + gap)
        label = c.get('label', '')
        text = c.get('text', '')
        # Detect icon type from label for accent color
        is_warning = any(kw in label for kw in ('⚠', '风险', '限制', 'risk'))
        accent_token = 'error' if is_warning else 'primary'
        # Card background
        card = create_card_shape(pptx_slide, spec, left, cy, width, card_h)
        try:
            card.fill.solid()
            card.fill.fore_color.rgb = get_color(spec, 'primary_container')
            card.line.fill.background()
        except Exception:
            pass
        # Left accent bar
        try:
            bar = pptx_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left), Inches(cy), Inches(accent_bar_w), Inches(card_h)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = get_color(spec, accent_token)
            bar.line.fill.background()
        except Exception:
            pass
        text_left = left + accent_bar_w + 0.10
        text_w = width - accent_bar_w - 0.18
        if label:
            # Render label (bold, with icon)
            tb_lbl = pptx_slide.shapes.add_textbox(
                Inches(text_left), Inches(cy + 0.06),
                Inches(text_w), Inches(0.22)
            )
            tf_lbl = tb_lbl.text_frame
            p_lbl = tf_lbl.paragraphs[0]
            run_lbl = p_lbl.add_run()
            run_lbl.text = label
            run_lbl.font.size = Pt(get_font_size(spec, 'label'))
            run_lbl.font.bold = True
            run_lbl.font.color.rgb = get_color(spec, accent_token)
            apply_font_to_run(run_lbl, spec)
            text_top = cy + 0.28
            text_h = card_h - 0.34
        else:
            text_top = cy + 0.08
            text_h = card_h - 0.16
        # Render text body
        tb = pptx_slide.shapes.add_textbox(
            Inches(text_left), Inches(text_top),
            Inches(text_w), Inches(text_h)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(get_font_size(spec, 'callout_text'))
        run.font.color.rgb = get_color(spec, 'on_primary_container')
        apply_font_to_run(run, spec)


def render_region_progression(pptx_slide: Any, data: Any, bounds: tuple, spec: Dict) -> None:
    """Render a simple timeline/progression using timeline_items."""
    left, top, width, height = bounds
    items = data or []
    if not isinstance(items, list) or not items:
        return
    n = len(items)
    step_w = width / max(n, 1)
    y = top + height / 2
    for i, it in enumerate(items):
        cx = left + i * step_w + step_w / 2
        # draw a small circle using rectangle as approximation
        circ = pptx_slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 0.08), Inches(y - 0.08), Inches(0.16), Inches(0.16))
        circ.fill.solid()
        circ.fill.fore_color.rgb = get_color(spec, 'primary')
        circ.line.fill.background()
        tb = pptx_slide.shapes.add_textbox(Inches(cx - step_w/2 + 0.02), Inches(y + 0.12), Inches(step_w - 0.04), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = it.get('phase', '')
        run.font.size = Pt(get_font_size(spec, 'label'))
        run.font.color.rgb = get_color(spec, 'on_surface')


def _classify_bullet_icon(text: str) -> tuple:
    """Classify a bullet text and return (icon_shape, color_token).

    Returns shape type + color for rendering a small icon shape beside the bullet.
    """
    t = text.lower()
    # Risk / warning
    if any(kw in t for kw in ('风险', '挑战', '限制', '注意', '回退', '⚠', 'risk', 'warn')):
        return MSO_SHAPE.ISOSCELES_TRIANGLE, 'error'
    # Recommendation / action
    if any(kw in t for kw in ('建议', '建議', '推荐', '优先', '应', '要', '✓', 'recommend')):
        return MSO_SHAPE.OVAL, 'success'
    # Data / metric
    if re.search(r'\d+[%亿万MBK]', t):
        return MSO_SHAPE.RECTANGLE, 'secondary'
    # Default
    return MSO_SHAPE.OVAL, 'primary'


def render_region_bullets(pptx_slide: Any, data: Any, bounds: tuple, spec: Dict) -> None:
    """Render structured bullets with icon markers and adaptive spacing."""
    left, top, width, height = bounds
    bullets = data or []
    if not isinstance(bullets, list) or not bullets:
        return
    n = min(len(bullets), 8)
    row_h = min(0.85, max(0.4, height / max(n, 1)))
    icon_size = 0.14
    icon_pad = 0.08
    for i, b in enumerate(bullets[:n]):
        b_str = str(b)
        icon_shape, icon_token = _classify_bullet_icon(b_str)
        bullet_y = top + i * row_h
        icon_y = bullet_y + 0.04
        # Draw icon shape
        try:
            icon = pptx_slide.shapes.add_shape(
                icon_shape,
                Inches(left + 0.05), Inches(icon_y),
                Inches(icon_size), Inches(icon_size)
            )
            icon.fill.solid()
            icon.fill.fore_color.rgb = get_color(spec, icon_token)
            icon.line.fill.background()
        except Exception:
            pass
        # Text (shifted right for icon)
        text_left = left + 0.05 + icon_size + icon_pad
        tb = pptx_slide.shapes.add_textbox(
            Inches(text_left), Inches(bullet_y),
            Inches(width - 0.1 - icon_size - icon_pad), Inches(row_h - 0.05)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = b_str
        run.font.size = Pt(get_font_size(spec, 'bullet_text'))
        run.font.color.rgb = get_color(spec, 'on_surface')
        apply_font_to_run(run, spec)


def _add_arrowhead(connector: Any) -> None:
    """Add a triangle arrowhead to the end of a connector via XML manipulation.

    python-pptx doesn't expose a high-level arrowhead API, so we
    insert the ``<a:tailEnd type="triangle"/>`` element directly.
    """
    try:
        ln = connector._element.find(qn('a:ln'))
        if ln is None:
            from lxml import etree
            sp_pr = connector._element.find(qn('a:spPr'))
            if sp_pr is None:
                sp_pr = connector._element
            ln = sp_pr.find(qn('a:ln'))
        if ln is None:
            return
        # Remove existing tailEnd if present
        for te in ln.findall(qn('a:tailEnd')):
            ln.remove(te)
        from lxml import etree
        tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'triangle')
        tail.set('w', 'med')
        tail.set('len', 'med')
    except Exception:
        pass


def apply_shape_style(shape: Any, spec: Dict, style: str = 'primary') -> None:
    """Apply basic MD3-inspired style to a shape.

    style: 'primary'|'secondary'|'tertiary'|'outline'
    """
    try:
        if style == 'outline':
            shape.fill.solid()
            # transparent fill
            shape.fill.fore_color.rgb = get_color(spec, 'surface')
            shape.line.color.rgb = get_color(spec, 'outline')
            shape.line.width = Emu(int(12700))  # ~0.125pt
        else:
            # filled container
            fill_token = {
                'primary': 'primary_container',
                'secondary': 'secondary_container',
                'tertiary': 'primary_container'
            }.get(style, 'primary_container')
            shape.fill.solid()
            shape.fill.fore_color.rgb = get_color(spec, fill_token)
            shape.line.fill.background()
        # Text styling
        try:
            tf = shape.text_frame
            p = tf.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(get_font_size(spec, 'label'))
            # text color for filled vs outline
            if style == 'outline':
                run.font.color.rgb = get_color(spec, 'on_surface')
            else:
                run.font.color.rgb = get_color(spec, 'on_primary')
        except Exception:
            pass
    except Exception:
        pass

def render_region_architecture(pptx_slide: Any, data: Any, bounds: tuple, spec: Dict) -> None:
    """Render an architecture diagram (nodes + edges) within bounds.

    Nodes may include absolute inches values for x,y,w,h OR fractional (0..1) relative to bounds.
    If coordinates missing, nodes are auto-laid out horizontally.
    Edges reference node ids via 'from'/'to'.
    """
    if not data or not isinstance(data, dict):
        return
    nodes = data.get('nodes', []) or []
    edges = data.get('edges', []) or []
    left, top, width, height = bounds

    # early exit
    if not nodes:
        return

    # compute node rectangles
    placements = {}
    # auto layout parameters
    auto_n = len(nodes)
    auto_w = min(2.0, width / max(1, auto_n) - 0.1)
    auto_h = min(1.0, height * 0.5)
    for i, node in enumerate(nodes):
        nid = node.get('id', f'node{i}')
        nw = node.get('w')
        nh = node.get('h')
        nx = node.get('x')
        ny = node.get('y')
        # interpret fractional (0..1) as relative to bounds, else inches
        if nw is None:
            w_n = auto_w
        else:
            w_n = nw if nw > 1 else max(0.2, nw * width)
        if nh is None:
            h_n = auto_h
        else:
            h_n = nh if nh > 1 else max(0.2, nh * height)
        if nx is None:
            x_n = left + i * (width / max(1, auto_n)) + 0.05
        else:
            x_n = (left + nx * width) if (0 <= nx <= 1) else (left + nx)
        if ny is None:
            y_n = top + (height - h_n) / 2
        else:
            y_n = (top + ny * height) if (0 <= ny <= 1) else (top + ny)
        placements[nid] = (x_n, y_n, w_n, h_n, node.get('label', ''))

    # draw node shapes
    for nid, (x_n, y_n, w_n, h_n, lbl) in placements.items():
        shp = pptx_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_n), Inches(y_n), Inches(w_n), Inches(h_n)
        )
        shp.text_frame.clear()
        p = shp.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = lbl
        apply_shape_style(shp, spec, node_style := next((n.get('style') for n in nodes if n.get('id') == nid), 'primary'))

    # draw edges
    for edge in edges:
        f = edge.get('from')
        t = edge.get('to')
        if f not in placements or t not in placements:
            continue
        x1, y1, w1, h1, _ = placements[f]
        x2, y2, w2, h2, _ = placements[t]
        # use center points
        cx1 = x1 + w1 / 2
        cy1 = y1 + h1 / 2
        cx2 = x2 + w2 / 2
        cy2 = y2 + h2 / 2
        conn = pptx_slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            Inches(cx1), Inches(cy1), Inches(cx2), Inches(cy2)
        )
        try:
            conn.line.color.rgb = get_color(spec, 'outline')
            conn.line.width = Emu(int(19050))  # ~1.5pt
            _add_arrowhead(conn)
        except Exception:
            pass


def _parse_mermaid_to_flow(mermaid_code: str) -> tuple:
    """Parse simple Mermaid graph definitions into (steps, transitions).

    Handles basic patterns:
      graph LR; A-->B[Label]-->C; C-->D
      graph TD; A[Label] --> B
    Returns (steps_list, transitions_list).
    """
    import re
    steps_dict: dict = {}  # id -> label
    transitions = []

    # Strip graph directive
    code = re.sub(r'^\s*graph\s+\w+\s*;?\s*', '', mermaid_code.strip(), flags=re.IGNORECASE)
    # Split on semicolons for multiple chains
    chains = [c.strip() for c in code.split(';') if c.strip()]

    for chain in chains:
        # Split chain on --> or ---> (arrow)
        parts = re.split(r'\s*-+>\s*', chain)
        prev_id = None
        for part in parts:
            part = part.strip()
            if not part:
                continue
            # Parse "NodeId[Label text]" or just "NodeId"
            m = re.match(r'^([A-Za-z0-9_]+)\s*(?:\[([^\]]*)\])?$', part)
            if m:
                node_id = m.group(1)
                label = m.group(2) if m.group(2) else node_id
                if node_id not in steps_dict:
                    steps_dict[node_id] = label
                if prev_id and prev_id != node_id:
                    transitions.append({'from': prev_id, 'to': node_id})
                prev_id = node_id
            else:
                # Fallback: use the raw text as both id and label
                node_id = re.sub(r'[^A-Za-z0-9_]', '_', part)[:20]
                if node_id and node_id not in steps_dict:
                    steps_dict[node_id] = part
                if prev_id and node_id:
                    transitions.append({'from': prev_id, 'to': node_id})
                prev_id = node_id

    steps = [{'id': nid, 'label': lbl, 'type': 'process'} for nid, lbl in steps_dict.items()]
    if steps:
        steps[0]['type'] = 'start'
        steps[-1]['type'] = 'end'
    return steps, transitions


def render_region_flow(pptx_slide: Any, data: Any, bounds: tuple, spec: Dict) -> None:
    """Render a flow/process diagram within bounds.

    Steps: list of {id,label,type,x?,y?,w?,h?,style?}
    Transitions: list of {from,to,label,condition,style}
    Supports auto horizontal layout when coords missing.
    Also handles mermaid_code by parsing it into steps/transitions.
    """
    if not data or not isinstance(data, dict):
        return
    steps = data.get('steps', []) or []
    transitions = data.get('transitions', []) or []
    left, top, width, height = bounds

    # Fallback: parse mermaid_code if no structured steps provided
    if not steps and data.get('mermaid_code'):
        steps, transitions = _parse_mermaid_to_flow(data['mermaid_code'])

    if not steps:
        return

    # map step type -> AutoShape
    shape_map = {
        'start': MSO_AUTO_SHAPE_TYPE.OVAL,
        'process': MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        'decision': MSO_AUTO_SHAPE_TYPE.DIAMOND,
        'end': MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    }

    placements = {}
    n = len(steps)
    # Fix-1: wider gap between nodes (0.35" min) to avoid overlap
    gap = 0.35
    step_w = min(1.8, (width - gap * max(1, n - 1)) / max(1, n))
    # Fix-2: shorter nodes for better proportions
    step_h = min(0.65, height * 0.25)
    slot_w = width / max(1, n)
    for i, s in enumerate(steps):
        sid = s.get('id', f's{i}')
        sw = s.get('w')
        sh = s.get('h')
        sx = s.get('x')
        sy = s.get('y')
        if sw is None:
            w_s = step_w
        else:
            w_s = sw if sw > 1 else max(0.3, sw * width)
        if sh is None:
            h_s = step_h
        else:
            h_s = sh if sh > 1 else max(0.3, sh * height)
        if sx is None:
            x_s = left + i * slot_w + (slot_w - w_s) / 2
        else:
            x_s = (left + sx * width) if (0 <= sx <= 1) else (left + sx)
        if sy is None:
            # Position nodes in upper-third of area for visual balance
            y_s = top + height * 0.35 - h_s / 2
        else:
            y_s = (top + sy * height) if (0 <= sy <= 1) else (top + sy)
        placements[sid] = (x_s, y_s, w_s, h_s, s.get('label', ''), s.get('type', 'process'))

    # draw steps
    for sid, (x_s, y_s, w_s, h_s, lbl, stype) in placements.items():
        sh_type = shape_map.get(stype, MSO_AUTO_SHAPE_TYPE.RECTANGLE)
        shp = pptx_slide.shapes.add_shape(
            sh_type, Inches(x_s), Inches(y_s), Inches(w_s), Inches(h_s)
        )
        # Fix-3: enable word_wrap + auto-fit so long labels are not truncated
        tf = shp.text_frame
        tf.clear()
        tf.word_wrap = True
        try:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception:
            pass
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = lbl
        run.font.size = Pt(get_font_size(spec, 'table_cell'))
        apply_shape_style(shp, spec, s.get('style', 'primary') if (s := next((it for it in steps if it.get('id') == sid), None)) else 'primary')

    # draw transitions
    for tr in transitions:
        f = tr.get('from')
        t = tr.get('to')
        if f not in placements or t not in placements:
            continue
        x1, y1, w1, h1, _lbl1, _ = placements[f]
        x2, y2, w2, h2, _lbl2, _ = placements[t]
        # Fix-4: connect from right edge of source to left edge of target
        # so the line is visible between nodes (not hidden behind them)
        ex1 = x1 + w1       # source right edge
        ey1 = y1 + h1 / 2   # source vertical center
        ex2 = x2             # target left edge
        ey2 = y2 + h2 / 2   # target vertical center
        conn = pptx_slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            Inches(ex1), Inches(ey1), Inches(ex2), Inches(ey2)
        )
        try:
            conn.line.color.rgb = get_color(spec, 'on_surface')
            conn.line.width = Emu(int(19050))  # 1.5pt
            # Add arrowhead at end
            _add_arrowhead(conn)
        except Exception:
            pass
        # label at midpoint
        lbl = tr.get('label') or tr.get('condition')
        if lbl:
            mx = (ex1 + ex2) / 2
            my = (ey1 + ey2) / 2
            tb = pptx_slide.shapes.add_textbox(Inches(mx - 0.6), Inches(my - 0.15), Inches(1.2), Inches(0.3))
            p = tb.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = lbl
            run.font.size = Pt(get_font_size(spec, 'label'))
            run.font.color.rgb = get_color(spec, 'on_surface')


def _strip_group_label_suffix(items: list) -> None:
    """Remove the common trailing token from all labels in a group (in-place).

    When comparison items are split by attribute groups, labels like
    "AWS 架构", "Dell 架构", "阿里云 架构" share the suffix "架构".
    Stripping it yields cleaner column headers: "AWS", "Dell", "阿里云".
    Only strips if *all* labels share the same last token and each label
    has at least 2 tokens (so we don't empty anything).
    """
    labels = [it.get('label', '') for it in items]
    if len(labels) < 2:
        return
    # Split each label by whitespace
    tokens_list = [l.split() for l in labels]
    if any(len(t) < 2 for t in tokens_list):
        return
    # Check if all labels share the same last token
    common_suffix = tokens_list[0][-1]
    if all(t[-1] == common_suffix for t in tokens_list):
        for it, tokens in zip(items, tokens_list):
            it['label'] = ' '.join(tokens[:-1])


def render_region_comparison_split(pptx_slide: Any, data: Any, bounds: tuple, spec: Dict) -> None:
    """Render two comparison tables stacked vertically within bounds.

    `data` expected: {'groups': [group_a_items, group_b_items]}
    Each group gets its own independent attribute keys and is rendered with
    half the available height (minus a small gap).
    """
    if not data or not isinstance(data, dict):
        return
    groups = data.get('groups', [])
    if not groups:
        return

    left, top, width, height = bounds
    gap = 0.15
    n_groups = len(groups)
    group_h = (height - gap * (n_groups - 1)) / max(n_groups, 1)

    if n_groups == 2:
        # Render two groups side-by-side for direct comparison
        left_w = (width - 0.04) / 2
        right_left = left + left_w + 0.04
        # strip suffixes within each group for concise headers
        _strip_group_label_suffix(groups[0])
        _strip_group_label_suffix(groups[1])
        # left group
        g0_bounds = (left, top, left_w, height)
        render_region_comparison(pptx_slide, groups[0], g0_bounds, spec)
        # right group
        g1_bounds = (right_left, top, left_w, height)
        render_region_comparison(pptx_slide, groups[1], g1_bounds, spec)
        # vertical separator
        try:
            sep = pptx_slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(left + left_w - 0.01), Inches(top), Inches(0.02), Inches(height)
            )
            sep.fill.solid()
            sep.fill.fore_color.rgb = get_color(spec, 'on_surface_variant')
            sep.line.fill.background()
        except Exception:
            pass
    else:
        for i, group_items in enumerate(groups):
            if not group_items:
                continue
            # Strip common label suffix within each group so headers are concise.
            # e.g. ["AWS 架构", "Dell 架构"] → ["AWS", "Dell"]
            _strip_group_label_suffix(group_items)
            g_top = top + i * (group_h + gap)
            g_bounds = (left, g_top, width, group_h)
            render_region_comparison(pptx_slide, group_items, g_bounds, spec)


# Register region renderers
REGION_RENDERERS = {
    'chart': render_region_chart,
    'comparison_table': render_region_comparison,
    'comparison_table_split': render_region_comparison_split,
    'kpi_row': render_region_kpi,
    'callout_stack': render_region_callout,
    'progression': render_region_progression,
    'bullet_list': render_region_bullets,
    'architecture': render_region_architecture,
    'flow': render_region_flow,
}


# Dispatcher of slide types (partial set)

RENDERERS = {
    'title': None,  # implemented below
    'section_divider': None,
    'bullet-list': None,
    'two-column': None,
    'comparison': None,
    'decision': None,
    'data-heavy': None,
}


def _resolve_bg_image(spec: Dict, slide_type: str, sd: Dict) -> str:
    img = sd.get('background_image', '')
    if img:
        return img
    stl = spec.get('slide_type_layouts', {})
    entry = stl.get(slide_type, stl.get('default', {}))
    return entry.get('background_image', '')


def find_section_for_slide(slide_id: int, sections: List[Dict]) -> Dict:
    result = {}
    for sec in sorted(sections, key=lambda s: s.get('start_slide', 0)):
        if slide_id >= sec.get('start_slide', 999):
            result = sec
    return result


def _render_background(slide: Any, sd: Dict, spec: Dict, grid: GridSystem, stype: str) -> None:
    """Render the slide background image or solid fill."""
    bg_token = spec.get('slide_type_layouts', {}).get(stype, spec.get('slide_type_layouts', {}).get('default', {})).get('background', 'surface')
    bg_image_path = _resolve_bg_image(spec, stype, sd)
    if bg_image_path and os.path.isfile(bg_image_path):
        pic = slide.shapes.add_picture(bg_image_path, Emu(0), Emu(0), Inches(grid.slide_w), Inches(grid.slide_h))
        sp_tree = slide.shapes._spTree
        sp_tree.remove(pic._element)
        sp_tree.insert(2, pic._element)
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Inches(grid.slide_w), Inches(grid.slide_h))
        fill = overlay.fill
        fill.solid()
        fill.fore_color.rgb = get_color(spec, bg_token)
        from pptx.oxml.ns import qn as _qn
        solid_fill_elem = fill._fill.find(_qn('a:solidFill'))
        if solid_fill_elem is not None:
            srgb = solid_fill_elem.find(_qn('a:srgbClr'))
            if srgb is None:
                srgb = solid_fill_elem[0] if len(solid_fill_elem) else None
            if srgb is not None:
                try:
                    from lxml import etree
                    alpha = etree.SubElement(srgb, _qn('a:alpha'))
                    alpha.set('val', '40000')
                except Exception:
                    logging.getLogger(__name__).warning(
                        'lxml not available or failed to set overlay alpha; skipping alpha modification'
                    )
        overlay.line.fill.background()
        sp_tree = slide.shapes._spTree
        sp_tree.remove(overlay._element)
        sp_tree.insert(3, overlay._element)
    else:
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = get_color(spec, bg_token)


def _render_bullets_fallback(slide: Any, sd: Dict, spec: Dict, grid: GridSystem, bar_h: float) -> None:
    """Render simple bullet-list fallback content."""
    top = bar_h + 0.12
    content_bullets = sd.get('content', [])
    bullet_y = top + 0.15
    for i, bullet in enumerate(content_bullets[:8]):
        tb = slide.shapes.add_textbox(Inches(grid.margin_h + 0.1), Inches(bullet_y + i * 0.48), Inches(grid.usable_w - 0.2), Inches(0.45))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.size = Pt(get_font_size(spec, 'bullet_text'))
        run.font.color.rgb = get_color(spec, 'on_surface')
        p.line_spacing = 1.5


def render_section_divider_slide(slide: Any, sd: Dict, spec: Dict, grid: GridSystem) -> None:
    """Render section_divider content: title + content text + optional callouts.

    Section divider slides have title_bar='none', so we render the section title
    prominently in the center, plus any content[] text and callouts.
    """
    title_text = sd.get('title', '')
    content = sd.get('content', [])
    callouts = sd.get('components', {}).get('callouts', [])

    # Section title (large, centered)
    tb_title = slide.shapes.add_textbox(
        Inches(grid.margin_h + 0.5), Inches(grid.slide_h * 0.25),
        Inches(grid.usable_w - 1.0), Inches(1.0)
    )
    tf = tb_title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(get_font_size(spec, 'headline_large'))
    run.font.bold = True
    run.font.color.rgb = get_color(spec, 'on_primary')
    apply_font_to_run(run, spec)

    # Content text (subtitle under title)
    y_offset = grid.slide_h * 0.42
    for i, text in enumerate(content[:3]):
        tb = slide.shapes.add_textbox(
            Inches(grid.margin_h + 0.8), Inches(y_offset + i * 0.5),
            Inches(grid.usable_w - 1.6), Inches(0.45)
        )
        tf2 = tb.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = text
        run2.font.size = Pt(get_font_size(spec, 'slide_subtitle'))
        run2.font.color.rgb = get_color(spec, 'on_primary')
        apply_font_to_run(run2, spec)
        y_offset += 0.05

    # Callouts (compact, below content)
    if callouts:
        callout_y = y_offset + len(content) * 0.5 + 0.3
        for i, c in enumerate(callouts[:3]):
            tb = slide.shapes.add_textbox(
                Inches(grid.margin_h + 0.8), Inches(callout_y + i * 0.55),
                Inches(grid.usable_w - 1.6), Inches(0.45)
            )
            tf3 = tb.text_frame
            tf3.word_wrap = True
            p3 = tf3.paragraphs[0]
            run3 = p3.add_run()
            label = c.get('label', '')
            text = c.get('text', '')
            run3.text = f"{label}: {text}" if label else text
            run3.font.size = Pt(get_font_size(spec, 'callout_text'))
            run3.font.color.rgb = get_color(spec, 'on_primary')
            apply_font_to_run(run3, spec)


def _render_components_fallback(slide: Any, sd: Dict, spec: Dict, grid: GridSystem, bar_h: float) -> float:
    """Render components (kpis, bullets, comparison_items, callouts) for v1 slides.

    Returns the y-offset after rendering all components (for subsequent visual placement).
    """
    comps = sd.get('components', {})
    if not comps:
        return bar_h

    lz = get_layout_zones(spec)
    left = grid.margin_h
    current_y = bar_h + lz.get('content_margin_top', 0.12) + 0.05
    width = grid.usable_w
    bottom_bar_h = lz.get('bottom_bar_h', 0.35)
    available_h = grid.slide_h - current_y - bottom_bar_h - 0.4

    # 1. KPIs (render as a row at the top of content area)
    kpis = comps.get('kpis', [])
    if kpis:
        try:
            kpi_h = render_kpis(slide, kpis, spec, grid, left, current_y, width)
            current_y += kpi_h + 0.05
        except Exception:
            pass

    # Count total content items to allocate vertical space proportionally
    comp_bullets = comps.get('bullets', [])
    comp_items = comps.get('comparison_items', [])
    callouts = comps.get('callouts', [])
    risks = comps.get('risks', [])

    # 2. Bullets from components (different from content[])
    if comp_bullets:
        try:
            n_bullets = min(len(comp_bullets), 8)
            bullet_h = min(n_bullets * 0.85, available_h * 0.5)
            bounds = (left + 0.05, current_y, width - 0.1, bullet_h)
            render_region_bullets(slide, comp_bullets, bounds, spec)
            current_y += bullet_h + 0.15
        except Exception:
            pass

    # 3. Comparison items — give them most of available space
    if comp_items:
        try:
            remaining_h = max(2.0, grid.slide_h - current_y - bottom_bar_h - 0.5)
            bounds = (left, current_y, width, remaining_h)
            render_region_comparison(slide, comp_items, bounds, spec)
            # estimate height used
            n_attrs = 0
            for it in comp_items:
                n_attrs = max(n_attrs, len(it.get('attributes', {})))
            header_h = 0.45
            data_row_h = max(0.45, remaining_h - header_h) / max(n_attrs, 1) if n_attrs else 0.5
            current_y += header_h + n_attrs * data_row_h + 0.15
        except Exception:
            pass

    # 4. Callouts
    if callouts:
        try:
            remaining_h = max(1.0, grid.slide_h - current_y - bottom_bar_h - 0.4)
            bounds = (left, current_y, width, remaining_h)
            render_region_callout(slide, callouts, bounds, spec)
            current_y += len(callouts[:6]) * 0.7 + 0.15
        except Exception:
            pass

    # 5. Risks — rendered as warning cards with amber background
    if risks:
        try:
            remaining_h = max(1.0, grid.slide_h - current_y - bottom_bar_h - 0.4)
            n_risks = min(len(risks), 6)
            card_gap = 0.1
            card_h = min(0.65, max(0.4, (remaining_h - (n_risks - 1) * card_gap) / n_risks))
            for i, r in enumerate(risks[:n_risks]):
                ry = current_y + i * (card_h + card_gap)
                desc = r if isinstance(r, str) else r.get('description', r.get('risk', str(r)))
                label = '' if isinstance(r, str) else r.get('label', '')
                # Warning card background
                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(left + 0.05), Inches(ry),
                    Inches(width - 0.1), Inches(card_h)
                )
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xCD)  # light amber
                card.line.color.rgb = RGBColor(0xF5, 0x9E, 0x0B)  # amber border
                card.line.width = Pt(1)
                # Risk text
                tb = slide.shapes.add_textbox(
                    Inches(left + 0.18), Inches(ry + 0.06),
                    Inches(width - 0.36), Inches(card_h - 0.12)
                )
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                prefix = f"⚠ {label}: " if label else "⚠ "
                run.text = prefix + desc
                run.font.size = Pt(get_font_size(spec, 'bullet_text'))
                run.font.color.rgb = RGBColor(0x78, 0x35, 0x0F)  # dark amber text
                apply_font_to_run(run, spec)
            current_y += n_risks * (card_h + card_gap) + 0.1
        except Exception:
            pass

    return current_y


def _render_footer_and_notes(slide: Any, sd: Dict, spec: Dict, grid: GridSystem, stype: str, sec_title: str, sec_accent_token: str, sec_index: int, sections_len: int, slide_num: int, total_slides: int) -> None:
    """Render bottom bar and speaker notes."""
    if stype not in ('title', 'section_divider'):
        render_bottom_bar(
            slide,
            spec,
            grid,
            sec_title,
            slide_num,
            total_slides,
            section_index=sec_index,
            total_sections=sections_len,
            accent_token=sec_accent_token,
        )
    render_speaker_notes(slide, sd.get('speaker_notes'))


def render_slide_v2(prs: Presentation, sd: Dict, spec: Dict, grid: GridSystem, sections: List[Dict], slide_num: int, total_slides: int) -> None:
    """Render a v2 slide using layout_intent.regions[] and REGION_RENDERERS."""
    stype = sd.get('slide_type', 'bullet-list')
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Section context
    slide_id = sd.get('slide_id', slide_num)
    section = find_section_for_slide(slide_id, sections)
    sec_id = section.get('id', 'A')
    sec_title = section.get('title', '')
    sec_accent_token = spec.get('section_accents', {}).get(sec_id, 'primary')
    sec_index = next((i for i, s in enumerate(sections) if s.get('id') == sec_id), 0)

    # Background
    _render_background(slide, sd, spec, grid, stype)

    # Title bar: assertion or label title
    tb_mode = get_title_bar_mode(spec, stype)
    bar_h = 0.0

    # Special-case: Title slides should always render title content even when title_bar='none'
    if stype == 'title':
        try:
            if tb_mode != 'none':
                section_label = f"Section {sec_id} · {sec_title}" if sec_title else ''
                if sd.get('assertion'):
                    bar_h = render_assertion_title(slide, spec, grid, sd.get('assertion'), sd.get('title', ''), slide_num, total_slides, section_label=section_label, accent_color_token=sec_accent_token, mode=tb_mode)
                else:
                    bar_h = render_title_bar(slide, spec, grid, sd.get('title', ''), slide_num, total_slides, section_label=section_label, accent_color_token=sec_accent_token, mode=tb_mode)
            render_title_slide(slide, sd, spec, grid, slide_num, total_slides)
        except Exception:
            pass
        _render_footer_and_notes(slide, sd, spec, grid, stype, sec_title, sec_accent_token, sec_index, len(sections), slide_num, total_slides)
        return

    # Special-case: Section divider slides
    if stype == 'section_divider':
        try:
            render_section_divider_slide(slide, sd, spec, grid)
        except Exception:
            pass
        _render_footer_and_notes(slide, sd, spec, grid, stype, sec_title, sec_accent_token, sec_index, len(sections), slide_num, total_slides)
        return

    if tb_mode != 'none':
        # Use EA-provided section label if available (from absorbed section dividers)
        ea_section_label = sd.get('_section_label', '')
        if ea_section_label:
            section_label = ea_section_label
        else:
            section_label = f"Section {sec_id} · {sec_title}" if sec_title else ''
        if sd.get('assertion'):
            bar_h = render_assertion_title(slide, spec, grid, sd.get('assertion'), sd.get('title', ''), slide_num, total_slides, section_label=section_label, accent_color_token=sec_accent_token, mode=tb_mode) or 0.0
        else:
            bar_h = render_title_bar(slide, spec, grid, sd.get('title', ''), slide_num, total_slides, section_label=section_label, accent_color_token=sec_accent_token, mode=tb_mode) or 0.0

    # Regions rendering
    # Reserve extra bottom space when insight bar is present
    bottom_reserve = 0.95 if sd.get('insight') else 0.5
    try:
        regions = sd.get('layout_intent', {}).get('regions', []) or []
        for region in regions:
            renderer_key = region.get('renderer')
            renderer = REGION_RENDERERS.get(renderer_key)
            if not renderer:
                # unsupported region renderer; skip
                continue
            position = region.get('position', 'full')
            bounds = compute_region_bounds(position, grid, bar_h, bottom_reserve)
            data = resolve_data_source(sd, region.get('data_source', ''))
            try:
                renderer(slide, data, bounds, spec)
            except Exception:
                # best-effort: continue to next region
                continue
    except Exception:
        pass

    # Insight bar
    if sd.get('insight'):
        try:
            render_insight_bar(slide, sd.get('insight'), spec, grid, sec_accent_token)
        except Exception:
            pass

    # Footer & notes
    _render_footer_and_notes(slide, sd, spec, grid, stype, sec_title, sec_accent_token, sec_index, len(sections), slide_num, total_slides)


def render_slide(prs: Presentation, sd: Dict, spec: Dict, grid: GridSystem, sections: List[Dict], slide_num: int, total_slides: int) -> None:
    """Render a single slide with full styling (delegates to helpers)."""
    # Detect schema version: v2 uses layout_intent and custom region renderers
    try:
        if detect_schema_version(sd) == 2:
            return render_slide_v2(prs, sd, spec, grid, sections, slide_num, total_slides)
    except Exception:
        pass

    stype = sd.get('slide_type', 'bullet-list')
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Section context
    slide_id = sd.get('slide_id', slide_num)
    section = find_section_for_slide(slide_id, sections)
    sec_id = section.get('id', 'A')
    sec_title = section.get('title', '')
    sec_accent_token = spec.get('section_accents', {}).get(sec_id, 'primary')
    sec_index = next((i for i, s in enumerate(sections) if s.get('id') == sec_id), 0)

    # Background
    _render_background(slide, sd, spec, grid, stype)

    # Title bar
    tb_mode = get_title_bar_mode(spec, stype)
    bar_h = 0.0

    # Special-case: Title slides
    if stype == 'title':
        try:
            if tb_mode != 'none':
                section_label = f"Section {sec_id} · {sec_title}" if sec_title else ''
                if sd.get('assertion'):
                    bar_h = render_assertion_title(slide, spec, grid, sd.get('assertion'), sd.get('title', ''), slide_num, total_slides, section_label=section_label, accent_color_token=sec_accent_token, mode=tb_mode)
                else:
                    bar_h = render_title_bar(slide, spec, grid, sd.get('title', ''), slide_num, total_slides, section_label=section_label, accent_color_token=sec_accent_token, mode=tb_mode)
            render_title_slide(slide, sd, spec, grid, slide_num, total_slides)
        except Exception:
            pass
        _render_footer_and_notes(slide, sd, spec, grid, stype, sec_title, sec_accent_token, sec_index, len(sections), slide_num, total_slides)
        return

    # Special-case: Section divider slides
    if stype == 'section_divider':
        try:
            render_section_divider_slide(slide, sd, spec, grid)
        except Exception:
            pass
        _render_footer_and_notes(slide, sd, spec, grid, stype, sec_title, sec_accent_token, sec_index, len(sections), slide_num, total_slides)
        return

    if tb_mode != 'none':
        section_label = f"Section {sec_id} · {sec_title}" if sec_title else ''
        if sd.get('assertion'):
            bar_h = render_assertion_title(slide, spec, grid, sd.get('assertion'), sd.get('title', ''), slide_num, total_slides, section_label=section_label, accent_color_token=sec_accent_token, mode=tb_mode) or 0.0
        else:
            bar_h = render_title_bar(slide, spec, grid, sd.get('title', ''), slide_num, total_slides, section_label=section_label, accent_color_token=sec_accent_token, mode=tb_mode) or 0.0

    # Content[] text fallback — only render if no structured components exist
    # (components.bullets supersedes content[]; avoid double-render)
    comps = sd.get('components', {})
    has_components = any(comps.get(k) for k in ('kpis', 'bullets', 'comparison_items', 'callouts', 'risks'))
    if not has_components:
        _render_bullets_fallback(slide, sd, spec, grid, bar_h)

    # Render components (kpis, bullets, comparison_items, callouts, risks)
    content_bottom = _render_components_fallback(slide, sd, spec, grid, bar_h)

    # Try to render the main visual if present (useful for data-heavy slides)
    try:
        vis = sd.get('visual')
        if isinstance(vis, dict) and vis.get('type') not in (None, 'none'):
            lz = get_layout_zones(spec)
            left = grid.margin_h
            # Place visual below components if any were rendered
            top = max(content_bottom, bar_h + lz['content_margin_top'])
            width = grid.usable_w
            bottom_bar_h = lz['bottom_bar_h']
            height = max(1.0, grid.slide_h - top - bottom_bar_h - 0.3)
            render_visual(slide, vis, spec, grid, left, top, width, height)
    except Exception:
        # best-effort: don't fail slide rendering if visual rendering errors
        pass

    # Insight bar (insert above bottom bar if present)
    if sd.get('insight'):
        try:
            render_insight_bar(slide, sd.get('insight'), spec, grid, sec_accent_token)
        except Exception:
            pass

    # Footer & notes
    _render_footer_and_notes(slide, sd, spec, grid, stype, sec_title, sec_accent_token, sec_index, len(sections), slide_num, total_slides)
