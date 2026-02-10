"""Native Gantt chart renderer using python-pptx shapes."""

from typing import Any, Dict, List
from datetime import datetime
from dateutil.relativedelta import relativedelta
import logging

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from ...protocols.visual_data_protocol import VisualDataProtocol, GanttData, GanttTask
from ...protocols.renderer_interface import IRenderer
from ..base import BaseRenderer
from ...helpers import get_color, get_font_size, apply_font_to_run


logger = logging.getLogger(__name__)


class NativeGanttRenderer(BaseRenderer):
    """Renders Gantt charts using python-pptx native shapes (rectangles, text).
    
    Features:
    - Horizontal timeline with month labels
    - Task bars with status-based coloring
    - Task names and durations
    - Supports 3-12 tasks comfortably
    
    Quality: Medium (70/100) - Simple but fast and zero-dependency
    """
    
    @property
    def name(self) -> str:
        return "native-gantt"
    
    @property
    def supported_types(self) -> List[str]:
        return ["gantt", "gantt_chart"]
    
    def is_available(self) -> bool:
        """Always available (native python-pptx)."""
        return True
    
    def estimate_quality(self, visual_data: VisualDataProtocol) -> int:
        """Estimate rendering quality based on data complexity.
        
        Returns:
            int: 70 for simple gantt (≤10 tasks), 40 for complex (>10 tasks)
        """
        try:
            if isinstance(visual_data.data, GanttData):
                task_count = len(visual_data.data.tasks)
                if task_count <= 10:
                    return 70  # Good for simple gantt
                elif task_count <= 15:
                    return 50  # Acceptable but crowded
                else:
                    return 40  # Recommend external renderer
            return 60  # Default
        except Exception:
            return 50
    
    def _do_render(
        self,
        slide: Any,
        visual_data: VisualDataProtocol,
        spec: Dict[str, Any],
        left: float,
        top: float,
        width: float,
        height: float
    ) -> bool:
        """Render Gantt chart."""
        if not self._validate_data(visual_data):
            return False
        
        # Extract gantt data
        gantt_data = visual_data.data
        if not isinstance(gantt_data, GanttData):
            # Try to extract from placeholder_data (backward compatibility)
            pd = visual_data.placeholder_data or {}
            gantt_dict = pd.get('gantt_data', {})
            if not gantt_dict:
                logger.error(f"{self.name}: No valid GanttData found")
                return False
            try:
                gantt_data = GanttData(**gantt_dict)
            except Exception as e:
                logger.error(f"{self.name}: Failed to parse gantt_data: {e}")
                return False
        
        # Render components
        try:
            # Title (if present)
            title_height = 0.0
            if visual_data.title:
                title_height = self._render_title(
                    slide, visual_data.title, spec, left, top, width
                )
                top += title_height + 0.1
                height -= (title_height + 0.1)
            
            # Layout: task names on left (25%), timeline + bars on right (75%)
            name_ratio = 0.25
            name_width = width * name_ratio
            bar_area_left = left + name_width
            bar_area_width = width - name_width

            # Timeline axis — only over the bar area, not the task name column
            timeline_height = 0.5
            self._render_timeline(
                slide, gantt_data, spec, bar_area_left, top, bar_area_width, timeline_height
            )
            
            # Task bars (with row backgrounds)
            tasks_top = top + timeline_height + 0.15
            tasks_height = height - timeline_height - 0.15
            self._render_tasks(
                slide, gantt_data, spec, left, tasks_top, width, tasks_height
            )
            
            return True
            
        except Exception as e:
            logger.error(f"{self.name}: Rendering failed: {e}", exc_info=True)
            return False
    
    def _render_title(
        self,
        slide: Any,
        title: str,
        spec: Dict[str, Any],
        left: float,
        top: float,
        width: float
    ) -> float:
        """Render chart title.
        
        Returns:
            float: Height consumed by title
        """
        title_h = 0.35
        tb = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(title_h)
        )
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        
        run = p.add_run()
        run.text = f"📊 {title}"
        run.font.size = Pt(get_font_size(spec, 'table_header'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_surface')
        apply_font_to_run(run, spec)
        
        return title_h
    
    def _render_timeline(
        self,
        slide: Any,
        gantt_data: GanttData,
        spec: Dict[str, Any],
        left: float,
        top: float,
        width: float,
        height: float
    ) -> None:
        """Render horizontal timeline with month labels."""
        timeline = gantt_data.timeline
        
        # Parse start/end dates
        start_date = self._parse_date(timeline.start)
        end_date = self._parse_date(timeline.end)
        
        # Calculate months span
        months_span = (end_date.year - start_date.year) * 12 + (end_date.month - start_date.month) + 1
        
        # Timeline background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = get_color(spec, 'surface_variant')
        bg.line.color.rgb = get_color(spec, 'outline')
        bg.line.width = Pt(1)
        
        # Month labels
        month_width = width / months_span
        current_date = start_date
        
        for i in range(months_span):
            x = left + i * month_width
            
            # Month label
            tb = slide.shapes.add_textbox(
                Inches(x + 0.05), Inches(top + 0.08),
                Inches(month_width - 0.1), Inches(0.25)
            )
            tf = tb.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            
            run = p.add_run()
            run.text = current_date.strftime('%Y-%m')
            run.font.size = Pt(9)
            run.font.color.rgb = get_color(spec, 'on_surface_variant')
            apply_font_to_run(run, spec)
            
            # Vertical separator
            if i > 0:
                sep = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(x), Inches(top),
                    Pt(1), Inches(height)
                )
                sep.fill.solid()
                sep.fill.fore_color.rgb = get_color(spec, 'outline')
                sep.line.fill.background()
            
            current_date += relativedelta(months=1)
    
    def _render_tasks(
        self,
        slide: Any,
        gantt_data: GanttData,
        spec: Dict[str, Any],
        left: float,
        top: float,
        width: float,
        height: float
    ) -> None:
        """Render task bars."""
        tasks = gantt_data.tasks
        if not tasks:
            return
        
        # Calculate timeline span
        timeline = gantt_data.timeline
        start_date = self._parse_date(timeline.start)
        end_date = self._parse_date(timeline.end)
        total_months = (end_date.year - start_date.year) * 12 + (end_date.month - start_date.month) + 1
        
        # Task layout - dynamic height based on available space
        task_count = len(tasks)
        # Expand task bars to fill available vertical space
        max_task_h = min(1.2, (height - 0.2) / max(task_count, 1))
        task_height = max(0.60, max_task_h)  # At least 0.60, up to 1.2
        bar_height = task_height * 0.55
        
        # 状态到颜色的映射（使用设计规范中的颜色）
        status_colors = {
            'completed': 'primary',        # 深蓝色 - 已完成
            'done': 'primary',
            'active': 'accent_2',          # 进行中
            'in_progress': 'accent_2',
            'planned': 'accent_3',          # 计划中
            'pending': 'surface_variant',   # 灰色 - 待定
        }
        # Per-task color cycling for visual differentiation (when status is same)
        task_color_cycle = ['primary', 'accent_2', 'accent_3', 'accent_4', 'accent_1']
        # Text color on each bar background
        task_text_on = {
            'primary': 'on_primary',
            'accent_1': 'on_primary',
            'accent_2': 'on_primary',
            'accent_3': 'on_primary',
            'accent_4': 'on_primary',
            'surface_variant': 'on_surface',
            'outline': 'on_surface',
            'primary_container': 'on_primary_container',
        }
        
        # Check if all tasks have the same status → use color cycling
        all_same_status = len(set(t.status for t in tasks)) <= 1

        # Alternating row background colors for visual distinction
        row_bg_tokens = ['surface_variant', 'surface_dim']
        
        for i, task in enumerate(tasks):
            y = top + i * task_height
            
            # Row background strip (full width, alternating colors)
            row_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left), Inches(y),
                Inches(width), Inches(task_height)
            )
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = get_color(spec, row_bg_tokens[i % 2])
            row_bg.line.fill.background()
            
            # Task name (left side)
            name_width = width * 0.25
            name_tb = slide.shapes.add_textbox(
                Inches(left), Inches(y + 0.05),
                Inches(name_width - 0.1), Inches(bar_height)
            )
            name_tf = name_tb.text_frame
            name_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            name_tf.word_wrap = True
            p = name_tf.paragraphs[0]
            
            run = p.add_run()
            run.text = task.name
            run.font.size = Pt(get_font_size(spec, 'body'))
            run.font.color.rgb = get_color(spec, 'on_surface')
            apply_font_to_run(run, spec)
            
            # Task bar
            bar_area_left = left + name_width
            bar_area_width = width - name_width
            
            # Calculate bar position and width
            bar_start_x = bar_area_left + (task.start_month / total_months) * bar_area_width
            bar_width = (task.duration_months / total_months) * bar_area_width
            
            # Draw bar — use per-task cycling when all same status
            if all_same_status:
                color_token = task_color_cycle[i % len(task_color_cycle)]
            else:
                color_token = status_colors.get(task.status, 'outline')
            bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(bar_start_x), Inches(y + 0.05),
                Inches(bar_width), Inches(bar_height)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = get_color(spec, color_token)
            bar.line.fill.background()
            
            # Duration label on bar (if space allows)
            if bar_width > 0.5:
                dur_tb = slide.shapes.add_textbox(
                    Inches(bar_start_x), Inches(y + 0.05),
                    Inches(bar_width), Inches(bar_height)
                )
                dur_tf = dur_tb.text_frame
                dur_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                dur_tf.word_wrap = False
                p2 = dur_tf.paragraphs[0]
                p2.alignment = PP_ALIGN.CENTER
                
                run2 = p2.add_run()
                run2.text = f"{task.duration_months}M"
                run2.font.size = Pt(10)
                run2.font.bold = True
                text_token = task_text_on.get(color_token, 'on_primary')
                run2.font.color.rgb = get_color(spec, text_token)
                apply_font_to_run(run2, spec)
    
    def _parse_date(self, date_str: str) -> datetime:
        """Parse date string in YYYY-MM or YYYY-MM-DD format."""
        parts = date_str.split('-')
        if len(parts) == 2:
            return datetime(int(parts[0]), int(parts[1]), 1)
        elif len(parts) == 3:
            return datetime(int(parts[0]), int(parts[1]), int(parts[2]))
        else:
            raise ValueError(f"Invalid date format: {date_str}")
