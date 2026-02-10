from pptx import Presentation
import json
import sys

pptx_path='docs/presentations/MFT-20260210/MFT.pptx'
semantic_path='docs/presentations/MFT-20260210/slides_semantic.json'
design_path='docs/presentations/MFT-20260210/design_spec.json'

prs = Presentation(pptx_path)
with open(semantic_path) as f:
    semantic = json.load(f)
with open(design_path) as f:
    spec = json.load(f)

errors = []
slides = semantic.get('slides', [])

# MR-1: Backgrounds
for i, slide in enumerate(prs.slides):
    if slide.background.fill.type is None:
        errors.append(f"MR-1: Slide {i+1} has no background fill")

# MR-5: Section dividers
dividers = sum(1 for s in slides if s.get('slide_type') == 'section_divider')
sections = len(semantic.get('sections', []))
if len(slides) >= 15 and dividers == 0 and sections > 0:
    errors.append(f"MR-5: {len(slides)} slides, {sections} sections, but 0 dividers")

# MR-6: Title slide
if len(prs.slides) > 0:
    title_shapes = [s for s in prs.slides[0].shapes if s.has_text_frame]
    if len(title_shapes) < 3:
        errors.append(f"MR-6: Title slide has {len(title_shapes)} text frames, need ≥3")

# MR-8: Bottom bar
slide_h = prs.slide_height / 914400.0
for i, slide in enumerate(prs.slides[1:], 2):
    stype = slides[i-1].get('slide_type', '') if i-1 < len(slides) else ''
    if stype in ('title', 'section_divider'):
        continue
    bottom = [s for s in slide.shapes if (s.top + s.height) / 914400.0 > slide_h - 0.35]
    if not bottom:
        errors.append(f"MR-8: Slide {i} has no bottom bar")

# MR-9: Components
for i, sd in enumerate(slides):
    comps = sd.get('components', {})
    has_comps = any(comps.get(k) for k in comps)
    if has_comps and i < len(prs.slides):
        shape_count = len(prs.slides[i].shapes)
        if shape_count < 8:
            errors.append(f"MR-9: Slide {i+1} has components but only {shape_count} shapes")

# MR-11: Component Key Flexibility
for i, sd in enumerate(slides):
    comps = sd.get('components', {})
    items = comps.get('comparison_items') or comps.get('decisions') or []
    for item in items:
        data_keys = [k for k in item if k not in ('label', 'icon', 'color')]
        if len(data_keys) == 0 and len(item) > 1:
            errors.append(f"MR-11: Slide {i+1} component item has keys {list(item.keys())} but no data keys detected")

# MR-12: Content Deduplication
for i, sd in enumerate(slides):
    comps = sd.get('components', {})
    content = sd.get('content', [])
    comp_labels = set()
    for key in ('decisions', 'comparison_items'):
        for item in (comps.get(key) or []):
            if item.get('label'):
                comp_labels.add(item['label'])
    if comp_labels and content:
        dupes = [c for c in content if c in comp_labels]
        if dupes:
            errors.append(f"MR-12: Slide {i+1} content duplicates component labels: {dupes}")

# MR-13: Content Zone Space Utilization — simple heuristic
lz = spec.get('layout_zones', {})
title_bar_h = lz.get('title_bar_height_default', 0.55)
bottom_bar_h = max(lz.get('bottom_bar_height', 0.25), 0.25)
content_zone_top = title_bar_h + lz.get('content_margin_top', 0.12)
content_zone_h = slide_h - content_zone_top - bottom_bar_h - lz.get('content_bottom_margin', 0.2)
for i, sd in enumerate(slides):
    stype = sd.get('slide_type', '')
    if stype in ('title', 'section_divider'):
        continue
    if i >= len(prs.slides):
        continue
    content_shapes = [s for s in prs.slides[i].shapes
                      if s.top / 914400.0 > content_zone_top - 0.1
                      and (s.top + s.height) / 914400.0 < slide_h - bottom_bar_h + 0.1]
    if content_shapes:
        max_bottom = max((s.top + s.height) / 914400.0 for s in content_shapes)
        used_h = max_bottom - content_zone_top
        fill_ratio = used_h / content_zone_h if content_zone_h > 0 else 1.0
        if fill_ratio < 0.55:
            errors.append(f"MR-13: Slide {i+1} ({stype}) content fills only {int(fill_ratio*100)}% of content zone — likely whitespace problem")

# MR-14 and MR-15 simplified checks
for i, sd in enumerate(slides):
    if sd.get('slide_type') == 'title':
        comps = sd.get('components', {})
        has_comps = any(comps.get(k) for k in comps)
        if has_comps:
            errors.append(f"MR-15: Slide {i+1} is title slide but has components {list(k for k in comps if comps.get(k))}")

report = {
    'errors': errors,
    'status': 'PASS' if not errors else 'FAIL',
    'count': len(errors)
}
with open('docs/presentations/MFT-20260210/qa_mr_report.json', 'w', encoding='utf-8') as f:
    import json
    json.dump(report, f, ensure_ascii=False, indent=2)
if errors:
    print('MR checks found issues; see qa_mr_report.json')
    sys.exit(2)
else:
    print('MR checks passed; report written to qa_mr_report.json')
    sys.exit(0)
