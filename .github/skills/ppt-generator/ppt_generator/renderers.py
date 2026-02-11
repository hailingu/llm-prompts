"""Renderers for slides and components.

This module reuses helpers and grid to provide all rendering functions
exported by the original monolithic script.
"""
from typing import Any, Dict, List
import json
import os
import re
import logging

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from .helpers import (
    get_color, get_color_hex, get_font_size, hex_to_rgb, px_to_inches
)
from .grid import GridSystem


# For brevity we copy selectively the rendering helpers from the original
# script, focusing on title bar, bottom bar, speaker notes, and the main
# per-slide dispatcher. Additional component renderers (kpis, bullets,
# tables, visuals) are also included as smaller functions.


def get_layout_zones(spec: Dict) -> Dict[str, float]:
    lz = spec.get('layout_zones', {})
    return {
        'title_bar_h': lz.get('title_bar_height_default', 0.55),
        'title_bar_h_narrow': lz.get('title_bar_height_narrow', 0.4),
        'bottom_bar_h': max(lz.get('bottom_bar_height', 0.25), 0.25),
        'content_margin_top': lz.get('content_margin_top', 0.12),
        'content_bottom_margin': lz.get('content_bottom_margin', 0.2),
    }


def detect_schema_version(slide: Dict) -> int:
    """Detects whether a slide is v1 or v2 schema.

    Returns 2 if `layout_intent.regions` key exists (even if empty), otherwise 1.
    """
    try:
        if isinstance(slide, dict) and 'layout_intent' in slide and isinstance(slide.get('layout_intent'), dict) and 'regions' in slide['layout_intent']:
            return 2
    except Exception:
        pass
    return 1


def resolve_data_source(slide: Dict, path: str):
    """Resolve a dot-separated data source path from a slide dict.

    Examples:
      - 'components.kpis' -> slide['components']['kpis']
      - 'visual' -> slide['visual']
      - 'content' -> slide['content']

    Returns None if the path cannot be resolved.
    """
    if not path or not isinstance(path, str):
        return None
    cur = slide
    parts = path.split('.')
    try:
        for p in parts:
            if isinstance(cur, dict) and p in cur:
                cur = cur[p]
            else:
                # Not found; return None
                return None
        return cur
    except Exception:
        return None


def compute_region_bounds(position: str, grid: GridSystem, bar_h: float) -> tuple:
    """Compute (left, top, width, height) in inches for a given position marker.

    Supported markers:
      - 'full' -> full usable width
      - 'left-<pct>' / 'right-<pct>' -> left or right region with percentage of usable width
      - 'center-<pct>' -> centered region with percentage width
      - 'col-<start>-<span>' -> column-based region using grid.col_span
      - 'top-<pct>' -> full-width but height is pct of available content height

    Best-effort: if parsing fails, returns full area bounds.
    """
    left = grid.margin_h
    top = bar_h + 0.12
    usable_w = grid.usable_w
    slide_h = grid.slide_h

    # Default full area height: reserve 0.5" for bottom bar / margins
    default_available_h = max(1.0, slide_h - bar_h - 0.5)
    width = usable_w
    height = default_available_h

    if not position or not isinstance(position, str):
        return (left, top, width, height)

    pos = position.strip().lower()
    if pos == 'full':
        return (left, top, width, height)

    import re
    m = re.match(r'^(left|right|center)-(\d{1,3})$', pos)
    if m:
        side, val = m.group(1), int(m.group(2))
        pct = max(1, min(100, val)) / 100.0
        w = usable_w * pct
        if side == 'left':
            return (left, top, w, height)
        elif side == 'right':
            l = left + (usable_w - w)
            return (l, top, w, height)
        else:  # center
            l = left + (usable_w - w) / 2
            return (l, top, w, height)

    m2 = re.match(r'^top-(\d{1,3})$', pos)
    if m2:
        pct = max(1, min(100, int(m2.group(1)))) / 100.0
        h = default_available_h * pct
        return (left, top, width, h)

    m3 = re.match(r'^col-(\d+)-(\d+)$', pos)
    if m3:
        start = int(m3.group(1))
        span = int(m3.group(2))
        start = max(0, min(grid.columns - 1, start))
        span = max(1, min(grid.columns - start, span))
        l, w = grid.col_span(span, start)
        return (l, top, w, height)

    # fallback: return full
    return (left, top, width, height)


def get_title_bar_mode(spec: Dict, slide_type: str) -> str:
    if slide_type in ('title', 'section_divider'):
        return 'none'
    stl = spec.get('slide_type_layouts', {})
    entry = stl.get(slide_type, stl.get('default', {}))
    return entry.get('title_bar', 'standard')


def apply_font_to_run(run: Any, spec: Dict) -> None:
    try:
        font_family = spec.get('typography_system', {}).get('font_family')
        cjk_font = spec.get('typography_system', {}).get('cjk_font_family')
        if font_family:
            run.font.name = font_family
            r_pr = getattr(run._element, 'rPr', None)
            if r_pr is not None:
                try:
                    r_pr.rFonts.set(qn('a:ascii'), font_family)
                    r_pr.rFonts.set(qn('a:hAnsi'), font_family)
                except Exception:
                    pass
        if cjk_font:
            r_pr = getattr(run._element, 'rPr', None)
            if r_pr is not None:
                try:
                    r_pr.rFonts.set(qn('a:ea'), cjk_font)
                except Exception:
                    pass
    except Exception:
        pass


def create_card_shape(slide: Any, spec: Dict, left_in: float, top_in: float, width_in: float, height_in: float):
    try:
        br = spec.get('component_library', {}).get('card', {}).get('border_radius')
    except Exception:
        br = None
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if br is None or br != 0 else MSO_SHAPE.RECTANGLE
    return slide.shapes.add_shape(shape_type, Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))


... (file continues)

def render_bottom_bar(slide: Any, spec: Dict, grid: GridSystem, section_name: str, slide_num: int, total_slides: int,
                      section_index: int = 0, total_sections: int = 6, accent_token: str = 'primary') -> None:
    lz = get_layout_zones(spec)
    bar_h = lz['bottom_bar_h']
    bar_top = grid.slide_h - bar_h

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(bar_top), Inches(grid.slide_w), Inches(bar_h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = get_color(spec, 'surface_variant')
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(grid.margin_h), Inches(bar_top + 0.03), Inches(3), Inches(bar_h - 0.06))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = section_name
    run.font.size = Pt(get_font_size(spec, 'label'))
    run.font.bold = True
    run.font.color.rgb = get_color(spec, accent_token)

    prog_left = grid.slide_w / 2 - 1.5
    seg_w = 3.0 / max(total_sections, 1)
    for i in range(total_sections):
        seg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(prog_left + i * seg_w + 0.02), Inches(bar_top + bar_h / 2 - 0.03),
            Inches(seg_w - 0.04), Inches(0.06)
        )
        seg.fill.solid()
        if i <= section_index:
            seg.fill.fore_color.rgb = get_color(spec, accent_token)
        else:
            seg.fill.fore_color.rgb = get_color(spec, 'outline')
        seg.line.fill.background()

    tb = slide.shapes.add_textbox(
        Inches(grid.slide_w - grid.margin_h - 1.5), Inches(bar_top + 0.03), Inches(1.5), Inches(bar_h - 0.06)
    )
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{slide_num} / {total_slides}"
    run.font.size = Pt(get_font_size(spec, 'page_number'))
    run.font.color.rgb = get_color(spec, 'on_surface_variant')


def render_speaker_notes(slide: Any, notes_text: Any) -> None:
    if not notes_text:
        return
    text = notes_text if isinstance(notes_text, str) else json.dumps(notes_text, ensure_ascii=False)
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# Minimal component renderers used by slide renderers

def render_kpis(slide: Any, kpis: List[Dict], spec: Dict, _grid: GridSystem, left: float, top: float, width: float) -> float:
    if not kpis:
        return 0
    n = len(kpis)
    card_gap = 0.15
    card_w = (width - card_gap * (n - 1)) / max(n, 1)
    card_h = 0.85
    for i, kpi in enumerate(kpis):
        cx = left + i * (card_w + card_gap)
        card = create_card_shape(slide, spec, cx, top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = get_color(spec, 'primary_container')
        card.line.fill.background()
        # Value
        tb = slide.shapes.add_textbox(Inches(cx + 0.12), Inches(top + 0.10), Inches(card_w - 0.24), Inches(0.40))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        val = kpi.get('value', '')
        trend = kpi.get('trend', '')
        trend_arrow = {'up': ' ↑', 'down': ' ↓', 'stable': ' →'}.get(trend, '')
        run.text = f"{val}{trend_arrow}"
        run.font.size = Pt(get_font_size(spec, 'kpi_value'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_primary_container')
        apply_font_to_run(run, spec)
        # Label
        tb2 = slide.shapes.add_textbox(Inches(cx + 0.12), Inches(top + 0.52), Inches(card_w - 0.24), Inches(0.25))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = kpi.get('label', '')
        run2.font.size = Pt(get_font_size(spec, 'kpi_label'))
        run2.font.color.rgb = get_color(spec, 'on_primary_container')
    return card_h + 0.15


# Optional: python-pptx native charts
try:
    from pptx.chart.data import CategoryChartData, XyChartData
    from pptx.enum.chart import XL_CHART_TYPE
    HAS_PPTX_CHARTS = True
except Exception:
    HAS_PPTX_CHARTS = False
    CategoryChartData = None
    XyChartData = None
    XL_CHART_TYPE = None

# Map visual types to XL chart types
NATIVE_CHART_TYPE_MAP = {
    'bar_chart': XL_CHART_TYPE.COLUMN_CLUSTERED if XL_CHART_TYPE is not None else None,
    'column_chart': XL_CHART_TYPE.COLUMN_CLUSTERED if XL_CHART_TYPE is not None else None,
    'horizontal_bar': XL_CHART_TYPE.BAR_CLUSTERED if XL_CHART_TYPE is not None else None,
    'line_chart': XL_CHART_TYPE.LINE_MARKERS if XL_CHART_TYPE is not None else None,
    'pie_chart': XL_CHART_TYPE.PIE if XL_CHART_TYPE is not None else None,
    'doughnut_chart': XL_CHART_TYPE.DOUGHNUT if XL_CHART_TYPE is not None else None,
    'radar_chart': XL_CHART_TYPE.RADAR if XL_CHART_TYPE is not None else None,
    'scatter_chart': XL_CHART_TYPE.XY_SCATTER if XL_CHART_TYPE is not None else None,
}
try:
    print(f"DEBUG package: HAS_PPTX_CHARTS={HAS_PPTX_CHARTS}", file=sys.stderr)
except Exception:
    pass


def apply_chart_theme(chart, spec, accent_token='primary'):
    """Minimal theme application: rotate series colors using spec['md3_palette'] if present."""
    try:
        palette = spec.get('md3_palette') or []
        if isinstance(palette, str):
            palette = [palette]
        if not palette:
            palette = ['#2563EB', '#10B981', '#F59E0B', '#A78BFA']
        for i, ser in enumerate(chart.series):
            color = palette[i % len(palette)]
            try:
                ser.format.fill.solid()
                ser.format.fill.fore_color.rgb = hex_to_rgb(color)
            except Exception:
                pass
    except Exception:
        pass


def render_native_chart(slide: Any, visual: Dict, spec: Dict, left: float, top: float, width: float, height: float, accent_token: str = 'primary') -> bool:
    """Attempt to render a chart using python-pptx native APIs. Returns True on success."""
    if not HAS_PPTX_CHARTS:
        return False
    pd = visual.get('placeholder_data', {})
    config = pd.get('chart_config', {})
    if not config:
        return False
    labels = config.get('labels') or config.get('x') or []
    series = config.get('series', [])
    if not series:
        return False
    chart_type = visual.get('type', 'line_chart')
    if chart_type != 'scatter_chart' and not labels:
        return False
    xl_type = NATIVE_CHART_TYPE_MAP.get(chart_type)
    if not xl_type:
        return False
    try:
        if xl_type == XL_CHART_TYPE.XY_SCATTER:
            try:
                data = XyChartData()
                all_x = None
                for s in series:
                    xy_x = s.get('x', [])
                    xy_y = s.get('y', [])
                    if not xy_x or not xy_y:
                        return False
                    if all_x is None:
                        all_x = xy_x
                    if len(xy_x) != len(all_x):
                        return False
                    data.add_series(s.get('name', 'Series'), list(zip(xy_x, xy_y)))
            except Exception:
                data = CategoryChartData()
                first_x = series[0].get('x', [])
                data.categories = [str(x) for x in first_x]
                for s in series:
                    data.add_series(s.get('name', 'Series'), s.get('y', []))
                xl_type = XL_CHART_TYPE.LINE_MARKERS
        else:
            data = CategoryChartData()
            data.categories = [str(l) for l in labels]
            for s in series:
                data.add_series(s.get('name', 'Series'), s.get('data', []))

        shape = slide.shapes.add_chart(xl_type, Inches(left), Inches(top), Inches(width), Inches(height), data)
        chart = shape.chart
        apply_chart_theme(chart, spec, accent_token)
        return True
    except Exception:
        return False


def render_visual(slide: Any, visual: Dict, spec: Dict, _grid: GridSystem, left: float, top: float, width: float, height: float) -> None:
    if not visual or visual.get('type') in (None, 'none'):
        return
    pd = visual.get('placeholder_data', {})
    # Prefer native chart rendering when possible
    if pd.get('chart_config'):
        try:
            if HAS_PPTX_CHARTS:
                try:
                    print(f"DEBUG package: attempting native chart type {visual.get('type')}", file=sys.stderr)
                except Exception:
                    pass
                ok = render_native_chart(slide, visual, spec, left, top, width, height)
                try:
                    print(f"DEBUG package: render_native_chart returned {ok} for {visual.get('type')}", file=sys.stderr)
                except Exception:
                    pass
                if ok:
                    return
        except Exception as e:
            try:
                print(f"DEBUG package: render_native_chart exception {e}", file=sys.stderr)
            except Exception:
                pass
        # Fallback: render as data table
        render_chart_table(slide, visual, spec, left, top, width, height)
    elif pd.get('mermaid_code'):
        render_mermaid_placeholder(slide, visual, spec, left, top, width, height)
    else:
        render_visual_placeholder(slide, visual, spec, left, top, width, height)


def render_chart_table(slide: Any, visual: Dict, spec: Dict, left: float, top: float, width: float, height: float) -> None:
    config = visual.get('placeholder_data', {}).get('chart_config', {})
    labels = config.get('labels', [])
    series = config.get('series', [])
    if not labels:
        return
    n_cols = len(labels)
    col_w = width / max(n_cols, 1)
    n_rows = len(series) + 1
    title_offset = 0.35 if visual.get('title') else 0
    usable_h = height - title_offset
    row_h = max(0.42, min(1.2, usable_h / max(n_rows, 1)))
    if visual.get('title'):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top - 0.30), Inches(width), Inches(0.28))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = visual['title']
        run.font.size = Pt(get_font_size(spec, 'table_header'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_surface')
    for j, label in enumerate(labels):
        tb = slide.shapes.add_textbox(Inches(left + j * col_w), Inches(top), Inches(col_w), Inches(row_h))
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = str(label)
        run.font.size = Pt(get_font_size(spec, 'table_header'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_surface')
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top + row_h - 0.02), Inches(width), Pt(2))
    sep.fill.solid()
    sep.fill.fore_color.rgb = get_color(spec, 'primary')
    sep.line.fill.background()
    for r, s in enumerate(series):
        ry = top + (r + 1) * row_h
        if r % 2 == 1:
            stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(ry), Inches(width), Inches(row_h))
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = get_color(spec, 'surface_variant')
            stripe.line.fill.background()
        data = s.get('data', [])
        for j, val in enumerate(data):
            if j >= n_cols:
                break
            tb = slide.shapes.add_textbox(Inches(left + j * col_w), Inches(ry), Inches(col_w), Inches(row_h))
            tf = tb.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(get_font_size(spec, 'table_cell'))
            run.font.color.rgb = get_color(spec, 'on_surface')
            if isinstance(val, (int, float)):
                p.alignment = PP_ALIGN.RIGHT


def render_mermaid_placeholder(slide: Any, visual: Dict, spec: Dict, left: float, top: float, width: float, height: float) -> None:
    mermaid = visual.get('placeholder_data', {}).get('mermaid_code', '')
    card = create_card_shape(slide, spec, left, top, width, min(height, 3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = get_color(spec, 'surface_variant')
    card.line.color.rgb = get_color(spec, 'outline')
    card.line.width = Pt(1)
    if visual.get('title'):
        tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.12), Inches(width - 0.4), Inches(0.28))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"📊 {visual['title']}"
        run.font.size = Pt(get_font_size(spec, 'table_header'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_surface')
        apply_font_to_run(run, spec)
    preview = '\n'.join(mermaid.strip().split('\n')[:8])
    if len(mermaid.strip().split('\n')) > 8:
        preview += '\n  ...'
    tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.50), Inches(width - 0.4), Inches(min(height - 0.7, 2.5)))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.font.size = Pt(11)
    run.font.color.rgb = get_color(spec, 'outline')
    apply_font_to_run(run, spec)
    run.text = preview


def render_visual_placeholder(slide: Any, visual: Dict, spec: Dict, left: float, top: float, width: float, height: float) -> None:
    card = create_card_shape(slide, spec, left, top, width, min(height, 2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = get_color(spec, 'surface_variant')
    card.line.color.rgb = get_color(spec, 'outline')
    card.line.width = Pt(1)
    label = visual.get('title', visual.get('type', 'Visual'))
    reqs = visual.get('content_requirements', [])
    text = f"[{label}]"
    if reqs:
        text += '\n' + '\n'.join(f"  • {r}" for r in reqs[:3])
    tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.3), Inches(width - 0.4), Inches(min(height - 0.6, 2.0)))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(get_font_size(spec, 'table_header'))
    run.font.color.rgb = get_color(spec, 'outline')
    apply_font_to_run(run, spec)


# Dispatcher of slide types (partial set)

RENDERERS = {
    'title': None,  # implemented below
    'section_divider': None,
    'bullet-list': None,
    'two-column': None,
    'comparison': None,
    'decision': None,
    'data-heavy': None,
}


def _resolve_bg_image(spec: Dict, slide_type: str, sd: Dict) -> str:
    img = sd.get('background_image', '')
    if img:
        return img
    stl = spec.get('slide_type_layouts', {})
    entry = stl.get(slide_type, stl.get('default', {}))
    return entry.get('background_image', '')


def find_section_for_slide(slide_id: int, sections: List[Dict]) -> Dict:
    result = {}
    for sec in sorted(sections, key=lambda s: s.get('start_slide', 0)):
        if slide_id >= sec.get('start_slide', 999):
            result = sec
    return result


def _render_background(slide: Any, sd: Dict, spec: Dict, grid: GridSystem, stype: str) -> None:
    """Render the slide background image or solid fill."""
    bg_token = spec.get('slide_type_layouts', {}).get(stype, spec.get('slide_type_layouts', {}).get('default', {})).get('background', 'surface')
    bg_image_path = _resolve_bg_image(spec, stype, sd)
    if bg_image_path and os.path.isfile(bg_image_path):
        pic = slide.shapes.add_picture(bg_image_path, Emu(0), Emu(0), Inches(grid.slide_w), Inches(grid.slide_h))
        sp_tree = slide.shapes._spTree
        sp_tree.remove(pic._element)
        sp_tree.insert(2, pic._element)
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Inches(grid.slide_w), Inches(grid.slide_h))
        fill = overlay.fill
        fill.solid()
        fill.fore_color.rgb = get_color(spec, bg_token)
        from pptx.oxml.ns import qn as _qn
        solid_fill_elem = fill._fill.find(_qn('a:solidFill'))
        if solid_fill_elem is not None:
            srgb = solid_fill_elem.find(_qn('a:srgbClr'))
            if srgb is None:
                srgb = solid_fill_elem[0] if len(solid_fill_elem) else None
            if srgb is not None:
                try:
                    from lxml import etree
                    alpha = etree.SubElement(srgb, _qn('a:alpha'))
                    alpha.set('val', '40000')
                except Exception:
                    logging.getLogger(__name__).warning(
                        'lxml not available or failed to set overlay alpha; skipping alpha modification'
                    )
        overlay.line.fill.background()
        sp_tree = slide.shapes._spTree
        sp_tree.remove(overlay._element)
        sp_tree.insert(3, overlay._element)
    else:
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = get_color(spec, bg_token)


def _render_bullets_fallback(slide: Any, sd: Dict, spec: Dict, grid: GridSystem, bar_h: float) -> None:
    """Render simple bullet-list fallback content."""
    top = bar_h + 0.12
    content_bullets = sd.get('content', [])
    bullet_y = top + 0.15
    for i, bullet in enumerate(content_bullets[:8]):
        tb = slide.shapes.add_textbox(Inches(grid.margin_h + 0.1), Inches(bullet_y + i * 0.48), Inches(grid.usable_w - 0.2), Inches(0.45))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.size = Pt(get_font_size(spec, 'bullet_text'))
        run.font.color.rgb = get_color(spec, 'on_surface')
        p.line_spacing = 1.5


def _render_footer_and_notes(slide: Any, sd: Dict, spec: Dict, grid: GridSystem, stype: str, sec_title: str, sec_accent_token: str, sec_index: int, sections_len: int, slide_num: int, total_slides: int) -> None:
    """Render bottom bar and speaker notes."""
    if stype not in ('title', 'section_divider'):
        render_bottom_bar(
            slide,
            spec,
            grid,
            sec_title,
            slide_num,
            total_slides,
            section_index=sec_index,
            total_sections=sections_len,
            accent_token=sec_accent_token,
        )
    render_speaker_notes(slide, sd.get('speaker_notes'))


def render_slide(prs: Presentation, sd: Dict, spec: Dict, grid: GridSystem, sections: List[Dict], slide_num: int, total_slides: int) -> None:
    """Render a single slide with full styling (delegates to helpers)."""
    stype = sd.get('slide_type', 'bullet-list')
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Section context
    slide_id = sd.get('slide_id', slide_num)
    section = find_section_for_slide(slide_id, sections)
    sec_id = section.get('id', 'A')
    sec_title = section.get('title', '')
    sec_accent_token = spec.get('section_accents', {}).get(sec_id, 'primary')
    sec_index = next((i for i, s in enumerate(sections) if s.get('id') == sec_id), 0)

    # Background
    _render_background(slide, sd, spec, grid, stype)

    # Title bar
    tb_mode = get_title_bar_mode(spec, stype)
    bar_h = 0.0
    if tb_mode != 'none':
        section_label = f"Section {sec_id} · {sec_title}" if sec_title else ''
        bar_h = render_title_bar(slide, spec, grid, sd.get('title', ''), slide_num, total_slides, section_label=section_label, accent_color_token=sec_accent_token, mode=tb_mode)

    # Content (fallback)
    _render_bullets_fallback(slide, sd, spec, grid, bar_h)

    # Try to render the main visual if present (useful for data-heavy slides)
    try:
        vis = sd.get('visual')
        if isinstance(vis, dict) and vis.get('type') not in (None, 'none'):
            lz = get_layout_zones(spec)
            left = grid.margin_h
            top = bar_h + lz['content_margin_top']
            width = grid.usable_w
            bottom_bar_h = lz['bottom_bar_h']
            height = max(1.0, grid.slide_h - top - bottom_bar_h - 0.3)
            render_visual(slide, vis, spec, grid, left, top, width, height)
    except Exception:
        # best-effort: don't fail slide rendering if visual rendering errors
        pass

    # Footer & notes
    _render_footer_and_notes(slide, sd, spec, grid, stype, sec_title, sec_accent_token, sec_index, len(sections), slide_num, total_slides)
