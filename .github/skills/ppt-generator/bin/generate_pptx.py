#!/usr/bin/env python3
"""Generate PPTX from slides_semantic.json + design_spec.json."""
import argparse
import math
import re
import sys
import os
import json
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# Optional: matplotlib for chart generation
try:
    import matplotlib
    matplotlib.use('Agg')  # Non-interactive backend
    import matplotlib.pyplot as plt
    import matplotlib.ticker as ticker
    HAS_MATPLOTLIB = True
except ImportError:
    HAS_MATPLOTLIB = False

# ════════════════════════════════════════════════════════════════
# §1  Token Helpers
# ════════════════════════════════════════════════════════════════

def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def get_color(spec: dict, token_name: str) -> RGBColor:
    # 1. Top-level color_system (legacy / flat structure)
    cs = spec.get('color_system') or spec.get('design_system', {}).get('color_system', {})
    val = cs.get(token_name)
    if val:
        return hex_to_rgb(val)
    # 2. tokens.colors (visual-designer's Material Design 3 structure)
    val = spec.get('tokens', {}).get('colors', {}).get(token_name)
    if val and isinstance(val, str):
        return hex_to_rgb(val)
    # 3. theme_tokens.palette (alternate nested structure)
    palette = spec.get('theme_tokens', {}).get('palette', {})
    val = palette.get(token_name)
    if val:
        return hex_to_rgb(val)
    # Last resort defaults for common tokens
    _DEFAULTS = {
        'primary': '#2563EB', 'secondary': '#10B981',
        'primary_container': '#E6F0FF', 'on_primary_container': '#0F172A',
        'on_primary': '#FFFFFF', 'surface': '#FFFFFF',
        'surface_variant': '#F3F4F6', 'surface_dim': '#F8FAFC',
        'on_surface': '#0F172A', 'on_surface_variant': '#6B7280',
        'outline': '#D1D5DB', 'muted': '#6B7280',
        'error': '#DC2626', 'warning': '#F59E0B', 'success': '#10B981',
        'accent_1': '#2563EB', 'accent_2': '#10B981',
        'accent_3': '#F59E0B', 'accent_4': '#A78BFA',
    }
    return hex_to_rgb(_DEFAULTS.get(token_name, '#1A1C1E'))


def get_color_hex(spec: dict, token_name: str) -> str:
    # 1. Top-level color_system
    cs = spec.get('color_system') or spec.get('design_system', {}).get('color_system', {})
    val = cs.get(token_name)
    if val:
        return val
    # 2. tokens.colors (Material Design 3 structure)
    val = spec.get('tokens', {}).get('colors', {}).get(token_name)
    if val and isinstance(val, str):
        return val
    # 3. theme_tokens.palette
    palette = spec.get('theme_tokens', {}).get('palette', {})
    val = palette.get(token_name)
    if val:
        return val
    _DEFAULTS = {
        'primary': '#2563EB', 'secondary': '#10B981',
        'primary_container': '#E6F0FF', 'on_primary_container': '#0F172A',
        'on_primary': '#FFFFFF', 'surface': '#FFFFFF',
        'surface_variant': '#F3F4F6', 'surface_dim': '#F8FAFC',
        'on_surface': '#0F172A', 'on_surface_variant': '#6B7280',
        'outline': '#D1D5DB', 'muted': '#6B7280',
        'error': '#DC2626', 'warning': '#F59E0B', 'success': '#10B981',
        'accent_1': '#2563EB', 'accent_2': '#10B981',
        'accent_3': '#F59E0B', 'accent_4': '#A78BFA',
    }
    return _DEFAULTS.get(token_name, '#1A1C1E')


def _is_light_color(hex_str: str) -> bool:
    """Check if a hex color is perceptually light (needs dark text)."""
    h = hex_str.lstrip('#')
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    luminance = 0.2126 * r + 0.7152 * g + 0.0722 * b
    return luminance > 160


def _text_color_for_bg(spec, bg_token):
    """Return appropriate text color token for a background color token.

    Light backgrounds (e.g., yellow/tertiary) get dark text;
    dark backgrounds (e.g., blue/primary) get white text.
    """
    hex_val = get_color_hex(spec, bg_token)
    if _is_light_color(hex_val):
        on_token = f'on_{bg_token}'
        cs = spec.get('color_system', {})
        if cs.get(on_token):
            return on_token
        return 'on_surface'
    return 'on_primary'


# Default font sizes for all roles used by the renderer.
# These ensure visual hierarchy even when design_spec omits some roles.
_FONT_SIZE_DEFAULTS = {
    'display_large': 40, 'headline_large': 28, 'title': 22,
    'slide_title': 22, 'slide_subtitle': 16, 'section_label': 10,
    'page_number': 10, 'label': 10, 'label_large': 12,
    'body': 11, 'body_text': 14, 'bullet_text': 14,
    'kpi_value': 20, 'kpi_label': 11,
    'table_header': 12, 'table_cell': 11,
    'callout_text': 13,
}


def get_font_size(spec: dict, role: str) -> int:
    # 1. Top-level typography_system (legacy / flat)
    ts = spec.get('typography_system') or spec.get('design_system', {}).get('typography_system', {})
    explicit = ts.get('explicit_sizes', {})
    if role in explicit:
        entry = explicit[role]
        if isinstance(entry, dict):
            return entry.get('size_pt', _FONT_SIZE_DEFAULTS.get(role, 16))
        return entry
    # 2. tokens.typography_system (Material Design 3 structure)
    token_ts = spec.get('tokens', {}).get('typography_system', {})
    if role in token_ts:
        entry = token_ts[role]
        if isinstance(entry, dict):
            return entry.get('size_pt', _FONT_SIZE_DEFAULTS.get(role, 16))
        return entry
    # 3. Fallback to type_scale in top-level typography_system
    scale = ts.get('type_scale', {})
    if role in scale:
        entry = scale[role]
        if isinstance(entry, dict):
            return entry.get('size_pt', 16)
        return entry
    return _FONT_SIZE_DEFAULTS.get(role, 14)


def px_to_inches(px):
    return px / 96.0


# ════════════════════════════════════════════════════════════════
# §2  Grid System
# ════════════════════════════════════════════════════════════════

class GridSystem:
    def __init__(self, spec: dict):
        grid = (spec.get('design_system', {}).get('grid_system')
                or spec.get('grid_system', {}))
        self.slide_w = grid.get('slide_width_inches', 13.333)
        self.slide_h = grid.get('slide_height_inches', 7.5)
        self.margin_h = px_to_inches(grid.get('margin_horizontal', 80))
        self.gutter = px_to_inches(grid.get('gutter', 24))
        self.columns = grid.get('columns', 12)
        self.usable_w = self.slide_w - 2 * self.margin_h
        self.col_w = (self.usable_w - self.gutter * (self.columns - 1)) / self.columns

    def col_span(self, n_cols: int, start_col: int = 0):
        left = self.margin_h + start_col * (self.col_w + self.gutter)
        width = n_cols * self.col_w + max(0, n_cols - 1) * self.gutter
        return left, width


# ════════════════════════════════════════════════════════════════
# §3  Layout Helpers
# ════════════════════════════════════════════════════════════════

def get_layout_zones(spec: dict) -> dict:
    lz = spec.get('layout_zones', {})
    return {
        'title_bar_h': lz.get('title_bar_height_default', 0.55),
        'title_bar_h_narrow': lz.get('title_bar_height_narrow', 0.4),
        'bottom_bar_h': max(lz.get('bottom_bar_height', 0.25), 0.25),
        'content_margin_top': lz.get('content_margin_top', 0.12),
        'content_bottom_margin': lz.get('content_bottom_margin', 0.2),
    }


def get_title_bar_mode(spec: dict, slide_type: str) -> str:
    # DEFENSIVE: title and section_divider are full-bleed slides — never show title bar
    if slide_type in ('title', 'section_divider'):
        return 'none'
    stl = spec.get('slide_type_layouts', {})
    entry = stl.get(slide_type, stl.get('default', {}))
    return entry.get('title_bar', 'standard')


# Helper: apply configured fonts to a text run (latin + CJK)
def apply_font_to_run(run, spec):
    try:
        font_family = spec.get('typography_system', {}).get('font_family')
        cjk_font = spec.get('typography_system', {}).get('cjk_font_family')
        if font_family:
            run.font.name = font_family
            # try to set rPr fonts for broader compatibility
            rPr = getattr(run._element, 'rPr', None)
            if rPr is not None:
                try:
                    rPr.rFonts.set(qn('a:ascii'), font_family)
                    rPr.rFonts.set(qn('a:hAnsi'), font_family)
                except Exception:
                    pass
        if cjk_font:
            rPr = getattr(run._element, 'rPr', None)
            if rPr is not None:
                try:
                    rPr.rFonts.set(qn('a:ea'), cjk_font)
                except Exception:
                    pass
    except Exception:
        pass


# Helper: create a card shape honoring component_library.card.border_radius
def create_card_shape(slide, spec, left_in, top_in, width_in, height_in):
    try:
        br = spec.get('component_library', {}).get('card', {}).get('border_radius')
    except Exception:
        br = None
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if br is None or br != 0 else MSO_SHAPE.RECTANGLE
    return slide.shapes.add_shape(shape_type, Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))


def get_bg_token(spec: dict, slide_type: str) -> str:
    stl = spec.get('slide_type_layouts', {})
    entry = stl.get(slide_type, stl.get('default', {}))
    bg = entry.get('background', 'surface')
    # DEFENSIVE: title and section_divider need dark/saturated background for white text
    if slide_type in ('title', 'section_divider'):
        if bg in ('surface', 'surface_variant', 'surface_dim', 'primary_container'):
            bg = 'primary'
    return bg


def get_content_fill(spec: dict, slide_type: str) -> str:
    """Get content fill strategy for a slide type.
    
    Returns 'expand' (default), 'center', or 'top-align'.
    When 'expand', components should grow to fill available vertical space.
    """
    stl = spec.get('slide_type_layouts', {})
    entry = stl.get(slide_type, stl.get('default', {}))
    return entry.get('content_fill', 'expand')


def get_section_accent(spec: dict, section_id: str) -> str:
    sa = spec.get('section_accents', {})
    return sa.get(section_id, 'primary')


def find_section_for_slide(slide_id: int, sections: list) -> dict:
    """Find which section a slide belongs to."""
    result = {}
    for sec in sorted(sections, key=lambda s: s.get('start_slide', 0)):
        if slide_id >= sec.get('start_slide', 999):
            result = sec
    return result


# ════════════════════════════════════════════════════════════════
# §3b  Bullet Dot Helper
# ════════════════════════════════════════════════════════════════

def _add_bullet_dot(slide, spec, left, center_y, accent_token='primary', size=0.08):
    """Add a small colored circle as a bullet marker.

    Args:
        slide: pptx slide object.
        spec: design spec dict.
        left: left edge x position in inches.
        center_y: vertical center of the bullet line in inches.
        accent_token: color token for the dot (default 'primary').
        size: diameter of the dot in inches.
    """
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(center_y - size / 2),
        Inches(size), Inches(size)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = get_color(spec, accent_token)
    dot.line.fill.background()
    return dot


def _render_highlighted_text(paragraph, text, spec, font_size_pt,
                             base_color_token='on_surface',
                             highlight_color_token='primary'):
    """Render text with automatic keyword highlighting.

    Detects numbers (e.g., 100TB, 3.5x, 40%), English technical terms (≥3 chars
    in an otherwise CJK context), and quoted text — rendering them in bold
    with the highlight color. The rest uses the base color.

    Args:
        paragraph: pptx paragraph object.
        text: full bullet text string.
        spec: design spec dict.
        font_size_pt: font size in points.
        base_color_token: color token for normal text.
        highlight_color_token: color token for highlighted keywords.
    """
    import re as _re
    # Pattern: numbers with units/percent, English words ≥3 chars, quoted text
    highlight_re = _re.compile(
        r'(\d+[\d,.]*\s*(?:%|TB|GB|PB|EB|MB|KB|x|X|倍|ms|μs|ns|us|IOPS|Gbps|Gbit|万|亿)?)'
        r'|([A-Za-z][A-Za-z0-9_./\-]{2,}(?:\s+[A-Za-z][A-Za-z0-9_./\-]{2,})*)'
        r'|("([^"]+)")'
        r"|('([^']+)')"
        r'|("([^"]+)")'
    )
    base_color = get_color(spec, base_color_token)
    hl_color = get_color(spec, highlight_color_token)

    last_end = 0
    has_highlights = False
    for m in highlight_re.finditer(text):
        start, end = m.start(), m.end()
        # Add normal text before the match
        if start > last_end:
            run = paragraph.add_run()
            run.text = text[last_end:start]
            run.font.size = Pt(font_size_pt)
            run.font.color.rgb = base_color
            apply_font_to_run(run, spec)
        # Add highlighted text
        run = paragraph.add_run()
        run.text = m.group(0)
        run.font.size = Pt(font_size_pt)
        run.font.bold = True
        run.font.color.rgb = hl_color
        apply_font_to_run(run, spec)
        last_end = end
        has_highlights = True

    # Remaining text after last match (or full text if no highlights)
    if last_end < len(text):
        run = paragraph.add_run()
        run.text = text[last_end:]
        run.font.size = Pt(font_size_pt)
        run.font.color.rgb = base_color
        apply_font_to_run(run, spec)
    elif not has_highlights:
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(font_size_pt)
        run.font.color.rgb = base_color
        apply_font_to_run(run, spec)


# ════════════════════════════════════════════════════════════════
# §4  Shared Renderers — Title bar, Bottom bar, Speaker Notes
# ════════════════════════════════════════════════════════════════

def render_title_bar(slide, spec, grid, title, slide_num, total_slides,
                     section_label='', accent_color_token='primary', mode='standard'):
    """Render title bar at top of slide."""
    if mode == 'none':
        return 0.0
    lz = get_layout_zones(spec)
    bar_h = lz['title_bar_h_narrow'] if mode == 'narrow' else lz['title_bar_h']

    # DEFENSIVE: When section_label is present, need extra height for label + title.
    # Without label: title starts at 0.10", needs ~0.35" for 22pt text → min 0.50"
    # With label: label occupies top 0.28", title needs ~0.35" → min 0.68"
    title_font_pt = get_font_size(spec, 'slide_title')
    title_text_h = title_font_pt / 72.0 + 0.08  # font height in inches + padding
    if section_label:
        min_bar_h = 0.28 + title_text_h  # label zone + title zone
    else:
        min_bar_h = 0.10 + title_text_h  # top padding + title zone
    bar_h = max(bar_h, min_bar_h)

    # Title bar background
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(grid.slide_w), Inches(bar_h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = get_color(spec, accent_color_token)
    bar.line.fill.background()

    # C1: Adaptive text color — light accents (yellow) get dark text
    text_tk = _text_color_for_bg(spec, accent_color_token)

    # Section label (small, top-left)
    if section_label:
        tb = slide.shapes.add_textbox(
            Inches(grid.margin_h), Inches(0.06),
            Inches(4), Inches(0.22)
        )
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = section_label
        run.font.size = Pt(get_font_size(spec, 'section_label'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, text_tk)
        apply_font_to_run(run, spec)

    # Title text
    title_top = 0.22 if section_label else 0.10
    tb = slide.shapes.add_textbox(
        Inches(grid.margin_h), Inches(title_top),
        Inches(grid.slide_w - 2 * grid.margin_h - 1.0), Inches(bar_h - title_top - 0.05)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(get_font_size(spec, 'slide_title'))
    run.font.bold = True
    run.font.color.rgb = get_color(spec, text_tk)
    apply_font_to_run(run, spec)

    # Page number
    tb_num = slide.shapes.add_textbox(
        Inches(grid.slide_w - 1.0), Inches(title_top),
        Inches(0.7), Inches(0.35)
    )
    tf_num = tb_num.text_frame
    p = tf_num.paragraphs[0]
    run = p.add_run()
    run.text = f"{slide_num}"
    run.font.size = Pt(get_font_size(spec, 'page_number'))
    run.font.bold = True
    run.font.color.rgb = get_color(spec, text_tk)
    apply_font_to_run(run, spec)
    p.alignment = PP_ALIGN.RIGHT

    # C3: Bottom separator line — color adapts to accent
    sep_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(bar_h - 0.025),
        Inches(grid.slide_w), Inches(0.025)
    )
    sep_line.fill.solid()
    sep_line.fill.fore_color.rgb = get_color(spec, text_tk)
    sep_line.line.fill.background()

    return bar_h


def render_bottom_bar(slide, spec, grid, section_name, slide_num, total_slides,
                      section_index=0, total_sections=6, accent_token='primary'):
    """Render bottom bar with section name, progress, and slide number."""
    lz = get_layout_zones(spec)
    bar_h = lz['bottom_bar_h']
    bar_top = grid.slide_h - bar_h

    # Background
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(bar_top),
        Inches(grid.slide_w), Inches(bar_h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = get_color(spec, 'surface_variant')
    bar.line.fill.background()

    # Section name (left)
    tb = slide.shapes.add_textbox(
        Inches(grid.margin_h), Inches(bar_top + 0.03),
        Inches(3), Inches(bar_h - 0.06)
    )
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = section_name
    run.font.size = Pt(get_font_size(spec, 'label'))
    run.font.bold = True
    run.font.color.rgb = get_color(spec, accent_token)

    # Progress bar (center) — filled rectangles for completed sections
    prog_left = grid.slide_w / 2 - 1.5
    seg_w = 3.0 / max(total_sections, 1)
    for i in range(total_sections):
        seg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(prog_left + i * seg_w + 0.02), Inches(bar_top + bar_h / 2 - 0.03),
            Inches(seg_w - 0.04), Inches(0.06)
        )
        seg.fill.solid()
        if i <= section_index:
            seg.fill.fore_color.rgb = get_color(spec, accent_token)
        else:
            seg.fill.fore_color.rgb = get_color(spec, 'outline')
        seg.line.fill.background()

    # Slide number (right)
    tb = slide.shapes.add_textbox(
        Inches(grid.slide_w - grid.margin_h - 1.5), Inches(bar_top + 0.03),
        Inches(1.5), Inches(bar_h - 0.06)
    )
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{slide_num} / {total_slides}"
    run.font.size = Pt(get_font_size(spec, 'page_number'))
    run.font.color.rgb = get_color(spec, 'on_surface_variant')


def render_speaker_notes(slide, notes_text):
    """Write speaker notes to slide."""
    if not notes_text:
        return
    text = notes_text if isinstance(notes_text, str) else json.dumps(notes_text, ensure_ascii=False)
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ════════════════════════════════════════════════════════════════
# §5  Component Renderers
# ════════════════════════════════════════════════════════════════

def _parse_kpi_number(value_str):
    """Extract numeric value from KPI string like '110 USD Bn' → 110."""
    if not value_str:
        return None
    m = re.search(r'[\d,.]+', str(value_str))
    if m:
        try:
            return float(m.group().replace(',', ''))
        except ValueError:
            return None
    return None


def render_kpis(slide, kpis, spec, grid, left, top, width):
    """Render KPI cards with proportion bars showing relative magnitude."""
    if not kpis:
        return 0
    n = len(kpis)
    card_gap = 0.15
    card_w = (width - card_gap * (n - 1)) / max(n, 1)
    card_h = 1.05  # slightly taller to fit proportion bar

    # Parse all numeric values and find max for scaling
    values = [_parse_kpi_number(kpi.get('value', '')) for kpi in kpis]
    max_val = max((v for v in values if v is not None), default=None)

    for i, kpi in enumerate(kpis):
        cx = left + i * (card_w + card_gap)
        card = create_card_shape(slide, spec, cx, top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = get_color(spec, 'primary_container')
        card.line.fill.background()
        add_shadow(card)
        # Value
        tb = slide.shapes.add_textbox(
            Inches(cx + 0.12), Inches(top + 0.08),
            Inches(card_w - 0.24), Inches(0.38)
        )
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        val = kpi.get('value', '')
        trend = kpi.get('trend', '')
        trend_arrow = {'up': ' ↑', 'down': ' ↓', 'stable': ' →'}.get(trend, '')
        run.text = f"{val}{trend_arrow}"
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'primary')
        apply_font_to_run(run, spec)
        # Label
        tb2 = slide.shapes.add_textbox(
            Inches(cx + 0.12), Inches(top + 0.48),
            Inches(card_w - 0.24), Inches(0.22)
        )
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = kpi.get('label', '')
        run2.font.size = Pt(13)
        run2.font.color.rgb = get_color(spec, 'on_primary_container')
        apply_font_to_run(run2, spec)
        # Proportion bar
        bar_top = top + 0.76
        bar_max_w = card_w - 0.24
        bar_h_px = 0.10
        # Background track
        track = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(cx + 0.12), Inches(bar_top),
            Inches(bar_max_w), Inches(bar_h_px)
        )
        track.fill.solid()
        track.fill.fore_color.rgb = get_color(spec, 'surface_variant')
        track.line.fill.background()
        # Filled proportion
        val_num = values[i]
        if val_num is not None and max_val and max_val > 0:
            ratio = val_num / max_val
            fill_w = max(0.1, bar_max_w * ratio)
            bar_fill = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(cx + 0.12), Inches(bar_top),
                Inches(fill_w), Inches(bar_h_px)
            )
            bar_fill.fill.solid()
            bar_fill.fill.fore_color.rgb = get_color(spec, 'primary')
            bar_fill.line.fill.background()
    return card_h + 0.15


def render_comparison_items(slide, items, spec, grid, left, top, width, avail_h=None):
    """Render comparison items as side-by-side cards.
    
    Handles any key structure in comparison_items:
    - Known keys: name/label, advantage, risk, recommendation, rationale, description
    - Generic keys: impact, feasibility, short_action, score, etc.
    - Numeric values shown with the key as label
    
    If avail_h is provided, cards expand to fill the available vertical space.
    """
    if not items:
        return 0
    n = len(items)
    card_gap = 0.15
    # M3: 2-column grid when 4+ items to prevent cards being too narrow for CJK
    if n >= 4:
        cols = 2
        rows = math.ceil(n / cols)
    else:
        cols = n
        rows = 1
    card_w = (width - card_gap * (cols - 1)) / max(cols, 1)
    row_gap = 0.15
    # Adaptive card height: distribute across rows
    if avail_h and avail_h > 0:
        card_h = max(1.6, (avail_h - row_gap * (rows - 1) - 0.15) / rows)
    else:
        card_h = 2.5 if rows == 1 else 1.6
    # Keys to skip (used for header, not body)
    skip_keys = {'name', 'label', 'id', 'slide_id'}
    # Keys with special formatting
    icon_map = {
        'advantage': '✅', 'risk': '⚠️', 'recommendation': '→',
        'short_action': '→', 'action': '→', 'next_step': '→',
    }
    # Comparison card color rotation using chart palette
    palette = get_chart_palette(spec)
    # Alternate card background tokens for visual variety
    _card_bg_tokens = ['surface_variant', 'primary_container', 'surface_dim']
    for i, item in enumerate(items):
        col = i % cols
        row = i // cols
        cx = left + col * (card_w + card_gap)
        cy = top + row * (card_h + row_gap)
        card = create_card_shape(slide, spec, cx, cy, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = get_color(spec, _card_bg_tokens[i % len(_card_bg_tokens)])
        card.line.fill.background()
        add_shadow(card)
        # Colored accent stripe at top (palette rotation)
        accent_color = palette[i % len(palette)]
        accent_stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(cx), Inches(cy), Inches(card_w), Inches(0.05)
        )
        accent_stripe.fill.solid()
        accent_stripe.fill.fore_color.rgb = hex_to_rgb(accent_color)
        accent_stripe.line.fill.background()
        # Name header
        tb = slide.shapes.add_textbox(
            Inches(cx + 0.12), Inches(cy + 0.12),
            Inches(card_w - 0.24), Inches(0.35)
        )
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = item.get('name', item.get('label', ''))
        run.font.size = Pt(18)  # Enlarged comparison card header
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'primary')
        apply_font_to_run(run, spec)
        # Separator line under header
        sep = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(cx + 0.12), Inches(cy + 0.50),
            Inches(card_w - 0.24), Inches(0.02)
        )
        sep.fill.solid()
        sep.fill.fore_color.rgb = hex_to_rgb(accent_color)
        sep.line.fill.background()
        # ensure header font follows style
        try:
            hdr_run = sep._element
        except Exception:
            hdr_run = None
        # Build details from ALL non-skip keys
        details = []
        # First: unpack 'attributes' dict (MO-0 schema: comparison_items[].attributes)
        attrs = item.get('attributes', {})
        if isinstance(attrs, dict) and attrs:
            for attr_key, attr_val in attrs.items():
                if attr_val is not None:
                    details.append(f"{attr_key}: {attr_val}")
        # Then: known formatted keys (flat structure fallback)
        for key in ('advantage', 'risk', 'recommendation', 'rationale', 'description',
                    'short_action', 'action', 'next_step'):
            if item.get(key):
                icon = icon_map.get(key, '')
                val = item[key]
                details.append(f"{icon} {val}" if icon else str(val))
        # Then: any remaining generic keys (numeric or string)
        for key, val in item.items():
            if key in skip_keys or key == 'attributes' or val is None:
                continue
            if key in ('advantage', 'risk', 'recommendation', 'rationale',
                       'description', 'short_action', 'action', 'next_step',
                       'highlight'):
                continue  # already handled or display-only
            if isinstance(val, (int, float)):
                pretty_key = key.replace('_', ' ').title()
                details.append(f"{pretty_key}: {val}")
            elif isinstance(val, str) and val:
                pretty_key = key.replace('_', ' ').title()
                details.append(f"{pretty_key}: {val}")
        detail_text = '\n'.join(details)
        if detail_text:
            tb2 = slide.shapes.add_textbox(
                Inches(cx + 0.12), Inches(cy + 0.55),
                Inches(card_w - 0.24), Inches(card_h - 0.65)
            )
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            run2 = p2.add_run()
            run2.text = detail_text
            run2.font.size = Pt(get_font_size(spec, 'body_text'))
            run2.font.color.rgb = get_color(spec, 'on_surface')
            p2.line_spacing = 1.5
    return rows * card_h + (rows - 1) * row_gap + 0.15


def render_comparison_horizontal(slide, items, spec, grid, left, top, width, avail_h=None):
    """方案2: Render comparison items as full-width horizontal bars.

    Each item is a horizontal card spanning the full width with:
    - Left accent color strip
    - Bold label on the left region
    - Key-value attributes flowing as inline text on the right
    This alternates visually with the vertical-card layout.
    """
    if not items:
        return 0
    n = len(items)
    gap = 0.12
    # Distribute bars evenly across available height
    if avail_h and avail_h > 0:
        bar_h = max(0.60, (avail_h - gap * (n - 1) - 0.15) / n)
    else:
        bar_h = 0.75

    palette = get_chart_palette(spec)
    _card_bg_tokens = ['surface_variant', 'primary_container', 'surface_dim']
    skip_keys = {'name', 'label', 'id', 'slide_id', 'highlight'}
    icon_map = {
        'advantage': '✅', 'risk': '⚠️', 'recommendation': '→',
        'short_action': '→', 'action': '→', 'next_step': '→',
    }
    cursor = top
    for i, item in enumerate(items):
        accent_color = palette[i % len(palette)]
        # Background card
        card = create_card_shape(slide, spec, left, cursor, width, bar_h)
        card.fill.solid()
        card.fill.fore_color.rgb = get_color(spec, _card_bg_tokens[i % len(_card_bg_tokens)])
        card.line.fill.background()
        add_shadow(card)
        # Left accent strip
        strip_w = 0.06
        strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(cursor),
            Inches(strip_w), Inches(bar_h)
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = hex_to_rgb(accent_color)
        strip.line.fill.background()
        # Label area (left 25%)
        label_w = min(width * 0.25, 2.5)
        tb = slide.shapes.add_textbox(
            Inches(left + strip_w + 0.12), Inches(cursor + 0.08),
            Inches(label_w - 0.24), Inches(bar_h - 0.16)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = item.get('name', item.get('label', ''))
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'primary')
        apply_font_to_run(run, spec)
        # Attributes area (right 75%)
        attr_left = left + strip_w + label_w + 0.10
        attr_w = width - strip_w - label_w - 0.22
        # Build details (same logic as vertical cards)
        details = []
        attrs = item.get('attributes', {})
        if isinstance(attrs, dict) and attrs:
            for attr_key, attr_val in attrs.items():
                if attr_val is not None:
                    details.append(f"{attr_key}: {attr_val}")
        for key in ('advantage', 'risk', 'recommendation', 'rationale',
                    'description', 'short_action', 'action', 'next_step'):
            if item.get(key):
                icon = icon_map.get(key, '')
                val = item[key]
                details.append(f"{icon} {val}" if icon else str(val))
        for key, val in item.items():
            if key in skip_keys or key == 'attributes' or val is None:
                continue
            if key in ('advantage', 'risk', 'recommendation', 'rationale',
                       'description', 'short_action', 'action', 'next_step'):
                continue
            if isinstance(val, (int, float, str)) and val:
                pretty_key = key.replace('_', ' ').title()
                details.append(f"{pretty_key}: {val}")

        if details:
            tb2 = slide.shapes.add_textbox(
                Inches(attr_left), Inches(cursor + 0.06),
                Inches(attr_w), Inches(bar_h - 0.12)
            )
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
            # Render each detail as a separate paragraph for clarity
            for j, detail in enumerate(details):
                if j == 0:
                    p2 = tf2.paragraphs[0]
                else:
                    p2 = tf2.add_paragraph()
                run2 = p2.add_run()
                run2.text = detail
                run2.font.size = Pt(get_font_size(spec, 'body_text'))
                run2.font.color.rgb = get_color(spec, 'on_surface')
                apply_font_to_run(run2, spec)
                p2.line_spacing = 1.3
        cursor += bar_h + gap
    return cursor - top


def render_decisions(slide, decisions, spec, grid, left, top, width, avail_h=None):
    """Render decision items as callout boxes.
    
    If avail_h is provided, cards are expanded to distribute evenly across the
    available vertical space, preventing large empty areas below.
    """
    if not decisions:
        return 0
    n = len(decisions)
    gap = 0.12
    # Adaptive card height: distribute available space evenly
    if avail_h and avail_h > 0:
        total_gaps = gap * (n - 1)
        h = max(0.65, (avail_h - total_gaps) / n)
    else:
        h = 0.65
    y = top
    for dec in decisions:
        # Background (respect border_radius token)
        bg = create_card_shape(slide, spec, left, y, width, h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = get_color(spec, 'primary_container')
        bg.line.fill.background()
        # Left accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(y), Inches(0.06), Inches(h)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = get_color(spec, 'primary')
        bar.line.fill.background()
        # Decision text
        tb = slide.shapes.add_textbox(
            Inches(left + 0.18), Inches(y + 0.08),
            Inches(width - 0.30), Inches(h - 0.16)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        # Read title from multiple possible field names (schema uses 'title')
        label = (dec.get('title') or dec.get('label')
                 or dec.get('decision') or '')
        desc = (dec.get('description') or dec.get('rationale') or '')
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        run.font.size = Pt(get_font_size(spec, 'body_text'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_primary_container')
        apply_font_to_run(run, spec)
        if desc:
            p2 = tf.add_paragraph()
            run2 = p2.add_run()
            run2.text = desc
            run2.font.size = Pt(get_font_size(spec, 'label'))
            run2.font.color.rgb = get_color(spec, 'on_surface_variant')
            apply_font_to_run(run2, spec)
        # Extra fields: budget, priority, timeline, status
        extras = []
        for fld, icon in [('budget', '💰'), ('priority', '🔹'),
                          ('timeline', '📅'), ('status', '📋')]:
            fval = dec.get(fld)
            if fval:
                extras.append(f"{icon} {fld.title()}: {fval}")
        if extras:
            p3 = tf.add_paragraph()
            run3 = p3.add_run()
            run3.text = '  |  '.join(extras)
            run3.font.size = Pt(get_font_size(spec, 'label'))
            run3.font.color.rgb = get_color(spec, 'on_surface_variant')
        y += h + gap
    return y - top


def render_risks_table(slide, risks, spec, grid, left, top, width):
    """Render risk items as colored table.

    Skips rendering entirely if all risk items lack meaningful data
    (e.g., only have label/description but no probability/impact/mitigation).
    Individual rows missing all detail fields are also skipped.
    """
    if not risks:
        return 0
    # Filter out incomplete risk items that would render as empty rows
    detail_fields = ('probability', 'impact', 'mitigation', 'risk', 'name')
    valid_risks = []
    for r in risks:
        has_detail = any(r.get(f) for f in detail_fields)
        if has_detail:
            valid_risks.append(r)
    if not valid_risks:
        # All risks lack table-renderable data; render as simple callouts instead
        y = top
        for r in risks:
            label = r.get('label', r.get('name', ''))
            desc = r.get('description', '')
            text = f"⚠️ {label}"
            if desc:
                text += f" — {desc}"
            if not text.strip() or text.strip() == '⚠️':
                continue
            h = 0.45
            bg = create_card_shape(slide, spec, left, y, width, h)
            bg.fill.solid()
            bg.fill.fore_color.rgb = get_color(spec, 'surface_variant')
            bg.line.fill.background()
            # Left accent bar (warning color)
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left), Inches(y), Inches(0.05), Inches(h)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = get_color(spec, 'warning')
            bar.line.fill.background()
            tb = slide.shapes.add_textbox(
                Inches(left + 0.15), Inches(y + 0.06),
                Inches(width - 0.25), Inches(h - 0.12)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(get_font_size(spec, 'body_text'))
            run.font.color.rgb = get_color(spec, 'on_surface')
            y += h + 0.08
        return y - top if y > top else 0
    risks = valid_risks
    headers = ['风险', '概率', '影响', '缓解措施']
    col_widths = [width * 0.28, width * 0.10, width * 0.10, width * 0.52]
    row_h = 0.38
    # Header row
    for j, hdr in enumerate(headers):
        cx = left + sum(col_widths[:j])
        tb = slide.shapes.add_textbox(
            Inches(cx), Inches(top), Inches(col_widths[j]), Inches(row_h)
        )
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = hdr
        run.font.size = Pt(get_font_size(spec, 'table_header'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_surface')
    # Header line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top + row_h - 0.02),
        Inches(width), Pt(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = get_color(spec, 'primary')
    line.line.fill.background()
    # Data rows
    prob_colors = {
        'H': 'error', 'M': 'tertiary', 'L': 'secondary',
        '高': 'error', '中': 'tertiary', '低': 'secondary'
    }
    for i, risk in enumerate(risks):
        ry = top + (i + 1) * row_h
        row_data = [
            '⚠️ ' + str(risk.get('risk', risk.get('name', ''))),
            risk.get('probability', ''),
            risk.get('impact', ''),
            risk.get('mitigation', '')
        ]
        # Zebra stripe
        if i % 2 == 1:
            stripe = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left), Inches(ry), Inches(width), Inches(row_h)
            )
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = get_color(spec, 'surface_variant')
            stripe.line.fill.background()
        for j, val in enumerate(row_data):
            cx = left + sum(col_widths[:j])
            tb = slide.shapes.add_textbox(
                Inches(cx), Inches(ry), Inches(col_widths[j]), Inches(row_h)
            )
            tf = tb.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(get_font_size(spec, 'table_cell'))
            # Color-code probability/impact
            if j in (1, 2) and val in prob_colors:
                run.font.color.rgb = get_color(spec, prob_colors[val])
                run.font.bold = True
            else:
                run.font.color.rgb = get_color(spec, 'on_surface')
    return (len(risks) + 1) * row_h + 0.1


def render_callouts(slide, callouts, spec, grid, left, top, width):
    """Render callout boxes."""
    if not callouts:
        return 0
    y = top
    for co in callouts:
        h = 0.55
        bg = create_card_shape(slide, spec, left, y, width, h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = get_color(spec, 'primary_container')
        bg.line.fill.background()
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(y), Inches(0.05), Inches(h)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = get_color(spec, 'primary')
        bar.line.fill.background()
        tb = slide.shapes.add_textbox(
            Inches(left + 0.15), Inches(y + 0.08),
            Inches(width - 0.25), Inches(h - 0.16)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        text = co if isinstance(co, str) else co.get('text', '')
        run.text = text
        run.font.size = Pt(get_font_size(spec, 'callout_text'))
        run.font.color.rgb = get_color(spec, 'on_primary_container')
        apply_font_to_run(run, spec)
        y += h + 0.08
    return y - top


def render_action_items(slide, items, spec, grid, left, top, width):
    """Render action items as a structured list."""
    if not items:
        return 0
    row_h = 0.40
    # Header
    headers = ['行动', '负责人', '截止日期']
    col_w = [width * 0.55, width * 0.2, width * 0.25]
    for j, hdr in enumerate(headers):
        cx = left + sum(col_w[:j])
        tb = slide.shapes.add_textbox(
            Inches(cx), Inches(top), Inches(col_w[j]), Inches(row_h)
        )
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = hdr
        run.font.size = Pt(get_font_size(spec, 'table_header'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_surface')
    for i, item in enumerate(items):
        ry = top + (i + 1) * row_h
        vals = [
            item.get('action', ''),
            item.get('owner', ''),
            item.get('deadline', '')
        ]
        for j, val in enumerate(vals):
            cx = left + sum(col_w[:j])
            tb = slide.shapes.add_textbox(
                Inches(cx), Inches(ry), Inches(col_w[j]), Inches(row_h)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(get_font_size(spec, 'table_cell'))
            run.font.color.rgb = get_color(spec, 'on_surface')
    return (len(items) + 1) * row_h + 0.1


def render_timeline_items(slide, items, spec, grid, left, top, width):
    """Render timeline as horizontal milestones."""
    if not items:
        return 0
    n = len(items)
    seg_w = width / max(n, 1)
    dot_r = 0.10
    line_y = top + dot_r
    # Connector line
    if n > 1:
        connector = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left + seg_w / 2), Inches(line_y - 0.015),
            Inches(width - seg_w), Inches(0.03)
        )
        connector.fill.solid()
        connector.fill.fore_color.rgb = get_color(spec, 'outline')
        connector.line.fill.background()
    status_colors = {
        'completed': 'secondary', 'done': 'secondary',
        'in_progress': 'primary', 'active': 'primary',
        'planned': 'outline', 'pending': 'outline',
    }
    for i, item in enumerate(items):
        cx = left + i * seg_w + seg_w / 2 - dot_r
        status = item.get('status', 'planned')
        color_token = status_colors.get(status, 'outline')
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(cx), Inches(top), Inches(dot_r * 2), Inches(dot_r * 2)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = get_color(spec, color_token)
        dot.line.fill.background()
        # Milestone text
        tb = slide.shapes.add_textbox(
            Inches(left + i * seg_w), Inches(top + dot_r * 2 + 0.08),
            Inches(seg_w), Inches(0.30)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = item.get('milestone', item.get('label', ''))
        run.font.size = Pt(get_font_size(spec, 'body_text'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_surface')
        # Date
        date_val = item.get('date', item.get('time', ''))
        if date_val:
            tb2 = slide.shapes.add_textbox(
                Inches(left + i * seg_w), Inches(top + dot_r * 2 + 0.38),
                Inches(seg_w), Inches(0.20)
            )
            tf2 = tb2.text_frame
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = date_val
            run2.font.size = Pt(get_font_size(spec, 'label'))
            run2.font.color.rgb = get_color(spec, 'on_surface_variant')
    return dot_r * 2 + 0.7


def render_table_data(slide, table_data, spec, grid, left, top, width):
    """Render table_data as Material-styled table.

    Supports flexible input formats:
    - dict: { headers: [...], rows: [[...],[...]] }
    - list: [ {label: 'Col1', values: [...]}, {label: 'Col2', values: [...]} ] (column-oriented)
    """
    if not table_data:
        return 0
    headers = []
    rows = []
    # Accept both dict and list formats
    if isinstance(table_data, dict):
        headers = table_data.get('headers', [])
        rows = table_data.get('rows', [])
    elif isinstance(table_data, list):
        # Convert column-oriented list into headers + rows
        cols = table_data
        headers = [c.get('label', '') for c in cols]
        n_rows = 0
        for c in cols:
            vals = c.get('values', [])
            if isinstance(vals, list):
                n_rows = max(n_rows, len(vals))
        for i in range(n_rows):
            row = []
            for c in cols:
                vals = c.get('values', [])
                row.append(vals[i] if i < len(vals) else '')
            rows.append(row)
    else:
        return 0
    if not headers and not rows:
        return 0
    n_cols = len(headers) if headers else (len(rows[0]) if rows else 0)
    col_w = width / max(n_cols, 1)
    row_h = 0.38
    # Headers
    if headers:
        for j, hdr in enumerate(headers):
            tb = slide.shapes.add_textbox(
                Inches(left + j * col_w), Inches(top),
                Inches(col_w), Inches(row_h)
            )
            tf = tb.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(hdr)
            run.font.size = Pt(get_font_size(spec, 'table_header'))
            run.font.bold = True
            run.font.color.rgb = get_color(spec, 'on_surface')
        # Header line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top + row_h - 0.02),
            Inches(width), Pt(2)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = get_color(spec, 'primary')
        line.line.fill.background()
    # Data rows
    for i, row in enumerate(rows):
        ry = top + (i + (1 if headers else 0)) * row_h
        if i % 2 == 1:
            stripe = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left), Inches(ry), Inches(width), Inches(row_h)
            )
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = get_color(spec, 'surface_variant')
            stripe.line.fill.background()
        for j, val in enumerate(row):
            if j >= n_cols:
                break
            tb = slide.shapes.add_textbox(
                Inches(left + j * col_w), Inches(ry),
                Inches(col_w), Inches(row_h)
            )
            tf = tb.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(get_font_size(spec, 'table_cell'))
            run.font.color.rgb = get_color(spec, 'on_surface')
    total_rows = len(rows) + (1 if headers else 0)
    return total_rows * row_h + 0.1


def render_bullets_component(slide, bullets, spec, grid, left, top, width,
                             accent_token='primary'):
    """Render components.bullets as a styled bullet list with colored dot markers.

    bullets is a list of {"text": "..."} dicts (schema format).
    """
    if not bullets:
        return 0
    y = top
    dot_indent = 0.20
    line_h = 0.45
    for item in bullets[:8]:
        text = item.get('text', '') if isinstance(item, dict) else str(item)
        if not text:
            continue
        _add_bullet_dot(slide, spec, left + 0.05, y + line_h / 2,
                        accent_token=accent_token, size=0.08)
        tb = slide.shapes.add_textbox(
            Inches(left + 0.05 + dot_indent), Inches(y),
            Inches(width - 0.10 - dot_indent), Inches(line_h)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _render_highlighted_text(p, text, spec, get_font_size(spec, 'bullet_text'))
        p.line_spacing = 1.5
        y += 0.48
    return y - top


def render_key_assertion(slide, assertion, spec, grid, left, top, width):
    """Render a key_assertion string as a prominent callout box."""
    if not assertion:
        return 0
    text = assertion if isinstance(assertion, str) else str(assertion)
    h = 0.70
    bg = create_card_shape(slide, spec, left, top, width, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = get_color(spec, 'primary_container')
    bg.line.fill.background()
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(0.06), Inches(h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = get_color(spec, 'primary')
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(
        Inches(left + 0.18), Inches(top + 0.10),
        Inches(width - 0.30), Inches(h - 0.20)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"💡 {text}"
    run.font.size = Pt(get_font_size(spec, 'body_text'))
    run.font.bold = True
    run.font.color.rgb = get_color(spec, 'on_primary_container')
    apply_font_to_run(run, spec)
    return h + 0.12


def render_components(slide, components, spec, grid, left, top, width, max_height):
    """Dispatch components rendering in priority order."""
    if not components:
        return 0
    cursor = 0
    priority_order = ['kpis', 'key_assertion', 'decisions', 'comparison_items',
                      'bullets', 'risks', 'callouts', 'action_items',
                      'timeline_items', 'table_data']
    renderers = {
        'kpis': render_kpis,
        'key_assertion': render_key_assertion,
        'comparison_items': render_comparison_items,
        'decisions': render_decisions,
        'bullets': render_bullets_component,
        'risks': render_risks_table,
        'callouts': render_callouts,
        'action_items': render_action_items,
        'timeline_items': render_timeline_items,
        'table_data': render_table_data,
    }
    for key in priority_order:
        data = components.get(key)
        if not data:
            continue
        if cursor > max_height - 0.5:
            # Overflow notice
            tb = slide.shapes.add_textbox(
                Inches(left), Inches(top + cursor),
                Inches(width), Inches(0.30)
            )
            p = tb.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = "… 详见 Speaker Notes"
            run.font.size = Pt(get_font_size(spec, 'label'))
            run.font.color.rgb = get_color(spec, 'on_surface_variant')
            run.font.italic = True
            break
        renderer = renderers.get(key)
        if renderer:
            used = renderer(slide, data, spec, grid, left, top + cursor, width)
            cursor += used
    return cursor


def add_shadow(shape, blur_pt=6, offset_pt=2):
    """Add Material elevation Level 1 shadow."""
    try:
        spPr = shape._element.spPr
        effectLst = spPr.makeelement(qn('a:effectLst'), {})
        outerShdw = effectLst.makeelement(qn('a:outerShdw'), {
            'blurRad': str(int(Pt(blur_pt))),
            'dist': str(int(Pt(offset_pt))),
            'dir': '5400000',
            'algn': 'bl',
            'rotWithShape': '0',
        })
        srgbClr = outerShdw.makeelement(qn('a:srgbClr'), {'val': '000000'})
        alpha = srgbClr.makeelement(qn('a:alpha'), {'val': '20000'})
        srgbClr.append(alpha)
        outerShdw.append(srgbClr)
        effectLst.append(outerShdw)
        spPr.append(effectLst)
    except Exception:
        pass  # Shadow is optional — don't fail the build


# ════════════════════════════════════════════════════════════════
# §6  Visual Renderers
# ════════════════════════════════════════════════════════════════

def get_chart_palette(spec: dict) -> list:
    """Get chart color palette from design_spec.

    Searches (in priority):
    1. color_system.chart_colors (recommended)
    2. component_library.chart_palette (legacy)
    3. tokens.accessibility.chart_color_palette
    4. Builds palette from accent_1..accent_4 + primary + secondary
    5. Fallback to a default colorblind-safe palette
    """
    # 1. color_system.chart_colors
    cc = spec.get('color_system', {}).get('chart_colors', [])
    if cc and isinstance(cc, list) and len(cc) >= 2:
        return cc
    # 2. component_library.chart_palette
    cp = spec.get('component_library', {}).get('chart_palette', [])
    if cp and isinstance(cp, list) and len(cp) >= 2:
        return cp
    # 3. tokens.accessibility.chart_color_palette
    ap = spec.get('tokens', {}).get('accessibility', {}).get('chart_color_palette', [])
    if ap and isinstance(ap, list) and len(ap) >= 2:
        return ap
    # 4. Build from accent tokens
    palette = []
    for token in ('primary', 'accent_1', 'accent_2', 'accent_3', 'accent_4', 'secondary'):
        hex_val = get_color_hex(spec, token)
        if hex_val and hex_val not in palette:
            palette.append(hex_val)
    if len(palette) >= 3:
        return palette
    # 5. Default colorblind-safe palette (Okabe-Ito inspired)
    return ['#0072B2', '#E69F00', '#009E73', '#CC79A7', '#56B4E9', '#D55E00', '#F0E442', '#000000']


def generate_chart_image(visual, spec, output_path, width_in, height_in, dpi=200):
    """Generate a chart image from chart_config using matplotlib.

    Supports: line_chart, bar_chart, bar_line_chart, composite_charts,
    stacked_bar, area_chart, and generic fallback.

    Returns the output file path on success, None on failure.
    """
    if not HAS_MATPLOTLIB:
        return None

    chart_type = visual.get('type', 'line_chart')
    pd = visual.get('placeholder_data', {})
    config = pd.get('chart_config', {})

    # Handle composite_charts (multiple sub-charts)
    if isinstance(config, dict) and not config.get('series') and not config.get('labels'):
        # Composite: config is a dict of sub-charts
        sub_charts = {k: v for k, v in config.items() if isinstance(v, dict) and v.get('series')}
        if not sub_charts:
            return None
        n_sub = len(sub_charts)
        fig, axes = plt.subplots(1, n_sub, figsize=(width_in, height_in), dpi=dpi)
        if n_sub == 1:
            axes = [axes]
        palette = get_chart_palette(spec)
        for idx, (sub_name, sub_config) in enumerate(sub_charts.items()):
            ax = axes[idx]
            labels = sub_config.get('labels', sub_config.get('x', []))
            series_list = sub_config.get('series', [])
            for si, s in enumerate(series_list):
                color = palette[si % len(palette)]
                data = s.get('data', [])
                ax.plot(labels[:len(data)], data, marker='o', color=color,
                        label=s.get('name', f'Series {si+1}'), linewidth=2, markersize=5)
            ax.set_title(sub_name.replace('_', ' ').title(),
                         fontsize=10, fontweight='bold', pad=8)
            ax.legend(fontsize=7, framealpha=0.8)
            ax.grid(True, alpha=0.3, linestyle='--')
            ax.tick_params(labelsize=8)
        plt.tight_layout(pad=1.5)
        fig.savefig(output_path, dpi=dpi, bbox_inches='tight',
                    facecolor='white', edgecolor='none')
        plt.close(fig)
        return output_path

    # Single chart
    labels = config.get('labels', config.get('x', []))
    series_list = config.get('series', [])
    if not labels or not series_list:
        return None

    palette = get_chart_palette(spec)
    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=dpi)

    # Determine effective chart type from visual.type
    effective_type = chart_type.lower().replace('-', '_').replace(' ', '_')

    if effective_type in ('bar_chart', 'column_chart', 'stacked_bar'):
        import numpy as np
        x_pos = np.arange(len(labels))
        bar_width = 0.8 / max(len(series_list), 1)
        for si, s in enumerate(series_list):
            color = palette[si % len(palette)]
            data = s.get('data', [])
            offset = (si - len(series_list) / 2 + 0.5) * bar_width
            bars = ax.bar(x_pos + offset, data[:len(labels)], bar_width,
                          label=s.get('name', f'Series {si+1}'), color=color,
                          edgecolor='white', linewidth=0.5)
            # Add value labels on bars
            for bar_item in bars:
                h = bar_item.get_height()
                if h > 0:
                    ax.text(bar_item.get_x() + bar_item.get_width() / 2., h,
                            f'{h:g}', ha='center', va='bottom', fontsize=7, fontweight='bold')
        ax.set_xticks(x_pos)
        ax.set_xticklabels([str(int(l)) if isinstance(l, (int, float)) and l == int(l) else str(l) for l in labels], fontsize=9)

    elif effective_type in ('bar_line_chart', 'combo_chart'):
        import numpy as np
        x_pos = np.arange(len(labels))
        # First series as bars, rest as lines
        if len(series_list) >= 1:
            s0 = series_list[0]
            ax.bar(x_pos, s0.get('data', [])[:len(labels)], 0.6,
                   label=s0.get('name', 'Series 1'), color=palette[0],
                   alpha=0.8, edgecolor='white', linewidth=0.5)
        ax2 = ax.twinx() if len(series_list) > 1 else None
        for si, s in enumerate(series_list[1:], 1):
            color = palette[si % len(palette)]
            data = s.get('data', [])
            target_ax = ax2 if ax2 else ax
            target_ax.plot(x_pos[:len(data)], data, marker='s', color=color,
                           label=s.get('name', f'Series {si+1}'), linewidth=2, markersize=6)
        ax.set_xticks(x_pos)
        ax.set_xticklabels([str(int(l)) if isinstance(l, (int, float)) and l == int(l) else str(l) for l in labels], fontsize=9)
        if ax2:
            ax2.legend(loc='upper left', fontsize=7, framealpha=0.8)

    elif effective_type in ('area_chart',):
        for si, s in enumerate(series_list):
            color = palette[si % len(palette)]
            data = s.get('data', [])
            ax.fill_between(labels[:len(data)], data, alpha=0.3, color=color)
            ax.plot(labels[:len(data)], data, color=color,
                    label=s.get('name', f'Series {si+1}'), linewidth=2)

    else:
        # Default: line chart
        for si, s in enumerate(series_list):
            color = palette[si % len(palette)]
            data = s.get('data', [])
            ax.plot(labels[:len(data)], data, marker='o', color=color,
                    label=s.get('name', f'Series {si+1}'), linewidth=2.5,
                    markersize=6, markerfacecolor='white', markeredgewidth=2,
                    markeredgecolor=color)
            # Add direct data labels on last point
            if data:
                ax.annotate(f'{data[-1]:g}',
                           xy=(labels[len(data)-1], data[-1]),
                           xytext=(5, 5), textcoords='offset points',
                           fontsize=8, fontweight='bold', color=color)

    # Style the chart
    ax.legend(fontsize=8, framealpha=0.9, edgecolor='#E0E0E0')
    ax.grid(True, alpha=0.2, linestyle='--', color='#9E9E9E')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#E0E0E0')
    ax.spines['bottom'].set_color('#E0E0E0')
    ax.tick_params(labelsize=9, colors='#616161')

    # Add source annotation if available
    source = config.get('source', '')
    if source:
        fig.text(0.99, 0.01, f'Source: {Path(source).stem}',
                 ha='right', va='bottom', fontsize=6, color='#9E9E9E', style='italic')

    plt.tight_layout(pad=1.2)
    fig.savefig(output_path, dpi=dpi, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close(fig)
    return output_path


def render_visual(slide, visual, spec, grid, left, top, width, height,
                  session_dir=''):
    """Route visual rendering based on placeholder_data type.

    Priority:
    1. Pre-rendered image (rendered_png_path / rendered_svg_path) — use if file exists
    2. Matplotlib chart generation from chart_config data — if matplotlib available
    3. Fallback to data table rendering for chart_config
    4. Mermaid placeholder for mermaid_code
    5. Generic visual placeholder
    """
    if not visual or visual.get('type') in (None, 'none'):
        return

    # 1. Check for pre-rendered images
    for img_key in ('rendered_png_path', 'rendered_svg_path'):
        img_path = visual.get(img_key, '')
        if img_path:
            # Try both relative to session_dir and as-is
            candidates = [img_path]
            if session_dir:
                candidates.insert(0, os.path.join(session_dir, img_path))
            for candidate in candidates:
                if os.path.isfile(candidate):
                    try:
                        pic = slide.shapes.add_picture(
                            candidate,
                            Inches(left), Inches(top),
                            Inches(width), Inches(height)
                        )
                        return
                    except Exception:
                        pass

    pd = visual.get('placeholder_data', {})

    # 2. Matplotlib chart generation from chart_config
    if pd.get('chart_config') and HAS_MATPLOTLIB:
        try:
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                tmp_path = tmp.name
            result = generate_chart_image(visual, spec, tmp_path, width, height)
            if result and os.path.isfile(tmp_path):
                pic = slide.shapes.add_picture(
                    tmp_path,
                    Inches(left), Inches(top),
                    Inches(width), Inches(height)
                )
                try:
                    os.unlink(tmp_path)
                except OSError:
                    pass
                return
        except Exception:
            pass  # Fall through to table rendering

    # 3. Fallback: chart_config as data table
    if pd.get('chart_config'):
        config = pd.get('chart_config', {})
        labels = config.get('labels', config.get('x', []))
        series = config.get('series', [])
        # Empty-chart guard: if no labels or no series, render an informative insight + placeholder
        if not labels or not series:
            cursor = top
            if visual.get('title'):
                strip_h = 0.45
                strip = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(left), Inches(cursor), Inches(width), Inches(strip_h)
                )
                strip.fill.solid()
                strip.fill.fore_color.rgb = get_color(spec, 'primary')
                strip.line.fill.background()
                tb = slide.shapes.add_textbox(
                    Inches(left + 0.12), Inches(cursor + 0.06),
                    Inches(width - 0.24), Inches(strip_h - 0.12)
                )
                tf = tb.text_frame
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = f"💡 暂无数据：{visual.get('title')}"
                run.font.size = Pt(get_font_size(spec, 'body_text'))
                run.font.bold = True
                run.font.color.rgb = get_color(spec, 'on_primary')
                apply_font_to_run(run, spec)
                cursor += strip_h + 0.10
            # Render placeholder card below the strip (or as sole placeholder)
            render_visual_placeholder(slide, visual, spec, left, cursor, width, max(1.0, height - (cursor - top)))
        else:
            render_chart_table(slide, visual, spec, left, top, width, height)
    elif pd.get('mermaid_code'):
        render_mermaid_placeholder(slide, visual, spec, left, top, width, height)
    else:
        render_visual_placeholder(slide, visual, spec, left, top, width, height)


def render_chart_table(slide, visual, spec, left, top, width, height):
    """Render chart_config as Material-styled data table.
    
    Row height scales to fill the available height so tables don't float
    in empty space.
    """
    config = visual.get('placeholder_data', {}).get('chart_config', {})
    labels = config.get('labels', [])
    series = config.get('series', [])
    if not labels:
        return

    n_cols = len(labels)
    col_w = width / max(n_cols, 1)
    # Adaptive row height: fill available height, clamped to reasonable range
    n_rows = len(series) + 1  # +1 for header
    title_offset = 0.35 if visual.get('title') else 0
    usable_h = height - title_offset
    row_h = max(0.42, min(1.2, usable_h / max(n_rows, 1)))

    # Title
    if visual.get('title'):
        tb = slide.shapes.add_textbox(
            Inches(left), Inches(top - 0.30),
            Inches(width), Inches(0.28)
        )
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = visual['title']
        run.font.size = Pt(get_font_size(spec, 'table_header'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_surface')

    # Header row
    for j, label in enumerate(labels):
        tb = slide.shapes.add_textbox(
            Inches(left + j * col_w), Inches(top),
            Inches(col_w), Inches(row_h)
        )
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = str(label)
        run.font.size = Pt(get_font_size(spec, 'table_header'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_surface')

    # Header separator
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top + row_h - 0.02),
        Inches(width), Pt(2)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = get_color(spec, 'primary')
    sep.line.fill.background()

    # Data rows
    for r, s in enumerate(series):
        ry = top + (r + 1) * row_h
        if r % 2 == 1:
            stripe = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left), Inches(ry), Inches(width), Inches(row_h)
            )
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = get_color(spec, 'surface_variant')
            stripe.line.fill.background()
        data = s.get('data', [])
        for j, val in enumerate(data):
            if j >= n_cols:
                break
            tb = slide.shapes.add_textbox(
                Inches(left + j * col_w), Inches(ry),
                Inches(col_w), Inches(row_h)
            )
            tf = tb.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(get_font_size(spec, 'table_cell'))
            run.font.color.rgb = get_color(spec, 'on_surface')
            if isinstance(val, (int, float)):
                p.alignment = PP_ALIGN.RIGHT


def render_mermaid_placeholder(slide, visual, spec, left, top, width, height):
    """Render mermaid code as simplified flowchart shapes.

    Parses basic mermaid graph syntax to extract node names and renders them
    as connected rounded-rectangle shapes with arrows.
    Falls back to styled text preview if parsing fails.
    """
    mermaid = visual.get('placeholder_data', {}).get('mermaid_code', '')
    if not mermaid.strip():
        return

    # Try to parse mermaid nodes from common patterns:
    # graph LR; A-->B-->C  or  graph TD\n  A[Label]-->B[Label]
    import re as _re
    lines = mermaid.strip().replace(';', '\n').split('\n')
    nodes = []
    seen = set()
    node_label_re = _re.compile(r'([A-Za-z0-9_]+)(?:\[([^\]]+)\])?')

    for line in lines:
        line = line.strip()
        if line.startswith('graph') or line.startswith('flowchart') or not line:
            continue
        # Split by --> or --- or ==> or -.- etc.
        parts = _re.split(r'-->|==>|---|-.->?|-\.->?|-->', line)
        for part in parts:
            part = part.strip()
            if not part:
                continue
            m = node_label_re.search(part)
            if m:
                node_id = m.group(1)
                node_label = m.group(2) or node_id
                if node_id not in seen:
                    seen.add(node_id)
                    nodes.append(node_label)

    if len(nodes) < 2:
        # Fallback: render as styled text preview
        _render_mermaid_text_fallback(slide, visual, spec, left, top, width, height)
        return

    # Render as horizontal flow of rounded rectangles with arrows
    n = len(nodes)
    max_nodes = min(n, 6)  # limit to 6 nodes
    arrow_w = 0.50
    total_arrows = max_nodes - 1
    node_w = (width - total_arrows * arrow_w - 0.3) / max_nodes
    node_h = min(0.70, height * 0.5)
    node_top = top + (height - node_h) / 2
    palette = get_chart_palette(spec)

    for i in range(max_nodes):
        nx = left + 0.15 + i * (node_w + arrow_w)
        # Node shape
        node_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(nx), Inches(node_top),
            Inches(node_w), Inches(node_h)
        )
        node_shape.fill.solid()
        node_color = palette[i % len(palette)]
        node_shape.fill.fore_color.rgb = hex_to_rgb(node_color)
        node_shape.line.color.rgb = hex_to_rgb(node_color)
        node_shape.line.width = Pt(1.5)
        add_shadow(node_shape)
        # Node label
        tf = node_shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = nodes[i]
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        apply_font_to_run(run, spec)

        # Arrow between nodes
        if i < max_nodes - 1:
            arrow_x = nx + node_w + 0.05
            arrow_y = node_top + node_h / 2
            # Arrow line
            arrow_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(arrow_x), Inches(arrow_y - 0.015),
                Inches(arrow_w - 0.15), Inches(0.03)
            )
            arrow_line.fill.solid()
            arrow_line.fill.fore_color.rgb = get_color(spec, 'on_surface_variant')
            arrow_line.line.fill.background()
            # Arrowhead (triangle)
            arrow_head = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE,
                Inches(arrow_x + arrow_w - 0.20), Inches(arrow_y - 0.06),
                Inches(0.12), Inches(0.12)
            )
            arrow_head.fill.solid()
            arrow_head.fill.fore_color.rgb = get_color(spec, 'on_surface_variant')
            arrow_head.line.fill.background()
            # Rotate triangle 90° to point right
            arrow_head.rotation = 90.0

    # Title above flow
    if visual.get('title'):
        tb = slide.shapes.add_textbox(
            Inches(left + 0.15), Inches(top + 0.08),
            Inches(width - 0.3), Inches(0.28)
        )
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"📊 {visual['title']}"
        run.font.size = Pt(get_font_size(spec, 'table_header'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_surface')
        apply_font_to_run(run, spec)


def _render_mermaid_text_fallback(slide, visual, spec, left, top, width, height):
    """Fallback: render mermaid code as styled text preview."""
    mermaid = visual.get('placeholder_data', {}).get('mermaid_code', '')
    card = create_card_shape(slide, spec, left, top, width, min(height, 3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = get_color(spec, 'surface_variant')
    card.line.color.rgb = get_color(spec, 'outline')
    card.line.width = Pt(1)
    if visual.get('title'):
        tb = slide.shapes.add_textbox(
            Inches(left + 0.2), Inches(top + 0.12),
            Inches(width - 0.4), Inches(0.28)
        )
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"📊 {visual['title']}"
        run.font.size = Pt(get_font_size(spec, 'table_header'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_surface')
        apply_font_to_run(run, spec)
    preview = '\n'.join(mermaid.strip().split('\n')[:8])
    if len(mermaid.strip().split('\n')) > 8:
        preview += '\n  ...'
    tb = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.50),
        Inches(width - 0.4), Inches(min(height - 0.7, 2.5))
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.font.size = Pt(11)
    run.font.color.rgb = get_color(spec, 'outline')
    apply_font_to_run(run, spec)
    run.text = preview
    run.font.size = Pt(11)
    run.font.color.rgb = get_color(spec, 'outline')


def render_visual_placeholder(slide, visual, spec, left, top, width, height):
    """Generic visual placeholder box."""
    card = create_card_shape(slide, spec, left, top, width, min(height, 2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = get_color(spec, 'surface_variant')
    card.line.color.rgb = get_color(spec, 'outline')
    card.line.width = Pt(1)
    label = visual.get('title', visual.get('type', 'Visual'))
    reqs = visual.get('content_requirements', [])
    text = f"[{label}]"
    if reqs:
        text += '\n' + '\n'.join(f"  • {r}" for r in reqs[:3])
    tb = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.3),
        Inches(width - 0.4), Inches(min(height - 0.6, 2.0))
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(get_font_size(spec, 'table_header'))
    run.font.color.rgb = get_color(spec, 'outline')
    apply_font_to_run(run, spec)


# ════════════════════════════════════════════════════════════════
# §7  Per-Slide-Type Renderers
# ════════════════════════════════════════════════════════════════

def render_slide_title(slide, sd, spec, grid, **ctx):
    """Title slide — full color bg, centered title, KPIs, author/date.

    Implements improved dynamic title height calculation for CJK/long titles
    and applies configured fonts.
    """
    title_text = sd.get('_deck_title', '') or sd.get('title', '')
    # Content layout anchors
    content_top = grid.slide_h * 0.30
    margin = grid.margin_h + 1.0
    cw = grid.slide_w - 2 * margin

    # Helper: detect CJK characters
    def is_cjk(ch):
        return any([
            '\u4e00' <= ch <= '\u9fff',  # CJK Unified Ideographs
            '\u3400' <= ch <= '\u4dbf',  # CJK Unified Ideographs Extension A
            '\u3000' <= ch <= '\u303f'   # punctuation
        ])

    # Helper to estimate text height more robustly (handles mixed CJK/Latin)
    def estimate_text_height(text, role, width_in):
        font_pt = get_font_size(spec, role)
        # per-char width multipliers: CJK nearly full em, Latin approx half-em
        total_w = 0.0
        for ch in text:
            if is_cjk(ch):
                total_w += (font_pt / 72.0) * 0.95
            elif ch.isspace():
                total_w += (font_pt / 72.0) * 0.35
            else:
                total_w += (font_pt / 72.0) * 0.55
        # prevent division by zero
        chars_per_line = max(1, int(width_in / max(0.0001, (total_w / max(1, len(text))))))
        # safe fallback: if text short, treat as single line
        if len(text) <= chars_per_line:
            lines = 1
        else:
            # approximate number of lines by dividing total width by available width
            lines = max(1, int((total_w / width_in) + 0.999))
        line_h = font_pt / 72.0 + 0.08
        return lines * line_h, font_pt

    # Estimate title and subtitle heights
    title_h, title_pt = estimate_text_height(title_text, 'display_large', cw)
    subtitles = sd.get('content', [])
    subtitle_line_h = get_font_size(spec, 'slide_subtitle') / 72.0 + 0.06
    subtitle_h = len(subtitles) * subtitle_line_h

    # KPI reservation
    comps = sd.get('components', {})
    kpi_h = 0
    if comps.get('kpis'):
        kpi_h = 0.85 + 0.15

    # Available height we can use for title+subtitle+kpis (conservative)
    available = max(1.0, grid.slide_h * 0.55 - content_top - 0.3)

    # If content overflows, reduce title size more aggressively; allow down to 16pt
    adjusted_pt = title_pt
    adjusted_title_h = title_h
    # First shrink to fit vertical available space
    while adjusted_title_h + subtitle_h + kpi_h > available and adjusted_pt > 16:
        adjusted_pt -= 2
        # recompute total width using new pt
        total_w = 0.0
        for ch in title_text:
            if is_cjk(ch):
                total_w += (adjusted_pt / 72.0) * 0.95
            elif ch.isspace():
                total_w += (adjusted_pt / 72.0) * 0.35
            else:
                total_w += (adjusted_pt / 72.0) * 0.55
        adjusted_lines = max(1, int((total_w / cw) + 0.999))
        adjusted_title_h = adjusted_lines * (adjusted_pt / 72.0 + 0.08)

    # Enforce a maximum of 2 lines for cover title; if still more, shrink more aggressively to fit 2 lines
    max_title_lines = 2
    # compute lines with current adjusted_pt
    def compute_lines_for_pt(pt):
        total_w = 0.0
        for ch in title_text:
            if is_cjk(ch):
                total_w += (pt / 72.0) * 0.95
            elif ch.isspace():
                total_w += (pt / 72.0) * 0.35
            else:
                total_w += (pt / 72.0) * 0.55
        return max(1, int((total_w / cw) + 0.999))

    lines_now = compute_lines_for_pt(adjusted_pt)
    while lines_now > max_title_lines and adjusted_pt > 14:
        adjusted_pt -= 1
        lines_now = compute_lines_for_pt(adjusted_pt)
    # If still too many lines at minimum font, insert a manual break near midpoint
    if lines_now > max_title_lines:
        mid = len(title_text) // 2
        # prefer splitting at nearest space around midpoint
        left = title_text.rfind(' ', 0, mid)
        right = title_text.find(' ', mid)
        if left != -1:
            split_pos = left
        elif right != -1:
            split_pos = right
        else:
            split_pos = mid
        title_text = title_text[:split_pos].strip() + '\n' + title_text[split_pos:].strip()
        # recompute height with the manual break
        lines_now = compute_lines_for_pt(adjusted_pt)
        adjusted_title_h = lines_now * (adjusted_pt / 72.0 + 0.08)

    # Extra conservative caps for CJK-heavy long titles (force smaller font to avoid wrapping unpredictably)
    if any(is_cjk(ch) for ch in title_text):
        cjk_count = sum(1 for ch in title_text if is_cjk(ch))
        txt_len = len(title_text.replace(' ',''))
        if txt_len > 24:
            cap = 24
        elif txt_len > 16:
            cap = 28
        else:
            cap = 32
        if adjusted_pt > cap:
            adjusted_pt = cap
            # recompute lines and adjusted_title_h with new pt
            total_w = 0.0
            for ch in title_text:
                if is_cjk(ch):
                    total_w += (adjusted_pt / 72.0) * 0.95
                elif ch.isspace():
                    total_w += (adjusted_pt / 72.0) * 0.35
                else:
                    total_w += (adjusted_pt / 72.0) * 0.55
            adjusted_lines = max(1, int((total_w / cw) + 0.999))
            adjusted_title_h = adjusted_lines * (adjusted_pt / 72.0 + 0.08)
            # If still exceeds 2 lines, try inserting break at punctuation or midpoint
            if adjusted_lines > max_title_lines:
                for sep in ['：', '，', ':', ',']:
                    if sep in title_text:
                        pos = title_text.find(sep)
                        title_text = title_text[:pos+1] + '\n' + title_text[pos+1:]
                        break
                else:
                    mid = len(title_text) // 2
                    title_text = title_text[:mid].strip() + '\n' + title_text[mid:].strip()
                # recompute with manual break
                adjusted_lines = compute_lines_for_pt(adjusted_pt)
                adjusted_title_h = adjusted_lines * (adjusted_pt / 72.0 + 0.08)

    # Add safety padding to title height (avoid tight fits)
    safety_padding = 0.30
    final_title_h = adjusted_title_h + safety_padding

    # Title textbox with dynamic height
    tb = slide.shapes.add_textbox(
        Inches(margin), Inches(content_top),
        Inches(cw), Inches(final_title_h)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(adjusted_pt)
    run.font.bold = True
    run.font.color.rgb = get_color(spec, 'on_primary')
    apply_font_to_run(run, spec)

    # Subtitle(s) positioned with larger safety gap
    subtitle_top = content_top + final_title_h + 0.25
    for i, bullet in enumerate(subtitles):
        tb2 = slide.shapes.add_textbox(
            Inches(margin), Inches(subtitle_top + i * (subtitle_line_h + 0.02)),
            Inches(cw), Inches(subtitle_line_h)
        )
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = bullet
        run2.font.size = Pt(get_font_size(spec, 'slide_subtitle'))
        run2.font.color.rgb = get_color(spec, 'on_primary')
        apply_font_to_run(run2, spec)

    # KPIs from components (below subtitle area)
    if comps.get('kpis'):
        kpi_top = subtitle_top + subtitle_h + 0.12
        render_kpis(slide, comps['kpis'], spec, grid,
                    grid.margin_h, kpi_top, grid.usable_w)


def render_section_divider(slide, sd, spec, grid, **ctx):
    """Section divider — accent-colored bg, centered title, callouts, progress."""
    meta = sd.get('metadata', {})
    sec_idx = meta.get('section_index', ctx.get('section_index', 0))
    total_sec = meta.get('total_sections', ctx.get('total_sections', 6))
    accent_tk = ctx.get('accent_token', 'primary')
    # M5: Text color adapts to background brightness
    text_tk = _text_color_for_bg(spec, accent_tk)
    # Title centered — use headline_large (28pt) for section dividers
    tb = slide.shapes.add_textbox(
        Inches(grid.margin_h + 1), Inches(grid.slide_h * 0.28),
        Inches(grid.usable_w - 2), Inches(1.0)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = sd.get('title', '')
    run.font.size = Pt(get_font_size(spec, 'headline_large'))
    run.font.bold = True
    run.font.color.rgb = get_color(spec, text_tk)
    apply_font_to_run(run, spec)
    # M2/M6: Decorative line only when content or callouts exist
    has_content = bool(sd.get('content'))
    has_callouts = bool(sd.get('components', {}).get('callouts'))
    # 方案3: If neither content nor callouts, extract first sentence from speaker_notes
    _auto_subtitle = ''
    if not has_content and not has_callouts:
        notes = sd.get('speaker_notes', '')
        if notes:
            # Extract first sentence (split on Chinese/English period or semicolon)
            for _delim in ['。', '；', '. ', '; ', '\n']:
                if _delim in notes:
                    _auto_subtitle = notes[:notes.index(_delim)].strip()
                    break
            if not _auto_subtitle:
                _auto_subtitle = notes.strip()[:80]  # fallback: first 80 chars
    if has_content or has_callouts or _auto_subtitle:
        deco_w = 2.0
        deco_left = (grid.slide_w - deco_w) / 2
        deco_top = grid.slide_h * 0.40
        deco_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(deco_left), Inches(deco_top),
            Inches(deco_w), Inches(0.03)
        )
        deco_line.fill.solid()
        deco_line.fill.fore_color.rgb = get_color(spec, text_tk)
        deco_line.line.fill.background()
    # Subtitle / content
    subtitle_y = grid.slide_h * 0.45
    for i, bullet in enumerate(sd.get('content', [])):
        tb2 = slide.shapes.add_textbox(
            Inches(grid.margin_h + 1.5), Inches(subtitle_y + i * 0.40),
            Inches(grid.usable_w - 3), Inches(0.38)
        )
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = bullet
        run2.font.size = Pt(get_font_size(spec, 'slide_subtitle'))
        run2.font.color.rgb = get_color(spec, text_tk)
        subtitle_y += 0.40
    # M1: Render callout components (e.g., "一行结论" callout in S0/S2)
    callouts = sd.get('components', {}).get('callouts', [])
    if callouts:
        callout_y = subtitle_y + 0.10 if has_content else grid.slide_h * 0.45
        for co in callouts:
            text = co if isinstance(co, str) else co.get('text', '')
            label = co.get('label', '') if isinstance(co, dict) else ''
            if not text:
                continue
            # Callout card with semi-transparent background
            cw = min(grid.usable_w - 2, 8.0)
            cl = (grid.slide_w - cw) / 2
            ch = 0.50
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(cl), Inches(callout_y),
                Inches(cw), Inches(ch)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = get_color(spec, 'surface')
            card.line.fill.background()
            # Callout text
            tb3 = slide.shapes.add_textbox(
                Inches(cl + 0.20), Inches(callout_y + 0.08),
                Inches(cw - 0.40), Inches(ch - 0.16)
            )
            tf3 = tb3.text_frame
            tf3.word_wrap = True
            p3 = tf3.paragraphs[0]
            p3.alignment = PP_ALIGN.CENTER
            display = f"{label}: {text}" if label else text
            run3 = p3.add_run()
            run3.text = display
            run3.font.size = Pt(get_font_size(spec, 'body_text'))
            run3.font.bold = True
            run3.font.color.rgb = get_color(spec, 'on_surface')
            apply_font_to_run(run3, spec)
            callout_y += ch + 0.10
    # 方案3: Render auto-subtitle from speaker_notes when no other content
    if _auto_subtitle and not has_content and not has_callouts:
        auto_y = grid.slide_h * 0.45
        tb_auto = slide.shapes.add_textbox(
            Inches(grid.margin_h + 1.5), Inches(auto_y),
            Inches(grid.usable_w - 3), Inches(0.50)
        )
        tf_auto = tb_auto.text_frame
        tf_auto.word_wrap = True
        p_auto = tf_auto.paragraphs[0]
        p_auto.alignment = PP_ALIGN.CENTER
        run_auto = p_auto.add_run()
        run_auto.text = _auto_subtitle
        run_auto.font.size = Pt(get_font_size(spec, 'slide_subtitle'))
        run_auto.font.color.rgb = get_color(spec, text_tk)
        run_auto.font.italic = True
        apply_font_to_run(run_auto, spec)
    # Progress indicator
    prog_w = 3.0
    prog_left = (grid.slide_w - prog_w) / 2
    seg_w = prog_w / max(total_sec, 1)
    prog_top = grid.slide_h * 0.78
    for i in range(total_sec):
        seg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(prog_left + i * seg_w + 0.03), Inches(prog_top),
            Inches(seg_w - 0.06), Inches(0.08)
        )
        seg.fill.solid()
        if i <= sec_idx:
            seg.fill.fore_color.rgb = get_color(spec, text_tk)
        else:
            seg.fill.fore_color.rgb = get_color(spec, 'surface_variant')
        seg.line.fill.background()


def _content_zone(grid, bar_h, content_h_est=None, center=False):
    """Calculate content zone below title bar.

    Args:
        grid: GridSystem instance.
        bar_h: Height of the title bar (inches).
        content_h_est: Estimated total content height (inches).
            When provided together with *center=True*, the returned
            top will be offset so that content is vertically centered
            in the available space.
        center: If True and content_h_est is provided, center content.

    Returns:
        (top, avail_h) — top edge of content zone and available height.
    """
    lz_top = bar_h + 0.12
    lz_h = grid.slide_h - lz_top - 0.35  # leave room for bottom bar
    if center and content_h_est is not None and content_h_est < lz_h:
        offset = (lz_h - content_h_est) / 2.0
        lz_top += offset
        lz_h = content_h_est
    return lz_top, lz_h


def _estimate_fill_ratio(sd):
    """Estimate how much of the content zone a slide's content will fill.

    Returns a ratio (0.0-1.0). Values below 0.5 indicate sparse content
    that may benefit from vertical centering.
    """
    content_items = len(sd.get('content', []))
    comps = sd.get('components', {})
    comp_items = sum(
        len(v) if isinstance(v, list) else (1 if v else 0)
        for v in comps.values()
    )
    has_visual = (sd.get('visual') and
                  sd['visual'].get('type') not in (None, 'none'))
    # Very rough height estimate (each bullet ≈ 0.45", each comp item ≈ 0.55")
    est_h = content_items * 0.45 + comp_items * 0.55 + (2.5 if has_visual else 0)
    max_h = 6.0  # approximate full content zone height
    return min(est_h / max_h, 1.0)


def render_slide_bullets(slide, sd, spec, grid, **ctx):
    """Bullet-list slide with optional visual."""
    bar_h = ctx.get('bar_h', 0.55)
    # Auto-center for sparse content
    fill_ratio = _estimate_fill_ratio(sd)
    use_center = fill_ratio < 0.45
    content_bullets_raw = sd.get('content', [])
    comps = sd.get('components', {})
    n_bullets = min(len(content_bullets_raw), 8)
    est_h = n_bullets * 0.48 + 0.30
    remaining_comps = {k: v for k, v in comps.items() if k != 'bullets' and v}
    if remaining_comps:
        est_h += 1.5  # approximate component height
    top, avail_h = _content_zone(grid, bar_h,
                                 content_h_est=est_h if use_center else None,
                                 center=use_center)
    has_visual = (sd.get('visual') and
                  sd['visual'].get('type') not in (None, 'none'))
    has_comps = any(comps.get(k) for k in comps)

    if has_visual:
        text_left, text_w = grid.col_span(7, 0)
        vis_left, vis_w = grid.col_span(5, 7)
    else:
        text_left, text_w = grid.col_span(12, 0)
        vis_left = vis_w = 0

    # Resolve bullet text: prefer content[], fall back to components.bullets
    content_bullets = sd.get('content', [])
    if not content_bullets and comps.get('bullets'):
        content_bullets = [
            (b.get('text', '') if isinstance(b, dict) else str(b))
            for b in comps['bullets']
        ]

    # Bullets
    accent_tk = ctx.get('accent_token', 'primary')
    dot_indent = 0.20
    n_bullets = len(content_bullets[:8])
    # Vertically center small bullet lists when a visual is present for better harmony
    if has_visual and n_bullets <= 2:
        bullet_total_h = n_bullets * 0.48
        bullet_y = top + max(0, (avail_h - bullet_total_h) / 2)
    else:
        bullet_y = top + 0.15
    for i, bullet in enumerate(content_bullets[:8]):
        by = bullet_y + i * 0.48
        _add_bullet_dot(slide, spec, text_left + 0.1, by + 0.225,
                        accent_token=accent_tk)
        tb = slide.shapes.add_textbox(
            Inches(text_left + 0.1 + dot_indent), Inches(by),
            Inches(text_w - 0.2 - dot_indent), Inches(0.45)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _render_highlighted_text(p, bullet, spec, get_font_size(spec, 'bullet_text'))
        p.line_spacing = 1.5

    # Components below bullets (skip 'bullets' since already rendered above)
    remaining_comps = {k: v for k, v in comps.items() if k != 'bullets' and v}
    if remaining_comps:
        comp_top = bullet_y + len(content_bullets[:8]) * 0.48 + 0.15
        render_components(slide, remaining_comps, spec, grid,
                          text_left + 0.1, comp_top, text_w - 0.2,
                          avail_h - (comp_top - top))

    # Visual on right
    if has_visual:
        render_visual(slide, sd['visual'], spec, grid,
                      vis_left, top + 0.15, vis_w, avail_h - 0.3,
                      session_dir=ctx.get('session_dir', ''))


def render_slide_two_column(slide, sd, spec, grid, **ctx):
    """Two-column or comparison layout."""
    bar_h = ctx.get('bar_h', 0.55)
    top, avail_h = _content_zone(grid, bar_h)
    comps = sd.get('components', {})
    has_comps = any(comps.get(k) for k in comps if k != 'bullets')

    l_left, l_w = grid.col_span(6, 0)
    r_left, r_w = grid.col_span(6, 6)

    # Resolve bullet text
    content_bullets = sd.get('content', [])
    if not content_bullets and comps.get('bullets'):
        content_bullets = [
            (b.get('text', '') if isinstance(b, dict) else str(b))
            for b in comps['bullets']
        ]

    # Left: bullets
    accent_tk = ctx.get('accent_token', 'primary')
    dot_indent = 0.20
    n_bullets = len(content_bullets[:6])
    if sd.get('visual') and sd['visual'].get('type') not in (None, 'none') and n_bullets <= 2:
        bullet_total_h = n_bullets * 0.48
        bullet_y = top + max(0, (avail_h - bullet_total_h) / 2)
    else:
        bullet_y = top + 0.15
    for i, bullet in enumerate(content_bullets[:6]):
        by = bullet_y + i * 0.48
        _add_bullet_dot(slide, spec, l_left + 0.1, by + 0.225,
                        accent_token=accent_tk)
        tb = slide.shapes.add_textbox(
            Inches(l_left + 0.1 + dot_indent), Inches(by),
            Inches(l_w - 0.2 - dot_indent), Inches(0.45)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _render_highlighted_text(p, bullet, spec, get_font_size(spec, 'bullet_text'))
        p.line_spacing = 1.4

    # Right: visual or components
    if sd.get('visual') and sd['visual'].get('type') not in (None, 'none'):
        render_visual(slide, sd['visual'], spec, grid,
                      r_left, top + 0.15, r_w, avail_h - 0.3,
                      session_dir=ctx.get('session_dir', ''))
    elif has_comps:
        render_components(slide, comps, spec, grid,
                          r_left, top + 0.15, r_w, avail_h - 0.3)


def render_slide_comparison(slide, sd, spec, grid, **ctx):
    """Comparison — cards + optional visual hybrid layout.
    
    layout_variant=0: vertical cards (side-by-side or 2×2 grid)
    layout_variant=1: horizontal bars (full-width, stacked vertically)
    """
    bar_h = ctx.get('bar_h', 0.55)
    top, avail_h = _content_zone(grid, bar_h)
    comps = sd.get('components', {})
    fill_mode = get_content_fill(spec, sd.get('slide_type', 'comparison'))
    variant = ctx.get('layout_variant', 0)
    has_visual = (sd.get('visual') and
                  sd['visual'].get('type') not in (None, 'none') and
                  sd['visual'].get('placeholder_data', {}))

    if comps.get('comparison_items'):
        # Determine layout: hybrid (cards + chart) or full-width cards
        if has_visual:
            cards_left, cards_w = grid.col_span(7, 0)
            vis_left, vis_w = grid.col_span(5, 7)
        else:
            cards_left = grid.margin_h
            cards_w = grid.usable_w
            vis_left = vis_w = 0

        content = sd.get('content', [])
        n_bullets = min(len(content), 4)
        bullet_h = n_bullets * 0.40 if n_bullets else 0
        callout_h = 0.8 if comps.get('callouts') else 0
        if fill_mode == 'expand':
            cards_avail = avail_h - 0.15 - bullet_h - callout_h
        else:
            cards_avail = None

        # 方案2: Alternate between vertical cards and horizontal bars
        if variant == 1:
            h = render_comparison_horizontal(slide, comps['comparison_items'],
                                             spec, grid, cards_left, top + 0.15,
                                             cards_w, avail_h=cards_avail)
        else:
            h = render_comparison_items(slide, comps['comparison_items'], spec, grid,
                                        cards_left, top + 0.15, cards_w,
                                        avail_h=cards_avail)
        cursor = top + 0.15 + h
        # Content bullets below cards (complementary context)
        if content:
            accent_tk = ctx.get('accent_token', 'primary')
            dot_indent = 0.20
            bullet_spacing = 0.40
            for i, bullet in enumerate(content[:4]):
                by = cursor + i * bullet_spacing
                _add_bullet_dot(slide, spec, cards_left + 0.1, by + 0.175,
                                accent_token=accent_tk)
                tb = slide.shapes.add_textbox(
                    Inches(cards_left + 0.1 + dot_indent), Inches(by),
                    Inches(cards_w - 0.2 - dot_indent), Inches(0.35)
                )
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                _render_highlighted_text(p, bullet, spec, get_font_size(spec, 'bullet_text'))
            cursor += n_bullets * bullet_spacing
        # Callouts below
        if comps.get('callouts'):
            render_callouts(slide, comps['callouts'], spec, grid,
                            cards_left, cursor, cards_w)
        # Visual on right (hybrid layout)
        if has_visual and vis_w > 0:
            render_visual(slide, sd['visual'], spec, grid,
                          vis_left, top + 0.15, vis_w, avail_h - 0.3,
                          session_dir=ctx.get('session_dir', ''))
    else:
        # No comparison_items — check if visual-only or fallback to two-column
        if has_visual:
            render_visual(slide, sd['visual'], spec, grid,
                          grid.margin_h, top + 0.15, grid.usable_w, avail_h - 0.3,
                          session_dir=ctx.get('session_dir', ''))
        else:
            render_slide_two_column(slide, sd, spec, grid, **ctx)


def render_slide_data_heavy(slide, sd, spec, grid, **ctx):
    """Data-heavy: components first (KPIs, decisions, etc), then visual.
    
    layout_variant=0: stacked (components top, visual bottom) or 6+6 split
    layout_variant=1: mirrored split (visual LEFT, components RIGHT)
    Patch A: when ≤1 bullet + visual + no heavy components → insight strip
    """
    bar_h = ctx.get('bar_h', 0.55)
    top, avail_h = _content_zone(grid, bar_h)
    comps = sd.get('components', {})
    has_visual = (sd.get('visual') and
                  sd['visual'].get('type') not in (None, 'none'))

    # --- Patch A: Insight strip for sparse data-heavy slides ---
    bullet_items = comps.get('bullets', [])
    n_bullets = len(bullet_items) if isinstance(bullet_items, list) else 0
    non_bullet_comps = {k: v for k, v in comps.items()
                        if k != 'bullets' and v}
    if has_visual and n_bullets <= 1 and not non_bullet_comps:
        cursor = top + 0.10
        accent_tk = ctx.get('accent_token', 'primary')
        if n_bullets == 1:
            strip_h = 0.45
            strip = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(grid.margin_h), Inches(cursor),
                Inches(grid.usable_w), Inches(strip_h)
            )
            strip.fill.solid()
            strip.fill.fore_color.rgb = get_color(spec, accent_tk)
            strip.line.fill.background()
            tb = slide.shapes.add_textbox(
                Inches(grid.margin_h + 0.20), Inches(cursor + 0.06),
                Inches(grid.usable_w - 0.40), Inches(strip_h - 0.12)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"💡 {bullet_items[0]}"
            run.font.size = Pt(get_font_size(spec, 'body_text'))
            run.font.bold = True
            run.font.color.rgb = get_color(spec, 'on_primary')
            apply_font_to_run(run, spec)
            cursor += strip_h + 0.10
        vis_h = max(1.5, avail_h - (cursor - top) - 0.10)
        render_visual(slide, sd['visual'], spec, grid,
                      grid.margin_h, cursor, grid.usable_w, vis_h,
                      session_dir=ctx.get('session_dir', ''))
        return
    # --- End Patch A ---

    variant = ctx.get('layout_variant', 0)

    # M4: Count total component items for layout decision
    _comp_count = 0
    for _k, _v in comps.items():
        if isinstance(_v, list):
            _comp_count += len(_v)
        elif _v:
            _comp_count += 1

    _stacked = False
    _mirrored = False
    if variant == 1 and has_visual and any(comps.get(k) for k in comps):
        # 方案1: Mirrored split — visual LEFT (6 cols), components RIGHT (6 cols)
        vis_left, vis_w = grid.col_span(6, 0)
        comp_left, comp_w = grid.col_span(6, 6)
        _mirrored = True
    elif has_visual and any(comps.get(k) for k in comps) and _comp_count <= 2:
        # M4: Stacked layout when few components + visual
        comp_left = grid.margin_h
        comp_w = grid.usable_w
        vis_left = grid.margin_h
        vis_w = grid.usable_w
        _stacked = True
    elif has_visual and any(comps.get(k) for k in comps):
        comp_left, comp_w = grid.col_span(6, 0)
        vis_left, vis_w = grid.col_span(6, 6)
    elif has_visual:
        comp_left = comp_w = 0
        vis_left = grid.margin_h
        vis_w = grid.usable_w
    else:
        comp_left = grid.margin_h
        comp_w = grid.usable_w
        vis_left = vis_w = 0

    cursor = top + 0.10

    # Render all components
    if any(comps.get(k) for k in comps) and comp_w > 0:
        h = render_components(slide, comps, spec, grid,
                              comp_left, cursor, comp_w,
                              avail_h - 0.10)
        cursor += h

    # Visual
    if has_visual:
        if _stacked:
            # M4: Visual below components, full width
            vis_top = cursor + 0.10
            vis_h = max(1.5, avail_h - (cursor - top) - 0.10)
        elif _mirrored:
            # 方案1: Visual on LEFT side, same vertical zone
            vis_top = top + 0.10
            vis_h = avail_h - 0.10
        else:
            vis_top = top + 0.10
            vis_h = avail_h - 0.10
        render_visual(slide, sd['visual'], spec, grid,
                      vis_left, vis_top, vis_w, vis_h,
                      session_dir=ctx.get('session_dir', ''))


def render_slide_matrix(slide, sd, spec, grid, **ctx):
    """Matrix layout — for risk/evaluation grids."""
    bar_h = ctx.get('bar_h', 0.55)
    top, avail_h = _content_zone(grid, bar_h)
    comps = sd.get('components', {})
    cursor = top + 0.10

    # Bullets
    accent_tk = ctx.get('accent_token', 'primary')
    dot_indent = 0.20
    for i, bullet in enumerate(sd.get('content', [])[:4]):
        by = cursor + i * 0.42
        _add_bullet_dot(slide, spec, grid.margin_h + 0.1, by + 0.19,
                        accent_token=accent_tk)
        tb = slide.shapes.add_textbox(
            Inches(grid.margin_h + 0.1 + dot_indent), Inches(by),
            Inches(grid.usable_w - 0.2 - dot_indent), Inches(0.38)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _render_highlighted_text(p, bullet, spec, get_font_size(spec, 'bullet_text'))
    cursor += len(sd.get('content', [])[:4]) * 0.42 + 0.15

    # Components (risks, table_data, etc)
    if any(comps.get(k) for k in comps):
        render_components(slide, comps, spec, grid,
                          grid.margin_h, cursor, grid.usable_w,
                          avail_h - (cursor - top))
    elif sd.get('visual'):
        render_visual(slide, sd['visual'], spec, grid,
                      grid.margin_h, cursor, grid.usable_w,
                      avail_h - (cursor - top),
                      session_dir=ctx.get('session_dir', ''))


def render_slide_timeline(slide, sd, spec, grid, **ctx):
    """Timeline/gantt layout."""
    bar_h = ctx.get('bar_h', 0.55)
    top, avail_h = _content_zone(grid, bar_h)
    comps = sd.get('components', {})
    cursor = top + 0.15

    if comps.get('timeline_items'):
        h = render_timeline_items(slide, comps['timeline_items'], spec, grid,
                                  grid.margin_h, cursor, grid.usable_w)
        cursor += h
    else:
        # Render bullets as milestones
        accent_tk = ctx.get('accent_token', 'primary')
        dot_indent = 0.20
        for i, bullet in enumerate(sd.get('content', [])):
            by = cursor + i * 0.45
            _add_bullet_dot(slide, spec, grid.margin_h + 0.1, by + 0.20,
                            accent_token=accent_tk)
            tb = slide.shapes.add_textbox(
                Inches(grid.margin_h + 0.1 + dot_indent), Inches(by),
                Inches(grid.usable_w - 0.2 - dot_indent), Inches(0.40)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            _render_highlighted_text(p, bullet, spec, get_font_size(spec, 'bullet_text'))
            cursor += 0.45

    # Remaining components
    remaining = {k: v for k, v in comps.items() if k != 'timeline_items' and v}
    if remaining:
        render_components(slide, remaining, spec, grid,
                          grid.margin_h, cursor, grid.usable_w,
                          avail_h - (cursor - top))

    # Visual
    if sd.get('visual') and sd['visual'].get('type') not in (None, 'none'):
        render_visual(slide, sd['visual'], spec, grid,
                      grid.margin_h, cursor + 0.10, grid.usable_w,
                      avail_h - (cursor - top) - 0.2,
                      session_dir=ctx.get('session_dir', ''))


def render_slide_flowchart(slide, sd, spec, grid, **ctx):
    """Flowchart layout — bullets + full-width visual."""
    bar_h = ctx.get('bar_h', 0.55)
    top, avail_h = _content_zone(grid, bar_h)
    comps = sd.get('components', {})
    cursor = top + 0.10

    # Resolve bullet text
    content_bullets = sd.get('content', [])
    if not content_bullets and comps.get('bullets'):
        content_bullets = [
            (b.get('text', '') if isinstance(b, dict) else str(b))
            for b in comps['bullets']
        ]

    # Bullets
    accent_tk = ctx.get('accent_token', 'primary')
    dot_indent = 0.20
    for i, bullet in enumerate(content_bullets[:3]):
        by = cursor + i * 0.40
        _add_bullet_dot(slide, spec, grid.margin_h + 0.1, by + 0.19,
                        accent_token=accent_tk)
        tb = slide.shapes.add_textbox(
            Inches(grid.margin_h + 0.1 + dot_indent), Inches(by),
            Inches(grid.usable_w - 0.2 - dot_indent), Inches(0.38)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _render_highlighted_text(p, bullet, spec, get_font_size(spec, 'bullet_text'))
    cursor += min(len(content_bullets), 3) * 0.40 + 0.15

    # Components (skip bullets since already rendered)
    remaining_comps = {k: v for k, v in comps.items() if k != 'bullets' and v}
    if remaining_comps:
        h = render_components(slide, remaining_comps, spec, grid,
                              grid.margin_h, cursor, grid.usable_w,
                              avail_h - (cursor - top))
        cursor += h

    # Visual
    if sd.get('visual'):
        render_visual(slide, sd['visual'], spec, grid,
                      grid.margin_h, cursor, grid.usable_w,
                      avail_h - (cursor - top),
                      session_dir=ctx.get('session_dir', ''))


def render_slide_cta(slide, sd, spec, grid, **ctx):
    """Call-to-action — centered prominent card."""
    bar_h = ctx.get('bar_h', 0.55)
    top, avail_h = _content_zone(grid, bar_h)
    comps = sd.get('components', {})

    cl, cw = grid.col_span(8, 2)
    card_top = top + 0.40
    card_h = avail_h - 0.80

    # Card background
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(cl), Inches(card_top), Inches(cw), Inches(card_h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = get_color(spec, 'primary_container')
    card.line.fill.background()
    add_shadow(card)

    # Content
    cursor = card_top + 0.30
    for i, bullet in enumerate(sd.get('content', [])[:6]):
        tb = slide.shapes.add_textbox(
            Inches(cl + 0.35), Inches(cursor + i * 0.55),
            Inches(cw - 0.70), Inches(0.50)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = bullet
        run.font.size = Pt(get_font_size(spec, 'slide_subtitle'))
        run.font.bold = True
        run.font.color.rgb = get_color(spec, 'on_primary_container')
    cursor += len(sd.get('content', [])[:6]) * 0.55 + 0.15

    # Components inside card
    if any(comps.get(k) for k in comps):
        render_components(slide, comps, spec, grid,
                          cl + 0.35, cursor, cw - 0.70,
                          card_h - (cursor - card_top) - 0.20)


def render_slide_decision(slide, sd, spec, grid, **ctx):
    """Decision slide — left decisions + right visual/chart.
    
    Decision cards and the right-side visual/table are expanded to fill
    the full content zone height, preventing large whitespace areas.
    """
    bar_h = ctx.get('bar_h', 0.55)
    top, avail_h = _content_zone(grid, bar_h)
    comps = sd.get('components', {})
    fill_mode = get_content_fill(spec, sd.get('slide_type', 'decision'))

    l_left, l_w = grid.col_span(5, 0)
    r_left, r_w = grid.col_span(7, 5)

    # Left: decisions cards — expand to fill available height
    cursor = top + 0.15
    if comps.get('decisions'):
        # Count extra bullets to reserve space for them
        decision_labels = {(d.get('title') or d.get('label') or '').strip()
                           for d in comps['decisions']}
        extra_bullets = [b for b in sd.get('content', [])
                         if b.strip() not in decision_labels]
        bullet_reserve = len(extra_bullets[:3]) * 0.50
        if fill_mode == 'expand':
            decisions_avail = avail_h - 0.15 - bullet_reserve
        else:
            decisions_avail = None
        h = render_decisions(slide, comps['decisions'], spec, grid,
                             l_left, cursor, l_w, avail_h=decisions_avail)
        cursor += h
        # When structured decisions exist, render content bullets only if
        # they add information beyond what's in the decision cards.
        decision_labels = {(d.get('title') or d.get('label') or '').strip()
                           for d in comps['decisions']}
        extra_bullets = [b for b in sd.get('content', [])
                         if b.strip() not in decision_labels]
        accent_tk = ctx.get('accent_token', 'primary')
        dot_indent = 0.20
        for i, bullet in enumerate(extra_bullets[:3]):
            by = cursor + i * 0.45
            _add_bullet_dot(slide, spec, l_left + 0.1, by + 0.21,
                            accent_token=accent_tk)
            tb = slide.shapes.add_textbox(
                Inches(l_left + 0.1 + dot_indent), Inches(by),
                Inches(l_w - 0.2 - dot_indent), Inches(0.42)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            _render_highlighted_text(p, bullet, spec, get_font_size(spec, 'bullet_text'))
    else:
        # No structured decisions — render all content as bullets
        accent_tk = ctx.get('accent_token', 'primary')
        dot_indent = 0.20
        for i, bullet in enumerate(sd.get('content', [])[:5]):
            by = cursor + i * 0.45
            _add_bullet_dot(slide, spec, l_left + 0.1, by + 0.21,
                            accent_token=accent_tk)
            tb = slide.shapes.add_textbox(
                Inches(l_left + 0.1 + dot_indent), Inches(by),
                Inches(l_w - 0.2 - dot_indent), Inches(0.42)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            _render_highlighted_text(p, bullet, spec, get_font_size(spec, 'bullet_text'))

    # Right: visual or remaining components
    if sd.get('visual') and sd['visual'].get('type') not in (None, 'none'):
        render_visual(slide, sd['visual'], spec, grid,
                      r_left, top + 0.15, r_w, avail_h - 0.3,
                      session_dir=ctx.get('session_dir', ''))
    else:
        remaining = {k: v for k, v in comps.items() if k != 'decisions' and v}
        if remaining:
            render_components(slide, remaining, spec, grid,
                              r_left, top + 0.15, r_w, avail_h - 0.3)


# ════════════════════════════════════════════════════════════════
# §8  Render Dispatcher
# ════════════════════════════════════════════════════════════════

RENDERERS = {
    'title':          render_slide_title,
    'section_divider': render_section_divider,
    'bullet-list':    render_slide_bullets,
    'two-column':     render_slide_two_column,
    'comparison':     render_slide_comparison,
    'decision':       render_slide_decision,
    'data-heavy':     render_slide_data_heavy,
    'matrix':         render_slide_matrix,
    'timeline':       render_slide_timeline,
    'gantt':          render_slide_timeline,
    'flowchart':      render_slide_flowchart,
    'diagram':        render_slide_flowchart,
    'sequence':       render_slide_flowchart,
    'radar':          render_slide_data_heavy,
    'technical':      render_slide_two_column,
    'process':        render_slide_flowchart,
    'call_to_action': render_slide_cta,
    'table':          render_slide_data_heavy,
    'waterfall':      render_slide_data_heavy,
    'kpi_dashboard':  render_slide_data_heavy,
}


def _resolve_bg_image(spec: dict, slide_type: str, sd: dict) -> str:
    """Resolve background image path for a slide.

    Checks (in priority order):
    1. sd['background_image'] — per-slide override from slides_semantic.json
    2. slide_type_layouts[type]['background_image'] — type-level default from design_spec
    Returns empty string if none found.
    """
    # Per-slide override
    img = sd.get('background_image', '')
    if img:
        return img
    # Type-level default from design_spec
    stl = spec.get('slide_type_layouts', {})
    entry = stl.get(slide_type, stl.get('default', {}))
    return entry.get('background_image', '')


def render_slide(prs, sd, spec, grid, sections, slide_num, total_slides,
                 session_dir='', layout_variant=0):
    """Render a single slide with full styling."""
    stype = sd.get('slide_type', 'bullet-list')
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Determine section context
    slide_id = sd.get('slide_id', slide_num)
    section = find_section_for_slide(slide_id, sections)
    sec_id = section.get('id', 'A')
    sec_title = section.get('title', '')
    sec_accent_token = get_section_accent(spec, sec_id)
    sec_index = next((i for i, s in enumerate(sections) if s.get('id') == sec_id), 0)

    # Background — support background_image for cover/divider slides
    bg_token = get_bg_token(spec, stype)
    # M5: Section dividers use section accent color for visual distinction
    if stype == 'section_divider':
        bg_token = sec_accent_token
    bg_image_path = _resolve_bg_image(spec, stype, sd)
    _use_bg_image = False
    if bg_image_path and os.path.isfile(bg_image_path):
        # C2: Verify file is actually a raster image (not SVG masquerading as .jpg)
        try:
            with open(bg_image_path, 'rb') as _imgf:
                _hdr = _imgf.read(16)
                if _hdr.startswith(b'<') or _hdr.startswith(b'<?xml') or b'<svg' in _hdr:
                    _use_bg_image = False  # SVG — skip
                else:
                    _use_bg_image = True
        except Exception:
            _use_bg_image = False
    if _use_bg_image:
        # Add full-bleed background image
        pic = slide.shapes.add_picture(
            bg_image_path, Emu(0), Emu(0),
            Inches(grid.slide_w), Inches(grid.slide_h)
        )
        # Send to back so text renders on top
        sp_tree = slide.shapes._spTree
        sp_tree.remove(pic._element)
        sp_tree.insert(2, pic._element)
        # Add semi-transparent overlay for text legibility
        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(0), Emu(0),
            Inches(grid.slide_w), Inches(grid.slide_h)
        )
        fill = overlay.fill
        fill.solid()
        fill.fore_color.rgb = get_color(spec, bg_token)
        # Set overlay transparency (40% opaque = 60% see-through)
        from pptx.oxml.ns import qn
        solid_fill_elem = fill._fill.find(qn('a:solidFill'))
        if solid_fill_elem is not None:
            srgb = solid_fill_elem.find(qn('a:srgbClr'))
            if srgb is None:
                srgb = solid_fill_elem[0] if len(solid_fill_elem) else None
            if srgb is not None:
                from lxml import etree
                alpha = etree.SubElement(srgb, qn('a:alpha'))
                alpha.set('val', '20000')  # 20% opacity — more see-through
        overlay.line.fill.background()
        sp_tree = slide.shapes._spTree
        sp_tree.remove(overlay._element)
        sp_tree.insert(3, overlay._element)
    else:
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = get_color(spec, bg_token)

    # Title bar
    tb_mode = get_title_bar_mode(spec, stype)
    bar_h = 0
    if tb_mode != 'none':
        section_label = f"Section {sec_id} · {sec_title}" if sec_title else ''
        bar_h = render_title_bar(slide, spec, grid,
                                 sd.get('title', ''), slide_num, total_slides,
                                 section_label=section_label,
                                 accent_color_token=sec_accent_token,
                                 mode=tb_mode)

    # Content (dispatch by slide_type)
    renderer = RENDERERS.get(stype, render_slide_bullets)
    renderer(slide, sd, spec, grid, bar_h=bar_h,
             section_index=sec_index, total_sections=len(sections),
             session_dir=session_dir,
             accent_token=sec_accent_token,
             layout_variant=layout_variant)

    # Bottom bar (skip for title and section_divider)
    if stype not in ('title', 'section_divider'):
        render_bottom_bar(slide, spec, grid, sec_title, slide_num, total_slides,
                          section_index=sec_index, total_sections=len(sections),
                          accent_token=sec_accent_token)

    # Speaker notes
    render_speaker_notes(slide, sd.get('speaker_notes'))


# ════════════════════════════════════════════════════════════════
# §9  Main Entry Point
# ════════════════════════════════════════════════════════════════

def generate_pptx(semantic_path, design_spec_path, output_path):
    """Main entry: load inputs, render all slides, save PPTX."""
    with open(semantic_path, encoding='utf-8') as f:
        semantic = json.load(f)
    with open(design_spec_path, encoding='utf-8') as f:
        spec = json.load(f)

    grid = GridSystem(spec)

    prs = Presentation()
    prs.slide_width = Inches(grid.slide_w)
    prs.slide_height = Inches(grid.slide_h)

    sections = semantic.get('sections', [])
    slides_data = semantic.get('slides', [])
    total = len(slides_data)
    deck_title = semantic.get('title', '')
    # Strip trailing duration suffix like "（30 分钟）" from deck title
    deck_title = re.sub(r'[（(]\s*\d+\s*分钟\s*[）)]\s*$', '', deck_title).strip()

    # Derive session_dir from semantic_path for resolving relative image paths
    session_dir = os.path.dirname(os.path.abspath(semantic_path))

    # 方案4: Track consecutive same-type slides for layout alternation
    _prev_type = None
    _consec_count = 0
    for i, sd in enumerate(slides_data, 1):
        # Inject deck-level title into title slides for cover rendering
        if sd.get('slide_type') == 'title':
            sd['_deck_title'] = deck_title
        cur_type = sd.get('slide_type', 'bullet-list')
        if cur_type == _prev_type:
            _consec_count += 1
        else:
            _consec_count = 0
        _prev_type = cur_type
        layout_variant = _consec_count % 2  # 0 = default, 1 = alternate
        render_slide(prs, sd, spec, grid, sections, i, total,
                     session_dir=session_dir,
                     layout_variant=layout_variant)

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    print(f"✅ PPTX saved: {output_path} ({total} slides)")
    return output_path


def main():
    parser = argparse.ArgumentParser(description='Generate PPTX from semantic JSON + design spec')
    parser.add_argument('--semantic', required=True, help='Path to slides_semantic.json')
    parser.add_argument('--design', required=True, help='Path to design_spec.json')
    parser.add_argument('--output', required=True, help='Output PPTX path')
    args = parser.parse_args()
    generate_pptx(args.semantic, args.design, args.output)


if __name__ == '__main__':
    main()
