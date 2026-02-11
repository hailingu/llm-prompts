from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

pptx_path = 'docs/presentations/storage-frontier-20260211/Storage-Frontier.pptx'
prs = Presentation(pptx_path)
slide = prs.slides[0]

# Subtitle
tb = slide.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(10.0), Inches(0.6))
tf = tb.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = '关键趋势、厂商案例与可落地建议'
run.font.size = Pt(20)
run.font.bold = False
p.alignment = PP_ALIGN.LEFT

# Author / Date
tb2 = slide.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(10.0), Inches(0.4))
tf2 = tb2.text_frame
p2 = tf2.paragraphs[0]
run2 = p2.add_run()
run2.text = 'Author: markdown-writer-specialist    Date: 2026-02-11'
run2.font.size = Pt(12)
run2.font.italic = True
p2.alignment = PP_ALIGN.LEFT

prs.save(pptx_path)
print('Added subtitle and metadata to title slide')