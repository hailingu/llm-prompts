#!/usr/bin/env python3
import importlib.util
import json
import os
import traceback

log_path = '/tmp/render_debug.log'
with open(log_path, 'w', encoding='utf-8') as log:
    try:
        spec = importlib.util.spec_from_file_location('gen', os.path.abspath('skills/ppt-generator/bin/generate_pptx.py'))
        gen = importlib.util.module_from_spec(spec)
        spec.loader.exec_module(gen)
        log.write(f"HAS_PPTX_CHARTS={getattr(gen,'HAS_PPTX_CHARTS',None)}\n")
        log.write(f"NATIVE_CHART_TYPE_MAP keys={list(getattr(gen,'NATIVE_CHART_TYPE_MAP',{}).keys())}\n")
        with open('docs/presentations/storage-frontier-20260211/slides_semantic.json','r',encoding='utf-8') as f:
            sem=json.load(f)
        slides_to_check=[5,6,13,17]
        for sid in slides_to_check:
            sd = next(s for s in sem['slides'] if s.get('id')==sid)
            vis = sd.get('visual')
            log.write(f"Trying slide {sid} type={vis.get('type')}\n")
            try:
                prs_mod = importlib.util.find_spec('pptx') is not None
                from pptx import Presentation
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                ok = gen.render_native_chart(slide, vis, json.load(open('docs/presentations/storage-frontier-20260211/design_spec.json')) , 1.0,1.0,6.0,3.0)
                log.write(f"ok: {ok}\n")
                log.write(f"shapes count: {len(slide.shapes)}\n")
                for sh in slide.shapes:
                    log.write(f" shape has_chart: {hasattr(sh,'chart')}\n")
            except Exception as e:
                log.write(f"exception for slide {sid}: {e}\n")
                log.write(traceback.format_exc())
    except Exception as e:
        log.write(f"fatal exception: {e}\n")
        log.write(traceback.format_exc())

print('Wrote debug log to', log_path)
