from pptx import Presentation
import json

pptx_path = 'docs/presentations/storage-frontier-20260211/Storage-Frontier.pptx'
semantic_path = 'docs/presentations/storage-frontier-20260211/slides_semantic.json'

prs = Presentation(pptx_path)
with open(semantic_path) as f:
    sd = json.load(f)

report = []
for idx, s in enumerate(sd['slides']):
    vis = s.get('visual', {})
    if vis.get('type','none') != 'none':
        slide_report = {'slide_id': s.get('id'), 'visual_type': vis.get('type'), 'chart': False, 'table': False, 'picture': False, 'shapes': len(prs.slides[idx].shapes)}
        shapes = prs.slides[idx].shapes
        for sh in shapes:
            # safe chart check
            try:
                _ = sh.chart
                slide_report['chart'] = True
            except Exception:
                pass
            try:
                _ = sh.table
                slide_report['table'] = True
            except Exception:
                pass
            try:
                if sh.shape_type == 13:
                    slide_report['picture'] = True
            except Exception:
                pass
        report.append(slide_report)

print(json.dumps(report, indent=2, ensure_ascii=False))
