#!/usr/bin/env python3
"""CI-ready PPTX generator script

Usage:
  python3 scripts/generate_pptx_ci.py --semantic output/MFT_slides_semantic.json \
      --design output/MFT_design_spec.json --output docs/presentations/mft-20260206/MFT.pptx

This script is intentionally conservative: it uses existing pre-rendered PNG/SVG assets from
`design_spec.visual_assets_manifest.assets` when available, otherwise it creates text placeholders.
It preserves speaker notes verbatim and applies basic layout & colors from design_spec.
"""
import argparse
import json
import os
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except Exception as e:
    print('Missing python-pptx. Please install requirements in CI (python-pptx, pillow).')
    raise


def hex_to_rgb_tuple(h):
    h = h.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def add_title_slide(prs, slide, spec, slide_data):
    layout = prs.slide_layouts[6]  # blank
    s = prs.slides.add_slide(layout)
    # background
    bg = s.background
    color_hex = spec.get('slide_type_layouts', {}).get('title', {}).get('background', 'primary')
    color = spec.get('color_system', {}).get(color_hex, '#2563EB')
    r, g, b = hex_to_rgb_tuple(color)
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)
    # title
    tx = s.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data.get('title', '')
    p.font.size = Pt(spec.get('typography_system', {}).get('explicit_sizes', {}).get('slide_title', 24))
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    # callouts/kpis
    comps = slide_data.get('components', {})
    if comps.get('kpis'):
        left = Inches(1)
        top = Inches(4)
        for kpi in comps['kpis']:
            tb = s.shapes.add_textbox(left, top, Inches(3.5), Inches(1))
            t = tb.text_frame
            t.text = f"{kpi.get('label')}: {kpi.get('value')}"
            t.paragraphs[0].font.size = Pt(spec.get('typography_system', {}).get('explicit_sizes', {}).get('kpi_value', 28))
            top += Inches(0.9)
    # notes
    notes = s.notes_slide.notes_text_frame
    notes.text = slide_data.get('speaker_notes', '')
    return s


def add_content_slide(prs, spec, slide_data, assets_map):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    # title bar
    title = slide_data.get('title', '')
    tx = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(spec.get('typography_system', {}).get('explicit_sizes', {}).get('slide_title', 24))
    p.font.bold = True
    # content: image or bullets
    visual = slide_data.get('visual', {})
    asset = None
    if visual:
        # find matching asset in assets_map by slide_id
        asset = assets_map.get(slide_data.get('slide_id'))
    if asset:
        img_path = os.path.join(spec.get('visual_assets_manifest', {}).get('output_folder', 'docs/presentations/mft-20260206/images'), Path(asset['path']).name)
        if os.path.exists(img_path):
            # add picture centered
            pic = s.shapes.add_picture(img_path, Inches(1), Inches(1.3), width=Inches(10.33))
        else:
            # fallback to code preview
            cb = s.shapes.add_textbox(Inches(1), Inches(1.3), Inches(10.33), Inches(5))
            cb.text_frame.text = visual.get('placeholder_data', {}).get('mermaid_code', '[Visual]')[:800]
    else:
        # bullets
        left = Inches(1)
        top = Inches(1.3)
        width = Inches(10.33)
        height = Inches(4.5)
        tb = s.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        contents = slide_data.get('content', [])
        if isinstance(contents, list):
            for i, c in enumerate(contents[:8]):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = c
                p.level = 0
                p.font.size = Pt(spec.get('typography_system', {}).get('explicit_sizes', {}).get('bullet_text', 16))
    # notes
    notes = s.notes_slide.notes_text_frame
    notes.text = slide_data.get('speaker_notes', '')
    return s


def build_assets_map(spec):
    assets = spec.get('visual_assets_manifest', {}).get('assets', [])
    m = {a['slide_id']: a for a in assets}
    return m


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--semantic', required=True)
    parser.add_argument('--design', required=True)
    parser.add_argument('--output', required=True)
    args = parser.parse_args()

    with open(args.semantic, encoding='utf-8') as f:
        semantic = json.load(f)
    with open(args.design, encoding='utf-8') as f:
        design = json.load(f)

    prs = Presentation()
    # set slide size
    slide_w = Inches(design.get('grid_system', {}).get('slide_width_inches', 13.333))
    slide_h = Inches(design.get('grid_system', {}).get('slide_height_inches', 7.5))
    prs.slide_width = slide_w
    prs.slide_height = slide_h

    assets_map = build_assets_map(design)

    for slide in semantic.get('slides', []):
        stype = slide.get('slide_type')
        if stype == 'title':
            add_title_slide(prs, None, design, slide)
        elif stype == 'section_divider':
            # simple divider
            s = prs.slides.add_slide(prs.slide_layouts[6])
            bg = s.background
            color_hex = design.get('color_system', {}).get(design.get('section_accents', {}).get(slide.get('metadata', {}).get('section_id','A'), 'primary'))
            r,g,b = hex_to_rgb_tuple(color_hex)
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(r,g,b)
            tx = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1))
            t = tx.text_frame
            t.text = slide.get('title','')
            t.paragraphs[0].font.size = Pt(design.get('typography_system', {}).get('explicit_sizes', {}).get('slide_subtitle', 18))
            t.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
            s.notes_slide.notes_text_frame.text = slide.get('speaker_notes','')
        else:
            add_content_slide(prs, design, slide, assets_map)

    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    prs.save(args.output)
    print('PPTX saved to', args.output)


if __name__ == '__main__':
    main()
