#!/usr/bin/env python3
"""Verify layout: content below title bar. Skip shapes inside the title bar."""
import json
from pptx import Presentation

sem = json.load(open('docs/presentations/storage-frontier-20260211/slides_semantic.json', encoding='utf-8'))
prs = Presentation('docs/presentations/storage-frontier-20260211/storage-frontier-v10-assertion-packaged.pptx')

for idx in range(len(prs.slides)):
    slide = prs.slides[idx]
    sd = sem['slides'][idx]
    stype = sd['slide_type']
    shapes = []
    for sh in slide.shapes:
        top = sh.top / 914400.0
        h = sh.height / 914400.0
        left = sh.left / 914400.0
        w = sh.width / 914400.0
        txt = ''
        try:
            if hasattr(sh, 'text') and sh.text:
                txt = sh.text[:50]
        except Exception:
            pass
        shapes.append((top, h, left, w, txt))

    # Find title bar (full-width rectangle at top=0, h > 0.3)
    bar_bottom = 0
    for top, h, left, w, txt in shapes:
        if top < 0.01 and w > 12 and h > 0.3 and h < 1.5:
            bar_bottom = top + h
            break

    # Content = shapes BELOW bar_bottom, not footer area, not full-width backgrounds
    content_shapes = [(top, h, txt) for top, h, left, w, txt in shapes 
                      if top >= bar_bottom - 0.01  # at or below bar bottom
                      and top < 7.0               # not footer
                      and w < 12                   # not full-width background
                      and h > 0.1                  # not tiny dots
                      and top > 0.5]               # skip title bar text

    if content_shapes:
        first_content_top = min(t for t, h, txt in content_shapes)
        overlap = "OVERLAP" if first_content_top < bar_bottom - 0.05 else "OK"
    else:
        first_content_top = 0
        overlap = "NO_CONTENT" if stype not in ('title', 'section_divider') else "SPECIAL"

    n_shapes = len(slide.shapes)
    flag = "!!" if overlap in ("OVERLAP", "NO_CONTENT") else "  "
    print(f"{flag} Slide {idx+1:2d} ({stype:18s}): bar_bot={bar_bottom:.2f}  content_start={first_content_top:.2f}  shapes={n_shapes:2d}  {overlap}")
