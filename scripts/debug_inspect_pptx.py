import os, json, copy, tempfile, subprocess
from pptx import Presentation

src_dir = os.path.join('docs', 'presentations', 'storage-frontier-20260211')
slides = 'slides_semantic.json'
spec = 'design_spec.json'
with tempfile.TemporaryDirectory() as td:
    tmp_sem = os.path.join(td, 'slides_semantic_no_images.json')
    with open(os.path.join(src_dir, slides),'r',encoding='utf-8') as f:
        sem = json.load(f)
    sem2 = copy.deepcopy(sem)
    for sd in sem2.get('slides',[]):
        vis = sd.get('visual')
        if isinstance(vis, dict):
            for rk in list(vis.keys()):
                if rk.startswith('rendered_'):
                    vis.pop(rk,None)
            pd = vis.get('placeholder_data') or {}
            if isinstance(pd, dict):
                for pkk,pv in list(pd.items()):
                    if pkk.startswith('rendered_'):
                        pd.pop(pkk,None)
                    elif isinstance(pv, dict):
                        for rk2 in list(pv.keys()):
                            if rk2.startswith('rendered_'):
                                pv.pop(rk2,None)
    with open(tmp_sem,'w',encoding='utf-8') as wf:
        json.dump(sem2,wf,ensure_ascii=False)
    out = os.path.join(td,'out.pptx')
    cmd = ['python3', os.path.join('..','..','..','.github','skills','ppt-generator','bin','generate_pptx.py'), '--semantic', tmp_sem, '--design', spec, '--output', out]
    # run in src_dir
    res = subprocess.run(cmd, cwd=src_dir, capture_output=True, text=True)
    print('returncode', res.returncode)
    print('stdout', res.stdout)
    print('stderr', res.stderr)
    if not os.path.isfile(out):
        print('no out file')
    else:
        prs = Presentation(out)
        for i, slide in enumerate(prs.slides, 1):
            print('Slide', i, 'shapes:', len(slide.shapes))
            for shape in slide.shapes:
                try:
                    t = shape.shape_type
                except Exception:
                    t = None
                print('  type', t, 'has_chart', hasattr(shape,'chart'))
