#!/usr/bin/env python3
import json, copy, os, subprocess, tempfile, shutil
from pathlib import Path

repo = Path(__file__).resolve().parents[1]
src_dir = repo / 'docs' / 'presentations' / 'storage-frontier-20260211'
with open(src_dir / 'slides_semantic.json', 'r', encoding='utf-8') as f:
    sem = json.load(f)
sem2 = copy.deepcopy(sem)
for sd in sem2.get('slides', []):
    vis = sd.get('visual')
    if isinstance(vis, dict):
        for rk in list(vis.keys()):
            if rk.startswith('rendered_'):
                vis.pop(rk, None)
        pd = vis.get('placeholder_data') or {}
        if isinstance(pd, dict):
            for pkk, pv in list(pd.items()):
                if pkk.startswith('rendered_'):
                    pd.pop(pkk, None)
                elif isinstance(pv, dict):
                    for rk2 in list(pv.keys()):
                        if rk2.startswith('rendered_'):
                            pv.pop(rk2, None)

with tempfile.TemporaryDirectory() as td:
    tmp_sem = Path(td) / 'slides_semantic_no_images.json'
    with open(tmp_sem, 'w', encoding='utf-8') as wf:
        json.dump(sem2, wf, ensure_ascii=False)
    out_pptx = Path(td) / 'storage-frontier-v9-native-chart-native.pptx'
    cmd = [
        'python3',
        str(repo / 'skills' / 'ppt-generator' / 'bin' / 'generate_pptx.py'),
        '--semantic', str(tmp_sem),
        '--design', str(src_dir / 'design_spec.json'),
        '--output', str(out_pptx),
    ]
    print('Running:', ' '.join(cmd))
    subprocess.run(cmd, cwd=str(src_dir), check=True)
    dest = src_dir / out_pptx.name
    shutil.copy2(out_pptx, dest)
    print('Copied generated PPTX to', dest)
    # Report chart shapes
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = Presentation(dest)
    chart_count = sum(1 for s in prs.slides for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.CHART)
    print('slides:', len(prs.slides), 'chart_shapes:', chart_count)
