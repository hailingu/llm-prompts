#!/usr/bin/env python3
import importlib.util
import json
import os
import traceback

log_path = '/tmp/render_debug_fixed.log'
with open(log_path, 'w', encoding='utf-8') as log:
    try:
        # Use explicit workspace root to avoid ambiguity in cwd/__file__ behavior in different environments
        repo_root = '/Users/guhailin/Git/llm-prompts'
        gen_path = os.path.join(repo_root, 'skills', 'ppt-generator', 'bin', 'generate_pptx.py')
        log.write(f"gen_path={gen_path}\n")
        spec = importlib.util.spec_from_file_location('gen', gen_path)
        gen = importlib.util.module_from_spec(spec)
        spec.loader.exec_module(gen)
        log.write(f"HAS_PPTX_CHARTS={getattr(gen,'HAS_PPTX_CHARTS',None)}\n")
        try:
            keys = list(getattr(gen,'NATIVE_CHART_TYPE_MAP',{}).keys())
        except Exception as e:
            keys = f"error:{e}"
        log.write(f"NATIVE_CHART_TYPE_MAP keys={keys}\n")
        slides_path = os.path.join(repo_root,'docs','presentations','storage-frontier-20260211')
        with open(os.path.join(slides_path,'slides_semantic.json'),'r',encoding='utf-8') as f:
            sem=json.load(f)
        slides_to_check=[5,6,7,13,17]
        design = json.load(open(os.path.join(slides_path,'design_spec.json'),'r',encoding='utf-8'))
        for sid in slides_to_check:
            sd = next(s for s in sem['slides'] if s.get('id')==sid)
            vis = sd.get('visual')
            log.write(f"Trying slide {sid} type={vis.get('type')}\n")
            log.write(f"vis: {json.dumps(vis, indent=2, ensure_ascii=False)}\n")
            try:
                from pptx import Presentation
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                ok = gen.render_native_chart(slide, vis, design, 1.0,1.0,6.0,3.0)
                log.write(f"ok: {ok}\n")
                log.write(f"shapes count: {len(slide.shapes)}\n")
                for idx,sh in enumerate(slide.shapes):
                    log.write(f" shape[{idx}] type={type(sh)} has_chart: {hasattr(sh,'chart')}\n")
            except Exception as e:
                log.write(f"exception for slide {sid}: {e}\n")
                log.write(traceback.format_exc())
    except Exception as e:
        log.write(f"fatal exception: {e}\n")
        log.write(traceback.format_exc())

print('Wrote debug log to', log_path)
