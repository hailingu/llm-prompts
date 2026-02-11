#!/usr/bin/env python3
"""Full diagnostic: dump every slide's shapes with positions and text."""
import json
from pptx import Presentation

sem = json.load(open('docs/presentations/storage-frontier-20260211/slides_semantic.json', encoding='utf-8'))
prs = Presentation('docs/presentations/storage-frontier-20260211/storage-frontier-v10-assertion-packaged.pptx')

for idx in range(len(prs.slides)):
    slide = prs.slides[idx]
    sd = sem['slides'][idx]
    stype = sd['slide_type']
    n = len(slide.shapes)
    
    # Find title bar bottom
    bar_bottom = 0
    for sh in slide.shapes:
        t = sh.top / 914400.0
        h = sh.height / 914400.0
        w = sh.width / 914400.0
        if t < 0.01 and w > 12 and h > 0.3 and h < 1.5:
            bar_bottom = t + h
            break

    # Find content shapes (below bar, above footer, not bg)
    content_shapes = []
    for sh in slide.shapes:
        t = sh.top / 914400.0
        h = sh.height / 914400.0
        w = sh.width / 914400.0
        if t > bar_bottom - 0.01 and t < 7.0 and w < 13 and t > 0.5:
            txt = ''
            try:
                if hasattr(sh, 'text') and sh.text:
                    txt = sh.text[:50]
            except:
                pass
            content_shapes.append((t, h, w, txt))
    
    if content_shapes:
        content_top = min(t for t, h, w, txt in content_shapes)
        content_bottom = max(t + h for t, h, w, txt in content_shapes)
        content_used = content_bottom - content_top
    else:
        content_top = content_bottom = content_used = 0

    avail = 7.25 - bar_bottom  # footer at ~7.25
    usage_pct = (content_used / avail * 100) if avail > 0 else 0
    
    # Detect issues
    issues = []
    if stype not in ('title', 'section_divider'):
        if usage_pct < 40:
            issues.append(f"LOW_FILL({usage_pct:.0f}%)")
        if content_top > 0 and content_top > bar_bottom + 0.5:
            issues.append("GAP_TOP")
    
    # Check for empty visual placeholders (large empty shapes)
    for sh in slide.shapes:
        t = sh.top / 914400.0
        h = sh.height / 914400.0
        w = sh.width / 914400.0
        txt = ''
        try:
            if hasattr(sh, 'text') and sh.text:
                txt = sh.text.strip()
        except:
            pass
        if h > 2.0 and w > 5.0 and not txt and t > 0.5 and t < 7.0:
            issues.append(f"EMPTY_RECT({w:.1f}x{h:.1f}@{t:.1f})")
    
    status = " ".join(issues) if issues else "OK"
    print(f"Slide {idx+1:2d} ({stype:18s}): bar={bar_bottom:.2f} content={content_top:.2f}-{content_bottom:.2f} used={content_used:.1f}/{avail:.1f}in ({usage_pct:.0f}%) shapes={n:2d} {status}")
