from pptx import Presentation
from pptx.util import Inches

pptx_path = 'docs/presentations/storage-frontier-20260211/Storage-Frontier.pptx'
prs = Presentation(pptx_path)

# Slide 7 is index 6 (1-based id 7)
slide_idx = 6
slide = prs.slides[slide_idx]

# Add two images side by side: left = 1.0in, right = 7.5in, width 5.5in each
left1 = Inches(1.0)
left2 = Inches(7.0)
top = Inches(2.0)
width = Inches(5.5)

slide.shapes.add_picture('docs/presentations/storage-frontier-20260211/images/figure_pmem.png', left1, top, width=width)
slide.shapes.add_picture('docs/presentations/storage-frontier-20260211/images/figure_media_share.png', left2, top, width=width)

prs.save(pptx_path)
print('Inserted composite images into slide 7 and saved PPTX')