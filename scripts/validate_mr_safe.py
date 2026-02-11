from pptx import Presentation
import json
import sys

pptx_path = 'docs/presentations/storage-frontier-20260211/Storage-Frontier.pptx'
semantic_path = 'docs/presentations/storage-frontier-20260211/slides_semantic.json'
design_path = 'docs/presentations/storage-frontier-20260211/design_spec.json'

prs = Presentation(pptx_path)
with open(semantic_path) as f:
    semantic = json.load(f)
with open(design_path) as f:
    spec = json.load(f)

errors = []
slides = semantic.get('slides', [])

# MR-1 Backgrounds
for i in range(len(prs.slides)):
    try:
        s = prs.slides[i]
        bg_type = None
        try:
            bg_type = s.background.fill.type
        except Exception as e:
            errors.append(f"MR-1: Slide {i+1} background check error: {e}")
            continue
        if bg_type is None:
            errors.append(f"MR-1: Slide {i+1} has no background fill")
    except Exception as e:
        errors.append(f"MR-1: Could not access slide {i+1}: {e}")

# MR-5: Section dividers
dividers = sum(1 for s in slides if s.get('slide_type') == 'section_divider')
sections = len(semantic.get('sections', []))
if len(slides) >= 15 and dividers == 0 and sections > 0:
    errors.append(f"MR-5: {len(slides)} slides, {sections} sections, but 0 dividers")

# MR-6: Title slide completeness
try:
    title_shapes = [s for s in prs.slides[0].shapes if s.has_text_frame]
    if len(title_shapes) < 3:
        errors.append(f"MR-6: Title slide has {len(title_shapes)} text frames, need ≥3")
except Exception as e:
    errors.append(f"MR-6: Error checking title slide: {e}")

# MR-8: Bottom bar presence
slide_h = prs.slide_height / 914400
for i in range(1, len(prs.slides)):
    try:
        s = prs.slides[i]
        stype = slides[i].get('slide_type','') if i < len(slides) else ''
        if stype in ('title','section_divider'):
            continue
        bottom = [sh for sh in s.shapes if ((sh.top + sh.height) / 914400) > (slide_h - 0.35)]
        if not bottom:
            errors.append(f"MR-8: Slide {i+1} has no bottom bar")
    except Exception as e:
        errors.append(f"MR-8: Error on slide {i+1}: {e}")

# MR-9: Components shape count
for i, sd in enumerate(slides):
    try:
        comps = sd.get('components', {})
        has_comps = any(comps.get(k) for k in comps)
        if has_comps and i < len(prs.slides):
            shape_count = len(prs.slides[i].shapes)
            if shape_count < 8:
                errors.append(f"MR-9: Slide {i+1} has components but only {shape_count} shapes")
    except Exception as e:
        errors.append(f"MR-9: Error on slide {i+1}: {e}")

# MR-11: Component key flexibility
for i, sd in enumerate(slides):
    try:
        comps = sd.get('components', {})
        items = comps.get('comparison_items') or comps.get('decisions') or []
        for item in items:
            data_keys = [k for k in item if k not in ('label','icon','color')]
            if len(data_keys) == 0 and len(item) > 1:
                errors.append(f"MR-11: Slide {i+1} component item has keys {list(item.keys())} but no data keys detected")
    except Exception as e:
        errors.append(f"MR-11: Error on slide {i+1}: {e}")

# MR-12: Content dedup
for i, sd in enumerate(slides):
    try:
        comps = sd.get('components', {})
        content = sd.get('content', [])
        comp_labels = set()
        for key in ('decisions','comparison_items'):
            for item in (comps.get(key) or []):
                if item.get('label'):
                    comp_labels.add(item['label'])
        if comp_labels and content:
            dupes = [c for c in content if c in comp_labels]
            if dupes:
                errors.append(f"MR-12: Slide {i+1} content duplicates component labels: {dupes}")
    except Exception as e:
        errors.append(f"MR-12: Error on slide {i+1}: {e}")

# MR-13: content zone utilization
lz = spec.get('layout_zones', {})
title_bar_h = lz.get('title_bar_height_default', 0.55)
bottom_bar_h = max(lz.get('bottom_bar_height', 0.25), 0.25)
content_zone_top = title_bar_h + lz.get('content_margin_top', 0.12)
content_zone_h = slide_h - content_zone_top - bottom_bar_h - lz.get('content_bottom_margin', 0.2)
for i, sd in enumerate(slides):
    try:
        stype = sd.get('slide_type','')
        if stype in ('title','section_divider'):
            continue
        if i >= len(prs.slides):
            continue
        s = prs.slides[i]
        content_shapes = [sh for sh in s.shapes if (sh.top / 914400) > (content_zone_top - 0.1) and ((sh.top + sh.height) / 914400) < (slide_h - bottom_bar_h + 0.1)]
        if content_shapes:
            max_bottom = max((sh.top + sh.height) / 914400 for sh in content_shapes)
            used_h = max_bottom - content_zone_top
            fill_ratio = used_h / content_zone_h if content_zone_h > 0 else 1.0
            if fill_ratio < 0.55:
                errors.append(f"MR-13: Slide {i+1} ({stype}) content fills only {fill_ratio:.0%} of content zone (used {used_h:.1f}in / {content_zone_h:.1f}in)")
    except Exception as e:
        errors.append(f"MR-13: Error on slide {i+1}: {e}")

# MR-14: components vs visual dedup
for i, sd in enumerate(slides):
    try:
        comps = sd.get('components', {})
        visual = sd.get('visual', {})
        vis_type = visual.get('type','none')
        has_comps = any(comps.get(k) for k in comps)
        has_visual_data = vis_type != 'none' and visual.get('placeholder_data', {})
        if has_comps and has_visual_data:
            comp_labels = set()
            for key in ('decisions','comparison_items','kpis'):
                for item in (comps.get(key) or []):
                    comp_labels.add(item.get('label', item.get('title','')))
            vis_data = visual.get('placeholder_data', {}).get('chart_config', {})
            vis_labels = set()
            for series in vis_data.get('series', []):
                vis_labels.add(series.get('name',''))
            overlap = comp_labels & vis_labels
            if overlap:
                errors.append(f"MR-14: Slide {i+1} has overlapping data in components and visual: {overlap}")
    except Exception as e:
        errors.append(f"MR-14: Error on slide {i+1}: {e}")

# MR-15: title slide component suppression
for i, sd in enumerate(slides):
    try:
        if sd.get('slide_type') == 'title':
            comps = sd.get('components', {})
            has_comps = any(comps.get(k) for k in comps)
            if has_comps:
                errors.append(f"MR-15: Slide {i+1} is title slide but has components {list(k for k in comps if comps.get(k))}")
    except Exception as e:
        errors.append(f"MR-15: Error on slide {i+1}: {e}")

if errors:
    print('❌ VALIDATION FAILED:')
    for e in errors:
        print('  ', e)
    sys.exit(1)
else:
    print(f"✅ All checks passed ({len(prs.slides)} slides)")
    sys.exit(0)
