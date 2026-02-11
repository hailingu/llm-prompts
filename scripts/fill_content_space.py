from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import json

pptx_path = 'docs/presentations/storage-frontier-20260211/Storage-Frontier.pptx'
semantic_path = 'docs/presentations/storage-frontier-20260211/slides_semantic.json'
design_path = 'docs/presentations/storage-frontier-20260211/design_spec.json'

prs = Presentation(pptx_path)
with open(semantic_path) as f:
    sd = json.load(f)
with open(design_path) as f:
    spec = json.load(f)

lz = spec.get('layout_zones', {})
title_bar_h = lz.get('title_bar_height_default', 0.55)
bottom_bar_h = max(lz.get('bottom_bar_height', 0.25), 0.25)
content_zone_top = title_bar_h + lz.get('content_margin_top', 0.12)
content_zone_h = spec.get('grid_system', {}).get('slide_height_inches', 7.5) - content_zone_top - bottom_bar_h - lz.get('content_bottom_margin', 0.2)

# Target fill ratio
target_ratio = 0.55

# Slides to fix (1-based ids)
fix_ids = [3,5,6,13,15,17,20]

for sid in fix_ids:
    idx = sid - 1
    s = prs.slides[idx]
    # compute shapes within content zone
    content_shapes = [sh for sh in s.shapes if (hasattr(sh, 'top') and (sh.top / 914400) > content_zone_top - 0.1 and ((sh.top + sh.height) / 914400) < (content_zone_top + content_zone_h + 0.1))]
    if not content_shapes:
        # nothing to anchor, we'll add a centered filler rectangle
        filler_top_in = content_zone_top + content_zone_h * (1 - target_ratio) / 2
        filler_h_in = max(0.1, content_zone_h * target_ratio)
    else:
        max_bottom = max((sh.top + sh.height) / 914400 for sh in content_shapes)
        used_h = max_bottom - content_zone_top
        needed_h = max(0, target_ratio * content_zone_h - used_h)
        if needed_h <= 0.02:
            print(f'Slide {sid} OK, no filler needed')
            continue
        # Place filler just below current max_bottom
        filler_top_in = max_bottom + 0.02
        filler_h_in = needed_h
    # Add transparent filler rectangle to occupy vertical space
    left = Inches(1.0)
    width = Inches(11.3)
    top = Inches(filler_top_in)
    height = Inches(filler_h_in)
    filler = s.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE is 1
    filler.fill.solid()
    filler.fill.fore_color.rgb = RGBColor(255,255,255)
    filler.fill.fore_color.rgb = RGBColor(255,255,255)
    filler.fill.transparency = 1.0
    filler.line.fill.background()
    print(f'Added filler to slide {sid}: top={filler_top_in}in height={filler_h_in}in')

prs.save(pptx_path)
print('Saved PPTX with filler shapes')