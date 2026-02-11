#!/usr/bin/env python3
"""Detailed shape dump for specific slides."""
import json
from pptx import Presentation

sem = json.load(open('docs/presentations/storage-frontier-20260211/slides_semantic.json', encoding='utf-8'))
prs = Presentation('docs/presentations/storage-frontier-20260211/storage-frontier-v10-assertion-packaged.pptx')

for idx in [2, 4, 8, 14]:  # slides 3, 5, 9, 15
    slide = prs.slides[idx]
    sd = sem['slides'][idx]
    print(f"\n=== Slide {idx+1} ({sd['slide_type']}) ===")
    for i, sh in enumerate(slide.shapes):
        top = sh.top / 914400.0
        h = sh.height / 914400.0
        left = sh.left / 914400.0
        w = sh.width / 914400.0
        txt = ''
        try:
            if hasattr(sh, 'text') and sh.text:
                txt = sh.text[:60]
        except Exception:
            pass
        print(f"  {i+1:2d}: top={top:.2f} h={h:.2f} left={left:.2f} w={w:.2f}  [{txt}]")
