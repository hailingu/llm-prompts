#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches
import json
import sys

pptx_path = sys.argv[1]
semantic_path = sys.argv[2]
design_path = sys.argv[3]
out_path = sys.argv[4]

with open(semantic_path, 'r', encoding='utf-8') as f:
    semantic = json.load(f)
with open(design_path, 'r', encoding='utf-8') as f:
    design = json.load(f)

prs = Presentation(pptx_path)
errors = []
slides = semantic.get('slides', [])

# MR-1 Background fills
for i, slide in enumerate(prs.slides):
    try:
        t = slide.background.fill.type
    except Exception:
        t = None
    if t is None:
        errors.append(f"MR-1: Slide {i+1} has no background fill")

# MR-5 Section dividers
dividers = sum(1 for s in slides if s.get('slide_type') == 'section_divider')
sections = len(semantic.get('sections', []))
if len(slides) >= 15 and dividers == 0 and sections > 0:
    errors.append(f"MR-5: {len(slides)} slides, {sections} sections, but 0 dividers")

# MR-6 Title slide completeness
title_shapes = [s for s in prs.slides[0].shapes if s.has_text_frame]
if len(title_shapes) < 3:
    errors.append(f"MR-6: Title slide has {len(title_shapes)} text frames, need ≥3")

# MR-8 Bottom bar presence
slide_h = prs.slide_height / 914400
for i, slide in enumerate(prs.slides[1:], 2):
    stype = slides[i-1].get('slide_type', '') if i-1 < len(slides) else ''
    if stype in ('title','section_divider'):
        continue
    bottom = [s for s in slide.shapes if (s.top + s.height) / 914400 > slide_h - 0.35]
    if not bottom:
        errors.append(f"MR-8: Slide {i} has no bottom bar")

# MR-9 Components rendering check
for i, sd in enumerate(slides):
    comps = sd.get('components', {})
    has_comps = any(comps.get(k) for k in comps)
    if has_comps and i < len(prs.slides):
        shape_count = len(prs.slides[i].shapes)
        if shape_count < 8:
            errors.append(f"MR-9: Slide {i+1} has components but only {shape_count} shapes")

# MR-11 Component key flexibility
for i, sd in enumerate(slides):
    comps = sd.get('components', {})
    items = comps.get('comparison_items') or comps.get('decisions') or []
    for item in items:
        data_keys = [k for k in item if k not in ('label', 'icon', 'color')]
        if len(data_keys) == 0 and len(item) > 1:
            errors.append(f"MR-11: Slide {i+1} component item has keys {list(item.keys())} but no data keys detected")

# MR-12 Content deduplication
for i, sd in enumerate(slides):
    comps = sd.get('components', {})
    content = sd.get('content', [])
    comp_labels = set()
    for key in ('decisions', 'comparison_items'):
        for item in (comps.get(key) or []):
            if item.get('label'):
                comp_labels.add(item['label'])
    if comp_labels and content:
        dupes = [c for c in content if c in comp_labels]
        if dupes:
            errors.append(f"MR-12: Slide {i+1} content duplicates component labels: {dupes}")

# MR-14 Components vs visual deduplication
for i, sd in enumerate(slides):
    comps = sd.get('components', {})
    visual = sd.get('visual', {})
    vis_type = visual.get('type', 'none')
    has_comps = any(comps.get(k) for k in comps)
    has_visual_data = vis_type != 'none' and visual.get('placeholder_data', {})
    if has_comps and has_visual_data:
        comp_labels = set()
        for key in ('decisions', 'comparison_items', 'kpis'):
            for item in (comps.get(key) or []):
                comp_labels.add(item.get('label', item.get('title', '')))
        vis_data = visual.get('placeholder_data', {}).get('chart_config', {})
        vis_labels = set()
        for series in vis_data.get('series', []):
            vis_labels.add(series.get('name', ''))
        overlap = comp_labels & vis_labels
        if overlap:
            errors.append(f"MR-14: Slide {i+1} has overlapping data in components and visual: {overlap}")

# MR-15 Title slide component suppression
for i, sd in enumerate(slides):
    if sd.get('slide_type') == 'title':
        comps = sd.get('components', {})
        has_comps = any(comps.get(k) for k in comps)
        if has_comps:
            errors.append(f"MR-15: Slide {i+1} is title slide but has components {list(k for k in comps if comps.get(k))}")

# Write report
report = {
    'meta': {'pptx': pptx_path, 'semantic': semantic_path, 'design': design_path},
    'errors': errors,
    'error_count': len(errors)
}
with open(out_path, 'w', encoding='utf-8') as f:
    json.dump(report, f, ensure_ascii=False, indent=2)
print('Wrote report to', out_path)
if errors:
    print('Errors:', len(errors))
else:
    print('No errors found')
