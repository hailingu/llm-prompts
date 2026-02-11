#!/usr/bin/env python3
import importlib.util
import json
import os

spec = importlib.util.spec_from_file_location('gen', os.path.abspath('.github/skills/ppt-generator/bin/generate_pptx.py'))
gen = importlib.util.module_from_spec(spec)
spec.loader.exec_module(gen)
print('HAS_PPTX_CHARTS', getattr(gen,'HAS_PPTX_CHARTS',None))
print('NATIVE_CHART_TYPE_MAP keys', list(getattr(gen,'NATIVE_CHART_TYPE_MAP',{}).keys()))
with open('docs/presentations/storage-frontier-20260211/slides_semantic.json','r',encoding='utf-8') as f:
    sem=json.load(f)
sd=next(s for s in sem['slides'] if s.get('id')==5)
vis=sd.get('visual')
print('vis type', vis.get('type'))
cc = (vis.get('placeholder_data') or {}).get('chart_config',{})
print('chart_config keys', list(cc.keys()))
print('labels', cc.get('labels'))
print('series count', len(cc.get('series',[])))

from pptx import Presentation
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
try:
    ok = gen.render_native_chart(slide, vis, json.load(open('docs/presentations/storage-frontier-20260211/design_spec.json')) , 1.0,1.0,6.0,3.0)
    print('render_native_chart ok:', ok)
    print('shapes:', len(slide.shapes))
    for sh in slide.shapes:
        print(' shape has_chart', hasattr(sh,'chart'))
except Exception as e:
    print('exception during render_native_chart:', e)
