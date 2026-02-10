#!/usr/bin/env python3
"""追踪 v16 生成时的坐标参数"""

import sys
sys.path.insert(0, 'skills/ppt-generator')

# Monkey patch to log coordinates
original_render_visual = None

def log_render_visual(slide, visual, spec, grid, left, top, width, height):
    print(f"🔍 render_visual 被调用:")
    print(f"   visual.type={visual.get('type')}")
    print(f"   left={left:.3f}, top={top:.3f}, width={width:.3f}, height={height:.3f}")
    print(f"   grid.slide_h={grid.slide_h:.3f}")
    return original_render_visual(slide, visual, spec, grid, left, top, width, height)

# Import and patch
import importlib
spec = importlib.util.spec_from_file_location("generate_pptx", "skills/ppt-generator/bin/generate_pptx.py")
generate_pptx = importlib.util.module_from_spec(spec)

# Patch before exec
import types
original_globals = generate_pptx.__dict__.copy()

print("=" * 70)
print("准备生成并追踪坐标...")
print("=" * 70)

# 只生成 Slide 29
import json
from pptx import Presentation
from pptx.util import Inches, Pt

with open('docs/presentations/MFT-20260210/slides_semantic.json') as f:
    slides_data = json.load(f)
with open('docs/presentations/MFT-20260210/design_spec.json') as f:
    design_spec = json.load(f)

slide29_data = [s for s in slides_data['slides'] if s['id'] == 29][0]

# 手动模拟渲染流程
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 构造 GridSystem
class GridSystem:
    def __init__(self):
        self.slide_w = 13.333
        self.slide_h = 7.5
        self.margin_h = 0.833
        self.margin_v = 0.5
        self.usable_w = 11.667
        self.usable_h = 6.5

grid = GridSystem()

print("\n计算坐标:")
bar_h = 0.55
lz_top = bar_h + 0.12
lz_h = grid.slide_h - lz_top - 0.35
print(f"  bar_h={bar_h}")
print(f"  lz_top={lz_top}")
print(f"  lz_h={lz_h}")

cursor = lz_top + 0.15
print(f"  cursor (初始)={cursor}")

# Slide 29 没有 timeline_items, bullets, components
print(f"  cursor (无内容后)={cursor}")

# render_visual 调用
visual_left = grid.margin_h
visual_top = cursor + 0.10
visual_width = grid.usable_w
visual_height = lz_h - (cursor - lz_top) - 0.2

print(f"\n预期传给 render_visual 的参数:")
print(f"  left={visual_left:.3f}")
print(f"  top={visual_top:.3f}")
print(f"  width={visual_width:.3f}")
print(f"  height={visual_height:.3f}")

print(f"\n✅ 如果 timeline 背景在 top={visual_top:.3f}，正常")
print(f"❌ 如果 timeline 背景在 top=7.250，说明坐标计算有问题")
