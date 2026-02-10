"""Deep analysis of v27: layout, overlap, colors, card sizing."""
from pptx import Presentation
from pptx.util import Emu, Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from collections import Counter, defaultdict
import json

prs = Presentation('/private/tmp/MFT-v27.pptx')
W, H = 13.33, 7.5
EMU_PER_INCH = 914400

def emu_to_in(v):
    return v / EMU_PER_INCH

print("=" * 70)
print("1. PER-SLIDE LAYOUT ANALYSIS")
print("=" * 70)

for idx, sl in enumerate(prs.slides):
    sn = idx + 1
    shapes = list(sl.shapes)
    if not shapes:
        continue
    
    # Categorize shapes
    ovals = []
    rounded_rects = []
    textboxes = []
    rects = []
    pictures = []
    
    for sh in shapes:
        l = emu_to_in(sh.left)
        t = emu_to_in(sh.top)
        w = emu_to_in(sh.width)
        h = emu_to_in(sh.height)
        info = {'l': l, 't': t, 'w': w, 'h': h, 'r': l+w, 'b': t+h}
        
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                st = str(sh.auto_shape_type)
                if 'OVAL' in st:
                    # Get text inside
                    txt = ''
                    if sh.has_text_frame:
                        txt = sh.text_frame.text
                    ovals.append({**info, 'text': txt})
                elif 'ROUNDED' in st:
                    rounded_rects.append(info)
                else:
                    rects.append(info)
            except:
                rects.append(info)
        elif sh.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            txt = sh.text_frame.text if sh.has_text_frame else ''
            textboxes.append({**info, 'text': txt[:60]})
        elif sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pictures.append(info)
    
    # Content zone analysis (exclude title bar area y<0.7 and page num area)
    content_shapes = [s for s in shapes if emu_to_in(s.top) > 0.6]
    if not content_shapes:
        continue
    content_min_y = min(emu_to_in(s.top) for s in content_shapes)
    content_max_y = max(emu_to_in(s.top + s.height) for s in content_shapes)
    content_fill = (content_max_y - content_min_y) / (H - 0.67) * 100
    
    # Check for overlaps between ovals and textboxes
    overlaps = []
    for oval in ovals:
        for tb in textboxes:
            # Check bounding box overlap
            if (oval['l'] < tb['r'] and oval['r'] > tb['l'] and 
                oval['t'] < tb['b'] and oval['b'] > tb['t']):
                # Compute overlap area
                ox = max(0, min(oval['r'], tb['r']) - max(oval['l'], tb['l']))
                oy = max(0, min(oval['b'], tb['b']) - max(oval['t'], tb['t']))
                if ox > 0.05 and oy > 0.05:
                    overlaps.append({
                        'oval_pos': f"({oval['l']:.2f},{oval['t']:.2f})",
                        'oval_text': oval.get('text', ''),
                        'tb_pos': f"({tb['l']:.2f},{tb['t']:.2f})",
                        'tb_text': tb.get('text', '')[:40],
                        'overlap': f"{ox:.2f}x{oy:.2f}"
                    })
    
    # Print slide info
    if rounded_rects or ovals or content_fill > 65 or overlaps:
        print(f"\n  S{sn:2d}: fill={content_fill:.0f}%  y=[{content_min_y:.2f}-{content_max_y:.2f}]")
        if rounded_rects:
            heights = [f"{r['h']:.2f}" for r in rounded_rects]
            print(f"       Cards: {len(rounded_rects)}x [{', '.join(heights)}]")
        if ovals:
            oval_info = [f"({o['l']:.2f},{o['t']:.2f}) sz={o['w']:.2f}" for o in ovals]
            print(f"       Ovals: {len(ovals)}x [{', '.join(oval_info)}]")
        if overlaps:
            print(f"       ⚠️ OVERLAPS ({len(overlaps)}):")
            for ov in overlaps:
                print(f"         oval@{ov['oval_pos']} ↔ text@{ov['tb_pos']} [{ov['tb_text']}] ({ov['overlap']})")

print()
print("=" * 70)
print("2. TEXT OVERFLOW DETECTION")
print("=" * 70)
print("   (Text boxes where text length relative to box width suggests overflow)")
for idx, sl in enumerate(prs.slides):
    sn = idx + 1
    for sh in sl.shapes:
        if sh.has_text_frame:
            text = sh.text_frame.text
            w_in = emu_to_in(sh.width)
            h_in = emu_to_in(sh.height)
            if not text or len(text) < 5:
                continue
            # Rough CJK-aware char width estimate
            cjk_count = sum(1 for c in text if '\u4e00' <= c <= '\u9fff' or '\u3000' <= c <= '\u303f' or '\uff00' <= c <= '\uffef')
            latin_count = len(text) - cjk_count
            est_width_in = (cjk_count * 0.18 + latin_count * 0.09)  # rough at ~14pt
            lines_needed = max(1, est_width_in / max(w_in - 0.2, 0.5))
            line_h = 0.25  # approximate line height
            needed_h = lines_needed * line_h
            if needed_h > h_in * 1.3 and len(text) > 15:
                print(f"  S{sn:2d}: box({w_in:.2f}x{h_in:.2f}) text='{text[:50]}...' est_lines={lines_needed:.1f} needed_h={needed_h:.2f}")

print()
print("=" * 70)
print("3. KPI CARD DETAIL (icon vs value positioning)")
print("=" * 70)
for idx, sl in enumerate(prs.slides):
    sn = idx + 1
    shapes = list(sl.shapes)
    # Find primary_container colored cards (KPI cards)
    kpi_cards = []
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                if 'ROUNDED' in str(sh.auto_shape_type):
                    c = str(sh.fill.fore_color.rgb) if sh.fill and sh.fill.type else ''
                    if c in ('C8A951', 'D4B062'):  # primary_container gold variants
                        kpi_cards.append({
                            'l': emu_to_in(sh.left), 't': emu_to_in(sh.top),
                            'w': emu_to_in(sh.width), 'h': emu_to_in(sh.height),
                            'color': c
                        })
            except:
                pass
    if kpi_cards:
        print(f"\n  S{sn:2d}: {len(kpi_cards)} KPI cards")
        for kc in kpi_cards:
            print(f"       card@({kc['l']:.2f},{kc['t']:.2f}) {kc['w']:.2f}x{kc['h']:.2f} color={kc['color']}")
        # Find ovals near KPI cards
        for oval in [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]:
            try:
                if 'OVAL' in str(oval.auto_shape_type):
                    ol = emu_to_in(oval.left)
                    ot = emu_to_in(oval.top)
                    ow = emu_to_in(oval.width)
                    print(f"       oval@({ol:.2f},{ot:.2f}) sz={ow:.2f}")
            except:
                pass

print()
print("=" * 70)
print("4. COLOR DISTRIBUTION (all fill colors)")
print("=" * 70)
color_map = Counter()
bg_colors = Counter()
for idx, sl in enumerate(prs.slides):
    bg = sl.background
    try:
        if bg.fill and bg.fill.type:
            bc = str(bg.fill.fore_color.rgb)
            bg_colors[bc] += 1
    except:
        pass
    for sh in sl.shapes:
        try:
            if sh.fill and sh.fill.type is not None:
                c = str(sh.fill.fore_color.rgb)
                color_map[c] += 1
        except:
            pass

print("  Shape fill colors:")
for c, cnt in color_map.most_common(15):
    print(f"    #{c}: {cnt}")
print(f"\n  Slide backgrounds: {dict(bg_colors)}")

print()
print("=" * 70)
print("5. CARD HEIGHT HISTOGRAM")
print("=" * 70)
all_card_h = []
for sl in prs.slides:
    for sh in sl.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                if 'ROUNDED' in str(sh.auto_shape_type):
                    all_card_h.append(emu_to_in(sh.height))
            except:
                pass
buckets = defaultdict(int)
for h in all_card_h:
    bucket = f"{int(h*2)/2:.1f}-{int(h*2)/2+0.5:.1f}"
    buckets[bucket] += 1
for b in sorted(buckets.keys()):
    bar = '█' * buckets[b]
    print(f"  {b}in: {bar} ({buckets[b]})")
