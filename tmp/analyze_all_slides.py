#!/usr/bin/env python3
"""分析v23所有slide的布局和尺寸"""
import json
from pptx import Presentation
from pptx.util import Inches, Emu

# 1. 分析semantic JSON结构
with open('docs/presentations/MFT-20260210/slides_semantic.json', 'r', encoding='utf-8') as f:
    data = json.load(f)

print("=" * 100)
print("Part 1: Semantic JSON 数据结构分析")
print("=" * 100)

for i, s in enumerate(data['slides']):
    visual = s.get('visual', {})
    bullets = s.get('bullets', [])
    kpis = s.get('kpis', [])
    decisions = s.get('decisions', [])
    sections = s.get('sections', [])
    items = s.get('items', [])
    columns = s.get('columns', [])
    rows = s.get('rows', [])
    
    content_parts = []
    if bullets: content_parts.append(f'bullets={len(bullets)}')
    if kpis: content_parts.append(f'kpis={len(kpis)}')
    if decisions: content_parts.append(f'decisions={len(decisions)}')
    if sections: content_parts.append(f'sections={len(sections)}')
    if items: content_parts.append(f'items={len(items)}')
    if columns: content_parts.append(f'columns={len(columns)}')
    if rows: content_parts.append(f'rows={len(rows)}')
    if visual:
        vtype = visual.get('type', '?')
        content_parts.append(f'visual={vtype}')
    
    title = s.get('title', '')[:50]
    content_str = ', '.join(content_parts) if content_parts else 'EMPTY'
    print(f'S{i+1:2d}: {content_str:60s} | {title}')

# 2. 分析v23 PPTX实际布局
print("\n" + "=" * 100)
print("Part 2: v23 PPTX 实际布局分析")
print("=" * 100)

prs = Presentation('/private/tmp/MFT-v23.pptx')
slide_w = prs.slide_width / 914400
slide_h = prs.slide_height / 914400
print(f"Slide尺寸: {slide_w:.1f} x {slide_h:.1f} inches\n")

for idx, slide in enumerate(prs.slides):
    shapes = list(slide.shapes)
    print(f"\n--- Slide {idx+1} ({len(shapes)} shapes) ---")
    
    for shape in sorted(shapes, key=lambda s: (s.top, s.left)):
        x = shape.left / 914400
        y = shape.top / 914400
        w = shape.width / 914400
        h = shape.height / 914400
        
        shape_info = f"  [{shape.shape_type}] {shape.name:30s}"
        pos_info = f"pos=({x:.2f},{y:.2f}) size=({w:.2f}x{h:.2f})"
        
        # 检查是否有文字
        text_preview = ""
        if shape.has_text_frame:
            text = shape.text_frame.text[:40].replace('\n', '|')
            text_preview = f' text="{text}"'
        
        # 检查是否是图片
        if hasattr(shape, 'image'):
            text_preview = f' [IMAGE {shape.image.content_type}]'
        
        print(f"{shape_info} {pos_info}{text_preview}")

# 3. 布局问题检测
print("\n" + "=" * 100)
print("Part 3: 布局问题检测")
print("=" * 100)

for idx, slide in enumerate(prs.slides):
    shapes = list(slide.shapes)
    issues = []
    
    for shape in shapes:
        x = shape.left / 914400
        y = shape.top / 914400
        w = shape.width / 914400
        h = shape.height / 914400
        right = x + w
        bottom = y + h
        
        # 检查溢出
        if right > slide_w + 0.1:
            issues.append(f"  ⚠️ {shape.name} 右侧溢出 ({right:.2f} > {slide_w:.1f})")
        if bottom > slide_h + 0.1:
            issues.append(f"  ⚠️ {shape.name} 底部溢出 ({bottom:.2f} > {slide_h:.1f})")
        
        # 检查过小元素
        if w < 0.3 and h < 0.3 and shape.has_text_frame:
            issues.append(f"  ⚠️ {shape.name} 过小 ({w:.2f}x{h:.2f})")
        
        # 检查文字框是否太窄
        if shape.has_text_frame and w < 1.0 and len(shape.text_frame.text) > 20:
            issues.append(f"  ⚠️ {shape.name} 文字框过窄 ({w:.2f}in for {len(shape.text_frame.text)} chars)")

    # 检查重叠
    for i_s, s1 in enumerate(shapes):
        for j_s in range(i_s + 1, len(shapes)):
            s2 = shapes[j_s]
            # 计算重叠
            x1, y1 = s1.left/914400, s1.top/914400
            w1, h1 = s1.width/914400, s1.height/914400
            x2, y2 = s2.left/914400, s2.top/914400
            w2, h2 = s2.width/914400, s2.height/914400
            
            overlap_x = max(0, min(x1+w1, x2+w2) - max(x1, x2))
            overlap_y = max(0, min(y1+h1, y2+h2) - max(y1, y2))
            overlap_area = overlap_x * overlap_y
            
            min_area = min(w1*h1, w2*h2)
            if min_area > 0 and overlap_area / min_area > 0.3:
                issues.append(f"  ⚠️ {s1.name} 与 {s2.name} 重叠 {overlap_area/min_area*100:.0f}%")
    
    # 检查空白率
    total_content_area = 0
    for shape in shapes:
        total_content_area += (shape.width/914400) * (shape.height/914400)
    fill_rate = total_content_area / (slide_w * slide_h) * 100
    if fill_rate < 20:
        issues.append(f"  ⚠️ 内容填充率过低: {fill_rate:.0f}%")
    
    if issues:
        print(f"\nSlide {idx+1}:")
        for issue in issues:
            print(issue)
