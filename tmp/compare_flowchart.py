#!/usr/bin/env python3
"""对比 v16 和 v17 的 Slide 19 flowchart"""

from pptx import Presentation

print("=" * 70)
print("📊 对比 Slide 19 (flowchart) 渲染效果")
print("=" * 70)

for version in ['v16', 'v17']:
    pptx_path = f'/private/tmp/MFT-{version}.pptx'
    prs = Presentation(pptx_path)
    slide19 = prs.slides[18]  # 0-based index
    
    print(f"\n{version.upper()}:")
    print(f"  形状数量: {len(slide19.shapes)}")
    
    # 查找图片
    pictures = [s for s in slide19.shapes if hasattr(s, 'image')]
    print(f"  图片数量: {len(pictures)}")
    
    if pictures:
        pic = pictures[0]
        print(f"  图片尺寸: {pic.width/914400:.2f} x {pic.height/914400:.2f} inches")
        print(f"  图片位置: left={pic.left/914400:.2f}, top={pic.top/914400:.2f}")
        
        # 图片文件大小（如果能获取）
        if hasattr(pic, 'image') and hasattr(pic.image, 'blob'):
            img_size = len(pic.image.blob)
            print(f"  图片大小: {img_size/1024:.1f} KB")

print("\n" + "=" * 70)
print("✅ 请打开 v17 查看颜色是否与 PPT 主题一致（深蓝+金色）")
print("=" * 70)
