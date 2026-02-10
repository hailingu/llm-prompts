"""Verify v28b KPI card sizing for S10, S22, S27."""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
E = 914400
def to_in(v): return v / E

prs = Presentation('/private/tmp/MFT-v28b.pptx')

for sn in [10, 16, 18, 22, 27]:
    sl = prs.slides[sn - 1]
    cards, ovals, texts = [], [], []
    for sh in sl.shapes:
        l, t, w, h = to_in(sh.left), to_in(sh.top), to_in(sh.width), to_in(sh.height)
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                st = str(sh.auto_shape_type)
                if 'ROUNDED' in st:
                    cards.append({'l': l, 't': t, 'w': w, 'h': h})
                elif 'OVAL' in st:
                    ovals.append({'l': l, 't': t, 'sz': w})
            except: pass
        elif sh.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and sh.has_text_frame:
            txt = sh.text_frame.text.strip()
            if txt and len(txt) > 1:
                texts.append({'l': l, 't': t, 'w': w, 'h': h, 'text': txt[:50]})
    
    print(f"\n{'='*60}")
    print(f"S{sn}:")
    if cards:
        for c in cards:
            # Find texts inside this card
            inner = [t for t in texts if t['l'] >= c['l'] and t['l'] < c['l']+c['w'] and t['t'] >= c['t'] and t['t'] < c['t']+c['h']]
            print(f"  Card({c['w']:.2f}x{c['h']:.2f}) @({c['l']:.2f},{c['t']:.2f})")
            for t in inner:
                print(f"    Text({t['w']:.2f}x{t['h']:.2f}): '{t['text']}'")
    if ovals:
        for o in ovals:
            print(f"  Oval sz={o['sz']:.2f} @({o['l']:.2f},{o['t']:.2f})")
