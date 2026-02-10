#!/usr/bin/env python3
"""测试 NativeGanttRenderer 的坐标使用"""

import sys
sys.path.insert(0, 'skills/ppt-generator')

import json
from pptx import Presentation
from pptx.util import Inches

from ppt_generator.renderers.native.gantt import NativeGanttRenderer
from ppt_generator.protocols.visual_data_protocol import VisualDataProtocol

# 加载设计规范
with open('docs/presentations/MFT-20260210/design_spec.json') as f:
    spec = json.load(f)

# 加载 Slide 29 数据
with open('docs/presentations/MFT-20260210/slides_semantic.json') as f:
    slides_data = json.load(f)
    slide29 = [s for s in slides_data['slides'] if s['id'] == 29][0]

# 创建测试 PPT
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# 构造 VisualDataProtocol
visual = slide29['visual']
pd = visual.get('placeholder_data', {})
vdp = VisualDataProtocol(
    type=visual.get('type'),
    title=visual.get('title'),
    data=pd.get('gantt_data'),
    placeholder_data=pd
)

# 预期坐标（从 render_slide_timeline 计算）
expected_left = 0.83
expected_top = 0.92
expected_width = 11.67
expected_height = 6.13

print("=" * 70)
print("测试 Native Gantt Renderer 坐标使用")
print("=" * 70)
print(f"\n传入坐标:")
print(f"  left={expected_left}, top={expected_top}")
print(f"  width={expected_width}, height={expected_height}")

# 渲染
renderer = NativeGanttRenderer()
success = renderer.render(
    slide, vdp, spec,
    expected_left, expected_top,
    expected_width, expected_height
)

print(f"\n渲染结果: {'成功' if success else '失败'}")

# 分析形状位置
print(f"\n生成的形状:")
for idx, shape in enumerate(slide.shapes):
    left_in = shape.left / 914400
    top_in = shape.top / 914400
    width_in = shape.width / 914400
    height_in = shape.height / 914400
    
    type_name = {1: 'RECTANGLE', 17: 'TEXT_BOX'}.get(shape.shape_type, f'TYPE_{shape.shape_type}')
    print(f"  [{idx}] {type_name}: ({left_in:.3f}, {top_in:.3f})in, {width_in:.3f}x{height_in:.3f}in")
    
    # 检查是否使用了传入的坐标
    if abs(left_in - expected_left) < 0.1 or abs(top_in - expected_top) < 0.5:
        print(f"       ✅ 使用了传入的坐标")
    else:
        if top_in > 6.0:
            print(f"       ❌ 位置异常：在底部边缘！")

# 保存测试 PPT
test_path = '/tmp/test_gantt_coords.pptx'
prs.save(test_path)
print(f"\n测试PPT已保存: {test_path}")
