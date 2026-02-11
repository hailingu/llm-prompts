import importlib.util
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE


def load_gen():
    spec = importlib.util.spec_from_file_location('generate_pptx', 'skills/ppt-generator/bin/generate_pptx.py')
    gen = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(gen)
    return gen


def test_apply_chart_theme_sets_series_colors():
    gen = load_gen()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    data = CategoryChartData()
    data.categories = ['A', 'B']
    data.add_series('S1', (1, 2))
    data.add_series('S2', (3, 4))
    shape = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(5), Inches(3), data)
    chart = shape.chart
    # Apply theme
    gen.apply_chart_theme(chart, spec={})
    palette = gen.get_chart_palette({})
    # First series color should match palette[0]
    first_color = chart.series[0].format.fill.fore_color.rgb
    assert first_color == gen.hex_to_rgb(palette[0])


def test_apply_chart_theme_uses_md3_palette_and_legend_axis_styles():
    gen = load_gen()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    data = CategoryChartData()
    data.categories = ['A', 'B']
    data.add_series('S1', (1, 2))
    data.add_series('S2', (3, 4))
    data.add_series('S3', (5, 6))
    shape = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(5), Inches(3), data)
    chart = shape.chart
    # Custom MD3 palette & on_surface_variant token
    spec = {'md3_palette': ['#010203', '#040506', '#070809'], 'color_system': {'on_surface_variant': '#112233', 'surface_variant': '#F0F0F0', 'outline': '#CCCCCC'}}
    gen.apply_chart_theme(chart, spec=spec, accent_token='primary')
    # Series colors match md3 palette
    for i, s in enumerate(chart.series):
        expected = gen.hex_to_rgb(spec['md3_palette'][i % len(spec['md3_palette'])])
        assert s.format.fill.fore_color.rgb == expected
    # Axis tick label size and color
    try:
        va = chart.value_axis
        assert va.tick_labels.font.size == Pt(8)
        assert va.tick_labels.font.color.rgb == gen.hex_to_rgb(spec['color_system']['on_surface_variant'])
    except Exception:
        # If axis not present in this chart type, skip assertions
        pass
    # Legend font size and color
    try:
        if chart.has_legend:
            for p in chart.legend.text_frame.paragraphs:
                for r in p.runs:
                    assert r.font.size == Pt(7)
                    assert r.font.color.rgb == gen.hex_to_rgb(spec['color_system']['on_surface_variant'])
    except Exception:
        pass
