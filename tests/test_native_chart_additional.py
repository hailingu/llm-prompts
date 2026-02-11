import importlib.util
from pptx import Presentation


def load_gen():
    spec = importlib.util.spec_from_file_location('generate_pptx', 'skills/ppt-generator/bin/generate_pptx.py')
    gen = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(gen)
    return gen


def test_render_native_chart_composite_first_child():
    gen = load_gen()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    visual = {
        'type': 'composite_charts',
        'placeholder_data': {
            'chart_config': {
                'pmem': {
                    'labels': [2023,2024,2025,2026],
                    'series': [{'name': 'PMem', 'data': [50,80,150,300]}]
                },
                'media_share': {
                    'labels': [2023,2024,2025,2026],
                    'series': [{'name': 'SSD', 'data': [45,48,52,58]}]
                }
            }
        }
    }
    spec = {}
    ok = gen.render_native_chart(slide, visual, spec, left=1.0, top=1.0, width=6.0, height=3.0)
    assert ok is True
    assert hasattr(slide.shapes[-1], 'chart')


def test_render_native_chart_bar_line_combo_renders_column():
    gen = load_gen()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    visual = {
        'type': 'bar_line_chart',
        'placeholder_data': {
            'chart_config': {
                'labels': [2023,2024,2025,2026],
                'series': [
                    {'name':'Training IOPS','data':[500000,1000000,2000000,4000000]},
                    {'name':'Training Bandwidth','data':[20,40,80,160]},
                ]
            }
        }
    }
    spec = {}
    ok = gen.render_native_chart(slide, visual, spec, left=1.0, top=1.0, width=6.0, height=3.0)
    assert ok is True
    assert hasattr(slide.shapes[-1], 'chart')
