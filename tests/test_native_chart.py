import importlib.util
from pptx import Presentation
from pptx.util import Inches


def load_gen():
    spec = importlib.util.spec_from_file_location('generate_pptx', 'skills/ppt-generator/bin/generate_pptx.py')
    gen = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(gen)
    return gen


def test_render_native_chart_bar_adds_chart():
    gen = load_gen()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    visual = {
        'type': 'bar_chart',
        'placeholder_data': {
            'chart_config': {
                'labels': ['A', 'B', 'C'],
                'series': [
                    {'name': 'S1', 'data': [1, 2, 3]},
                    {'name': 'S2', 'data': [4, 5, 6]},
                ]
            }
        }
    }
    spec = {}
    ok = gen.render_native_chart(slide, visual, spec, left=1.0, top=1.0, width=6.0, height=3.0)
    assert ok is True
    shape = slide.shapes[-1]
    assert hasattr(shape, 'chart')


def test_render_native_chart_empty_returns_false():
    gen = load_gen()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    visual = {'type': 'bar_chart', 'placeholder_data': {'chart_config': {'labels': [], 'series': []}}}
    spec = {}
    ok = gen.render_native_chart(slide, visual, spec, left=1.0, top=1.0, width=6.0, height=3.0)
    assert ok is False


def test_render_native_chart_scatter():
    gen = load_gen()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    visual = {
        'type': 'scatter_chart',
        'placeholder_data': {
            'chart_config': {
                'series': [
                    {'name': 'S1', 'x': [1, 2, 3], 'y': [4, 5, 6]},
                ]
            }
        }
    }
    spec = {}
    ok = gen.render_native_chart(slide, visual, spec, left=1.0, top=1.0, width=6.0, height=3.0)
    assert ok is True
    shape = slide.shapes[-1]
    assert hasattr(shape, 'chart')
