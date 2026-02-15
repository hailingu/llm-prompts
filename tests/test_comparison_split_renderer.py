import os, sys
ROOT = os.getcwd()
PKG_PATH = os.path.join(ROOT, 'skills', 'ppt-generator')
if ROOT not in sys.path:
    sys.path.insert(0, ROOT)
if PKG_PATH not in sys.path:
    sys.path.insert(0, PKG_PATH)

from pptx import Presentation
from pptx.util import Inches
from ppt_generator.renderers import render_region_comparison_split
from ppt_generator.grid import GridSystem


def make_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    spec = {}
    grid = GridSystem(spec)
    return prs, slide, spec, grid


def test_comparison_split_two_groups_renders():
    prs, slide, spec, grid = make_slide()
    left, top, width, height = grid.margin_h, 1.0, grid.usable_w, 3.0
    groups = [
        [{'label': 'A 架构', 'attributes': {'x': 1}}, {'label': 'B 架构', 'attributes': {'x': 2}}],
        [{'label': 'C 试点', 'attributes': {'y': 3}}, {'label': 'D 试点', 'attributes': {'y': 4}}],
    ]
    render_region_comparison_split(slide, {'groups': groups}, (left, top, width, height), spec)
    assert len(slide.shapes) > 0


def test_comparison_split_merges_into_single_table_and_rows_align():
    prs, slide, spec, grid = make_slide()
    left, top, width, height = grid.margin_h, 1.0, grid.usable_w, 4.0
    groups = [
        [{'label': '並行訓練', 'attributes': {'推薦協議': 'NVMe/TCP / RDMA', '網絡要點': '100/200GbE leaf-spine', '驗證指標': 'IO 並發、隊列延遲、丟包率'}},
         {'label': '低延遲 DB', 'attributes': {'推薦協議': 'RDMA', '驗證指標': 'P99 延遲、QDepth', '回退策略': '本地 NVMe/實例'}}],
        [{'label': 'PMem (持久內存)', 'attributes': {'適用': '元數據/日志加速', '優點': '低 P99 延遲', '局限': '高成本、容量有限'}},
         {'label': 'CXL (遠內存 pooling)', 'attributes': {'適用': '大模型內存彈性', '優點': '內存共享/擴展', '局限': '生態成熟度、隔離/權限'}}],
    ]
    render_region_comparison_split(slide, {'groups': groups}, (left, top, width, height), spec)
    # find header label top (any of the four labels)
    label_texts = {it['label'] for grp in groups for it in grp}
    label_top = None
    for sh in slide.shapes:
        try:
            txt = sh.text.strip()
        except Exception:
            txt = ''
        if txt in label_texts:
            label_top = sh.top
            break
    assert label_top is not None

    # Check that all labels are in the same header row (within tolerance), indicating a single merged table
    tol = 2000
    label_row = [sh for sh in slide.shapes if getattr(sh, 'top', 0) and abs(getattr(sh, 'top', 0) - label_top) <= tol and sh.text.strip() in label_texts]
    expected_cols = len(groups[0]) + len(groups[1])
    assert len(label_row) == expected_cols

    # As a weaker check, ensure there are at least 3 attribute rows (rows below header)
    bottoms = [sh for sh in slide.shapes if getattr(sh, 'top', 0) > label_top and getattr(sh, 'width', 0) > 0]
    rows_set = set((sh.top // 10000) for sh in bottoms)  # coarse bucketing
    assert len(rows_set) >= 3
