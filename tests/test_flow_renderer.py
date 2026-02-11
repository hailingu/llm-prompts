import os
import sys
ROOT = os.getcwd()
PKG_PATH = os.path.join(ROOT, 'skills', 'ppt-generator')
if ROOT not in sys.path:
    sys.path.insert(0, ROOT)
if PKG_PATH not in sys.path:
    sys.path.insert(0, PKG_PATH)

from pptx import Presentation
from ppt_generator.renderers import render_region_flow


def test_render_region_flow_basic():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    data = {
        'steps': [
            {'id': 's1', 'label': 'Start', 'type': 'start'},
            {'id': 's2', 'label': 'Process', 'type': 'process'},
            {'id': 's3', 'label': 'Decision', 'type': 'decision'},
        ],
        'transitions': [
            {'from': 's1', 'to': 's2', 'label': ''},
            {'from': 's2', 'to': 's3', 'label': 'if yes'},
        ]
    }
    bounds = (1.0, 1.0, 10.0, 4.0)
    before = len(slide.shapes)
    render_region_flow(slide, data, bounds, spec={})
    after = len(slide.shapes)
    assert after - before >= 5  # 3 steps + 2 connectors (labels may be separate)

    texts = [s.text_frame.text for s in slide.shapes if hasattr(s, 'text_frame') and s.text_frame is not None and s.text_frame.text]
    assert 'Start' in texts
    assert 'Process' in texts
    assert 'Decision' in texts
    assert 'if yes' in texts


if __name__ == '__main__':
    print('Run: python3 -m pytest tests/test_flow_renderer.py')