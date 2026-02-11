import sys
import os
sys.path.insert(0, os.path.abspath('skills/ppt-generator'))
from pptx import Presentation
from ppt_generator import renderers as r


def test_render_assertion_title_inserts_headline_and_subtitle():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    spec = {
        'color_system': {
            'primary': '#2563EB',
            'on_primary': '#FFFFFF',
            'on_surface_variant': '#6B7280'
        }
    }
    grid = r.GridSystem(spec)
    assertion = 'We must invest in NVMe today'
    title = 'NVMe adoption trends 2023-2026'
    bar_h = r.render_assertion_title(slide, spec, grid, assertion, title, slide_num=5, _total_slides=23, section_label='Section A', accent_color_token='primary')
    assert bar_h >= 0.85
    texts = [sh.text.strip() for sh in slide.shapes if hasattr(sh, 'text') and sh.text and sh.text.strip()]
    # Expect assertion and title to be present among shape texts
    assert any(assertion in t for t in texts), f"Assertion not found in shapes: {texts}"
    assert any(title in t for t in texts), f"Subtitle not found in shapes: {texts}"
