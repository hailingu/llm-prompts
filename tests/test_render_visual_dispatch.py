import os
import subprocess
import tempfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def test_generate_pptx_creates_chart_shapes():
    """Integration: run generator on storage-frontier sample and assert chart shapes exist."""
    src_dir = os.path.join('docs', 'presentations', 'storage-frontier-20260211')
    slides = 'slides_semantic.json'
    spec = 'design_spec.json'
    with tempfile.TemporaryDirectory() as td:
        out_pptx = os.path.join(td, 'storage-frontier-test.pptx')

        # Create a temporary semantic JSON with pre-rendered images removed so
        # native chart path is exercised by the generator.
        import json
        orig_sem = os.path.join(src_dir, slides)
        with open(orig_sem, 'r', encoding='utf-8') as f:
            semantic = json.load(f)
        import copy
        sem2 = copy.deepcopy(semantic)
        for sd in sem2.get('slides', []):
            vis = sd.get('visual')
            if isinstance(vis, dict):
                for rk in list(vis.keys()):
                    if rk.startswith('rendered_'):
                        vis.pop(rk, None)
                pd = vis.get('placeholder_data') or {}
                if isinstance(pd, dict):
                    for pkk, pv in list(pd.items()):
                        if pkk.startswith('rendered_'):
                            pd.pop(pkk, None)
                        elif isinstance(pv, dict):
                            for rk2 in list(pv.keys()):
                                if rk2.startswith('rendered_'):
                                    pv.pop(rk2, None)
        tmp_sem_path = os.path.join(td, 'slides_semantic_no_images.json')
        with open(tmp_sem_path, 'w', encoding='utf-8') as wf:
            json.dump(sem2, wf, ensure_ascii=False)

        cmd = [
            'python3',
            os.path.join('..', '..', '..', 'skills', 'ppt-generator', 'bin', 'generate_pptx.py'),
            '--semantic', tmp_sem_path,
            '--design', spec,
            '--output', out_pptx,
        ]
        # Run generator in the presentation sample directory
        subprocess.run(cmd, cwd=src_dir, check=True)

        assert os.path.isfile(out_pptx), 'Output PPTX was not created'

        # Opening the PPTX to ensure it was written successfully (visuals may
        # be rendered as images or placeholders depending on pre-rendered assets
        # and pipeline configuration). We primarily assert generation succeeded.
        prs = Presentation(out_pptx)
        assert len(prs.slides) > 0

        # Ensure the generated PPTX contains at least one native chart shape
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        chart_count = sum(1 for s in prs.slides for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.CHART)
        assert chart_count > 0, f"Expected at least one chart in generated PPTX, found {chart_count}"

        # Integration sanity: explicitly exercise the package's native chart
        # renderer function to ensure it can produce a Chart shape (unit-level
        # check in an integration context).
        import importlib.util
        spec = importlib.util.spec_from_file_location('generate_pptx', os.path.join('skills', 'ppt-generator', 'bin', 'generate_pptx.py'))
        gen = importlib.util.module_from_spec(spec)
        spec.loader.exec_module(gen)
        prs2 = Presentation()
        slide2 = prs2.slides.add_slide(prs2.slide_layouts[6])
        visual = {
            'type': 'bar_chart',
            'placeholder_data': {
                'chart_config': {
                    'labels': ['A', 'B', 'C'],
                    'series': [
                        {'name': 'S1', 'data': [1, 2, 3]},
                        {'name': 'S2', 'data': [4, 5, 6]},
                    ]
                }
            }
        }
        ok = gen.render_native_chart(slide2, visual, {}, left=1.0, top=1.0, width=6.0, height=3.0)
        assert ok is True
        shape = slide2.shapes[-1]
        assert hasattr(shape, 'chart')
