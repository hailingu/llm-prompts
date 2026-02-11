import sys
import os
sys.path.insert(0, os.path.abspath('skills/ppt-generator'))
from pptx import Presentation
from ppt_generator import renderers as r


def test_render_insight_bar_renders_bar_and_text():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    spec = {'color_system': {'primary': '#2563EB', 'on_primary': '#FFFFFF'}}
    grid = r.GridSystem(spec)
    insight = 'Key takeaway: focus on NVMe deployment.'
    # Call the function directly
    r.render_insight_bar(slide, insight, spec, grid, accent_token='primary')
    texts = [sh.text.strip() for sh in slide.shapes if hasattr(sh, 'text') and sh.text and sh.text.strip()]
    assert any('💡' in t for t in texts), f"No emoji found in shapes: {texts}"
    assert any('Key takeaway' in t for t in texts), f"Insight text missing: {texts}"
