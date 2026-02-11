import os
import sys
ROOT = os.getcwd()
PKG_PATH = os.path.join(ROOT, 'skills', 'ppt-generator')
if ROOT not in sys.path:
    sys.path.insert(0, ROOT)
if PKG_PATH not in sys.path:
    sys.path.insert(0, PKG_PATH)

from pptx import Presentation
from ppt_generator.grid import GridSystem
from ppt_generator.renderers import render_slide, render_slide_v2, detect_schema_version


def make_prs_and_spec():
    prs = Presentation()
    spec = {}
    grid = GridSystem(spec)
    sections = []
    return prs, spec, grid, sections


def test_detect_schema_version_v2():
    sd = {'layout_intent': {'regions': []}}
    assert detect_schema_version(sd) == 2


def test_render_slide_v2_regions_invoked():
    prs, spec, grid, sections = make_prs_and_spec()
    sd = {
        'slide_id': 1,
        'title': 'V2 Slide',
        'slide_type': 'two-column',
        'layout_intent': {
            'template': 'two-column',
            'regions': [
                {'id': 'r1', 'position': 'left-50', 'renderer': 'bullet_list', 'data_source': 'components.bullets'},
                {'id': 'r2', 'position': 'right-50', 'renderer': 'kpi_row', 'data_source': 'components.kpis'},
            ]
        },
        'components': {
            'bullets': ['a','b'],
            'kpis': [{'label': 'K1', 'value': '10'}]
        }
    }
    render_slide(prs, sd, spec, grid, sections, slide_num=1, total_slides=1)
    # After render, there should be shapes on the slide
    s = prs.slides[-1]
    assert len(s.shapes) > 0


def test_render_slide_v2_mixed_ok():
    prs, spec, grid, sections = make_prs_and_spec()
    sd = {
        'slide_id': 2,
        'title': 'V2 Mixed',
        'layout_intent': {'regions': [{'id': 'r1', 'position': 'full', 'renderer': 'comparison_table', 'data_source': 'components.comparison_items'}]},
        'components': {'comparison_items': [{'label': 'A', 'attributes': {'x': 1}}, {'label': 'B', 'attributes': {'x': 2}}]}
    }
    render_slide_v2(prs, sd, spec, grid, sections, slide_num=2, total_slides=2)
    s = prs.slides[-1]
    assert len(s.shapes) > 0
