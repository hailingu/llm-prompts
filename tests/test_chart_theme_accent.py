import importlib.util
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE


def load_gen():
    spec = importlib.util.spec_from_file_location('generate_pptx', 'skills/ppt-generator/bin/generate_pptx.py')
    gen = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(gen)
    return gen


def test_apply_chart_theme_respects_section_accents():
    gen = load_gen()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    data = CategoryChartData()
    data.categories = ['A', 'B']
    data.add_series('S1', (1, 2))
    data.add_series('S2', (3, 4))
    shape = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(5), Inches(3), data)
    chart = shape.chart
    spec = {'section_accents': {'accent_x': ['#AA0000', '#00AA00']}}
    gen.apply_chart_theme(chart, spec=spec, accent_token='accent_x')
    # Series colors should match palette from section_accents
    assert chart.series[0].format.fill.fore_color.rgb == gen.hex_to_rgb('#AA0000')
    assert chart.series[1].format.fill.fore_color.rgb == gen.hex_to_rgb('#00AA00')
