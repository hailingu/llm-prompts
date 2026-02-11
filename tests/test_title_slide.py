import os
import sys
import json
ROOT = os.getcwd()
PKG = os.path.join(ROOT, 'skills', 'ppt-generator')
if ROOT not in sys.path:
    sys.path.insert(0, ROOT)
if PKG not in sys.path:
    sys.path.insert(0, PKG)

from pptx import Presentation
from ppt_generator.grid import GridSystem
from ppt_generator.renderers import render_slide


def test_title_slide_renders_title_and_meta(tmp_path):
    sem = json.load(open('docs/presentations/storage-frontier-20260211/slides_semantic.json', encoding='utf-8'))
    design = json.load(open('docs/presentations/storage-frontier-20260211/design_spec.json', encoding='utf-8'))
    slides = sem.get('slides', [])
    assert slides
    sd = slides[0]
    prs = Presentation()
    grid = GridSystem(design)
    total = len(slides)
    # render
    render_slide(prs, sd, design, grid, sem.get('sections', []), 1, total)
    s = prs.slides[-1]
    # Expect at least 3 text-containing shapes (title, subtitle, metadata)
    text_shapes = [sh for sh in s.shapes if hasattr(sh, 'text') and (sh.text or '').strip()]
    assert len(text_shapes) >= 3, f"expected >=3 text shapes on title slide, got {len(text_shapes)}"
    # one of them must contain the deck title
    titles = [sh.text for sh in text_shapes]
    assert any('存储行业前沿' in t for t in titles), f"title text missing in {titles}"
