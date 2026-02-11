import importlib.util
import os
import sys

ROOT = os.getcwd()
PKG_PATH = os.path.join(ROOT, 'skills', 'ppt-generator')
if ROOT not in sys.path:
    sys.path.insert(0, ROOT)
if PKG_PATH not in sys.path:
    sys.path.insert(0, PKG_PATH)

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE

from ppt_generator.renderers import (
    render_region_chart, render_region_comparison, render_region_kpi,
    render_region_callout, render_region_progression, render_region_bullets,
    REGION_RENDERERS
)
from ppt_generator.grid import GridSystem


def make_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    spec = {}
    grid = GridSystem(spec)
    return prs, slide, spec, grid


def test_each_region_renderer_basic_runs():
    prs, slide, spec, grid = make_slide()
    # chart
    chart_cfg = {'labels': ['A', 'B', 'C'], 'series': [{'name': 's', 'data': [1,2,3]}]}
    render_region_chart(slide, chart_cfg, (grid.margin_h, 1.0, grid.usable_w, 2.5), spec)
    # comparison
    comp = [{'label': 'A', 'attributes': {'x': 1}}, {'label': 'B', 'attributes': {'x': 2}}]
    render_region_comparison(slide, comp, (grid.margin_h, 3.8, grid.usable_w, 2.0), spec)
    # kpis
    kpis = [{'label': 'K1', 'value': '10'}, {'label': 'K2', 'value': '20'}]
    render_region_kpi(slide, kpis, (grid.margin_h, 6.0, grid.usable_w, 0.9), spec)
    # callouts
    calls = [{'text': 'Callout 1'}, {'text': 'Callout 2'}]
    render_region_callout(slide, calls, (grid.margin_h, 7.2, grid.usable_w, 1.6), spec)
    # progression
    timeline = [{'phase': 'P1'}, {'phase': 'P2'}, {'phase': 'P3'}]
    render_region_progression(slide, timeline, (grid.margin_h, 9.0, grid.usable_w, 1.4), spec)
    # bullets
    bullets = ['a', 'b', 'c']
    render_region_bullets(slide, bullets, (grid.margin_h, 10.6, grid.usable_w, 2.0), spec)

    # assert some shapes were added
    assert len(slide.shapes) > 0


def test_region_renderers_registry_has_entries():
    expected = {'chart', 'comparison_table', 'kpi_row', 'callout_stack', 'progression', 'bullet_list'}
    assert expected.issubset(set(REGION_RENDERERS.keys()))


def test_two_regions_do_not_overlap():
    prs, slide, spec, grid = make_slide()
    # render two regions side by side: left-50 and right-50
    bounds_l = (grid.margin_h, 1.0, grid.usable_w * 0.5 - 0.02, 2.0)
    bounds_r = (grid.margin_h + grid.usable_w * 0.5 + 0.02, 1.0, grid.usable_w * 0.5 - 0.02, 2.0)
    render_region_bullets(slide, ['one', 'two'], bounds_l, spec)
    render_region_bullets(slide, ['three'], bounds_r, spec)
    # collect left positions of first textbox in each region
    lefts = [sh.left for sh in slide.shapes if hasattr(sh, 'left')]
    # ensure there are at least two distinct lefts (non-overlapping by position)
    assert len(set(lefts)) >= 2
