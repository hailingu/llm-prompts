import json
import importlib.util
from pptx import Presentation

from pathlib import Path

# Import packaged renderers
import sys
sys.path.insert(0, str(Path('skills/ppt-generator').resolve()))
from ppt_generator import renderers


def test_p1_integration_assertion_and_insight_rendered(tmp_path):
    semantic = json.load(open('docs/presentations/storage-frontier-20260211/slides_semantic.json', encoding='utf-8'))
    design = json.load(open('docs/presentations/storage-frontier-20260211/design_spec.json', encoding='utf-8'))

    grid = renderers.GridSystem(design)
    prs = Presentation()
    prs.slide_width = renderers.Inches(grid.slide_w)
    prs.slide_height = renderers.Inches(grid.slide_h)

    sections = semantic.get('sections', [])
    slides = semantic.get('slides', [])

    for i, sd in enumerate(slides, 1):
        if sd.get('slide_type') == 'title':
            sd['_deck_title'] = semantic.get('title', '')
        renderers.render_slide(prs, sd, design, grid, sections, i, len(slides))

    out = tmp_path / 'p1_demo.pptx'
    prs.save(str(out))

    p = Presentation(str(out))

    # Verify each slide with assertion/insight has corresponding shapes
    for i, sd in enumerate(slides, 1):
        if sd.get('assertion'):
            slide = p.slides[i - 1]
            texts = [sh.text.strip() for sh in slide.shapes if hasattr(sh, 'text') and sh.text and sh.text.strip()]
            assert any(sd.get('assertion') in t for t in texts), f"Assertion not found on slide {i}: {texts}"
        if sd.get('insight'):
            slide = p.slides[i - 1]
            texts = [sh.text.strip() for sh in slide.shapes if hasattr(sh, 'text') and sh.text and sh.text.strip()]
            assert any('💡' in t and sd.get('insight') in t for t in texts), f"Insight not found on slide {i}: {texts}"
