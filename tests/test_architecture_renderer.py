import os
import sys
ROOT = os.getcwd()
PKG_PATH = os.path.join(ROOT, 'skills', 'ppt-generator')
if ROOT not in sys.path:
    sys.path.insert(0, ROOT)
if PKG_PATH not in sys.path:
    sys.path.insert(0, PKG_PATH)

from pptx import Presentation
from pptx.util import Inches
from ppt_generator.renderers import render_region_architecture


def test_render_region_architecture_basic():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    data = {
        'nodes': [
            {'id': 'n1', 'label': 'Service A', 'x': 0.05, 'y': 0.1, 'w': 1.8, 'h': 0.7, 'style': 'primary'},
            {'id': 'n2', 'label': 'Service B', 'x': 0.5, 'y': 0.1, 'w': 1.8, 'h': 0.7, 'style': 'secondary'},
        ],
        'edges': [
            {'from': 'n1', 'to': 'n2', 'label': 'calls'}
        ]
    }
    bounds = (1.0, 1.0, 10.0, 4.0)
    before = len(slide.shapes)
    render_region_architecture(slide, data, bounds, spec={})
    after = len(slide.shapes)
    assert after - before >= 3  # two nodes + one connector

    # verify node labels present
    texts = [s.text_frame.text for s in slide.shapes if hasattr(s, 'text_frame') and s.text_frame is not None and s.text_frame.text]
    assert 'Service A' in texts
    assert 'Service B' in texts


if __name__ == '__main__':
    print('Run: python3 -m pytest tests/test_architecture_renderer.py')